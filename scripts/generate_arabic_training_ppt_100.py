from __future__ import annotations

import zipfile
import json
import re
from datetime import date
from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
ASSETS_DIR = ROOT / "attached_assets"
OUTPUT_DIR = ROOT / "artifacts" / "xray-academy"
OUTPUT_FILE = OUTPUT_DIR / "XRay_Training_Program_AR_100plus.pptx"
KNOWLEDGE_PROFILE_FILE = OUTPUT_DIR / "training_kb_profile.json"
LOGO_DIR = ROOT / "artifacts" / "xray-academy" / "logo_candidates"
LOGO_FILE = LOGO_DIR / "rapiscan_ase_logo.png"
REFERENCE_PPT = ASSETS_DIR / "LZBV_Full_Service_Course_Rapiscan_7-2-19_1784135178108.pptx"


PRIMARY = RGBColor(0x0D, 0x3B, 0x66)
SECONDARY = RGBColor(0x18, 0x72, 0xB4)
ACCENT = RGBColor(0xF2, 0x5C, 0x2A)
TEXT = RGBColor(0x12, 0x1A, 0x24)
MUTED = RGBColor(0x55, 0x63, 0x73)
BG = RGBColor(0xF6, 0xFA, 0xFD)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


MODULE_BLUEPRINTS = [
    {
        "name": "أساسيات النظام والبنية التقنية",
        "keywords": ["system", "architecture", "component", "subsystem", "overview", "controller", "module", "modes"],
        "defaults": [
            "التعريف بنظام فحص المركبات ودوره في الأمن التشغيلي",
            "المكونات الأساسية: المصدر، الكاشف، منصة التحكم، وحدات الأمان",
            "دورة الإشارة من التقاط البيانات حتى عرض الصورة",
            "العوامل المؤثرة على اعتمادية النظام في الموقع",
            "المتطلبات المسبقة للتشغيل الآمن قبل بدء الوردية",
            "نقاط الفشل المحتملة في البنية الكهربائية والميكانيكية",
            "أفضل ممارسات التوثيق الفني لسجل النظام",
            "مؤشرات الأداء الأساسية KPI للتشغيل اليومي",
        ],
        "case": "حالة: تذبذب مفاجئ في قراءة إحدى قنوات الكشف أثناء بداية الوردية",
        "quiz": "سؤال: ما أول مكوّن يتم عزله عند ظهور خطأ إشارة متكرر؟",
    },
    {
        "name": "مبادئ الفيزياء الإشعاعية والتصوير",
        "keywords": ["radiation", "x-ray", "image", "attenuation", "backscatter", "density", "energy", "detector"],
        "defaults": [
            "فكرة التأين وتأثير طاقة الأشعة السينية على المواد",
            "الفرق بين الامتصاص، المرور، والتشتت في الصور",
            "العلاقة بين الكثافة والعدد الذري وتباين الألوان",
            "تأثير سمك الجسم على دقة الاكتشاف",
            "فهم مناطق فقد البيانات داخل الصورة",
            "اختلاف استجابة المواد العضوية وغير العضوية",
            "أسباب التشويش البصري وكيفية تقليله",
            "تفسير الصور المركبة متعددة الطبقات",
        ],
        "case": "حالة: جسم مختلط المادة يعطي نتائج غير حاسمة بين النقل والارتداد",
        "quiz": "سؤال: متى نستخدم إعادة المسح بدل اتخاذ قرار فوري؟",
    },
    {
        "name": "بروتوكولات السلامة والإجراءات التنظيمية",
        "keywords": ["safety", "hazard", "warning", "caution", "protocol", "dose", "boundary", "compliance"],
        "defaults": [
            "حدود الجرعات المسموح بها ومبدأ ALARA",
            "إدارة مناطق التحذير والعزل ومنع الوصول",
            "التحقق من سلامة إشارات المرور الضوئية",
            "تسلسل الإيقاف الطارئ وإعادة التهيئة",
            "إجراءات التدقيق الداخلي في الامتثال",
            "السجلات الإلزامية للحوادث وملاحظات السلامة",
            "أدوار فريق التشغيل مقابل فريق الصيانة",
            "ممارسات الحد من التعرض غير المقصود",
        ],
        "case": "حالة: إنذار إشعاعي غير طبيعي أثناء مرور مركبة في النفق",
        "quiz": "سؤال: ما الخطوة الصحيحة قبل استئناف التشغيل بعد الإنذار؟",
    },
    {
        "name": "التشغيل القياسي ومراقبة الجودة",
        "keywords": ["operation", "scan", "start", "traffic", "quality", "workflow", "procedure", "control"],
        "defaults": [
            "قائمة التحقق قبل التشغيل خطوة بخطوة",
            "معايرة بداية اليوم وتأكيد الثبات",
            "إدارة حركة المركبات داخل المسار بأمان",
            "إطلاق المسح ومراقبة مؤشرات النظام الحية",
            "معايير جودة الصورة المقبولة تشغيليًا",
            "التعامل مع فقد الإطار أو تأخر المعالجة",
            "إنهاء المهمة وحفظ نتائج الفحص",
            "مراجعة نهاية الوردية وتحسين الإعدادات",
        ],
        "case": "حالة: جودة الصورة تهبط تدريجيًا مع زيادة سرعة المرور",
        "quiz": "سؤال: ما المؤشر الأول الذي تراجعه عند انخفاض الوضوح؟",
    },
    {
        "name": "تفسير الصور وتقييم المخاطر التشغيلية",
        "keywords": ["interpret", "analysis", "threat", "risk", "pattern", "organic", "inorganic", "decision"],
        "defaults": [
            "قراءة الأنماط العامة للأجسام عالية الكثافة",
            "تمييز الإشارات الكاذبة عن المؤشرات الحقيقية",
            "التقييم السياقي للصورة مع بيانات المركبة",
            "استخدام التكبير والمعالجة لتحسين القراءة",
            "الحدود الفنية للنظام في الكشف المعقد",
            "بناء قرار تشغيلي آمن عند عدم اليقين",
            "تصنيف الحالات إلى منخفضة ومتوسطة وعالية الخطورة",
            "تحويل الحالات الحساسة للمراجعة المتخصصة",
        ],
        "case": "حالة: صورة تحمل نمطًا ملتبسًا في منطقة محمية هيكليًا",
        "quiz": "سؤال: متى تتحول الحالة مباشرة إلى تحقيق أمني؟",
    },
    {
        "name": "الأعطال المتكررة والتحقق الأولي",
        "keywords": ["fault", "alarm", "error", "failure", "diagnostic", "code", "status", "reset"],
        "defaults": [
            "أعطال المستشعرات وطرق العزل السريع",
            "انخفاض أداء وحدة الجهد العالي وأعراضه",
            "مشكلات الاتصال بين وحدات النظام",
            "انقطاع التخزين أو تأخر الكتابة في السجل",
            "فشل المعايرة ومتى نعيد الإجراء",
            "إجراءات الاسترجاع الآمن بعد توقف مفاجئ",
            "مصفوفة قرار: إصلاح موقعي أم تصعيد",
            "توثيق سبب الجذر لكل عطل متكرر",
        ],
        "case": "حالة: رسالة Fault متقطعة تختفي بعد إعادة التشغيل",
        "quiz": "سؤال: ما الخطر من تكرار إعادة التشغيل بدون تحليل سبب الجذر؟",
    },
    {
        "name": "الصيانة الوقائية والجداول الدورية",
        "keywords": ["maintenance", "preventive", "inspection", "schedule", "service", "checklist", "test", "repair"],
        "defaults": [
            "خطة الصيانة اليومية والأسبوعية والشهرية",
            "فحص التبريد والتهوية وحالة المرشحات",
            "اختبار الكابلات ونقاط التأريض الحرجة",
            "فحص ميكانيكي للأجزاء المتحركة والبوابات",
            "مراجعة مؤشرات الصحة العامة للنظام",
            "إدارة قطع الغيار ذات الأولوية العالية",
            "معايير قبول الأداء بعد كل صيانة",
            "التنسيق بين الصيانة والتشغيل لتقليل التوقف",
        ],
        "case": "حالة: ارتفاع حرارة مستمر بعد عملية تنظيف دورية",
        "quiz": "سؤال: ما الحد الذي يستوجب إيقاف النظام للصيانة الفورية؟",
    },
    {
        "name": "الاستجابة للطوارئ وخطط الاستمرارية",
        "keywords": ["emergency", "incident", "evacuation", "shutdown", "contingency", "crisis", "response", "recovery"],
        "defaults": [
            "خريطة الطوارئ التشغيلية ومسؤوليات الفريق",
            "سيناريو فقد الطاقة أثناء عملية المسح",
            "سيناريو إنذار إشعاعي حرج وإخلاء المنطقة",
            "التواصل مع فرق الأمن والسلامة المؤسسية",
            "استمرارية الخدمة في وضع التشغيل الاحتياطي",
            "إدارة الرسائل الإعلامية الداخلية أثناء الحوادث",
            "توثيق الدروس المستفادة بعد كل حدث",
            "اختبارات دورية لخطة الطوارئ",
        ],
        "case": "حالة: تعطل مصدر رئيسي أثناء ذروة تشغيلية",
        "quiz": "سؤال: ما ترتيب الأولويات في أول 3 دقائق للطوارئ؟",
    },
    {
        "name": "سيناريوهات ميدانية ودراسات حالة",
        "keywords": ["scenario", "field", "practical", "exercise", "demonstration", "laboratory", "hands-on", "case"],
        "defaults": [
            "بناء سيناريو تدريبي واقعي من بيانات حقيقية",
            "توزيع الأدوار بين المشغل والمشرف والفني",
            "التنفيذ تحت ضغط الوقت مع الحفاظ على السلامة",
            "تحليل القرار المتخذ أثناء الحالة",
            "نقاط القوة والأخطاء الشائعة لدى المتدربين",
            "إعادة التمرين مع تحسين الأداء",
            "توحيد منهج الاستجابة داخل الفريق",
            "تحويل الملاحظات إلى إجراءات تشغيل معيارية",
        ],
        "case": "حالة: مركبة عالية التعقيد مع إشارات متضاربة في الصورة",
        "quiz": "سؤال: كيف توازن بين سرعة القرار ودقته تحت الضغط؟",
    },
    {
        "name": "التقييم النهائي وخطة التحسين",
        "keywords": ["assessment", "quiz", "exam", "evaluation", "objective", "improvement", "score", "pass"],
        "defaults": [
            "هيكل الاختبار النظري ومعايير النجاح",
            "روبرك تقييم الأداء العملي في الموقع",
            "قياس الكفاءة قبل وبعد البرنامج",
            "تحليل الفجوات المعرفية لكل متدرب",
            "خطة علاج فردية للحالات المتأخرة",
            "خطة تطوير الفريق خلال 90 يومًا",
            "مؤشرات متابعة ما بعد التدريب",
            "ربط التدريب بمخرجات الأداء المؤسسي",
        ],
        "case": "حالة: متدرب يجتاز النظري ويفشل في التطبيق الميداني",
        "quiz": "سؤال: ما الإجراء التطويري الأكثر أثرًا خلال أول شهر؟",
    },
]


def ensure_logo() -> Path | None:
    LOGO_DIR.mkdir(parents=True, exist_ok=True)
    if LOGO_FILE.exists():
        return LOGO_FILE

    if not REFERENCE_PPT.exists():
        return None

    with zipfile.ZipFile(REFERENCE_PPT) as zf:
        candidate = "ppt/media/image2.png"
        if candidate in zf.namelist():
            LOGO_FILE.write_bytes(zf.read(candidate))
            return LOGO_FILE

    return None


def pick_training_images(max_count: int = 40) -> list[Path]:
    imgs = sorted(ASSETS_DIR.glob("image_*.png"))
    return [p for p in imgs if p.is_file()][:max_count]


def _read_text_sources() -> list[str]:
    texts: list[str] = []
    patterns = ["Pasted-*.txt", "*.md"]
    for pattern in patterns:
        for path in sorted(ASSETS_DIR.glob(pattern))[:40]:
            try:
                content = path.read_text(encoding="utf-8", errors="ignore").strip()
            except Exception:
                continue
            if content:
                texts.append(content)
    return texts


def _read_reference_ppt_text() -> list[str]:
    texts: list[str] = []
    ppt_candidates = [
        REFERENCE_PPT,
        ASSETS_DIR / "ASE_9200_1_Rv_A_ZBV__FSM_instructor_1784134504895.pptx",
    ]
    for ppt in ppt_candidates:
        if not ppt.exists():
            continue
        try:
            with zipfile.ZipFile(ppt) as zf:
                for name in zf.namelist():
                    if not re.match(r"^ppt/slides/slide\\d+\\.xml$", name):
                        continue
                    raw = zf.read(name).decode("utf-8", errors="ignore")
                    chunks = re.findall(r"<a:t>(.*?)</a:t>", raw)
                    line = " ".join(c.strip() for c in chunks if c.strip())
                    if line:
                        texts.append(line)
        except Exception:
            continue
    return texts


def _to_lines(corpus: list[str]) -> list[str]:
    lines: list[str] = []
    for block in corpus:
        for part in re.split(r"[\n\r\.\!\?؛:]+", block):
            txt = re.sub(r"\s+", " ", part).strip(" -•\t")
            if 18 <= len(txt) <= 180:
                lines.append(txt)
    # Preserve order while deduplicating.
    return list(dict.fromkeys(lines))


def _arabize_topic(text: str) -> str:
    replacements = {
        "safety": "السلامة",
        "radiation": "الإشعاع",
        "maintenance": "الصيانة",
        "operation": "التشغيل",
        "troubleshooting": "استكشاف الأعطال",
        "quality": "الجودة",
        "training": "التدريب",
        "system": "النظام",
        "component": "المكوّن",
        "procedure": "الإجراء",
    }
    out = text
    for en, ar in replacements.items():
        out = re.sub(rf"\b{en}\b", ar, out, flags=re.IGNORECASE)
    return out


def build_modules_from_knowledge_base() -> tuple[list[dict[str, object]], dict[str, object]]:
    corpus = _read_text_sources() + _read_reference_ppt_text()
    lines = _to_lines(corpus)

    modules: list[dict[str, object]] = []
    for blueprint in MODULE_BLUEPRINTS:
        matched: list[str] = []
        for line in lines:
            low = line.lower()
            if any(k in low for k in blueprint["keywords"]):
                candidate = _arabize_topic(line)
                if candidate not in matched:
                    matched.append(candidate)
            if len(matched) >= 8:
                break

        # Fill from defaults if KB did not provide enough direct lines.
        for default in blueprint["defaults"]:
            if len(matched) >= 8:
                break
            if default not in matched:
                matched.append(default)

        modules.append(
            {
                "name": blueprint["name"],
                "topics": matched[:8],
                "case": blueprint["case"],
                "quiz": blueprint["quiz"],
            }
        )

    profile = {
        "corpus_sources_count": len(corpus),
        "candidate_lines_count": len(lines),
        "modules_count": len(modules),
        "generated_at": date.today().isoformat(),
        "module_topic_counts": {m["name"]: len(m["topics"]) for m in modules},
    }
    return modules, profile


def set_text(shape, text: str, size: int, color: RGBColor, bold: bool = False, align: PP_ALIGN = PP_ALIGN.RIGHT) -> None:
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.name = "Arial"
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold


def add_background(slide) -> None:
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)


def add_header(slide, title: str, subtitle: str = "", logo_path: Path | None = None) -> None:
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.95))
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY
    header.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(3.2), Inches(0.08), Inches(9.7), Inches(0.45))
    set_text(title_box, title, 22, WHITE, bold=True)

    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(3.2), Inches(0.50), Inches(9.7), Inches(0.3))
        set_text(subtitle_box, subtitle, 12, WHITE)

    if logo_path and logo_path.exists():
        slide.shapes.add_picture(str(logo_path), Inches(0.25), Inches(0.10), Inches(2.7), Inches(0.70))


def add_footer(slide, page_num: int, total: int) -> None:
    footer = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.12), Inches(13.33), Inches(0.38))
    footer.fill.solid()
    footer.fill.fore_color.rgb = RGBColor(0xE8, 0xF0, 0xF7)
    footer.line.fill.background()

    right = slide.shapes.add_textbox(Inches(7.0), Inches(7.17), Inches(6.0), Inches(0.2))
    set_text(right, "برنامج تدريبي احترافي | تشغيل وصيانة أنظمة فحص المركبات", 10, MUTED)

    left = slide.shapes.add_textbox(Inches(0.2), Inches(7.17), Inches(2.0), Inches(0.2))
    set_text(left, f"{page_num}/{total}", 10, MUTED, align=PP_ALIGN.LEFT)


def add_bullets(slide, lines: Iterable[str], title: str) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.65), Inches(1.15), Inches(7.45), Inches(0.65))
    set_text(title_box, title, 28, SECONDARY, bold=True)

    body = slide.shapes.add_textbox(Inches(0.8), Inches(1.95), Inches(7.3), Inches(4.9))
    tf = body.text_frame
    tf.clear()
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"• {line}"
        p.alignment = PP_ALIGN.RIGHT
        p.level = 0
        run = p.runs[0]
        run.font.name = "Arial"
        run.font.size = Pt(22 if idx == 0 else 20)
        run.font.color.rgb = TEXT


def add_image(slide, image_path: Path, caption: str) -> None:
    frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.35), Inches(1.55), Inches(4.6), Inches(4.95))
    frame.fill.solid()
    frame.fill.fore_color.rgb = WHITE
    frame.line.color.rgb = RGBColor(0xD2, 0xDE, 0xE8)

    slide.shapes.add_picture(str(image_path), Inches(8.55), Inches(1.75), Inches(4.2), Inches(4.25))

    cap = slide.shapes.add_textbox(Inches(8.45), Inches(6.10), Inches(4.4), Inches(0.35))
    set_text(cap, caption, 11, MUTED, align=PP_ALIGN.CENTER)


def add_trainer_notes(slide, note_text: str) -> None:
    notes = slide.notes_slide.notes_text_frame
    notes.clear()
    notes.text = "شرح المدرب:\n" + note_text


def add_cover(prs: Presentation, logo_path: Path | None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = PRIMARY
    bg.line.fill.background()

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(5.45), Inches(7.2), Inches(0.24))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()

    title = slide.shapes.add_textbox(Inches(0.9), Inches(1.65), Inches(11.8), Inches(1.5))
    set_text(title, "الدبلوم التطبيقي الشامل", 44, WHITE, bold=True)

    subtitle = slide.shapes.add_textbox(Inches(0.9), Inches(3.0), Inches(11.8), Inches(1.2))
    set_text(subtitle, "تشغيل وصيانة وسلامة أنظمة فحص المركبات بالأشعة السينية", 22, WHITE)

    meta = slide.shapes.add_textbox(Inches(0.9), Inches(6.1), Inches(11.8), Inches(0.5))
    set_text(meta, f"نسخة المدرب | تاريخ الإصدار: {date.today().isoformat()} | الحد الأدنى 100 شريحة", 13, WHITE)

    if logo_path and logo_path.exists():
        slide.shapes.add_picture(str(logo_path), Inches(0.95), Inches(0.25), Inches(3.3), Inches(0.85))


def add_outline(prs: Presentation, logo_path: Path | None, total: int, modules: list[dict[str, object]]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_header(slide, "الخطة العامة للبرنامج", "10 محاور تدريبية تفصيلية", logo_path)
    add_bullets(
        slide,
        [
            *[
                f"المحور {idx}: {str(module['name'])}"
                for idx, module in enumerate(modules, start=1)
            ],
        ],
        "خارطة المحتوى",
    )
    add_footer(slide, 2, total)
    add_trainer_notes(slide, "ابدأ بشرح ناتج التعلم النهائي، ثم اربط كل محور بالممارسة الميدانية المتوقعة من المتدرب.")


def build_slides(prs: Presentation, logo_path: Path | None, images: list[Path], total: int, modules: list[dict[str, object]]) -> None:
    image_count = len(images)
    page = 3

    for module_index, module in enumerate(modules, start=1):
        module_name = str(module["name"])
        topics = list(module["topics"])

        intro = prs.slides.add_slide(prs.slide_layouts[6])
        add_background(intro)
        add_header(intro, f"المحور {module_index}: {module_name}", "شريحة تمهيدية", logo_path)
        add_bullets(
            intro,
            [
                "مخرجات التعلم لهذا المحور موضوعة على مستوى المهارة التطبيقية.",
                "سيتم ربط كل مفهوم بنموذج واقعي من بيئة التشغيل.",
                "التركيز على السلامة والانضباط الإجرائي وجودة القرار.",
                "في نهاية المحور يوجد تقييم قصير سريع.",
            ],
            "أهداف المحور",
        )
        if image_count:
            add_image(intro, images[(page - 1) % image_count], "تمهيد بصري للمحور")
        add_footer(intro, page, total)
        add_trainer_notes(intro, f"قدّم هذا المحور على أنه سلسلة قرارات تشغيلية مترابطة وليس معلومات منفصلة. المحور: {module_name}")
        page += 1

        for topic_idx, topic in enumerate(topics, start=1):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            add_background(slide)
            add_header(slide, f"المحور {module_index}: {module_name}", f"الشريحة التفصيلية {topic_idx}/8", logo_path)
            add_bullets(
                slide,
                [
                    topic,
                    "الشرح العملي: كيف ينعكس هذا المفهوم على قرار المشغل داخل غرفة التحكم.",
                    "نقطة حرجة: الخطأ الأكثر شيوعًا وكيفية تجنبه قبل أن يتحول إلى حادثة.",
                    "مؤشر قياس: ما الدليل الذي يثبت أن الإجراء نُفِّذ بشكل صحيح.",
                ],
                f"تفصيل مهني: {topic}",
            )
            if image_count:
                add_image(slide, images[(page - 1) % image_count], f"دعم بصري - {module_name}")
            add_footer(slide, page, total)
            add_trainer_notes(
                slide,
                "ابدأ بالمفهوم ثم اعرض مثالًا ميدانيًا، بعدها اطلب من المتدرب وصف القرار الصحيح خلال 30 ثانية.",
            )
            page += 1

        case_slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_background(case_slide)
        add_header(case_slide, f"المحور {module_index}: {module_name}", "دراسة حالة تطبيقية", logo_path)
        add_bullets(
            case_slide,
            [
                str(module["case"]),
                "المطلوب من الفريق: تحليل السبب المحتمل خلال دقيقتين.",
                "الإجراء: صياغة خطة استجابة من 5 خطوات واضحة.",
                "المخرجات: تحديد القرار النهائي ومعيار التحقق بعد التنفيذ.",
            ],
            "حالة ميدانية",
        )
        if image_count:
            add_image(case_slide, images[(page - 1) % image_count], "مشهد داعم لدراسة الحالة")
        add_footer(case_slide, page, total)
        add_trainer_notes(case_slide, "قسّم المتدربين إلى مجموعات صغيرة، واطلب من كل مجموعة تقديم خطة مختلفة ثم ناقش الأفضلية وفق معيار السلامة.")
        page += 1

        quiz_slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_background(quiz_slide)
        add_header(quiz_slide, f"المحور {module_index}: {module_name}", "اختبار سريع", logo_path)
        add_bullets(
            quiz_slide,
            [
                str(module["quiz"]),
                "A) إجراء فوري غير موثق",
                "B) تحليل السبب ثم قرار آمن موثق",
                "C) الاستمرار دون إيقاف",
                "D) تحويل الحالة دون تحقق",
            ],
            "تحقق من الفهم",
        )
        if image_count:
            add_image(quiz_slide, images[(page - 1) % image_count], "صورة مرتبطة بسؤال الاختبار")
        add_footer(quiz_slide, page, total)
        add_trainer_notes(quiz_slide, "اقرأ السؤال بصوت واضح، امنح 20 ثانية للتفكير، ثم ناقش سبب صحة الإجابة وليس الإجابة فقط.")
        page += 1


def add_closing(prs: Presentation, logo_path: Path | None, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_header(slide, "الخاتمة وخطة ما بعد التدريب", "تعزيز استدامة الأداء", logo_path)
    add_bullets(
        slide,
        [
            "اعتماد خطة متابعة 30/60/90 يوم لقياس أثر التدريب.",
            "تنفيذ إعادة تقييم عملي دوري للحالات الحرجة.",
            "تحديث إجراءات التشغيل بناءً على الدروس المستفادة.",
            "تطوير مسار متقدم للمتفوقين ومسار دعم للمتعثرين.",
        ],
        "الخطوات التالية",
    )
    add_footer(slide, total, total)
    add_trainer_notes(slide, "اختتم البرنامج برسالة واضحة: السلامة والانضباط الإجرائي هما معيار الجودة الحقيقي.")


def main() -> None:
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    logo_path = ensure_logo()
    images = pick_training_images(max_count=60)
    modules, profile = build_modules_from_knowledge_base()

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    total_slides = 1 + 1 + (len(modules) * 11) + 1
    add_cover(prs, logo_path)
    add_outline(prs, logo_path, total_slides, modules)
    build_slides(prs, logo_path, images, total_slides, modules)
    add_closing(prs, logo_path, total_slides)

    target_file = OUTPUT_FILE
    try:
        prs.save(str(target_file))
    except PermissionError:
        stamped = date.today().strftime("%Y%m%d")
        target_file = OUTPUT_DIR / f"XRay_Training_Program_AR_100plus_{stamped}.pptx"
        prs.save(str(target_file))

    print(f"Saved: {target_file}")
    print(f"Slides: {len(prs.slides)}")
    print(f"Logo used: {logo_path if logo_path else 'None'}")
    print(f"Images used: {len(images)}")
    KNOWLEDGE_PROFILE_FILE.write_text(
        json.dumps(profile, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    print(f"Knowledge profile: {KNOWLEDGE_PROFILE_FILE}")


if __name__ == "__main__":
    main()
