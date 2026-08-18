from __future__ import annotations

from datetime import date
from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
ASSETS_DIR = ROOT / "attached_assets"
OUTPUT_DIR = ROOT / "artifacts" / "xray-academy"
OUTPUT_FILE = OUTPUT_DIR / "XRay_Training_Program_AR.pptx"


PRIMARY = RGBColor(0x0D, 0x3B, 0x66)
ACCENT = RGBColor(0xF9, 0x57, 0x38)
SOFT = RGBColor(0xF2, 0xF7, 0xFA)
TEXT = RGBColor(0x12, 0x1A, 0x24)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def pick_images(max_count: int = 12) -> list[Path]:
    preferred = [
        p
        for p in sorted(ASSETS_DIR.glob("image_*.png"))
        if p.is_file()
    ]
    return preferred[:max_count]


def set_ar_title(shape, text: str) -> None:
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.RIGHT
    run = p.runs[0]
    run.font.name = "Arial"
    run.font.size = Pt(34)
    run.font.bold = True
    run.font.color.rgb = WHITE


def set_ar_subtitle(shape, text: str) -> None:
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.RIGHT
    run = p.runs[0]
    run.font.name = "Arial"
    run.font.size = Pt(18)
    run.font.color.rgb = WHITE


def add_bg_rect(slide, color: RGBColor) -> None:
    rect = slide.shapes.add_shape(
        1,
        Inches(0),
        Inches(0),
        Inches(13.33),
        Inches(7.5),
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    slide.shapes._spTree.remove(rect._element)
    slide.shapes._spTree.insert(2, rect._element)


def add_section_band(slide, text: str) -> None:
    band = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(0.85))
    band.fill.solid()
    band.fill.fore_color.rgb = PRIMARY
    band.line.fill.background()

    tf = band.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.RIGHT
    run = p.runs[0]
    run.font.name = "Arial"
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = WHITE


def add_ar_bullets(slide, title: str, bullets: Iterable[str], footer: str = "") -> None:
    add_section_band(slide, title)

    body = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(7.4), Inches(5.8))
    tf = body.text_frame
    tf.clear()

    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {line}"
        p.alignment = PP_ALIGN.RIGHT
        p.level = 0
        run = p.runs[0]
        run.font.name = "Arial"
        run.font.size = Pt(24 if i == 0 else 21)
        run.font.color.rgb = TEXT

    if footer:
        f = slide.shapes.add_textbox(Inches(0.8), Inches(6.9), Inches(12.0), Inches(0.4))
        ft = f.text_frame
        ft.clear()
        p = ft.paragraphs[0]
        p.text = footer
        p.alignment = PP_ALIGN.RIGHT
        run = p.runs[0]
        run.font.name = "Arial"
        run.font.size = Pt(12)
        run.font.color.rgb = RGBColor(0x55, 0x63, 0x73)


def add_image_block(slide, image_path: Path, caption: str) -> None:
    frame = slide.shapes.add_shape(1, Inches(8.5), Inches(1.4), Inches(4.4), Inches(4.8))
    frame.fill.solid()
    frame.fill.fore_color.rgb = SOFT
    frame.line.color.rgb = RGBColor(0xD2, 0xDE, 0xE8)

    slide.shapes.add_picture(str(image_path), Inches(8.7), Inches(1.6), Inches(4.0), Inches(4.2))

    cap = slide.shapes.add_textbox(Inches(8.6), Inches(5.95), Inches(4.1), Inches(0.6))
    ctf = cap.text_frame
    ctf.clear()
    p = ctf.paragraphs[0]
    p.text = caption
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.name = "Arial"
    run.font.size = Pt(12)
    run.font.color.rgb = RGBColor(0x2A, 0x36, 0x44)


def make_title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide, PRIMARY)

    accent = slide.shapes.add_shape(1, Inches(0.8), Inches(5.4), Inches(5.5), Inches(0.35))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.8), Inches(2.0))
    set_ar_title(title_box, "البرنامج التدريبي الشامل")

    subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.3), Inches(11.8), Inches(1.2))
    set_ar_subtitle(
        subtitle_box,
        "تشغيل وصيانة نظام فحص المركبات بالأشعة السينية | نسخة المدرب",
    )

    meta = slide.shapes.add_textbox(Inches(0.8), Inches(6.2), Inches(11.8), Inches(0.6))
    tf = meta.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = f"تاريخ الإصدار: {date.today().isoformat()}  |  مدة البرنامج: 6 ساعات"
    p.alignment = PP_ALIGN.RIGHT
    run = p.runs[0]
    run.font.name = "Arial"
    run.font.size = Pt(14)
    run.font.color.rgb = WHITE


def make_content(prs: Presentation, images: list[Path]) -> None:
    slides_data = [
        (
            "أهداف البرنامج",
            [
                "فهم بنية نظام الفحص ووظائف المكونات الرئيسية",
                "تطبيق إجراءات السلامة الإشعاعية قبل وأثناء وبعد التشغيل",
                "تنفيذ خطوات التشغيل القياسية وتحسين جودة الصور",
                "تشخيص الأعطال الشائعة والتصعيد الفني الصحيح",
            ],
        ),
        (
            "المتطلبات قبل التشغيل",
            [
                "التحقق من جاهزية الموقع ومنطقة الأمان",
                "فحص مصدر الطاقة ووحدة التحكم وملحقات الحماية",
                "مراجعة حالة الإنذارات والسجل اليومي",
                "تأكيد توفر معدات الوقاية الشخصية للفريق",
            ],
        ),
        (
            "مبادئ التصوير بالأشعة السينية",
            [
                "العلاقة بين كثافة المادة وسلوك الامتصاص/المرور",
                "التمييز بين المواد العضوية وغير العضوية عبر التباين",
                "تأثير زاوية المسح وسرعة المرور على جودة الصورة",
                "متى نعيد المسح ومتى نرفع الحالة للتقييم",
            ],
        ),
        (
            "السلامة الإشعاعية",
            [
                "الالتزام بحدود الجرعة والسياسات التنظيمية",
                "إدارة منطقة العزل ومنع الوصول غير المصرح",
                "الاستجابة الفورية لأي إنذار إشعاعي غير طبيعي",
                "توثيق كل حادثة في سجل السلامة المعتمد",
            ],
        ),
        (
            "تسلسل التشغيل القياسي",
            [
                "تهيئة النظام واختبار الاستجابة الذاتية",
                "تفعيل وضع المسح ومراقبة حركة المركبة",
                "مراجعة الصورة الأولية والتحقق من اكتمال البيانات",
                "إنهاء المهمة وحفظ السجل التشغيلي",
            ],
        ),
        (
            "مراقبة جودة الصورة",
            [
                "الوضوح: تمييز الحدود الدقيقة للأجسام",
                "التباين: فصل الطبقات المادية بوضوح",
                "الضجيج: تخفيض التشويش عبر الإعداد الصحيح",
                "الثبات: الحفاظ على جودة ثابتة عبر النوبات",
            ],
        ),
        (
            "الإنذارات والأعطال المتكررة",
            [
                "انقطاع إشارة المستشعر: تحقق من الأسلاك والاتصال",
                "انخفاض جودة الصورة: راجع المعايرة وسرعة المرور",
                "تنبيه وحدة الجهد العالي: أوقف التشغيل فورًا",
                "فشل التخزين: تحقق من المساحة وسلامة قاعدة البيانات",
            ],
        ),
        (
            "الصيانة الوقائية",
            [
                "فحص يومي: نظافة العناصر الحساسة وحالة التبريد",
                "فحص أسبوعي: اختبار الكابلات ونقاط التثبيت",
                "فحص شهري: مراجعة الأداء ومعايير الجرعة",
                "إدارة قطع الغيار الحرجة ضمن خطة الاستمرارية",
            ],
        ),
        (
            "دراسة حالة تشغيلية",
            [
                "الحالة: تكرار إنذار أثناء المسح في الذروة",
                "التحليل: ربط الحدث بقراءات المستشعر",
                "الإجراء: تحويل النظام إلى Safe State",
                "النتيجة: استعادة التشغيل بعد تحقق السلامة",
            ],
        ),
        (
            "اختبار المعرفة",
            [
                "ما الإجراء الأول عند ظهور إنذار إشعاعي؟",
                "متى يتم إيقاف المسح وإعادة التهيئة؟",
                "كيف تفرق بين خطأ تشغيلي وخطأ مكون تقني؟",
                "ما عناصر التقرير النهائي بعد كل حادثة؟",
            ],
        ),
        (
            "تقييم عملي للمتدرب",
            [
                "تطبيق قائمة فحص ما قبل التشغيل بالكامل",
                "تنفيذ مسح فعلي ضمن زمن تشغيلي قياسي",
                "قراءة الإنذارات واتخاذ قرار صحيح خلال 60 ثانية",
                "توثيق الحادثة وفق النموذج المؤسسي",
            ],
        ),
        (
            "الخلاصة وخطة التحسين",
            [
                "الالتزام بالسلامة هو معيار النجاح الأول",
                "التحسين المستمر يعتمد على جودة البيانات الميدانية",
                "المراجعة الدورية تقلل الأعطال وتزيد الجاهزية",
                "التدريب العملي المنتظم يرفع كفاءة الفريق",
            ],
        ),
    ]

    img_cycle = images if images else []
    for idx, (title, bullets) in enumerate(slides_data):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_ar_bullets(
            slide,
            title,
            bullets,
            footer=f"المحور {idx + 1} من {len(slides_data)} | برنامج المدرب",
        )
        if img_cycle:
            pic = img_cycle[idx % len(img_cycle)]
            add_image_block(slide, pic, f"صورة توضيحية تدريبية: {pic.name}")


def main() -> None:
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    images = pick_images(12)
    make_title_slide(prs)
    make_content(prs, images)

    prs.save(str(OUTPUT_FILE))
    print(f"Saved: {OUTPUT_FILE}")
    print(f"Slides: {len(prs.slides)}")
    print(f"Images used: {len(images)}")


if __name__ == "__main__":
    main()
