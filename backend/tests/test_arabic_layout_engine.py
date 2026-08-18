"""Tests for the Arabic Layout Transformation Engine."""

import io
import sys

sys.path.insert(0, ".")

from pptx import Presentation
from pptx.util import Inches, Emu

from api.utils.arabic_pptx_profile import ARABIC_PROFILE
from api.utils.arabic_layout_engine import transform_slide_layout


def _blank_deck():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])   # blank layout
    return prs, slide


def _title_deck():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])   # "Title Only" layout
    slide.shapes.title.text_frame.text = "نظرة عامة على النظام"
    return prs, slide


# ── Mirroring rules ───────────────────────────────────────────────────────────

def test_left_box_mirrors_to_right():
    prs, slide = _blank_deck()
    box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(2), Inches(1))
    box.text_frame.text = "محتوى"
    orig_left, w = box.left, box.width
    sw = prs.slide_width

    transform_slide_layout(slide, sw, prs.slide_height, ARABIC_PROFILE)
    assert box.left == sw - orig_left - w, "left box must mirror to the right"


def test_symmetric_full_width_box_stays():
    """Equal margins → mirror is the identity → no move."""
    prs, slide = _blank_deck()
    sw = prs.slide_width
    box = slide.shapes.add_textbox(Emu(int(sw * 0.05)), Inches(2),
                                   Emu(sw - 2 * int(sw * 0.05)), Inches(1))
    box.text_frame.text = "محتوى عريض"
    orig_left = box.left

    transform_slide_layout(slide, sw, prs.slide_height, ARABIC_PROFILE)
    assert box.left == orig_left, "symmetric shapes must not move"


def test_wide_asymmetric_box_swaps_margins():
    """A wide box with unequal margins (e.g. text beside a gutter of bullet
    chips) must swap margins so the whole composition mirrors coherently."""
    prs, slide = _blank_deck()
    sw = prs.slide_width
    left, w = int(sw * 0.11), int(sw * 0.84)
    box = slide.shapes.add_textbox(Emu(left), Inches(2), Emu(w), Inches(1))
    box.text_frame.text = "صف نصي بجانب دائرة مرقمة"

    transform_slide_layout(slide, sw, prs.slide_height, ARABIC_PROFILE)
    assert box.left == sw - left - w, "asymmetric wide box must swap margins"


def test_right_arrow_flips_direction():
    """Directional autoshapes must flip (flipH), not just move."""
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.oxml.ns import qn
    prs, slide = _blank_deck()
    sw = prs.slide_width
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, Inches(1), Inches(2), Inches(2), Inches(1))
    orig_left, orig_w = arrow.left, arrow.width

    transform_slide_layout(slide, sw, prs.slide_height, ARABIC_PROFILE)
    xfrm = arrow._element.find(qn("p:spPr") + "/" + qn("a:xfrm"))
    assert xfrm.get("flipH") == "1", "arrow must point the RTL way"
    assert arrow.left == sw - orig_left - orig_w, "arrow must also margin-swap"


def test_connector_flips():
    """Connectors are directional — flipH must toggle."""
    from pptx.enum.shapes import MSO_CONNECTOR
    from pptx.oxml.ns import qn
    prs, slide = _blank_deck()
    sw = prs.slide_width
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(1), Inches(1), Inches(4), Inches(2))

    transform_slide_layout(slide, sw, prs.slide_height, ARABIC_PROFILE)
    xfrm = conn._element.find(qn("p:spPr") + "/" + qn("a:xfrm"))
    assert xfrm is not None and xfrm.get("flipH") == "1"


def test_group_children_mirror_in_group_space():
    """A group's internal arrangement must flip, not just its position."""
    from pptx.oxml.ns import qn
    prs, slide = _blank_deck()
    sw = prs.slide_width
    grp = slide.shapes.add_group_shape()
    a = grp.shapes.add_textbox(Inches(0), Inches(0), Inches(1), Inches(1))
    b = grp.shapes.add_textbox(Inches(3), Inches(0), Inches(2), Inches(1))
    a.text_frame.text = "أ"
    b.text_frame.text = "ب"
    gx = grp._element.find(qn("p:grpSpPr") + "/" + qn("a:xfrm"))
    ox = int(gx.find(qn("a:chOff")).get("x"))
    ext = int(gx.find(qn("a:chExt")).get("cx"))
    exp_a = 2 * ox + ext - a.left - a.width
    exp_b = 2 * ox + ext - b.left - b.width

    transform_slide_layout(slide, sw, prs.slide_height, ARABIC_PROFILE)
    assert (a.left, b.left) == (exp_a, exp_b), \
        "children must swap within group child-coordinate space"


def test_failure_rolls_slide_back(monkeypatch):
    """A mid-slide exception must restore pre-transform geometry but keep text."""
    import api.utils.arabic_layout_engine as ale
    prs, slide = _blank_deck()
    sw = prs.slide_width
    box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(2), Inches(1))
    box.text_frame.text = "نص مترجم"
    orig_left = box.left

    calls = {"n": 0}
    def boom(left, width, slide_w):
        calls["n"] += 1
        if calls["n"] >= 2:  # raise AFTER shape.left was already mutated
            raise RuntimeError("boom")
        return 0
    monkeypatch.setattr(ale, "_h_overflow", boom)

    warnings = transform_slide_layout(slide, sw, prs.slide_height, ARABIC_PROFILE)
    assert any("restored" in w for w in warnings)
    restored = slide.shapes[0]
    assert restored.left == orig_left, "geometry must roll back"
    assert restored.text_frame.text == "نص مترجم", "translation must survive rollback"


def test_centered_box_not_mirrored():
    prs, slide = _blank_deck()
    sw = prs.slide_width
    w = int(sw * 0.3)
    left = (sw - w) // 2
    box = slide.shapes.add_textbox(Emu(left), Inches(2), Emu(w), Inches(1))
    box.text_frame.text = "مركزي"

    transform_slide_layout(slide, sw, prs.slide_height, ARABIC_PROFILE)
    assert box.left == left, "centered shapes must stay put"


def test_mirror_preserves_offcanvas_bleed():
    """A shape bleeding off the left edge mirrors to bleed off the right by
    the same amount — no warnings, no clamping."""
    prs, slide = _blank_deck()
    sw = prs.slide_width
    box = slide.shapes.add_textbox(Emu(-100000), Inches(2), Inches(2), Inches(1))
    box.text_frame.text = "نزيف"
    w = box.width

    warns = transform_slide_layout(slide, sw, prs.slide_height, ARABIC_PROFILE)
    assert box.left == sw - (-100000) - w
    assert warns == [], f"mirror must not change overflow: {warns}"


# ── Title placeholder preservation ───────────────────────────────────────────

def test_title_placeholder_preserved_and_mirrored_rtl():
    from api.utils.arabic_layout_engine import _detect_top_banner_rect

    prs, slide = _title_deck()
    sw, sh = prs.slide_width, prs.slide_height
    t = slide.shapes.title
    orig_left, orig_top, orig_w, orig_h = t.left, t.top, t.width, t.height
    orig_margin_top = t.text_frame.margin_top
    orig_margin_right = t.text_frame.margin_right
    orig_ph_type = str(t.placeholder_format.type)

    transform_slide_layout(slide, sw, sh, ARABIC_PROFILE)
    t = slide.shapes.title
    _, banner_top, _, banner_h = _detect_top_banner_rect(slide, sw, sh)
    banner_bottom = banner_top + banner_h
    assert t.left == sw - orig_left - orig_w
    assert t.top <= orig_top
    assert t.top >= banner_top
    assert t.width == orig_w
    assert t.height == orig_h
    assert t.top + t.height <= banner_bottom
    assert str(t.placeholder_format.type) == orig_ph_type
    assert t.text_frame.margin_top >= orig_margin_top
    assert t.text_frame.margin_right == orig_margin_right

    # Right-aligned paragraphs and vertical centering.
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.text import MSO_VERTICAL_ANCHOR
    assert slide.shapes.title.text_frame.paragraphs[0].alignment == PP_ALIGN.RIGHT
    assert slide.shapes.title.text_frame.vertical_anchor == MSO_VERTICAL_ANCHOR.MIDDLE


def test_title_box_expands_height_before_aggressive_shrink():
    """Long Arabic titles should gain height first to avoid top glyph clipping."""
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    sw, sh = prs.slide_width, prs.slide_height
    t = slide.shapes.title
    t.left = Inches(0.6)
    t.width = Inches(6.0)
    t.top = Inches(0.25)
    t.height = Inches(0.70)
    t.text_frame.text = (
        "تدريب نظام EDAQ ومعايرة الكسب والتشخيص في التشغيل الميداني"
    )

    orig_center = t.top + (t.height // 2)
    orig_h = t.height

    warns = transform_slide_layout(slide, sw, sh, ARABIC_PROFILE)

    t = slide.shapes.title
    new_center = t.top + (t.height // 2)
    assert t.left == sw - int(Inches(0.6)) - int(Inches(6.0))
    assert t.text_frame.paragraphs[0].alignment == PP_ALIGN.RIGHT
    assert t.height >= orig_h
    # Preserve visual vertical position (center stays effectively the same).
    assert abs(new_center - orig_center) <= Emu(91440)  # <= 0.1 inch
    assert not any("top-glyph clearance failed" in w for w in warns)


def test_cover_slide_layout_is_mirrored():
    """CENTER_TITLE cover slides keep their art-directed geometry."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])   # Title Slide (ctrTitle)
    slide.shapes.title.text_frame.text = "عنوان الغلاف"
    box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(2), Inches(1))
    box.text_frame.text = "زخرفة"
    t = slide.shapes.title
    t_left, t_top, t_w, t_h = t.left, t.top, t.width, t.height
    transform_slide_layout(slide, prs.slide_width, prs.slide_height, ARABIC_PROFILE)
    assert slide.shapes.title.left == t_left
    assert slide.shapes.title.top == t_top
    assert slide.shapes.title.width == t_w
    assert slide.shapes.title.height == t_h


# ── End-to-end gating through rebuild_pptx ────────────────────────────────────

def _deck_bytes(prs):
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def test_ltr_source_deck_gets_transformed():
    """English-source deck: title stays in original placeholder geometry."""
    from api.utils.doc_rebuilder import rebuild_pptx, _detect_top_banner_rect

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text_frame.text = "System Overview"
    orig_title = slide.shapes.title
    orig_left, orig_top, orig_w, orig_h = (
        orig_title.left,
        orig_title.top,
        orig_title.width,
        orig_title.height,
    )
    orig_count = len(slide.shapes)
    src = _deck_bytes(prs)

    segments = [{
        "target": "نظرة عامة على النظام",
        "source": "System Overview",
        "loc": {"format": "pptx", "slide_idx": 0, "shape_idx": 0, "para_idx": 0},
    }]
    out = rebuild_pptx(src, segments, "ar")

    prs2 = Presentation(io.BytesIO(out))
    t = prs2.slides[0].shapes.title
    _, banner_top, _, banner_h = _detect_top_banner_rect(prs2.slides[0], prs2.slide_width, prs2.slide_height)
    banner_bottom = banner_top + banner_h
    assert len(prs2.slides[0].shapes) == orig_count
    assert t.left == prs2.slide_width - orig_left - orig_w
    assert t.top <= orig_top
    assert t.top >= banner_top
    assert t.width == orig_w
    assert t.height == orig_h
    assert t.top + t.height <= banner_bottom
    assert "نظرة عامة" in t.text_frame.text


def test_arabic_source_deck_geometry_fully_preserved():
    """Arabic-source deck: the engine must not run — geometry identical."""
    from api.utils.doc_rebuilder import rebuild_pptx

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text_frame.text = "كورس المشغل"
    box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(2), Inches(1))
    box.text_frame.text = "نص جانبي"
    src = _deck_bytes(prs)

    orig_geo = [(sh.shape_id, sh.left, sh.top, sh.width, sh.height)
                for sh in Presentation(io.BytesIO(src)).slides[0].shapes]

    segments = [{
        "target": "دورة المشغل",
        "source": "كورس المشغل",
        "loc": {"format": "pptx", "slide_idx": 0, "shape_idx": 0, "para_idx": 0},
    }]
    out = rebuild_pptx(src, segments, "ar")
    new_geo = [(sh.shape_id, sh.left, sh.top, sh.width, sh.height)
               for sh in Presentation(io.BytesIO(out)).slides[0].shapes]
    assert orig_geo == new_geo, "Arabic-source decks must keep exact geometry"


def test_rebuild_mirrors_title_placeholder_position_for_rtl_default_mode():
    """Default Arabic rebuild must mirror title placeholder in place."""
    from api.utils.doc_rebuilder import rebuild_pptx
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text_frame.text = "System Configuration"

    # Force an asymmetric title position to verify true mirror behavior.
    title.left = Inches(0.7)
    title.width = Inches(6.2)
    title.top = Inches(0.25)
    title.height = Inches(0.9)

    src = _deck_bytes(prs)
    src_title = Presentation(io.BytesIO(src)).slides[0].shapes.title
    sw = prs.slide_width
    exp_left = sw - src_title.left - src_title.width

    segments = [{
        "target": "تهيئة النظام",
        "source": "System Configuration",
        "loc": {"format": "pptx", "slide_idx": 0, "shape_idx": 0, "para_idx": 0},
    }]
    out = rebuild_pptx(src, segments, "ar")

    prs2 = Presentation(io.BytesIO(out))
    t = prs2.slides[0].shapes.title
    assert t.left == exp_left
    assert t.top == src_title.top
    assert t.height == src_title.height
    assert t.width == src_title.width
    assert t.text_frame.paragraphs[0].alignment == PP_ALIGN.RIGHT
    ppr = t.text_frame.paragraphs[0]._p.get_or_add_pPr()
    assert ppr.get("rtl") == "1"


def test_title_placeholder_mirrors_even_with_nearby_logo_obstacle():
    """A nearby top-right image must not prevent title placeholder mirroring."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    sw = prs.slide_width
    t = slide.shapes.title
    t.text_frame.text = "System Overview"
    orig_left, orig_w = t.left, t.width

    # Add a picture-like obstacle proxy in the top-right neighborhood.
    # (Rectangle with no text and sizable area is treated as a soft obstacle.)
    logo = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Emu(int(sw * 0.78)),
        Emu(int(prs.slide_height * 0.01)),
        Emu(int(sw * 0.18)),
        Emu(int(prs.slide_height * 0.10)),
    )
    logo.text_frame.text = ""

    transform_slide_layout(slide, sw, prs.slide_height, ARABIC_PROFILE)
    assert t.left == sw - orig_left - orig_w


def test_top_banner_textbox_detected_as_title_and_mirrored():
    """Top heading textboxes (non-placeholders) must mirror as titles."""
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Pt

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sw, sh = prs.slide_width, prs.slide_height

    # Simulated top banner strip.
    banner = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Emu(0),
        Emu(0),
        Emu(sw),
        Emu(int(0.16 * sh)),
    )
    banner.text_frame.text = ""

    # Title-like textbox (not a placeholder).
    title_box = slide.shapes.add_textbox(
        Emu(int(0.06 * sw)),
        Emu(int(0.03 * sh)),
        Emu(int(0.62 * sw)),
        Emu(int(0.10 * sh)),
    )
    title_box.text_frame.text = "System Detector Overview"
    title_box.text_frame.paragraphs[0].runs[0].font.size = Pt(36)

    # Top-right logo obstacle proxy.
    logo = slide.shapes.add_shape(
        1,
        Emu(int(0.80 * sw)),
        Emu(int(0.02 * sh)),
        Emu(int(0.16 * sw)),
        Emu(int(0.10 * sh)),
    )
    logo.text_frame.text = ""

    orig_left = title_box.left
    orig_width = title_box.width
    expected_mirror_left = sw - orig_left - orig_width
    safe_gap = int(0.015 * sw)

    transform_slide_layout(slide, sw, sh, ARABIC_PROFILE)

    # Title may shift further left than pure mirror to avoid top-right overlay.
    assert title_box.left <= expected_mirror_left
    assert title_box.left + title_box.width <= logo.left - safe_gap
    assert title_box.text_frame.paragraphs[0].alignment == PP_ALIGN.RIGHT
    ppr = title_box.text_frame.paragraphs[0]._p.get_or_add_pPr()
    assert ppr.get("rtl") == "1"


def test_rebuild_mirrors_non_placeholder_top_title_textbox():
    """Default rebuild must mirror top heading textbox titles, not placeholders only."""
    from api.utils.doc_rebuilder import rebuild_pptx
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Pt

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sw, sh = prs.slide_width, prs.slide_height

    # Banner + title-like textbox.
    b = slide.shapes.add_shape(1, Emu(0), Emu(0), Emu(sw), Emu(int(0.16 * sh)))
    b.text_frame.text = ""
    tb = slide.shapes.add_textbox(
        Emu(int(0.05 * sw)),
        Emu(int(0.03 * sh)),
        Emu(int(0.60 * sw)),
        Emu(int(0.10 * sh)),
    )
    tb.text_frame.text = "Detector Training Intro"
    tb.text_frame.paragraphs[0].runs[0].font.size = Pt(34)

    src = _deck_bytes(prs)
    src_prs = Presentation(io.BytesIO(src))
    src_tb = src_prs.slides[0].shapes[1]
    exp_left = sw - src_tb.left - src_tb.width

    segments = [{
        "target": "مقدمة تدريب الكاشف",
        "source": "Detector Training Intro",
        "loc": {"format": "pptx", "slide_idx": 0, "shape_idx": 1, "para_idx": 0},
    }]
    out = rebuild_pptx(src, segments, "ar")

    prs2 = Presentation(io.BytesIO(out))
    out_tb = prs2.slides[0].shapes[1]
    assert out_tb.left == exp_left
    assert out_tb.text_frame.paragraphs[0].alignment == PP_ALIGN.RIGHT
    ppr = out_tb.text_frame.paragraphs[0]._p.get_or_add_pPr()
    assert ppr.get("rtl") == "1"
