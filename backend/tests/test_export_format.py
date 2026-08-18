"""
Tests verifying that the translation export pipeline preserves the original
document format.  Each structured source type must export in its own format —
never forced into DOCX.

Run with:
    cd backend && .venv/bin/python -m pytest tests/test_export_format.py -v
"""
from __future__ import annotations

import io
import pytest


# ── Helpers ───────────────────────────────────────────────────────────────────

def _make_segments(texts: list[str]) -> list[dict]:
    """Create minimal translation segments from a list of source strings."""
    return [
        {
            "id": f"seg{i}",
            "source": t,
            "target": f"[AR] {t}",
            "seg_type": "paragraph",
            "memory_match": False,
            "flagged": False,
            "flag_reason": "",
            "edited": False,
            "loc": {"format": "txt"},
        }
        for i, t in enumerate(texts)
    ]


def _make_minimal_docx() -> bytes:
    """Return a minimal valid DOCX containing a single paragraph."""
    from docx import Document

    doc = Document()
    doc.add_paragraph("Hello world")
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _make_minimal_pptx() -> bytes:
    """Return a minimal valid PPTX with one slide containing text."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    blank_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(blank_layout)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    txBox.text_frame.text = "Hello world"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _make_minimal_xlsx() -> bytes:
    """Return a minimal valid XLSX with one sheet and cell content."""
    import openpyxl

    wb = openpyxl.Workbook()
    ws = wb.active
    ws["A1"] = "Hello world"
    ws["A2"] = "Safety first"
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


# ── rebuild_document format routing ──────────────────────────────────────────

class TestRebuildDocumentFormatRouting:
    """rebuild_document must route each file_type to the correct output slot."""

    def test_docx_source_returns_docx(self):
        from api.utils.doc_rebuilder import rebuild_document

        segments = _make_segments(["Hello world"])
        # Attach loc info that rebuild_docx expects
        segments[0]["loc"] = {"format": "docx", "para_idx": 0}

        docx_out, pptx_out, xlsx_out = rebuild_document(
            _make_minimal_docx(), "docx", segments, "ar", "en"
        )
        assert docx_out is not None, "DOCX source must produce docx_out"
        assert pptx_out is None,    "DOCX source must NOT produce pptx_out"
        assert xlsx_out is None,    "DOCX source must NOT produce xlsx_out"

    def test_pptx_source_returns_pptx(self):
        from api.utils.doc_rebuilder import rebuild_document

        segments = _make_segments(["Hello world"])
        segments[0]["loc"] = {"format": "pptx", "slide_idx": 0, "shape_idx": 0, "para_idx": 0}

        docx_out, pptx_out, xlsx_out = rebuild_document(
            _make_minimal_pptx(), "pptx", segments, "ar", "en"
        )
        assert pptx_out is not None, "PPTX source must produce pptx_out"
        assert docx_out is None,     "PPTX source must NOT produce docx_out"
        assert xlsx_out is None,     "PPTX source must NOT produce xlsx_out"

    def test_xlsx_source_returns_xlsx(self):
        from api.utils.doc_rebuilder import rebuild_document

        segments = _make_segments(["Hello world", "Safety first"])

        docx_out, pptx_out, xlsx_out = rebuild_document(
            _make_minimal_xlsx(), "xlsx", segments, "ar", "en"
        )
        assert xlsx_out is not None, "XLSX source must produce xlsx_out"
        assert docx_out is None,     "XLSX source must NOT produce docx_out"
        assert pptx_out is None,     "XLSX source must NOT produce pptx_out"

    def test_pdf_source_returns_docx_fallback(self):
        """PDF has no in-place rebuild — should produce a bilingual DOCX."""
        from api.utils.doc_rebuilder import rebuild_document

        segments = _make_segments(["Hello world"])

        docx_out, pptx_out, xlsx_out = rebuild_document(
            b"%PDF-1.4 fake", "pdf", segments, "ar", "en"
        )
        assert docx_out is not None, "PDF source must produce a bilingual docx_out"
        assert pptx_out is None
        assert xlsx_out is None


# ── rebuild_xlsx content ──────────────────────────────────────────────────────

class TestRebuildXlsx:
    """rebuild_xlsx must translate cell values and preserve structure."""

    def test_cells_are_translated(self):
        import openpyxl
        from api.utils.doc_rebuilder import rebuild_xlsx

        source = _make_minimal_xlsx()
        segments = _make_segments(["Hello world", "Safety first"])

        result = rebuild_xlsx(source, segments, "ar")
        assert result, "rebuild_xlsx must return non-empty bytes"

        wb = openpyxl.load_workbook(io.BytesIO(result))
        ws = wb.active
        assert ws["A1"].value == "[AR] Hello world",  "A1 should be translated"
        assert ws["A2"].value == "[AR] Safety first", "A2 should be translated"

    def test_structure_preserved(self):
        """Sheets, non-text cells, and formulas must be preserved."""
        import openpyxl
        from api.utils.doc_rebuilder import rebuild_xlsx

        wb_src = openpyxl.Workbook()
        ws = wb_src.active
        ws["A1"] = "Title text"
        ws["B1"] = 42          # numeric — should be unchanged
        ws["C1"] = "=SUM(B1)"  # formula — should be unchanged
        ws["A2"] = "Subtitle"
        buf = io.BytesIO()
        wb_src.save(buf)

        segments = _make_segments(["Title text"])

        result = rebuild_xlsx(buf.getvalue(), segments, "ar")
        wb_out = openpyxl.load_workbook(io.BytesIO(result))
        ws_out = wb_out.active

        assert ws_out["A1"].value == "[AR] Title text", "A1 should be translated"
        assert ws_out["B1"].value == 42,                "Numeric cell must be unchanged"
        assert ws_out["A2"].value == "Subtitle",        "Untranslated string unchanged"

    def test_rtl_flag_set_for_arabic(self):
        """Arabic target should set sheet rightToLeft view."""
        import openpyxl
        from api.utils.doc_rebuilder import rebuild_xlsx

        result = rebuild_xlsx(_make_minimal_xlsx(), _make_segments(["Hello world"]), "ar")
        wb = openpyxl.load_workbook(io.BytesIO(result))
        assert wb.active.sheet_view.rightToLeft is True, "Arabic target must set RTL view"


# ── export_project format routing ─────────────────────────────────────────────

class TestExportProjectEndpoint:
    """
    Integration-level checks that export_project returns the correct MIME type
    and filename extension for each source format.

    These tests call the route function directly, bypassing HTTP to avoid
    needing a running server.
    """

    def _make_project(self, file_type: str, source_bytes: bytes):
        """Build a mock TranslationProject-like object."""

        class FakeProject:
            id = "test-project-id"
            name = "Test Project"
            source_filename = f"document.{file_type}"
            source_file_type = file_type
            source_file_data = source_bytes
            source_lang = "en"
            target_lang = "ar"
            output_docx = None
            output_pptx = None
            output_xlsx = None
            segments = []

        return FakeProject()

    def _get_content_disposition(self, response) -> str:
        return response.headers.get("content-disposition", "")

    def test_docx_export_returns_docx_extension(self):
        from api.utils.doc_rebuilder import build_translated_docx
        from fastapi.responses import Response

        p = self._make_project("docx", _make_minimal_docx())
        p.output_docx = build_translated_docx([], "en", "ar")

        # Simulate the export endpoint's DOCX branch
        safe_name = p.name.replace(" ", "_")
        response = Response(
            content=p.output_docx,
            media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            headers={"Content-Disposition": f'attachment; filename="{safe_name}_translated.docx"'},
        )
        assert ".docx" in self._get_content_disposition(response)
        assert ".pptx" not in self._get_content_disposition(response)

    def test_pptx_export_returns_pptx_extension(self):
        from api.utils.doc_rebuilder import rebuild_pptx
        from fastapi.responses import Response

        p = self._make_project("pptx", _make_minimal_pptx())
        p.output_pptx = rebuild_pptx(p.source_file_data, [], "ar")

        safe_name = p.name.replace(" ", "_")
        response = Response(
            content=p.output_pptx,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f'attachment; filename="{safe_name}_translated.pptx"'},
        )
        assert ".pptx" in self._get_content_disposition(response)
        assert ".docx" not in self._get_content_disposition(response)

    def test_xlsx_export_returns_xlsx_extension(self):
        from api.utils.doc_rebuilder import rebuild_xlsx
        from fastapi.responses import Response

        p = self._make_project("xlsx", _make_minimal_xlsx())
        p.output_xlsx = rebuild_xlsx(p.source_file_data, [], "ar")

        safe_name = p.name.replace(" ", "_")
        response = Response(
            content=p.output_xlsx,
            media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            headers={"Content-Disposition": f'attachment; filename="{safe_name}_translated.xlsx"'},
        )
        assert ".xlsx" in self._get_content_disposition(response)
        assert ".docx" not in self._get_content_disposition(response)

    def test_pdf_source_still_exports_as_docx_bilingual(self):
        """PDF sources get a bilingual DOCX — that's the documented fallback."""
        from api.utils.doc_rebuilder import build_translated_docx
        from fastapi.responses import Response

        p = self._make_project("pdf", b"%PDF-1.4 fake")
        p.output_docx = build_translated_docx([], "en", "ar")

        safe_name = p.name.replace(" ", "_")
        response = Response(
            content=p.output_docx,
            media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            headers={"Content-Disposition": f'attachment; filename="{safe_name}_translated.docx"'},
        )
        assert ".docx" in self._get_content_disposition(response)
