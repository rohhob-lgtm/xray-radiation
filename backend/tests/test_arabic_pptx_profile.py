"""
Tests for the Arabic PPTX reference profile and the rebuild_pptx Arabic formatting.

Run with:
    cd backend && .venv/bin/python -m pytest tests/test_arabic_pptx_profile.py -v
"""
import io
import pytest


# ── arabic_pptx_profile module ────────────────────────────────────────────────

class TestArabicPptxProfile:
    def test_profile_has_required_keys(self):
        from api.utils.arabic_pptx_profile import ARABIC_PROFILE
        assert "fonts" in ARABIC_PROFILE
        assert "rtl_rules" in ARABIC_PROFILE
        assert "semantic_styles" in ARABIC_PROFILE
        assert "line_spacing" in ARABIC_PROFILE
        assert "text_frame_margins_emu" in ARABIC_PROFILE
        assert "score_weights" in ARABIC_PROFILE

    def test_primary_font_is_simplified_arabic(self):
        from api.utils.arabic_pptx_profile import ARABIC_PROFILE
        assert ARABIC_PROFILE["fonts"]["primary"] == "Simplified Arabic"

    def test_rtl_rules_set_rtl_on_all_paragraphs(self):
        from api.utils.arabic_pptx_profile import ARABIC_PROFILE
        assert ARABIC_PROFILE["rtl_rules"]["set_rtl_on_all_paragraphs"] is True

    def test_default_body_alignment_is_right(self):
        from api.utils.arabic_pptx_profile import ARABIC_PROFILE
        assert ARABIC_PROFILE["rtl_rules"]["default_body_alignment"] == "right"

    def test_table_tblpr_rtl_true(self):
        from api.utils.arabic_pptx_profile import ARABIC_PROFILE
        assert ARABIC_PROFILE["rtl_rules"]["table_tblPr_rtl"] is True

    def test_score_weights_sum_to_one(self):
        from api.utils.arabic_pptx_profile import ARABIC_PROFILE
        total = sum(ARABIC_PROFILE["score_weights"].values())
        assert abs(total - 1.0) < 1e-9, f"Score weights sum to {total}, expected 1.0"

    def test_margin_values_match_reference(self):
        from api.utils.arabic_pptx_profile import ARABIC_PROFILE
        m = ARABIC_PROFILE["text_frame_margins_emu"]
        assert m["left"]   == 91_440
        assert m["right"]  == 91_440
        assert m["top"]    == 45_720
        assert m["bottom"] == 45_720

    def test_line_spacing_dominant_is_150(self):
        from api.utils.arabic_pptx_profile import ARABIC_PROFILE
        assert ARABIC_PROFILE["line_spacing"]["dominant_pct"] == 150


class TestIsArabicCapable:
    def test_simplified_arabic_capable(self):
        from api.utils.arabic_pptx_profile import is_arabic_capable
        assert is_arabic_capable("Simplified Arabic") is True

    def test_traditional_arabic_capable(self):
        from api.utils.arabic_pptx_profile import is_arabic_capable
        assert is_arabic_capable("Traditional Arabic") is True

    def test_arial_capable(self):
        from api.utils.arabic_pptx_profile import is_arabic_capable
        assert is_arabic_capable("Arial") is True

    def test_tahoma_capable(self):
        from api.utils.arabic_pptx_profile import is_arabic_capable
        assert is_arabic_capable("Tahoma") is True

    def test_calibri_capable(self):
        from api.utils.arabic_pptx_profile import is_arabic_capable
        assert is_arabic_capable("Calibri") is True

    def test_montserrat_not_capable(self):
        from api.utils.arabic_pptx_profile import is_arabic_capable
        assert is_arabic_capable("Montserrat") is False

    def test_roboto_not_capable(self):
        from api.utils.arabic_pptx_profile import is_arabic_capable
        assert is_arabic_capable("Roboto") is False

    def test_none_returns_true(self):
        """None means inherited — treated as capable to avoid unnecessary substitutions."""
        from api.utils.arabic_pptx_profile import is_arabic_capable
        assert is_arabic_capable(None) is True

    def test_empty_string_returns_true(self):
        from api.utils.arabic_pptx_profile import is_arabic_capable
        assert is_arabic_capable("") is True

    def test_case_insensitive(self):
        from api.utils.arabic_pptx_profile import is_arabic_capable
        assert is_arabic_capable("ARIAL") is True
        assert is_arabic_capable("arial") is True

    def test_unknown_font_not_capable(self):
        from api.utils.arabic_pptx_profile import is_arabic_capable
        assert is_arabic_capable("SomeCorporateFontXYZ") is False


class TestGetArabicFont:
    def test_capable_font_returned_unchanged(self):
        from api.utils.arabic_pptx_profile import get_arabic_font
        font, substituted = get_arabic_font("Arial")
        assert font == "Arial"
        assert substituted is False

    def test_incapable_font_substituted(self):
        from api.utils.arabic_pptx_profile import get_arabic_font
        # Spec §9 preferred fallback: Arial first
        font, substituted = get_arabic_font("Montserrat")
        assert font == "Arial"
        assert substituted is True

    def test_none_font_unchanged(self):
        from api.utils.arabic_pptx_profile import get_arabic_font
        font, substituted = get_arabic_font(None)
        assert font is None
        assert substituted is False

    def test_title_role_uses_correct_substitute(self):
        from api.utils.arabic_pptx_profile import get_arabic_font
        # Spec §9 preferred fallback: Arial first for titles
        font, substituted = get_arabic_font("Raleway", "title")
        assert font == "Arial"
        assert substituted is True


# ── rebuild_pptx Arabic integration ──────────────────────────────────────────

def _make_minimal_pptx() -> bytes:
    """Create an in-memory PPTX with one slide, one title, one body paragraph."""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    blank_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(blank_layout)

    # Title
    title = slide.shapes.title
    title.text = "Hello World"

    # Body
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "First bullet"
    p = tf.add_paragraph()
    p.text = "Second bullet"

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _make_segments(slide_idx=0):
    """Build minimal pptx segments matching the shape indices of _make_minimal_pptx."""
    return [
        {
            "id": "s1",
            "source": "Hello World",
            "target": "مرحبا بالعالم",
            "seg_type": "slide_title",
            "loc": {"format": "pptx", "slide_idx": slide_idx, "shape_idx": 0, "para_idx": 0},
        },
        {
            "id": "s2",
            "source": "First bullet",
            "target": "الرصاصة الأولى",
            "seg_type": "slide_body",
            "loc": {"format": "pptx", "slide_idx": slide_idx, "shape_idx": 1, "para_idx": 0},
        },
        {
            "id": "s3",
            "source": "Second bullet",
            "target": "الرصاصة الثانية",
            "seg_type": "slide_body",
            "loc": {"format": "pptx", "slide_idx": slide_idx, "shape_idx": 1, "para_idx": 1},
        },
    ]


class TestRebuildPptxArabic:
    def test_output_is_valid_pptx(self):
        from api.utils.doc_rebuilder import rebuild_pptx, validate_pptx_bytes
        pptx_in = _make_minimal_pptx()
        result = rebuild_pptx(pptx_in, _make_segments(), "ar")
        ok, reason = validate_pptx_bytes(result)
        assert ok, f"Invalid PPTX: {reason}"

    def test_translated_text_present(self):
        from pptx import Presentation
        from api.utils.doc_rebuilder import rebuild_pptx
        pptx_in = _make_minimal_pptx()
        result = rebuild_pptx(pptx_in, _make_segments(), "ar")
        prs = Presentation(io.BytesIO(result))
        slide = prs.slides[0]
        all_text = " ".join(
            para.text for shape in slide.shapes
            if shape.has_text_frame
            for para in shape.text_frame.paragraphs
        )
        assert "مرحبا بالعالم" in all_text
        assert "الرصاصة الأولى" in all_text
        assert "الرصاصة الثانية" in all_text

    def test_rtl_set_on_translated_paragraphs(self):
        """Every translated paragraph must have <a:pPr rtl="1"/>."""
        import lxml.etree as ET
        from pptx import Presentation
        from api.utils.doc_rebuilder import rebuild_pptx

        NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
        pptx_in = _make_minimal_pptx()
        result = rebuild_pptx(pptx_in, _make_segments(), "ar")
        prs = Presentation(io.BytesIO(result))
        slide = prs.slides[0]

        translated_texts = {"مرحبا بالعالم", "الرصاصة الأولى", "الرصاصة الثانية"}
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                if para.text.strip() in translated_texts:
                    pPr = para._p.find(f"{{{NS}}}pPr")
                    assert pPr is not None, f"No pPr on para: {para.text!r}"
                    assert pPr.get("rtl") == "1", \
                        f"rtl != 1 on translated para: {para.text!r}"

    def test_non_arabic_lang_no_rtl(self):
        """English translation must NOT get rtl=1."""
        import lxml.etree as ET
        from pptx import Presentation
        from api.utils.doc_rebuilder import rebuild_pptx

        NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
        segs = [
            {
                "id": "x1",
                "source": "Hello World",
                "target": "Bonjour le monde",
                "seg_type": "slide_title",
                "loc": {"format": "pptx", "slide_idx": 0, "shape_idx": 0, "para_idx": 0},
            },
        ]
        pptx_in = _make_minimal_pptx()
        result = rebuild_pptx(pptx_in, segs, "fr")   # French, not RTL
        prs = Presentation(io.BytesIO(result))
        slide = prs.slides[0]
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                if "Bonjour" in (para.text or ""):
                    pPr = para._p.find(f"{{{NS}}}pPr")
                    if pPr is not None:
                        assert pPr.get("rtl") != "1", \
                            "French paragraph should not have rtl=1"

    def test_output_is_non_empty(self):
        from api.utils.doc_rebuilder import rebuild_pptx
        pptx_in = _make_minimal_pptx()
        result = rebuild_pptx(pptx_in, _make_segments(), "ar")
        assert len(result) > 1000  # must be a real PPTX, not empty bytes


class TestDocExtractorPptxTables:
    def test_table_cells_extracted(self):
        """Table cells in PPTX slides are now extracted as table_cell segments."""
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from api.utils.doc_extractor import extract_pptx

        # Build a PPTX with one slide containing a 2×2 table
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank
        tbl_shape = slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(6), Inches(2))
        tbl = tbl_shape.table
        tbl.cell(0, 0).text = "Header A"
        tbl.cell(0, 1).text = "Header B"
        tbl.cell(1, 0).text = "Value one"
        tbl.cell(1, 1).text = "Value two"

        buf = io.BytesIO()
        prs.save(buf)
        pptx_bytes = buf.getvalue()

        segs = extract_pptx(pptx_bytes)
        table_segs = [s for s in segs if s.get("seg_type") == "table_cell"]
        texts = {s["source"] for s in table_segs}

        assert "Header A" in texts
        assert "Header B" in texts
        assert "Value one" in texts
        assert "Value two" in texts

    def test_table_seg_has_correct_loc_fields(self):
        from pptx import Presentation
        from pptx.util import Inches
        from api.utils.doc_extractor import extract_pptx

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        tbl_shape = slide.shapes.add_table(1, 2, Inches(1), Inches(1), Inches(6), Inches(2))
        tbl_shape.table.cell(0, 0).text = "CellText"
        tbl_shape.table.cell(0, 1).text = "OtherCell"

        buf = io.BytesIO()
        prs.save(buf)
        segs = extract_pptx(buf.getvalue())
        cell_segs = [s for s in segs if s.get("seg_type") == "table_cell"]

        assert cell_segs, "Expected table_cell segments"
        loc = cell_segs[0]["loc"]
        assert loc["format"] == "pptx"
        assert loc["table_cell"] is True
        assert "slide_idx" in loc
        assert "shape_idx" in loc
        assert "row_idx" in loc
        assert "col_idx" in loc


# ── Layout-preservation regression tests ─────────────────────────────────────

def _mini_deck_bytes(texts, table=False):
    """Build a minimal in-memory PPTX (blank layout); returns bytes."""
    import io
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if table:
        gfx = slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(6), Inches(2))
        for r in range(2):
            for c in range(2):
                gfx.table.cell(r, c).text = texts[(r * 2 + c) % len(texts)]
    else:
        for i, t in enumerate(texts):
            box = slide.shapes.add_textbox(Inches(1), Inches(1 + i), Inches(7), Inches(0.8))
            box.text_frame.text = t
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _slide1_xml(pptx_bytes):
    import io, zipfile
    with zipfile.ZipFile(io.BytesIO(pptx_bytes)) as z:
        return z.read("ppt/slides/slide1.xml").decode("utf-8")


def test_mixed_script_ltr_majority_source_still_gets_rtl():
    """English-majority source containing one Arabic term is an LTR source:
    the translated paragraph must still receive rtl=1 + right alignment."""
    from api.utils.doc_rebuilder import rebuild_pptx

    src_text = "Press the زر button to start the scanner now"
    src = _mini_deck_bytes([src_text])
    segments = [{
        "target": "اضغط الزر لبدء تشغيل الماسح الضوئي",
        "source": src_text,
        "loc": {"format": "pptx", "slide_idx": 0, "shape_idx": 0, "para_idx": 0},
    }]
    xml = _slide1_xml(rebuild_pptx(src, segments, "ar"))
    assert 'rtl="1"' in xml, "LTR-majority source must still get rtl=1"
    assert 'algn="r"' in xml, "LTR-majority source body must be right-aligned"


def test_arabic_source_paragraph_layout_untouched():
    """Arabic-source deck: translated paragraphs keep original direction and
    alignment attributes exactly (no rtl, no algn added)."""
    from api.utils.doc_rebuilder import rebuild_pptx

    src = _mini_deck_bytes(["كورس المشغل"])
    segments = [{
        "target": "دورة المشغل",
        "source": "كورس المشغل",
        "loc": {"format": "pptx", "slide_idx": 0, "shape_idx": 0, "para_idx": 0},
    }]
    xml = _slide1_xml(rebuild_pptx(src, segments, "ar"))
    assert 'rtl="1"' not in xml, "Arabic-source paragraph must not gain rtl attr"
    assert 'algn=' not in xml, "Arabic-source paragraph must not gain algn attr"
    assert "دورة المشغل" in xml, "text must still be translated"


def test_arabic_source_table_tblPr_unchanged():
    """Arabic-source deck with a table: tblPr@rtl must remain unset."""
    import io
    import lxml.etree as ET
    from api.utils.doc_rebuilder import rebuild_pptx

    cells = ["الخلية الأولى", "الخلية الثانية", "الخلية الثالثة", "الخلية الرابعة"]
    src = _mini_deck_bytes(cells, table=True)
    segments = [{
        "target": "الخلية المترجمة",
        "source": cells[0],
        "loc": {"format": "pptx", "slide_idx": 0, "shape_idx": 0,
                "table_cell": True, "row_idx": 0, "col_idx": 0},
    }]
    out = rebuild_pptx(src, segments, "ar")
    root = ET.fromstring(_slide1_xml(out).encode("utf-8"))
    ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    tblPrs = root.findall(".//a:tbl/a:tblPr", ns)
    assert tblPrs, "table should still exist"
    for tblPr in tblPrs:
        assert tblPr.get("rtl") is None, "Arabic-source table must not gain tblPr@rtl"
