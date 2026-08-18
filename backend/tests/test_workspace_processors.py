"""File processor unit tests — each processor against a small in-memory fixture."""
import io
import os
import sys

sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))

from api.services.workspace_processors.registry import process_file
from api.services.workspace_processors.csv_processor import process_csv
from api.services.workspace_processors.text_processor import process_text
from api.services.workspace_processors.code_processor import process_code


def _docx_bytes(paragraphs: list[str]) -> bytes:
    from docx import Document
    doc = Document()
    for p in paragraphs:
        doc.add_paragraph(p)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _pptx_bytes(slide_texts: list[str]) -> bytes:
    from pptx import Presentation
    prs = Presentation()
    layout = prs.slide_layouts[1]
    for text in slide_texts:
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = text
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _xlsx_bytes(rows: list[list]) -> bytes:
    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    for row in rows:
        ws.append(row)
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _pdf_bytes(text: str) -> bytes:
    import fitz
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72, 72), text)
    data = doc.tobytes()
    doc.close()
    return data


def test_word_processor_extracts_text():
    data = _docx_bytes(["Hello from Word", "Second paragraph"])
    result = process_file("doc.docx", data, "docx")
    assert not result["errors"]
    assert "Hello from Word" in result["text"]
    assert result["metadata"]["paragraph_count"] == 2


def test_pptx_processor_extracts_slides():
    data = _pptx_bytes(["Slide One", "Slide Two"])
    result = process_file("deck.pptx", data, "pptx")
    assert not result["errors"]
    assert result["metadata"]["slide_count"] == 2
    assert any("Slide One" in s["text"] for s in result["slides"])


def test_excel_processor_extracts_sheets_and_tables():
    data = _xlsx_bytes([["Name", "Value"], ["a", 1], ["b", 2]])
    result = process_file("sheet.xlsx", data, "xlsx")
    assert not result["errors"]
    assert result["metadata"]["sheet_count"] == 1
    assert result["tables"][0]["header"] == ["Name", "Value"]
    assert len(result["tables"][0]["rows"]) == 2


def test_pdf_processor_extracts_text():
    data = _pdf_bytes("Hello PDF world")
    result = process_file("doc.pdf", data, "pdf")
    assert not result["errors"]
    assert "Hello PDF" in result["text"]
    assert result["metadata"]["page_count"] == 1


def test_csv_processor():
    data = b"name,value\na,1\nb,2\n"
    result = process_csv("data.csv", data)
    assert not result["errors"]
    assert result["tables"][0]["header"] == ["name", "value"]
    assert len(result["tables"][0]["rows"]) == 2


def test_text_processor_redacts_secret():
    data = b'api_key = "sk-THIS-LOOKS-LIKE-A-REAL-SECRET-VALUE"\nnormal text follows'
    result = process_text("config.txt", data)
    assert "REDACTED" in result["text"]
    assert "sk-THIS-LOOKS-LIKE-A-REAL-SECRET-VALUE" not in result["text"]
    assert any("redacted" in w for w in result["warnings"])


def test_text_processor_env_filename_fully_redacted():
    data = b"SECRET_TOKEN=abc123\nDB_PASSWORD=hunter2"
    result = process_text(".env", data, filename=".env")
    assert "abc123" not in result["text"]
    assert "hunter2" not in result["text"]
    assert "REDACTED" in result["text"]


def test_code_processor_tags_language():
    data = b"def foo():\n    return 1\n"
    result = process_code("script.py", data, extension="py", filename="script.py")
    assert result["metadata"]["language"] == "python"
    assert "def foo" in result["text"]


def test_unsupported_extension_returns_warning_not_error():
    result = process_file("model.glb", b"\x00\x01", "glb")
    assert not result["errors"]
    assert result["warnings"]
