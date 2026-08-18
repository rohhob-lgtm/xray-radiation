"""PowerPoint (.pptx) processor — python-pptx, per-slide text + speaker notes."""
from __future__ import annotations

import io

from .base import make_result


def process_pptx(file_path: str, data: bytes) -> dict:
    warnings: list[str] = []
    errors: list[str] = []
    slides_out: list[dict] = []

    try:
        from pptx import Presentation

        prs = Presentation(io.BytesIO(data))
        for idx, slide in enumerate(prs.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text and shape.text.strip():
                    texts.append(shape.text.strip())
            notes = ""
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = (slide.notes_slide.notes_text_frame.text or "").strip()
            slides_out.append({"index": idx, "text": "\n".join(texts), "notes": notes})

        if not any(s["text"] or s["notes"] for s in slides_out):
            warnings.append("No extractable text found in this PowerPoint file.")
    except Exception as exc:
        errors.append(f"Failed to read PowerPoint file: {exc}")

    full_text = "\n\n".join(f"[Slide {s['index']}]\n{s['text']}" for s in slides_out if s["text"])
    return make_result(
        file_path, "pptx",
        text=full_text,
        slides=slides_out,
        metadata={"slide_count": len(slides_out)},
        warnings=warnings, errors=errors,
    )
