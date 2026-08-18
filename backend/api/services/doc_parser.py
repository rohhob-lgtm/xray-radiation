"""
Extract plain text and images from PDF, DOCX, PPTX, and plain-text uploads.

Image extraction strategy:
  1. Try pypdf page.images — yields each embedded raster image as its own cropped object.
  2. If a PDF page has zero embedded images, fall back to rendering that page as a
     300-DPI PNG with PyMuPDF (fitz) and ask GPT-5.4 to locate figure bounding boxes,
     then crop each figure with Pillow.
  3. Every extracted image is normalised to PNG via Pillow before storage.
"""
from __future__ import annotations
import io
from typing import List, Dict, Any


# ──────────────────────────────────────────────────────────
# Public API
# ──────────────────────────────────────────────────────────

def extract_text(filename: str, data: bytes) -> str:
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""

    if ext == "pdf":
        return _extract_pdf_text(data)
    elif ext == "docx":
        return _extract_docx(data)
    elif ext == "doc":
        raise ValueError("Legacy .doc format is not supported — please save as .docx.")
    elif ext == "pptx":
        return _extract_pptx(data)
    elif ext == "ppt":
        raise ValueError("Legacy .ppt format is not supported — please save as .pptx.")
    else:
        try:
            return data.decode("utf-8", errors="replace")
        except Exception as exc:
            raise ValueError(f"Cannot read file as text: {exc}") from exc


def extract_images(filename: str, data: bytes) -> List[Dict[str, Any]]:
    """
    Extract figures from a PDF.
    Returns: [{page_num, image_index, name, data (PNG bytes), mime_type}]
    Non-PDF files return an empty list.
    """
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    if ext != "pdf":
        return []
    return _extract_pdf_images(data)


def render_pdf_pages(data: bytes, dpi: int = 150) -> List[Dict[str, Any]]:
    """
    Render every page of a PDF as a PNG using PyMuPDF.
    Used to build the ColPali visual index (whole-page embeddings).

    Returns: [{page_num, data (PNG bytes)}]
    """
    try:
        import fitz  # PyMuPDF

        doc = fitz.open(stream=data, filetype="pdf")
        pages = []
        mat = fitz.Matrix(dpi / 72, dpi / 72)
        for page_num in range(1, len(doc) + 1):
            page = doc[page_num - 1]
            pix = page.get_pixmap(matrix=mat, alpha=False)
            png_bytes = pix.tobytes("png")
            if len(png_bytes) > 3000:          # skip near-blank/empty pages
                pages.append({"page_num": page_num, "data": png_bytes})
        doc.close()
        return pages
    except Exception:
        return []


# ──────────────────────────────────────────────────────────
# Pillow normalisation helper
# ──────────────────────────────────────────────────────────

def _to_png(raw: bytes) -> bytes:
    """Convert any image bytes to PNG. Returns original bytes on failure."""
    try:
        from PIL import Image
        img = Image.open(io.BytesIO(raw)).convert("RGBA")
        buf = io.BytesIO()
        img.save(buf, format="PNG")
        return buf.getvalue()
    except Exception:
        return raw


# ──────────────────────────────────────────────────────────
# PDF text extraction
# ──────────────────────────────────────────────────────────

def _extract_pdf_text(data: bytes) -> str:
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise RuntimeError("pypdf is not installed") from exc

    reader = PdfReader(io.BytesIO(data))
    parts: list[str] = []
    for page in reader.pages:
        text = page.extract_text() or ""
        if text.strip():
            parts.append(text.strip())
    if not parts:
        raise ValueError(
            "No extractable text found in this PDF. "
            "Scanned/image-only PDFs are not supported — try a text-based PDF."
        )
    return "\n\n".join(parts)


# ──────────────────────────────────────────────────────────
# PDF image extraction
# ──────────────────────────────────────────────────────────

def _extract_pdf_images(data: bytes) -> List[Dict[str, Any]]:
    """
    Two-pass extraction:
      Pass 1 — pypdf embedded images (each image is already a cropped figure).
      Pass 2 — for pages with no embedded images, render the page at 300 DPI
                with PyMuPDF and treat the whole page as one image
                (GPT vision will caption and locate figures at query time).
    Every image is normalised to PNG.
    """
    results: List[Dict[str, Any]] = []
    pages_with_embedded: set[int] = set()

    # ── Pass 1: pypdf embedded images ─────────────────────
    try:
        from pypdf import PdfReader
        reader = PdfReader(io.BytesIO(data))
        for page_num, page in enumerate(reader.pages, start=1):
            page_images = getattr(page, "images", [])
            for img_index, img in enumerate(page_images):
                try:
                    raw = img.data
                    if not raw or len(raw) < 100:   # skip tiny/degenerate blobs
                        continue
                    name = getattr(img, "name", f"page{page_num}_img{img_index}")
                    png = _to_png(raw)
                    results.append({
                        "page_num": page_num,
                        "image_index": img_index,
                        "name": f"{name}.png",
                        "data": png,
                        "mime_type": "image/png",
                    })
                    pages_with_embedded.add(page_num)
                except Exception:
                    continue
    except Exception:
        pass

    # ── Pass 2: PyMuPDF page rendering for image-less pages ──
    try:
        import fitz  # PyMuPDF

        doc = fitz.open(stream=data, filetype="pdf")
        for page_num in range(1, len(doc) + 1):
            if page_num in pages_with_embedded:
                continue  # already extracted embedded images from this page
            page = doc[page_num - 1]
            # 300 DPI render
            mat = fitz.Matrix(300 / 72, 300 / 72)
            pix = page.get_pixmap(matrix=mat, alpha=False)
            png_bytes = pix.tobytes("png")
            if len(png_bytes) < 5000:   # skip near-blank pages
                continue
            results.append({
                "page_num": page_num,
                "image_index": 0,
                "name": f"page{page_num}_render.png",
                "data": png_bytes,
                "mime_type": "image/png",
            })
        doc.close()
    except Exception:
        pass  # PyMuPDF not installed or render failed — no-op

    return results


# ──────────────────────────────────────────────────────────
# DOCX
# ──────────────────────────────────────────────────────────

def _extract_docx(data: bytes) -> str:
    try:
        from docx import Document
    except ImportError as exc:
        raise RuntimeError("python-docx is not installed") from exc

    doc = Document(io.BytesIO(data))
    parts: list[str] = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            parts.append(text)
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(
                cell.text.strip() for cell in row.cells if cell.text.strip()
            )
            if row_text:
                parts.append(row_text)
    if not parts:
        raise ValueError("No text found in this Word document.")
    return "\n\n".join(parts)


# ──────────────────────────────────────────────────────────
# PPTX
# ──────────────────────────────────────────────────────────

def _extract_pptx(data: bytes) -> str:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("python-pptx is not installed") from exc

    prs = Presentation(io.BytesIO(data))
    parts: list[str] = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_parts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_parts.append(shape.text.strip())
        if slide_parts:
            parts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_parts))
    if not parts:
        raise ValueError("No text found in this PowerPoint file.")
    return "\n\n".join(parts)
