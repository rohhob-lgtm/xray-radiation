"""
Layout style learner — extracts PPTX presentation style profiles without API calls.

When a PPTX file is uploaded, python-pptx extracts:
  - Slide dimensions and aspect ratio
  - Theme colours (up to 8)
  - Title font (name, size, bold, colour)
  - Body font (name, size, colour)
  - Background colour
  - Placeholder positions for all layouts
  - Table style (border colour, header fill)
  - is_rtl detection from slide content
  - Paragraph spacing from slide master

The profile is stored in LayoutStyle.properties and can be applied to future
translation jobs to produce on-brand Arabic PPTX output.
"""
from __future__ import annotations
import logging
from typing import Optional

from sqlalchemy.orm import Session

log = logging.getLogger(__name__)


# ── Safe property helpers ──────────────────────────────────────────────────────

def _safe(fn, default=None):
    """Call fn(); return default on any exception."""
    try:
        return fn()
    except Exception:
        return default


def _font_props(run) -> dict:
    """Safely read font size, name, bold, and colour from a run."""
    font = _safe(lambda: run.font)
    if font is None:
        return {}
    return {
        "size":  _safe(lambda: font.size),
        "name":  _safe(lambda: font.name),
        "bold":  _safe(lambda: font.bold),
        "color": _safe(lambda: str(font.color.rgb) if font.color and font.color.rgb else None),
    }


def _text_is_rtl(text: str) -> bool:
    """Return True when >30% of characters are in an RTL script."""
    rtl_count = sum(1 for ch in text if "\u0600" <= ch <= "\u06FF" or "\u0590" <= ch <= "\u05FF")
    return len(text) > 0 and rtl_count / len(text) > 0.30


def extract_layout(filename: str, pptx_bytes: bytes) -> Optional[dict]:
    """
    Extract a style profile from a PPTX file.
    Returns a dict suitable for storing in LayoutStyle.properties, or None on
    failure.  Individual extraction steps fail silently so a partial result is
    always returned if the file can be opened at all.
    """
    try:
        import io
        from pptx import Presentation

        prs = Presentation(io.BytesIO(pptx_bytes))
    except Exception as exc:
        log.warning("Cannot open PPTX %s: %s", filename, exc)
        return None

    # ── Slide dimensions ──────────────────────────────────────────────────────
    w_in = _safe(lambda: round(prs.slide_width  / 914400, 2), 10.0)
    h_in = _safe(lambda: round(prs.slide_height / 914400, 2),  7.5)
    slide_count = _safe(lambda: len(prs.slides), 0)

    # ── Theme colours ─────────────────────────────────────────────────────────
    theme_colors: list[str] = []
    try:
        from pptx.oxml.ns import qn
        for el in prs.slide_master.element.iter(qn("a:srgbClr")):
            val = _safe(lambda _el=el: _el.get("val", ""), "")
            if val and val not in theme_colors:
                theme_colors.append(val.upper())
            if len(theme_colors) >= 8:
                break
    except Exception:
        pass

    # ── Title / body fonts from first few slides ───────────────────────────────
    title_font_name  = "Calibri"
    title_font_size  = 28
    title_bold       = True
    title_color      = "1A1A1A"
    body_font_name   = "Calibri"
    body_font_size   = 18
    body_color       = "FFFFFF"
    bg_color         = "141414"

    slides = _safe(lambda: list(prs.slides[:5]), [])
    for slide in slides:
        shapes = _safe(lambda s=slide: list(s.shapes), [])
        for shape in shapes:
            has_tf = _safe(lambda sh=shape: sh.has_text_frame, False)
            if not has_tf:
                continue
            stype = _safe(lambda sh=shape: sh.shape_type, None)
            if stype == 13:
                continue
            tf = _safe(lambda sh=shape: sh.text_frame, None)
            if tf is None:
                continue
            paras = _safe(lambda t=tf: list(t.paragraphs), [])
            for para in paras:
                runs = _safe(lambda p=para: list(p.runs), [])
                for run in runs:
                    fp = _font_props(run)
                    fsize = fp.get("size")
                    fname = fp.get("name")
                    fbold = fp.get("bold")
                    fcolor = fp.get("color")
                    if fsize and fsize > 700000:   # > ~19pt → likely a title
                        if fname:  title_font_name = fname
                        title_font_size = round(fsize / 12700)
                        if fbold is not None: title_bold = bool(fbold)
                        if fcolor: title_color = fcolor
                    else:
                        if fname:  body_font_name = fname
                        if fsize:  body_font_size = round(fsize / 12700)
                        if fcolor: body_color = fcolor

        # Background colour
        try:
            bg = slide.background.fill
            if _safe(lambda: bg.type) is not None:
                bg.solid()
                rgb = _safe(lambda: str(bg.fore_color.rgb) if bg.fore_color and bg.fore_color.rgb else None)
                if rgb:
                    bg_color = rgb
        except Exception:
            pass

    # ── Placeholder positions — all layouts ───────────────────────────────────
    all_layouts: list[dict] = []
    try:
        for layout_idx, layout in enumerate(_safe(lambda: list(prs.slide_layouts), [])):
            layout_name = _safe(lambda l=layout: l.name, f"Layout {layout_idx}")
            phs = []
            for ph in _safe(lambda l=layout: list(l.placeholders), []):
                phs.append({
                    "idx":      _safe(lambda p=ph: p.placeholder_format.idx),
                    "type":     _safe(lambda p=ph: str(p.placeholder_format.type)),
                    "left_in":  _safe(lambda p=ph: round(p.left  / 914400, 3) if p.left  is not None else None),
                    "top_in":   _safe(lambda p=ph: round(p.top   / 914400, 3) if p.top   is not None else None),
                    "width_in": _safe(lambda p=ph: round(p.width / 914400, 3) if p.width is not None else None),
                    "height_in":_safe(lambda p=ph: round(p.height/ 914400, 3) if p.height is not None else None),
                })
            all_layouts.append({"name": layout_name, "placeholders": phs})
            if len(all_layouts) >= 12:
                break
    except Exception:
        pass

    # ── Paragraph spacing from slide master ───────────────────────────────────
    para_spacing_before = None
    para_spacing_after  = None
    try:
        from pptx.oxml.ns import qn as _qn
        master = prs.slide_master
        for ph in _safe(lambda: list(master.placeholders), []):
            tf = _safe(lambda p=ph: p.text_frame, None)
            if tf:
                paras = _safe(lambda t=tf: list(t.paragraphs), [])
                for para in paras[:3]:
                    pf = para._p.find(_qn("a:pPr"))
                    if pf is not None:
                        spcBef = pf.find(_qn("a:spcBef"))
                        spcAft = pf.find(_qn("a:spcAft"))
                        if spcBef is not None:
                            spcPts = spcBef.find(_qn("a:spcPts"))
                            if spcPts is not None:
                                val = spcPts.get("val")
                                if val:
                                    para_spacing_before = int(val) // 100  # hundredths of pt → pt
                        if spcAft is not None:
                            spcPts = spcAft.find(_qn("a:spcPts"))
                            if spcPts is not None:
                                val = spcPts.get("val")
                                if val:
                                    para_spacing_after = int(val) // 100
                    break
    except Exception:
        pass

    # ── is_rtl detection — sample up to 20 slides ────────────────────────────
    rtl_slide_count = 0
    sampled = _safe(lambda: list(prs.slides[:20]), [])
    for slide in sampled:
        all_text = " ".join(
            _safe(lambda sh=shape: sh.text_frame.text, "") or ""
            for shape in _safe(lambda s=slide: list(s.shapes), [])
            if _safe(lambda sh=shape: sh.has_text_frame, False)
        )
        if _text_is_rtl(all_text):
            rtl_slide_count += 1
    source_is_rtl = len(sampled) > 0 and rtl_slide_count / len(sampled) > 0.5

    # ── First-layout placeholder positions (kept for backward compat) ─────────
    placeholder_positions: list[dict] = []
    if all_layouts:
        placeholder_positions = all_layouts[0].get("placeholders", [])

    aspect = _safe(lambda: f"{w_in:.1f}x{h_in:.1f}", f"{w_in}x{h_in}")

    return {
        # Dimensions
        "slide_width_in":       w_in,
        "slide_height_in":      h_in,
        "slide_count":          slide_count,
        "aspect_ratio":         aspect,
        # Colors
        "theme_colors":         theme_colors[:8],
        "bg_color":             bg_color,
        # Title typography
        "title_font_name":      title_font_name,
        "title_font_size":      title_font_size,
        "title_bold":           title_bold,
        "title_color":          title_color,
        # Body typography
        "body_font_name":       body_font_name,
        "body_font_size":       body_font_size,
        "body_color":           body_color,
        # Spacing
        "para_spacing_before":  para_spacing_before,
        "para_spacing_after":   para_spacing_after,
        # Layouts
        "placeholder_positions": placeholder_positions,
        "all_layouts":          all_layouts,
        # RTL
        "source_is_rtl":        source_is_rtl,
        # Metadata (user-editable, stored here to avoid schema changes)
        "organisation":         "",
        "department":           "",
        "language":             "",
    }


def learn_and_store(db: Session, doc_id: str, filename: str, pptx_bytes: bytes) -> Optional[str]:
    """
    Extract and persist a LayoutStyle. Returns the new style ID, or None on failure.
    Idempotent — if a style for this filename already exists, return its ID.
    """
    from api.db.models import LayoutStyle

    existing = db.query(LayoutStyle).filter(
        LayoutStyle.source_filename == filename,
        LayoutStyle.source_doc_id == doc_id,
    ).first()
    if existing:
        log.debug("Layout style already exists for %s", filename)
        return existing.id

    props = extract_layout(filename, pptx_bytes)
    if not props:
        return None

    name = filename.replace(".pptx", "").replace("_", " ").replace("-", " ")[:80]

    style = LayoutStyle(
        name=name,
        source_filename=filename,
        source_doc_id=doc_id,
        properties=props,
        is_default=False,
    )
    db.add(style)
    try:
        db.commit()
        log.info("Saved layout style '%s' from %s", name, filename)
        return style.id
    except Exception as exc:
        db.rollback()
        log.warning("Layout style save failed for %s: %s", filename, exc)
        return None


def list_styles(db: Session) -> list[dict]:
    """List all stored layout styles with their enriched profile."""
    from api.db.models import LayoutStyle
    rows = db.query(LayoutStyle).order_by(LayoutStyle.created_at.desc()).all()
    return [_style_to_dict(s) for s in rows]


def get_style(db: Session, style_id: str) -> Optional[dict]:
    """Return a single style by ID, or None."""
    from api.db.models import LayoutStyle
    s = db.query(LayoutStyle).filter(LayoutStyle.id == style_id).first()
    return _style_to_dict(s) if s else None


def _style_to_dict(s) -> dict:
    """Serialize a LayoutStyle ORM row."""
    props = s.properties or {}
    return {
        "id":              s.id,
        "name":            s.name,
        "source_filename": s.source_filename,
        "source_doc_id":   s.source_doc_id,
        "is_default":      s.is_default,
        "created_at":      s.created_at.isoformat() if s.created_at else None,
        # Flatten user-editable metadata for easy frontend access
        "organisation":    props.get("organisation", ""),
        "department":      props.get("department", ""),
        "language":        props.get("language", ""),
        # Rich profile for display
        "style_profile":   props,
        # Legacy field kept for existing callers
        "properties":      props,
    }


def set_default_style(db: Session, style_id: str) -> bool:
    """Set a layout style as the default (clears others)."""
    from api.db.models import LayoutStyle
    db.query(LayoutStyle).update({LayoutStyle.is_default: False})
    style = db.query(LayoutStyle).filter(LayoutStyle.id == style_id).first()
    if not style:
        return False
    style.is_default = True
    db.commit()
    return True


# Backward-compat alias used by pptx_gen.py / study.py
set_default = set_default_style
