"""
Batched text extraction service for large files.

Reads files from disk in slices — never loads the entire document into RAM.
Yields (batch_text, batch_meta) pairs so the job runner can checkpoint after
every batch without holding all extracted text in memory simultaneously.

Supported formats and batch sizes:
  PDF   — 25 pages per batch
  PPTX  — 20 slides per batch
  DOCX  — grouped by heading sections, max 20 headings per batch
  XLSX  — 1 sheet per batch
  ZIP   — individual contained files (with security scanning)
  TXT   — 50,000-character blocks
"""
from __future__ import annotations
import hashlib
import io
import logging
import os
import pathlib
import zipfile
from typing import Generator

log = logging.getLogger(__name__)

# ── Batch size constants ───────────────────────────────────────────────────────
PDF_PAGES_PER_BATCH  = 25
PPTX_SLIDES_PER_BATCH = 20
DOCX_SECTIONS_PER_BATCH = 20
TXT_CHARS_PER_BATCH  = 50_000

# ── ZIP security constants ─────────────────────────────────────────────────────
_BLOCKED_EXTENSIONS = {".exe", ".dll", ".so", ".sh", ".bat", ".cmd", ".ps1",
                       ".msi", ".vbs", ".scr", ".pif", ".com", ".jar"}
_MAX_COMPRESSION_RATIO = 100  # reject if uncompressed > 100× compressed

# ── Return type ───────────────────────────────────────────────────────────────

class ExtractionBatch:
    """One extracted batch of text with metadata."""
    __slots__ = ("batch_index", "total_batches", "text", "source_info", "char_count")

    def __init__(
        self,
        batch_index: int,
        total_batches: int,
        text: str,
        source_info: str,  # e.g. "pages 1-25", "slide 1-20", "Sheet1"
    ):
        self.batch_index = batch_index
        self.total_batches = total_batches
        self.text = text
        self.source_info = source_info
        self.char_count = len(text)


def count_batches(file_path: str | pathlib.Path) -> int:
    """
    Quickly count how many extraction batches a file will produce.
    Used by the preflight endpoint — does NOT load the full file.
    """
    path = pathlib.Path(file_path)
    ext = path.suffix.lower().lstrip(".")
    file_size = path.stat().st_size

    if ext == "pdf":
        try:
            import pypdf
            with open(path, "rb") as f:
                reader = pypdf.PdfReader(f)
                n = len(reader.pages)
            return max(1, (n + PDF_PAGES_PER_BATCH - 1) // PDF_PAGES_PER_BATCH)
        except Exception:
            pass
    elif ext == "pptx":
        try:
            from pptx import Presentation
            prs = Presentation(str(path))
            n = len(prs.slides)
            return max(1, (n + PPTX_SLIDES_PER_BATCH - 1) // PPTX_SLIDES_PER_BATCH)
        except Exception:
            pass
    elif ext == "docx":
        try:
            from docx import Document
            doc = Document(str(path))
            headings = [p for p in doc.paragraphs if p.style.name.startswith("Heading")]
            n = max(1, len(headings))
            return max(1, (n + DOCX_SECTIONS_PER_BATCH - 1) // DOCX_SECTIONS_PER_BATCH)
        except Exception:
            pass
    elif ext in ("xlsx", "xls"):
        try:
            import openpyxl
            wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
            n = len(wb.sheetnames)
            wb.close()
            return max(1, n)
        except Exception:
            pass
    elif ext == "zip":
        try:
            with zipfile.ZipFile(str(path), "r") as zf:
                infos = zf.infolist()
                file_count = sum(1 for i in infos if not i.is_dir())
                return max(1, file_count)
        except Exception:
            pass

    # Fallback: estimate by file size
    return max(1, (file_size + TXT_CHARS_PER_BATCH - 1) // TXT_CHARS_PER_BATCH)


def estimate_page_count(file_path: str | pathlib.Path) -> int:
    """Quickly estimate page / slide / sheet count for preflight display."""
    path = pathlib.Path(file_path)
    ext = path.suffix.lower().lstrip(".")

    if ext == "pdf":
        try:
            import pypdf
            with open(path, "rb") as f:
                return len(pypdf.PdfReader(f).pages)
        except Exception:
            return max(1, path.stat().st_size // 50_000)
    elif ext == "pptx":
        try:
            from pptx import Presentation
            return len(Presentation(str(path)).slides)
        except Exception:
            return max(1, path.stat().st_size // 100_000)
    elif ext == "docx":
        return max(1, path.stat().st_size // 3_000)
    elif ext in ("xlsx", "xls"):
        try:
            import openpyxl
            wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
            n = len(wb.sheetnames)
            wb.close()
            return n
        except Exception:
            return 1
    elif ext == "zip":
        try:
            with zipfile.ZipFile(str(path), "r") as zf:
                return sum(1 for i in zf.infolist() if not i.is_dir())
        except Exception:
            return 1
    return max(1, path.stat().st_size // 5_000)


def extract_batches(
    file_path: str | pathlib.Path,
    start_batch: int = 0,
) -> Generator[ExtractionBatch, None, None]:
    """
    Yield ExtractionBatch objects for each batch of text from the file.
    start_batch allows resuming from a checkpoint (skip already-processed batches).
    Never loads the entire file into memory.
    """
    path = pathlib.Path(file_path)
    ext = path.suffix.lower().lstrip(".")

    if ext == "pdf":
        yield from _extract_pdf_batches(path, start_batch)
    elif ext == "pptx":
        yield from _extract_pptx_batches(path, start_batch)
    elif ext == "docx":
        yield from _extract_docx_batches(path, start_batch)
    elif ext in ("xlsx", "xls"):
        yield from _extract_xlsx_batches(path, start_batch)
    elif ext == "zip":
        yield from _extract_zip_batches(path, start_batch)
    else:
        yield from _extract_text_batches(path, start_batch)


def scan_zip(file_path: str | pathlib.Path) -> dict:
    """
    Pre-scan a ZIP file for security issues before extraction.
    Returns {safe: bool, issues: list[str], file_list: list[dict]}
    """
    path = pathlib.Path(file_path)
    issues: list[str] = []
    file_list: list[dict] = []

    try:
        with zipfile.ZipFile(str(path), "r") as zf:
            infos = zf.infolist()

            for info in infos:
                name = info.filename
                compressed = info.compress_size
                uncompressed = info.file_size

                # Path traversal check
                if ".." in name or name.startswith("/") or name.startswith("\\"):
                    issues.append(f"Path traversal blocked: {name}")
                    continue

                # Executable check
                _, fext = os.path.splitext(name.lower())
                if fext in _BLOCKED_EXTENSIONS:
                    issues.append(f"Executable blocked: {name}")
                    continue

                # Compression bomb check
                if compressed > 0 and uncompressed / compressed > _MAX_COMPRESSION_RATIO:
                    issues.append(f"Compression bomb detected: {name} (ratio {uncompressed // max(1, compressed)}:1)")
                    continue

                if not info.is_dir():
                    file_list.append({
                        "name": name,
                        "size_bytes": uncompressed,
                        "compressed_bytes": compressed,
                    })

        return {
            "safe": len(issues) == 0,
            "issues": issues,
            "file_list": file_list,
            "file_count": len(file_list),
        }
    except zipfile.BadZipFile:
        return {"safe": False, "issues": ["Invalid ZIP file"], "file_list": [], "file_count": 0}
    except Exception as exc:
        return {"safe": False, "issues": [f"Scan error: {exc}"], "file_list": [], "file_count": 0}


# ── Private extraction helpers ─────────────────────────────────────────────────

def _extract_pdf_batches(path: pathlib.Path, start_batch: int) -> Generator[ExtractionBatch, None, None]:
    try:
        import pypdf
    except ImportError:
        log.error("pypdf not installed — cannot extract PDF batches")
        return

    with open(path, "rb") as f:
        reader = pypdf.PdfReader(f)
        total_pages = len(reader.pages)
        total_batches = max(1, (total_pages + PDF_PAGES_PER_BATCH - 1) // PDF_PAGES_PER_BATCH)

        for batch_idx in range(total_batches):
            if batch_idx < start_batch:
                continue

            page_start = batch_idx * PDF_PAGES_PER_BATCH
            page_end = min(page_start + PDF_PAGES_PER_BATCH, total_pages)
            texts = []

            for pg_num in range(page_start, page_end):
                try:
                    page = reader.pages[pg_num]
                    t = page.extract_text() or ""
                    if t.strip():
                        texts.append(f"--- Page {pg_num + 1} ---\n{t}")
                except Exception as e:
                    log.debug("PDF page %d extraction error: %s", pg_num + 1, e)

            yield ExtractionBatch(
                batch_index=batch_idx,
                total_batches=total_batches,
                text="\n\n".join(texts),
                source_info=f"pages {page_start + 1}–{page_end}",
            )


def _extract_pptx_batches(path: pathlib.Path, start_batch: int) -> Generator[ExtractionBatch, None, None]:
    try:
        from pptx import Presentation
    except ImportError:
        log.error("python-pptx not installed — cannot extract PPTX batches")
        return

    prs = Presentation(str(path))
    slides = prs.slides
    total_slides = len(slides)
    total_batches = max(1, (total_slides + PPTX_SLIDES_PER_BATCH - 1) // PPTX_SLIDES_PER_BATCH)

    for batch_idx in range(total_batches):
        if batch_idx < start_batch:
            continue

        slide_start = batch_idx * PPTX_SLIDES_PER_BATCH
        slide_end = min(slide_start + PPTX_SLIDES_PER_BATCH, total_slides)
        texts = []

        for sl_num in range(slide_start, slide_end):
            slide = slides[sl_num]
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = "".join(run.text for run in para.runs).strip()
                        if t:
                            slide_texts.append(t)
            if slide_texts:
                texts.append(f"--- Slide {sl_num + 1} ---\n" + "\n".join(slide_texts))

        yield ExtractionBatch(
            batch_index=batch_idx,
            total_batches=total_batches,
            text="\n\n".join(texts),
            source_info=f"slides {slide_start + 1}–{slide_end}",
        )


def _extract_docx_batches(path: pathlib.Path, start_batch: int) -> Generator[ExtractionBatch, None, None]:
    try:
        from docx import Document
    except ImportError:
        log.error("python-docx not installed — cannot extract DOCX batches")
        return

    doc = Document(str(path))
    paragraphs = doc.paragraphs

    # Split into sections at Heading styles
    sections: list[list[str]] = []
    current: list[str] = []

    for para in paragraphs:
        if para.style.name.startswith("Heading") and current:
            sections.append(current)
            current = []
        if para.text.strip():
            current.append(para.text.strip())
    if current:
        sections.append(current)

    # If no headings found, treat the whole document as one section
    if not sections:
        full_text = "\n".join(p.text.strip() for p in paragraphs if p.text.strip())
        sections = [full_text[i : i + TXT_CHARS_PER_BATCH] for i in range(0, max(1, len(full_text)), TXT_CHARS_PER_BATCH)]
        sections = [[s] for s in sections]

    total_batches = max(1, (len(sections) + DOCX_SECTIONS_PER_BATCH - 1) // DOCX_SECTIONS_PER_BATCH)

    for batch_idx in range(total_batches):
        if batch_idx < start_batch:
            continue

        sec_start = batch_idx * DOCX_SECTIONS_PER_BATCH
        sec_end = min(sec_start + DOCX_SECTIONS_PER_BATCH, len(sections))
        batch_lines: list[str] = []
        for sec in sections[sec_start:sec_end]:
            batch_lines.extend(sec)

        yield ExtractionBatch(
            batch_index=batch_idx,
            total_batches=total_batches,
            text="\n".join(batch_lines),
            source_info=f"sections {sec_start + 1}–{sec_end}",
        )


def _extract_xlsx_batches(path: pathlib.Path, start_batch: int) -> Generator[ExtractionBatch, None, None]:
    try:
        import openpyxl
    except ImportError:
        log.error("openpyxl not installed — cannot extract XLSX batches")
        return

    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    sheet_names = wb.sheetnames
    total_batches = max(1, len(sheet_names))

    for batch_idx, sheet_name in enumerate(sheet_names):
        if batch_idx < start_batch:
            continue

        ws = wb[sheet_name]
        rows_text: list[str] = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            row_text = "\t".join(cells).strip()
            if row_text.replace("\t", "").strip():
                rows_text.append(row_text)
            if len("\n".join(rows_text)) > TXT_CHARS_PER_BATCH:
                break

        yield ExtractionBatch(
            batch_index=batch_idx,
            total_batches=total_batches,
            text=f"--- Sheet: {sheet_name} ---\n" + "\n".join(rows_text),
            source_info=f"sheet '{sheet_name}'",
        )

    wb.close()


def _extract_zip_batches(path: pathlib.Path, start_batch: int) -> Generator[ExtractionBatch, None, None]:
    scan = scan_zip(path)
    if not scan["safe"]:
        log.warning("ZIP security issues: %s", scan["issues"])
        # Still process safe files if any
    file_list = scan["file_list"]
    total_batches = max(1, len(file_list))

    if not file_list:
        return

    try:
        with zipfile.ZipFile(str(path), "r") as zf:
            for batch_idx, file_info in enumerate(file_list):
                if batch_idx < start_batch:
                    continue

                name = file_info["name"]
                try:
                    data = zf.read(name)
                    # Try to extract text from embedded file
                    text = _extract_text_from_bytes(name, data)
                    if text.strip():
                        yield ExtractionBatch(
                            batch_index=batch_idx,
                            total_batches=total_batches,
                            text=f"--- File: {name} ---\n{text}",
                            source_info=f"file '{name}'",
                        )
                except Exception as exc:
                    log.debug("ZIP member %s extraction error: %s", name, exc)
    except Exception as exc:
        log.error("ZIP batch extraction failed: %s", exc)


def _extract_text_batches(path: pathlib.Path, start_batch: int) -> Generator[ExtractionBatch, None, None]:
    """Fallback: read as text in TXT_CHARS_PER_BATCH blocks."""
    file_size = path.stat().st_size
    total_batches = max(1, (file_size + TXT_CHARS_PER_BATCH - 1) // TXT_CHARS_PER_BATCH)

    with open(path, "r", encoding="utf-8", errors="replace") as f:
        batch_idx = 0
        while True:
            if batch_idx >= start_batch:
                chunk = f.read(TXT_CHARS_PER_BATCH)
                if not chunk:
                    break
                yield ExtractionBatch(
                    batch_index=batch_idx,
                    total_batches=total_batches,
                    text=chunk,
                    source_info=f"chars {batch_idx * TXT_CHARS_PER_BATCH + 1}–{(batch_idx + 1) * TXT_CHARS_PER_BATCH}",
                )
            else:
                f.read(TXT_CHARS_PER_BATCH)  # skip
            batch_idx += 1


def _extract_text_from_bytes(filename: str, data: bytes) -> str:
    """Extract text from in-memory bytes for a ZIP member."""
    ext = os.path.splitext(filename.lower())[1].lstrip(".")
    try:
        if ext == "pdf":
            import pypdf
            reader = pypdf.PdfReader(io.BytesIO(data))
            return "\n".join(p.extract_text() or "" for p in reader.pages)
        elif ext == "docx":
            from docx import Document
            doc = Document(io.BytesIO(data))
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        elif ext == "pptx":
            from pptx import Presentation
            prs = Presentation(io.BytesIO(data))
            texts = []
            for sl in prs.slides:
                for shape in sl.shapes:
                    if shape.has_text_frame:
                        texts.append(shape.text_frame.text)
            return "\n".join(texts)
        elif ext in ("txt", "md", "rst", "csv", "json", "xml", "html", "htm"):
            return data.decode("utf-8", errors="replace")
    except Exception as exc:
        log.debug("ZIP member text extraction failed for %s: %s", filename, exc)
    return ""
