"""PPTX knowledge-base indexing and retrieval for training generation.

This module adds targeted, slide-level reference retrieval from previously uploaded
PowerPoint files while keeping the uploaded manual as the authoritative source.
"""
from __future__ import annotations

import io
import math
import re
from collections import defaultdict
from datetime import datetime, timezone
from typing import Any

from sqlalchemy.orm import Session

from api.db.models import PptxPresentationIndex, PptxSlideIndex


_EXCLUDE_STATUSES = {"obsolete", "do_not_use"}


def _now() -> datetime:
    return datetime.now(timezone.utc)


def _as_aware_utc(dt: datetime | None) -> datetime:
    """Normalize DB timestamps to timezone-aware UTC for safe datetime arithmetic."""
    if dt is None:
        return _now()
    if dt.tzinfo is None:
        return dt.replace(tzinfo=timezone.utc)
    return dt.astimezone(timezone.utc)


def _norm(value: Any) -> str:
    return re.sub(r"\s+", " ", str(value or "")).strip()


def _tokens(value: str) -> set[str]:
    return set(re.findall(r"[a-z0-9]{3,}", value.lower()))


def _language_hint(text: str) -> str:
    t = _norm(text)
    if not t:
        return "unknown"
    arabic_chars = sum(1 for ch in t if "\u0600" <= ch <= "\u06ff")
    ratio = arabic_chars / max(1, len(t))
    if ratio > 0.25:
        return "arabic"
    return "english"


def _is_heading_text(text: str) -> bool:
    t = _norm(text)
    if not t:
        return False
    if len(t) > 160:
        return False
    words = t.split()
    return 1 <= len(words) <= 18


def _shape_text(shape: Any) -> str:
    try:
        if getattr(shape, "has_text_frame", False) and shape.text_frame:
            lines = []
            for p in shape.text_frame.paragraphs:
                line = _norm(getattr(p, "text", ""))
                if line:
                    lines.append(line)
            return "\n".join(lines)
    except Exception:
        return ""
    return ""


def _table_text(shape: Any) -> str:
    try:
        if not getattr(shape, "has_table", False):
            return ""
        rows = []
        for row in shape.table.rows:
            cells = [_norm(cell.text) for cell in row.cells if _norm(cell.text)]
            if cells:
                rows.append(" | ".join(cells))
        return "\n".join(rows)
    except Exception:
        return ""


def _extract_title(slide: Any) -> tuple[str, Any | None]:
    title_shape = None
    title_text = ""
    try:
        if slide.shapes.title is not None:
            title_shape = slide.shapes.title
            title_text = _norm(slide.shapes.title.text)
    except Exception:
        title_shape = None
        title_text = ""

    if title_text:
        return title_text, title_shape

    for shape in slide.shapes:
        txt = _shape_text(shape)
        if _is_heading_text(txt):
            return txt[:240], shape
    return "Untitled Slide", None


def _layout_meta(slide: Any, title_shape: Any | None, slide_w: int, slide_h: int) -> dict:
    image_count = 0
    text_boxes = 0
    table_count = 0
    diagram_like = 0
    for shape in slide.shapes:
        st = str(getattr(shape, "shape_type", ""))
        name = _norm(getattr(shape, "name", ""))
        if "PICTURE" in st:
            image_count += 1
        if getattr(shape, "has_table", False):
            table_count += 1
        if getattr(shape, "has_text_frame", False):
            text_boxes += 1
        if "ARROW" in name.upper() or "FLOW" in name.upper() or "DIAGRAM" in name.upper():
            diagram_like += 1

    title_position = {}
    if title_shape is not None:
        try:
            title_position = {
                "left_ratio": round(float(title_shape.left) / max(1, slide_w), 4),
                "top_ratio": round(float(title_shape.top) / max(1, slide_h), 4),
                "width_ratio": round(float(title_shape.width) / max(1, slide_w), 4),
                "height_ratio": round(float(title_shape.height) / max(1, slide_h), 4),
            }
        except Exception:
            title_position = {}

    return {
        "title_position": title_position,
        "content_blocks": max(0, text_boxes - (1 if title_shape is not None else 0)),
        "image_count": image_count,
        "table_count": table_count,
        "text_boxes": text_boxes,
        "diagram_like_count": diagram_like,
    }


def _infer_manufacturer(text: str) -> str:
    lower = text.lower()
    known = [
        "rapiscan", "smiths detection", "nuctech", "astrophysics", "leidos", "votix", "l3", "canberra",
    ]
    for k in known:
        if k in lower:
            return k.title()
    m = re.search(r"manufacturer\s*[:\-]\s*([^\n]{2,80})", text, re.IGNORECASE)
    return _norm(m.group(1)) if m else ""


def _infer_model(text: str) -> str:
    m = re.search(r"\b(?:model|equipment model)\s*[:#\-]?\s*([A-Z0-9][A-Z0-9\-]{1,32})", text, re.IGNORECASE)
    return _norm(m.group(1)) if m else ""


def _infer_training_category(text: str) -> str:
    t = text.lower()
    if "maintenance" in t or "preventive" in t:
        return "maintenance"
    if "operator" in t or "operation" in t:
        return "operator"
    if "troubleshoot" in t or "fault" in t:
        return "troubleshooting"
    if "installation" in t or "commission" in t:
        return "installation"
    if "safety" in t or "radiation" in t:
        return "safety"
    return "technical"


def extract_pptx_index(filename: str, pptx_bytes: bytes) -> dict:
    """Extract presentation-level and slide-level metadata from a PPTX file."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(pptx_bytes))
    slide_w = int(prs.slide_width)
    slide_h = int(prs.slide_height)

    slides: list[dict] = []
    deck_text_blocks: list[str] = []
    topic_counter: defaultdict[str, int] = defaultdict(int)

    for idx, slide in enumerate(prs.slides, start=1):
        title, title_shape = _extract_title(slide)

        text_parts: list[str] = []
        table_parts: list[str] = []
        image_captions: list[str] = []
        diagram_labels: list[str] = []

        for shape in slide.shapes:
            txt = _shape_text(shape)
            if txt:
                text_parts.append(txt)
            tbl = _table_text(shape)
            if tbl:
                table_parts.append(tbl)

            name = _norm(getattr(shape, "name", ""))
            desc = ""
            try:
                desc = _norm(shape.alt_text)
            except Exception:
                desc = ""
            if name:
                if any(k in name.lower() for k in ["image", "picture", "photo", "fig", "figure"]):
                    image_captions.append(name)
                if any(k in name.lower() for k in ["diagram", "flow", "arrow", "chart", "process"]):
                    diagram_labels.append(name)
            if desc:
                image_captions.append(desc)

        notes = ""
        try:
            if slide.has_notes_slide and slide.notes_slide and slide.notes_slide.notes_text_frame:
                notes = _norm(slide.notes_slide.notes_text_frame.text)
        except Exception:
            notes = ""

        slide_text = "\n".join([p for p in text_parts if p])
        table_text = "\n".join([p for p in table_parts if p])
        image_text = "\n".join(sorted(set(image_captions)))
        diagram_text = "\n".join(sorted(set(diagram_labels)))

        combined = "\n".join([title, slide_text, table_text, notes, image_text, diagram_text])
        deck_text_blocks.append(combined)

        training_category = _infer_training_category(combined)
        topic_counter[training_category] += 1

        layout = _layout_meta(slide, title_shape, slide_w, slide_h)
        lang = _language_hint(combined)
        quality = 0.0
        quality += 0.35 if slide_text else 0.0
        quality += 0.20 if notes else 0.0
        quality += 0.20 if table_text else 0.0
        quality += 0.15 if layout.get("diagram_like_count", 0) > 0 else 0.0
        quality += 0.10 if layout.get("image_count", 0) > 0 else 0.0

        slides.append({
            "slide_number": idx,
            "slide_title": title,
            "slide_text": slide_text,
            "speaker_notes": notes,
            "table_content": table_text,
            "image_captions": image_text,
            "diagram_labels": diagram_text,
            "language": lang,
            "training_category": training_category,
            "visual_layout_metadata": layout,
            "quality_score": round(min(1.0, quality), 3),
        })

    deck_text = "\n".join(deck_text_blocks)
    manufacturer = _infer_manufacturer(deck_text)
    model = _infer_model(deck_text)
    language = _language_hint(deck_text)

    topics = [k for k, _ in sorted(topic_counter.items(), key=lambda kv: kv[1], reverse=True)][:12]
    top_title = slides[0]["slide_title"] if slides else filename.rsplit(".", 1)[0]

    return {
        "presentation": {
            "filename": filename,
            "course_title": top_title,
            "course_type": _infer_training_category(deck_text),
            "target_audience": "",
            "equipment_family": "",
            "equipment_name": model,
            "equipment_model": model,
            "manufacturer": manufacturer,
            "language": language,
            "slide_count": len(slides),
            "main_topics": topics,
        },
        "slides": slides,
    }


def store_pptx_index(db: Session, doc_id: str, filename: str, extracted: dict) -> dict:
    """Persist PPTX presentation and slide index rows for a knowledge-base file."""
    db.query(PptxSlideIndex).filter(PptxSlideIndex.doc_id == doc_id).delete()
    db.query(PptxPresentationIndex).filter(PptxPresentationIndex.doc_id == doc_id).delete()
    db.flush()

    pres = extracted.get("presentation") or {}
    presentation = PptxPresentationIndex(
        doc_id=doc_id,
        filename=filename,
        course_title=_norm(pres.get("course_title")),
        course_type=_norm(pres.get("course_type")),
        target_audience=_norm(pres.get("target_audience")),
        equipment_family=_norm(pres.get("equipment_family")),
        equipment_name=_norm(pres.get("equipment_name")),
        equipment_model=_norm(pres.get("equipment_model")),
        manufacturer=_norm(pres.get("manufacturer")),
        language=_norm(pres.get("language") or "unknown"),
        slide_count=int(pres.get("slide_count") or 0),
        main_topics=pres.get("main_topics") or [],
        source_status="unverified",
        uploaded_at=_now(),
    )
    db.add(presentation)
    db.flush()

    rows = extracted.get("slides") or []
    for row in rows:
        db.add(PptxSlideIndex(
            doc_id=doc_id,
            presentation_id=presentation.id,
            filename=filename,
            slide_number=int(row.get("slide_number") or 1),
            slide_title=_norm(row.get("slide_title")),
            slide_text=_norm(row.get("slide_text")),
            speaker_notes=_norm(row.get("speaker_notes")),
            table_content=_norm(row.get("table_content")),
            image_captions=_norm(row.get("image_captions")),
            diagram_labels=_norm(row.get("diagram_labels")),
            equipment_name=_norm(presentation.equipment_name),
            equipment_model=_norm(presentation.equipment_model),
            manufacturer=_norm(presentation.manufacturer),
            training_category=_norm(row.get("training_category")),
            language=_norm(row.get("language") or presentation.language),
            visual_layout_metadata=row.get("visual_layout_metadata") or {},
            source_status=presentation.source_status,
            quality_score=float(row.get("quality_score") or 0.0),
            uploaded_at=_now(),
        ))

    db.commit()
    return {
        "presentation_id": presentation.id,
        "slide_count": len(rows),
        "filename": filename,
    }


def update_source_control(db: Session, doc_id: str, updates: dict) -> dict:
    """Update source-control flags for a PPTX reference set."""
    pres = db.query(PptxPresentationIndex).filter(PptxPresentationIndex.doc_id == doc_id).first()
    if not pres:
        raise ValueError("PPTX reference index not found")

    allowed = {
        "source_status",
        "trusted",
        "manufacturer_approved",
        "internal_training_reference",
        "visual_template",
        "arabic_formatting_example",
        "obsolete",
        "do_not_use",
        "course_type",
        "target_audience",
        "equipment_family",
        "manufacturer",
        "equipment_model",
        "language",
    }
    for key, value in updates.items():
        if key in allowed:
            setattr(pres, key, value)

    if pres.obsolete:
        pres.source_status = "obsolete"
    elif pres.do_not_use:
        pres.source_status = "do_not_use"

    db.query(PptxSlideIndex).filter(PptxSlideIndex.doc_id == doc_id).update({
        "source_status": pres.source_status,
        "manufacturer": pres.manufacturer,
        "equipment_model": pres.equipment_model,
        "language": pres.language,
    })
    db.commit()
    db.refresh(pres)

    return {
        "doc_id": doc_id,
        "source_status": pres.source_status,
        "trusted": bool(pres.trusted),
        "manufacturer_approved": bool(pres.manufacturer_approved),
        "obsolete": bool(pres.obsolete),
        "do_not_use": bool(pres.do_not_use),
    }


def _score_topic_similarity(query: str, haystack: str) -> float:
    q = _tokens(query)
    h = _tokens(haystack)
    if not q or not h:
        return 0.0
    inter = len(q.intersection(h))
    return inter / max(1.0, math.sqrt(len(q) * len(h)))


def _classify_reference_category(slide: PptxSlideIndex, use_opts: dict, is_arabic_course: bool) -> str:
    text_blob = " ".join([
        slide.slide_title or "",
        slide.slide_text or "",
        slide.speaker_notes or "",
        slide.table_content or "",
        slide.diagram_labels or "",
    ]).lower()
    if is_arabic_course and use_opts.get("use_arabic_formatting_examples", True):
        if slide.language == "arabic":
            return "arabic_formatting"
    if use_opts.get("use_powerpoint_terminology", True):
        if "terminology" in text_blob or "definition" in text_blob or "glossary" in text_blob:
            return "terminology"
    if use_opts.get("use_powerpoint_layout_inspiration", True):
        if slide.visual_layout_metadata.get("image_count", 0) > 0 or slide.visual_layout_metadata.get("diagram_like_count", 0) > 0:
            return "visual"
    if use_opts.get("use_powerpoint_technical_support", True):
        if any(k in text_blob for k in ["kv", "ma", "fault", "alarm", "calibration", "procedure", "specification"]):
            return "technical"
    if use_opts.get("use_powerpoint_exercises_assessments", True):
        if any(k in text_blob for k in ["exercise", "quiz", "assessment", "knowledge check", "lab"]):
            return "instructional"
    return "instructional"


def detect_manual_conflicts(manual_index: dict, selected_refs: list[dict]) -> list[dict]:
    """Detect likely conflicts where technical values in PPT refs diverge from manual values."""
    manual_values: set[str] = set()
    manual_parts: set[str] = set()
    manual_faults: set[str] = set()

    for items in (manual_index.get("categories") or {}).values():
        for row in items:
            for v in row.get("technical_values") or []:
                manual_values.add(_norm(v).lower())
            for p in row.get("part_numbers") or []:
                manual_parts.add(_norm(p).lower())
            txt = _norm(row.get("text"))
            for fc in re.findall(r"\b[A-Z]{1,3}[- ]?\d{2,4}\b", txt):
                manual_faults.add(fc.lower())

    conflicts: list[dict] = []
    for ref in selected_refs:
        if ref.get("reference_category") != "technical":
            continue
        blob = _norm(ref.get("slide_text") or "")
        vals = {v.lower() for v in re.findall(r"\b\d+(?:\.\d+)?\s?(?:kv|ma|v|a|hz|mm|cm|kg|msv|usv|%)\b", blob, re.IGNORECASE)}
        parts = {p.lower() for p in re.findall(r"\b(?:PN|P/N)\s*[:#-]?\s*([A-Z0-9\-]{3,})\b", blob, re.IGNORECASE)}
        faults = {f.lower() for f in re.findall(r"\b[A-Z]{1,3}[- ]?\d{2,4}\b", blob)}

        mismatch_values = sorted(v for v in vals if manual_values and v not in manual_values)[:6]
        mismatch_parts = sorted(p for p in parts if manual_parts and p not in manual_parts)[:6]
        mismatch_faults = sorted(f for f in faults if manual_faults and f not in manual_faults)[:6]

        if mismatch_values or mismatch_parts or mismatch_faults:
            conflicts.append({
                "reference_file": ref.get("reference_file"),
                "reference_slide": ref.get("reference_slide"),
                "reference_category": ref.get("reference_category"),
                "manual_authority": "Uploaded manual values are authoritative",
                "ppt_values": {
                    "technical_values": mismatch_values,
                    "part_numbers": mismatch_parts,
                    "fault_codes": mismatch_faults,
                },
                "review_action": "Use manual values and flag for instructor review",
            })
    return conflicts


def retrieve_pptx_references(
    db: Session,
    *,
    equipment_name: str,
    equipment_model: str,
    manufacturer: str,
    course_type: str,
    language: str,
    topics: list[str],
    strictness: str = "balanced",
    use_options: dict | None = None,
    max_per_topic: int = 5,
) -> dict:
    """Retrieve targeted slide references from indexed PPTX files."""
    opts = use_options or {}
    query = (
        db.query(PptxPresentationIndex)
        .filter(PptxPresentationIndex.obsolete.is_(False))
        .filter(PptxPresentationIndex.do_not_use.is_(False))
    )

    presentations = query.all()
    if opts.get("use_manufacturer_approved_only"):
        presentations = [p for p in presentations if bool(p.manufacturer_approved)]
    searched_files = len(presentations)
    if not presentations:
        return {
            "searched_files": 0,
            "candidate_slides_found": 0,
            "selected_references": [],
            "selected_by_topic": {},
            "message": "No suitable PowerPoint reference found",
        }

    model_l = _norm(equipment_model).lower()
    mfr_l = _norm(manufacturer).lower()
    lang_l = _norm(language).lower()

    allowed_presentation_ids: list[str] = []
    for p in presentations:
        p_model = _norm(p.equipment_model).lower()
        p_mfr = _norm(p.manufacturer).lower()
        if strictness == "strict":
            if model_l and p_model and model_l != p_model:
                continue
            if mfr_l and p_mfr and mfr_l != p_mfr:
                continue
        elif strictness == "balanced":
            if mfr_l and p_mfr and mfr_l != p_mfr:
                if model_l and p_model and model_l in p_model:
                    pass
                else:
                    continue
        allowed_presentation_ids.append(p.id)

    if not allowed_presentation_ids:
        return {
            "searched_files": searched_files,
            "candidate_slides_found": 0,
            "selected_references": [],
            "selected_by_topic": {},
            "message": "No suitable PowerPoint reference found",
        }

    slides = (
        db.query(PptxSlideIndex)
        .filter(PptxSlideIndex.presentation_id.in_(allowed_presentation_ids))
        .all()
    )

    is_arabic_course = lang_l.startswith("arab")
    selected_by_topic: dict[str, list[dict]] = {}
    selected_all: list[dict] = []

    topic_list = [t for t in topics if _norm(t)]
    if not topic_list:
        topic_list = [course_type or "general"]

    for topic in topic_list:
        scored: list[tuple[float, dict]] = []
        for s in slides:
            if not _norm(s.slide_text) and not _norm(s.speaker_notes):
                continue
            blob = " ".join([
                _norm(s.slide_title), _norm(s.slide_text), _norm(s.speaker_notes),
                _norm(s.table_content), _norm(s.diagram_labels), _norm(s.image_captions),
            ])
            sim = _score_topic_similarity(topic, blob)
            if sim <= 0.05:
                continue

            p = next((x for x in presentations if x.id == s.presentation_id), None)
            if not p:
                continue

            eq_score = 1.0 if model_l and model_l == _norm(s.equipment_model).lower() else (0.6 if model_l and model_l in _norm(blob).lower() else 0.0)
            mfr_score = 1.0 if mfr_l and mfr_l == _norm(s.manufacturer).lower() else 0.0
            course_score = _score_topic_similarity(course_type, _norm(p.course_type))
            lang_score = 1.0 if lang_l and lang_l.startswith(_norm(s.language).lower()[:3]) else 0.2
            approval_score = 1.0 if p.manufacturer_approved else (0.8 if p.trusted else (0.7 if p.internal_training_reference else 0.3))
            uploaded_at_utc = _as_aware_utc(p.uploaded_at)
            recency_days = max(1.0, (_now() - uploaded_at_utc).days + 1)
            recency_score = 1.0 / math.log10(10 + recency_days)
            source_quality = max(0.1, min(1.0, float(s.quality_score or 0.0)))

            score = (
                0.23 * eq_score
                + 0.20 * mfr_score
                + 0.22 * sim
                + 0.08 * course_score
                + 0.06 * lang_score
                + 0.12 * approval_score
                + 0.05 * recency_score
                + 0.04 * source_quality
            )

            ref_category = _classify_reference_category(s, opts, is_arabic_course)
            scored.append((score, {
                "reference_file": s.filename,
                "reference_slide": s.slide_number,
                "reference_title": s.slide_title,
                "presentation_id": s.presentation_id,
                "presentation_total_slides": int(getattr(p, "slide_count", getattr(p, "total_slides", 0)) or 0),
                "reference_category": ref_category,
                "how_used": "retrieved_for_topic",
                "relevance_score": round(score, 4),
                "slide_text": s.slide_text,
                "speaker_notes": s.speaker_notes,
                "table_content": s.table_content,
                "diagram_labels": s.diagram_labels,
                "layout_metadata": s.visual_layout_metadata or {},
                "manufacturer": s.manufacturer,
                "equipment_model": s.equipment_model,
                "language": s.language,
                "source_status": s.source_status,
                "is_manufacturer_approved": bool(p.manufacturer_approved),
                "is_trusted": bool(p.trusted),
            }))

        scored.sort(key=lambda x: x[0], reverse=True)
        picks: list[dict] = []
        seen = set()
        for _, row in scored:
            dedupe_key = f"{row['reference_file']}::{row['reference_slide']}::{_norm(row['reference_title']).lower()}"
            if dedupe_key in seen:
                continue
            seen.add(dedupe_key)
            picks.append(row)
            if len(picks) >= max(3, min(8, max_per_topic)):
                break

        selected_by_topic[topic] = picks
        selected_all.extend(picks)

    # Cross-topic dedupe while preserving score preference
    deduped = {}
    for r in sorted(selected_all, key=lambda x: x.get("relevance_score", 0), reverse=True):
        key = f"{r['reference_file']}::{r['reference_slide']}"
        if key not in deduped:
            deduped[key] = r

    final_refs = list(deduped.values())
    final_refs = final_refs[:64]

    benchmark_max_slides = max([
        int(getattr(p, "slide_count", getattr(p, "total_slides", 0)) or 0)
        for p in presentations
    ] or [0])

    return {
        "searched_files": searched_files,
        "candidate_slides_found": len(slides),
        "selected_references": final_refs,
        "selected_by_topic": selected_by_topic,
        "manufacturer_approved_used": sum(1 for r in final_refs if r.get("is_manufacturer_approved")),
        "trusted_used": sum(1 for r in final_refs if r.get("is_trusted")),
        "arabic_formatting_refs": sum(1 for r in final_refs if r.get("reference_category") == "arabic_formatting"),
        "benchmark_reference_max_slides": benchmark_max_slides,
        "message": "ok" if final_refs else "No suitable PowerPoint reference found",
    }
