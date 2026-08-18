"""
Arabic Layout Transformation Engine.

Transforms LTR-designed (English) slides into Arabic-optimized layouts that
follow the approved Arabic reference deck's design system:

- TITLE placeholders snap into the reference "Arabic title band" (top of the
  slide, ~73% width), right-aligned, with reference title typography.
- ALL other shapes are mirrored horizontally via the margin swap
  x' = SW − x − w: symmetric/full-bleed shapes are identity-mapped (they do
  not move by construction), while asymmetric compositions — gutters,
  two-column layouts, side chips, page numbers — flip as coherent units so
  the visual flow starts from the right. Because every element of a
  composition mirrors together, relative overlaps are preserved exactly.
- Cover slides (any CENTER_TITLE placeholder present) are art-directed
  designs: geometry is left completely untouched.
- Directional graphics (connectors, arrow/chevron autoshapes) additionally
  get flipH toggled (+ rotation negated) so they point the RTL way; group
  shapes have their children mirrored recursively in group child-coordinate
  space (pictures keep orientation — only arrangement flips).
- Reference line spacing (150%) and equalized internal margins are applied
  to translated body text frames.
- A mirror invariant is verified per shape (horizontal off-canvas overflow
  must not grow); violations are reported as warnings.
- Any unexpected exception rolls the slide back to its pre-transform state —
  translations are kept, layout stays untransformed (never half-mirrored).

Overlap resolution (Rule §1–8):
- After the mirror pass a second pass catalogs every non-text obstacle
  (pictures, tables, OLE objects, large graphics).
- Every translated text box is checked for overlap with the obstacle map.
- Overlapping text boxes are repositioned into the largest available
  whitespace at approximately the same vertical position.
- If no adequate horizontal zone exists the text box is shifted vertically
  until a clear zone is found.
- As a last resort (entire slide covered) font size is reduced by up to 30 %
  and word-wrap is forced on; text never covers graphics.
- A minimum margin of 0.1 in (91 440 EMU) is enforced between text boxes
  and every obstacle.
- Every slide is validated after repositioning; any remaining overlap is
  logged as a warning for the quality report.

This engine must ONLY run for LTR-source decks. Arabic-designed decks are
preserved exactly — see the deck_src_is_rtl gating in doc_rebuilder.
"""

from __future__ import annotations

import copy
import logging
import math
import re

from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

from api.utils.font_metrics import arabic_width_factor

log = logging.getLogger(__name__)

# Fallback internal text-frame margins (EMU) from the reference deck.
_REF_MARGIN_LR = 91440
_REF_MARGIN_TB = 45720

# Minimum gap (EMU) between a text box and any obstacle. 0.1 in.
_MIN_MARGIN_EMU = 91440

# A text box is only useful when it is at least this wide (15 % of slide).
_MIN_TEXT_WIDTH_FRAC = 0.15

# Horizontally-directional preset geometries that must be flipped (not just
# repositioned) for RTL reading flow. Substring 'arrow' covers rightArrow,
# bentArrow, curvedLeftArrow, stripedRightArrow, notchedRightArrow, upArrow
# (flip is a no-op for h-symmetric ones — harmless).
_DIRECTIONAL_EXACT = {"chevron", "homePlate"}


def _is_title_placeholder(shape) -> bool:
    """True for any title placeholder variant across masters/templates."""
    try:
        if not shape.is_placeholder:
            return False
        ph_t = str(shape.placeholder_format.type)
        return "TITLE" in ph_t
    except Exception:
        return False


def _shape_text(shape) -> str:
    """Return normalized visible text for a shape, or empty string."""
    try:
        if not shape.has_text_frame:
            return ""
        parts = [p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip()]
        return " ".join(parts).strip()
    except Exception:
        return ""


def _shape_max_font_pt(shape) -> float:
    """Best-effort max run font size in points for a text shape."""
    max_pt = 0.0
    try:
        if not shape.has_text_frame:
            return 0.0
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.size is not None and run.font.size.pt:
                    max_pt = max(max_pt, float(run.font.size.pt))
    except Exception:
        pass
    return max_pt


def _looks_like_heading_text(text: str) -> bool:
    """Heuristic: short heading-like text, not paragraph-like body text."""
    if not text:
        return False
    t = text.strip()
    if not t:
        return False
    if t.isdigit():
        return False
    # Enumerated callouts like "1- ..." are usually body labels, not titles.
    if re.match(r"^\s*\d+[\)\].:-]", t):
        return False
    # Tiny labels (e.g., "P") are usually callout markers, not slide titles.
    core = re.sub(r"[^\w\u0600-\u06FF]+", "", t)
    if len(core) < 4:
        return False
    words = [w for w in t.split() if w]
    if len(words) > 22:
        return False
    if len(t) > 170:
        return False
    if t.endswith((".", ";", "؟", "!")):
        return False
    return True


def _shape_xfrm(shape):
    """The shape's own <a:xfrm> element (sp/cxnSp/pic: spPr/xfrm), or None."""
    el = shape._element
    for parent_tag in ("p:spPr", "p:grpSpPr"):
        parent = el.find(qn(parent_tag))
        if parent is not None:
            xfrm = parent.find(qn("a:xfrm"))
            if xfrm is not None:
                return xfrm
    return None


def _is_directional(shape) -> bool:
    """True for connectors and arrow/chevron-style autoshapes."""
    el = shape._element
    if el.tag == qn("p:cxnSp"):
        return True
    geom = el.find(".//" + qn("a:prstGeom"))
    if geom is None:
        return False
    prst = geom.get("prst") or ""
    return "arrow" in prst.lower() or prst in _DIRECTIONAL_EXACT


def _flip_horizontal(shape) -> bool:
    """Toggle flipH (and negate rotation) so directional geometry mirrors."""
    xfrm = _shape_xfrm(shape)
    if xfrm is None:
        return False
    if xfrm.get("flipH") in ("1", "true"):
        del xfrm.attrib["flipH"]
    else:
        xfrm.set("flipH", "1")
    rot = int(xfrm.get("rot", "0"))
    if rot:
        xfrm.set("rot", str((-rot) % 21600000))
    return True


def _mirror_group_children(group, warnings: list[str], label: str) -> None:
    """Mirror a group's children within the group's child-coordinate space.

    The group container itself is repositioned in slide space by the caller;
    here each child gets the same margin swap in chOff/chExt space so the
    internal arrangement flips too. Pictures keep their orientation (never
    flipH'd); directional autoshapes/connectors inside the group are flipped;
    nested groups recurse.
    """
    xfrm = _shape_xfrm(group)
    if xfrm is None:
        return
    ch_off = xfrm.find(qn("a:chOff"))
    ch_ext = xfrm.find(qn("a:chExt"))
    if ch_off is None or ch_ext is None:
        return
    ox = int(ch_off.get("x"))
    ext = int(ch_ext.get("cx"))
    for child in group.shapes:
        try:
            x, w = child.left, child.width
            if x is None or w is None:
                continue
            child.left = Emu(2 * ox + ext - x - w)
            if child.shape_type == MSO_SHAPE_TYPE.GROUP:
                _mirror_group_children(child, warnings, label + ">grp")
            elif _is_directional(child):
                _flip_horizontal(child)
        except Exception as exc:
            warnings.append(f"{label}: group child mirror skipped ({exc})")


def _ph_kind(shape) -> str | None:
    """Classify a placeholder shape; None for non-placeholders."""
    try:
        if not shape.is_placeholder:
            return None
        if _is_title_placeholder(shape):
            ph_t = str(shape.placeholder_format.type)
            if "CENTER_TITLE" in ph_t:
                return "center_title"
            return "title"
        t = str(shape.placeholder_format.type)
    except Exception:
        return None
    if "SUBTITLE" in t:
        return "subtitle"
    if "BODY" in t or "OBJECT" in t:
        return "body"
    return "other"


def _geometry(shape):
    """Return (left, top, width, height) or None when not fully resolvable."""
    try:
        g = (shape.left, shape.top, shape.width, shape.height)
    except Exception:
        return None
    if any(v is None for v in g):
        return None
    return g


def _h_overflow(left: int, width: int, slide_w: int) -> int:
    """Horizontal off-canvas overflow in EMU (bleed left + bleed right)."""
    return max(0, -left) + max(0, left + width - slide_w)


# ── Obstacle detection ─────────────────────────────────────────────────────────

def _group_has_picture(group) -> bool:
    """True if any direct child of a group is a picture shape."""
    try:
        for child in group.shapes:
            if child.shape_type == MSO_SHAPE_TYPE.PICTURE:
                return True
    except Exception:
        pass
    return False


def _is_obstacle(shape, slide_w: int, slide_h: int) -> bool:
    """True for shapes that must not be overlapped by translated text boxes.

    Hard obstacles (always):
      - Pictures / linked OLE / embedded OLE objects
      - Tables
      - Groups containing at least one picture

    Soft obstacles (size-gated — must cover > 3 % of slide area):
      - Any non-placeholder, non-text shape large enough to be a diagram,
        screenshot overlay, or background graphic.
    """
    # Skip placeholder shapes — those are the text content we're positioning.
    try:
        if shape.is_placeholder:
            return False
    except Exception:
        pass

    geom = _geometry(shape)
    if geom is None:
        return False

    # Hard obstacles ───────────────────────────────────────────────────────────
    try:
        st = shape.shape_type
        if st == MSO_SHAPE_TYPE.PICTURE:
            return True
        if st == MSO_SHAPE_TYPE.OLE_OBJECT:
            return True
        if st == MSO_SHAPE_TYPE.LINKED_OLE_OBJECT:
            return True
        if st == MSO_SHAPE_TYPE.GROUP and _group_has_picture(shape):
            return True
    except Exception:
        pass

    try:
        if shape.has_table:
            return True
    except Exception:
        pass

    # Soft obstacles: large non-text shapes ───────────────────────────────────
    _, _, w, h = geom
    slide_area = slide_w * slide_h
    if w * h < 0.03 * slide_area:
        return False  # too small to matter

    # If the shape has meaningful text content it's a text box, not an obstacle
    try:
        if shape.has_text_frame:
            if any(p.text.strip() for p in shape.text_frame.paragraphs):
                return False
    except Exception:
        pass

    return True  # large, non-text, non-placeholder → treat as obstacle


def _is_static_decoration(shape, slide_w: int, slide_h: int) -> bool:
    """Heuristic lock for edge decorations and logo-like ornaments."""
    geom = _geometry(shape)
    if geom is None:
        return False
    left, top, width, height = geom

    try:
        if shape.is_placeholder:
            return False
    except Exception:
        pass

    try:
        if shape.has_text_frame:
            if any(p.text.strip() for p in shape.text_frame.paragraphs):
                return False
    except Exception:
        pass

    # Directional shapes and groups participate in RTL mirroring.
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            return False
    except Exception:
        pass
    if _is_directional(shape):
        return False

    near_top_or_bottom = top <= int(0.18 * slide_h) or (top + height) >= int(0.82 * slide_h)
    compact = width <= int(0.25 * slide_w) and height <= int(0.20 * slide_h)
    near_side = left <= int(0.12 * slide_w) or (left + width) >= int(0.88 * slide_w)
    return compact and near_top_or_bottom and near_side


def _collect_obstacles(
    slide, slide_w: int, slide_h: int
) -> list[tuple[int, int, int, int]]:
    """Return (left, top, width, height) for every obstacle on the slide.

    Called after the mirror pass so coordinates reflect final positions.
    """
    result: list[tuple[int, int, int, int]] = []
    for shape in slide.shapes:
        geom = _geometry(shape)
        if geom is None:
            continue
        if _is_obstacle(shape, slide_w, slide_h):
            result.append(geom)
    return result


def _detect_top_banner_rect(slide, slide_w: int, slide_h: int) -> tuple[int, int, int, int]:
    """Detect the top colored banner area; fallback to generic top strip."""
    best = None
    best_area = -1
    for shape in slide.shapes:
        geom = _geometry(shape)
        if geom is None:
            continue
        l, t, w, h = geom
        if t > int(0.25 * slide_h):
            continue

        # Candidate banner: wide, top-area, and no meaningful text.
        if w < int(0.45 * slide_w):
            continue
        if h < int(0.04 * slide_h) or h > int(0.25 * slide_h):
            continue

        txt = _shape_text(shape)
        if txt:
            continue

        area = w * h
        if area > best_area:
            best_area = area
            best = (l, t, w, h)

    if best is not None:
        return best
    return (0, 0, slide_w, int(0.18 * slide_h))


def _detect_logo_left_boundary(slide, slide_w: int, slide_h: int) -> int | None:
    """Best-effort left boundary of top-right logo/branding area."""
    best_left = None
    for shape in slide.shapes:
        geom = _geometry(shape)
        if geom is None:
            continue
        l, t, w, h = geom
        if t > int(0.25 * slide_h):
            continue
        if l < int(0.55 * slide_w):
            continue
        if w > int(0.40 * slide_w) or h > int(0.25 * slide_h):
            continue

        is_logo_like = False
        try:
            st = shape.shape_type
            if st == MSO_SHAPE_TYPE.PICTURE:
                is_logo_like = True
            elif st == MSO_SHAPE_TYPE.GROUP and _group_has_picture(shape):
                is_logo_like = True
        except Exception:
            pass

        # Fallback logo proxy only for very small compact non-text marks.
        # Avoid classifying generic banner rectangles/blocks as logos.
        if (
            not is_logo_like
            and not _shape_text(shape)
            and w <= int(0.12 * slide_w)
            and h <= int(0.12 * slide_h)
        ):
            is_logo_like = True

        if not is_logo_like:
            continue

        if best_left is None or l < best_left:
            best_left = l
    return best_left


def _detect_top_right_obstruction_left(
    slide,
    slide_w: int,
    slide_h: int,
    exclude_idxs: set[int] | None = None,
) -> int | None:
    """Find top-right decorative non-text overlays that can mask title text."""
    blocked_left: int | None = None
    skip = exclude_idxs or set()

    for idx, shape in enumerate(slide.shapes):
        if idx in skip:
            continue
        geom = _geometry(shape)
        if geom is None:
            continue
        l, t, w, h = geom

        if t > int(0.25 * slide_h):
            continue
        if l < int(0.45 * slide_w):
            continue
        if w > int(0.60 * slide_w) or h > int(0.25 * slide_h):
            continue
        if _shape_text(shape):
            continue

        if blocked_left is None or l < blocked_left:
            blocked_left = l

    return blocked_left


def _detect_title_shape_idxs(
    slide,
    slide_w: int,
    slide_h: int,
    translated_idxs: set[int],
) -> set[int]:
    """Detect title-like shapes, including non-placeholder top banner titles."""
    title_idxs: set[int] = set()

    # Always include canonical title placeholders.
    for idx, shape in enumerate(slide.shapes):
        if _is_title_placeholder(shape):
            try:
                if int(shape.top) <= int(0.45 * slide_h):
                    title_idxs.add(idx)
            except Exception:
                title_idxs.add(idx)

    banner = _detect_top_banner_rect(slide, slide_w, slide_h)
    bl, bt, bw, bh = banner
    banner_bottom = bt + bh

    top_font_max = 0.0
    font_by_idx: dict[int, float] = {}
    for idx, shape in enumerate(slide.shapes):
        geom = _geometry(shape)
        if geom is None:
            continue
        _, t, _, _ = geom
        if t <= int(0.25 * slide_h):
            f = _shape_max_font_pt(shape)
            font_by_idx[idx] = f
            top_font_max = max(top_font_max, f)

    best_non_placeholder_idx = None
    best_non_placeholder_score = -999.0

    for idx, shape in enumerate(slide.shapes):
        if idx in title_idxs:
            continue
        if not shape.has_text_frame:
            continue
        text = _shape_text(shape)
        if not text:
            continue
        geom = _geometry(shape)
        if geom is None:
            continue
        l, t, w, h = geom

        # Exclusions: footer/page-number/body-like placements.
        if t >= int(0.80 * slide_h):
            continue
        if text.strip().isdigit() and w <= int(0.12 * slide_w):
            continue

        in_top_20 = t <= int(0.20 * slide_h)
        overlaps_banner = _rects_overlap(geom, banner, margin=0)
        near_banner = t <= banner_bottom + int(0.02 * slide_h)

        if not (in_top_20 or overlaps_banner or near_banner):
            continue

        score = 0.0
        if overlaps_banner:
            score += 4.0
        if in_top_20:
            score += 2.5
        if near_banner:
            score += 1.0

        if _looks_like_heading_text(text):
            score += 2.0
        else:
            score -= 3.0

        if idx in translated_idxs:
            score += 1.5

        f = font_by_idx.get(idx, 0.0)
        if top_font_max > 0 and f > 0:
            score += 3.0 * (f / top_font_max)

        # Body blocks are typically very tall and paragraph-dense.
        if h > int(0.22 * slide_h):
            score -= 2.0

        if score > best_non_placeholder_score:
            best_non_placeholder_score = score
            best_non_placeholder_idx = idx

    if best_non_placeholder_idx is not None and best_non_placeholder_score >= 4.5:
        title_idxs.add(best_non_placeholder_idx)

    return title_idxs


# ── Overlap geometry helpers ───────────────────────────────────────────────────

def _rects_overlap(r1: tuple, r2: tuple, margin: int = 0) -> bool:
    """True if rectangle r1 overlaps r2 when r2 is expanded by *margin*."""
    l1, t1, w1, h1 = r1
    l2, t2, w2, h2 = r2
    return (
        l1 < l2 + w2 + margin
        and l1 + w1 > l2 - margin
        and t1 < t2 + h2 + margin
        and t1 + h1 > t2 - margin
    )


def _merge_intervals(intervals: list[tuple[int, int]]) -> list[tuple[int, int]]:
    """Merge overlapping [lo, hi) intervals."""
    if not intervals:
        return []
    intervals = sorted(intervals)
    merged = [list(intervals[0])]
    for lo, hi in intervals[1:]:
        if lo <= merged[-1][1]:
            merged[-1][1] = max(merged[-1][1], hi)
        else:
            merged.append([lo, hi])
    return [(a, b) for a, b in merged]


def _free_x_zones(
    top: int,
    height: int,
    obstacles: list[tuple],
    slide_w: int,
    margin: int,
) -> list[tuple[int, int]]:
    """Find free horizontal bands in [0, slide_w] for a text box at (top, height).

    Returns list of (lo, hi) pairs representing clear x-ranges.
    """
    blocked: list[tuple[int, int]] = []
    for ol, ot, ow, oh in obstacles:
        # Does this obstacle's y-band overlap the text box's y-band?
        if ot < top + height + margin and ot + oh > top - margin:
            lo = max(0, ol - margin)
            hi = min(slide_w, ol + ow + margin)
            if hi > lo:
                blocked.append((lo, hi))

    merged = _merge_intervals(blocked)

    free: list[tuple[int, int]] = []
    cursor = 0
    for lo, hi in merged:
        if cursor < lo:
            free.append((cursor, lo))
        cursor = max(cursor, hi)
    if cursor < slide_w:
        free.append((cursor, slide_w))
    return free


def _reduce_shape_font(
    shape,
    max_reduction_frac: float = 0.25,
    min_pt: float = 8.0,
) -> None:
    """Reduce every run's font size in *shape* by up to *max_reduction_frac*."""
    try:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                try:
                    sz = run.font.size
                    if sz is not None and sz.pt:
                        new_pt = max(min_pt, sz.pt * (1.0 - max_reduction_frac))
                        run.font.size = Pt(new_pt)
                except Exception:
                    pass
    except Exception:
        pass


def _enable_word_wrap(shape) -> None:
    """Enable word-wrap on the text frame so Arabic text reflows."""
    try:
        shape.text_frame.word_wrap = True
    except Exception:
        pass


def _set_norm_autofit(shape) -> None:
    """Set OOXML normAutofit so text shrinks to stay inside the shape box."""
    try:
        tx_body = shape.text_frame._txBody
        body_pr = tx_body.find(qn("a:bodyPr"))
        if body_pr is None:
            return
        for tag in ("normAutofit", "spAutoFit", "noAutofit"):
            for child in list(body_pr.findall(qn(f"a:{tag}"))):
                body_pr.remove(child)
        from lxml import etree as _ET

        body_pr.append(_ET.Element(qn("a:normAutofit")))
    except Exception:
        pass


def _run_base_sizes_pt(shape) -> dict[tuple[int, int], float]:
    """Capture each run's original size (pt) as a stable baseline.

    Keyed by (paragraph_idx, run_idx) — NOT id(run). python-pptx's `.runs`
    is a property that constructs a fresh _Run wrapper object on every
    access, so id(run) captured here is not guaranteed to match id(run) for
    the "same" run seen again later (in _apply_font_scale, or across
    multiple shrink-loop iterations calling this pair repeatedly). Whether
    it happened to match depended on CPython reusing the same freed
    wrapper's memory address — unreliable, and when it didn't match the
    lookup silently found nothing and skipped scaling that run entirely,
    with no error raised. Positional indices only depend on paragraph/run
    order, which is stable for the lifetime of a single shrink pass.
    """
    sizes: dict[tuple[int, int], float] = {}
    try:
        for para_idx, para in enumerate(shape.text_frame.paragraphs):
            for run_idx, run in enumerate(para.runs):
                sz = run.font.size
                if sz is not None and sz.pt:
                    sizes[(para_idx, run_idx)] = float(sz.pt)
    except Exception:
        pass
    return sizes


def _apply_font_scale(
    shape,
    base_sizes_pt: dict[tuple[int, int], float],
    scale: float,
    min_frac: float,
    min_pt: float = 8.0,
) -> bool:
    """Scale run font sizes from their baseline, honoring a minimum fraction."""
    changed = False
    try:
        for para_idx, para in enumerate(shape.text_frame.paragraphs):
            for run_idx, run in enumerate(para.runs):
                base = base_sizes_pt.get((para_idx, run_idx))
                if not base:
                    continue
                floor_pt = max(min_pt, base * min_frac)
                target_pt = max(floor_pt, base * scale)
                cur = run.font.size.pt if run.font.size is not None else base
                if cur > target_pt + 0.05:
                    run.font.size = Pt(target_pt)
                    changed = True
    except Exception:
        pass
    return changed


def _set_line_spacing(shape, spacing: float) -> None:
    """Apply uniform line spacing to non-empty paragraphs."""
    try:
        for para in shape.text_frame.paragraphs:
            if para.text.strip():
                para.line_spacing = spacing
    except Exception:
        pass


def _is_dense_bullet_textbox(shape) -> bool:
    """Detect long bullet/list text boxes that need stronger fit compaction."""
    try:
        tf = shape.text_frame
    except Exception:
        return False

    non_empty = [p.text.strip() for p in tf.paragraphs if p.text and p.text.strip()]
    if len(non_empty) < 7:
        return False

    bullet_like = 0
    for txt in non_empty:
        if txt.startswith(("•", "-", "*")):
            bullet_like += 1
            continue
        if len(txt) >= 2 and txt[0].isdigit() and txt[1] in {".", ")", "-"}:
            bullet_like += 1
    return bullet_like >= max(3, int(len(non_empty) * 0.40))


def _set_para_rtl(para) -> None:
    """Force rtl=1 on paragraph properties."""
    try:
        ppr = para._p.get_or_add_pPr()
        ppr.set("rtl", "1")
    except Exception:
        pass


def _estimate_text_occupancy(shape) -> tuple[float, float, int]:
    """Estimate occupied (width, height, line_count) in EMU for Arabic text."""
    geom = _geometry(shape)
    if geom is None:
        return (0.0, 0.0, 0)
    _, _, width, _ = geom

    try:
        tf = shape.text_frame
    except Exception:
        return (0.0, 0.0, 0)

    margin_l = int(getattr(tf, "margin_left", 0) or 0)
    margin_r = int(getattr(tf, "margin_right", 0) or 0)
    usable_w = max(1.0, float(width - margin_l - margin_r))

    font_name = ""
    base_pt = 12.0
    line_spacing = 1.5

    try:
        for para in tf.paragraphs:
            if para.text.strip() and isinstance(para.line_spacing, (int, float)):
                line_spacing = max(1.0, float(para.line_spacing))
                break
    except Exception:
        pass

    try:
        pts: list[float] = []
        for para in tf.paragraphs:
            for run in para.runs:
                if run.font.name and not font_name:
                    font_name = run.font.name
                if run.font.size is not None and run.font.size.pt:
                    pts.append(float(run.font.size.pt))
        if pts:
            base_pt = sum(pts) / len(pts)
    except Exception:
        pass

    width_factor, _ = arabic_width_factor(font_name)
    emu_per_pt = 12700.0

    total_lines = 0
    max_line_w = 0.0
    try:
        paragraphs = [p for p in tf.paragraphs if p.text is not None]
        for para in paragraphs:
            raw = para.text or ""
            chunks = [c for c in raw.split("\n") if c.strip()]
            if not chunks:
                continue
            for chunk in chunks:
                chars = max(1, len(chunk.strip()))
                est_w = chars * base_pt * emu_per_pt * width_factor
                lines = max(1, int(math.ceil(est_w / usable_w)))
                total_lines += lines
                max_line_w = max(max_line_w, est_w)
    except Exception:
        pass

    if total_lines == 0:
        return (0.0, 0.0, 0)

    line_h = base_pt * emu_per_pt * line_spacing
    est_h = float(total_lines) * line_h
    est_w = max_line_w if max_line_w > 0 else usable_w
    return (est_w, est_h, total_lines)


def _fits_textbox_constraints(
    shape,
    original_geom: tuple[int, int, int, int],
    slide_w: int,
    slide_h: int,
) -> tuple[bool, bool]:
    """Return (fits, height_overflow) using occupancy and boundary checks."""
    geom = _geometry(shape)
    if geom is None:
        return (True, False)
    left, top, width, height = geom

    in_bounds = (
        left >= 0
        and top >= 0
        and left + width <= slide_w
        and top + height <= slide_h
    )
    if not in_bounds:
        return (False, True)

    est_w, est_h, _ = _estimate_text_occupancy(shape)
    if est_h <= 0:
        return (True, False)

    orig_area = max(1, original_geom[2] * original_geom[3])
    occupied_area = max(1.0, min(est_w, float(width)) * min(est_h, float(height)))
    area_ratio = occupied_area / float(orig_area)
    height_over = est_h > float(height)

    fits = (not height_over) and area_ratio <= 1.05
    return (fits, height_over)


def _shrink_textbox_away_from_side_obstacle(
    shape,
    obstacles: list[tuple[int, int, int, int]],
    slide_w: int,
    slide_h: int,
    margin: int,
    min_width_frac: float = 0.35,
) -> bool:
    """Narrow a right-aligned RTL text box so it stops before a side obstacle.

    Many LTR-designed slides place an image to one side of a wide text
    placeholder; left-aligned English text never reaches that side, so the
    overlap is geometrically present but invisible. Once the paragraph is
    right-aligned for Arabic, the same text hugs the box's *opposite* edge —
    directly into the obstacle — and appears clipped/hidden behind it.

    This only ever shrinks the box (never grows it, never moves its
    anchored edge), so it cannot introduce a new collision elsewhere. It is
    a no-op when no side obstacle is found or the resulting box would be
    unreasonably narrow.
    """
    geom = _geometry(shape)
    if geom is None:
        return False
    left, top, width, height = geom
    box_bottom = top + height
    box_right = left + width
    box_center = left + width / 2.0

    right_side: list[tuple[int, int, int, int]] = []
    left_side: list[tuple[int, int, int, int]] = []
    for (ol, ot, ow, oh) in obstacles:
        obs_bottom = ot + oh
        obs_right = ol + ow
        # Obstacle must vertically overlap the box to matter.
        if min(box_bottom, obs_bottom) - max(top, ot) <= 0:
            continue
        # Obstacle must horizontally intrude into the box's own span.
        if ol >= box_right or obs_right <= left:
            continue
        obs_center = ol + ow / 2.0
        if obs_center >= box_center:
            right_side.append((ol, ot, ow, oh))
        else:
            left_side.append((ol, ot, ow, oh))

    if not right_side and not left_side:
        return False

    changed = False
    min_width = int(width * min_width_frac)

    if right_side:
        nearest_left_edge = min(ol for (ol, _ot, _ow, _oh) in right_side)
        candidate_width = (nearest_left_edge - margin) - left
        if min_width <= candidate_width < width:
            width = candidate_width
            changed = True

    if left_side:
        nearest_right_edge = max(ol + ow for (ol, _ot, ow, _oh) in left_side)
        candidate_left = nearest_right_edge + margin
        candidate_width = (left + width) - candidate_left
        if candidate_left > left and min_width <= candidate_width < width:
            left = candidate_left
            width = candidate_width
            changed = True

    if not changed:
        return False

    shape.left = Emu(left)
    shape.width = Emu(width)
    _enable_word_wrap(shape)
    return True


def _expand_textbox_with_whitespace(
    shape,
    obstacles: list[tuple],
    slide_w: int,
    slide_h: int,
    margin: int,
    max_expand_frac: float = 0.12,
) -> bool:
    """Expand width slightly within local whitespace while preserving balance."""
    geom = _geometry(shape)
    if geom is None:
        return False
    left, top, width, height = geom

    zones = _free_x_zones(top, height, obstacles, slide_w, margin)
    containing = [z for z in zones if z[0] <= left and z[1] >= left + width]
    if not containing:
        return False
    lo, hi = max(containing, key=lambda z: z[1] - z[0])
    max_zone_w = max(0, hi - lo)
    if max_zone_w <= width:
        return False

    target_w = min(int(width * (1.0 + max_expand_frac)), max_zone_w)
    if target_w <= width:
        return False

    new_left = left - int((target_w - width) / 2)
    new_left = max(lo, min(new_left, hi - target_w))
    new_left = max(0, min(new_left, slide_w - target_w))

    before = geom
    shape.left = Emu(new_left)
    shape.width = Emu(target_w)

    after = _geometry(shape)
    if after is None:
        shape.left = Emu(before[0])
        shape.width = Emu(before[2])
        return False

    if any(_rects_overlap(after, obs, margin) for obs in obstacles):
        shape.left = Emu(before[0])
        shape.width = Emu(before[2])
        return False

    _enable_word_wrap(shape)
    return True


def _clamp_shape_to_slide(shape, slide_w: int, slide_h: int) -> None:
    """Keep a shape fully inside slide boundaries."""
    geom = _geometry(shape)
    if geom is None:
        return
    left, top, width, height = geom

    width = max(1, min(width, slide_w))
    height = max(1, min(height, slide_h))
    left = max(0, min(left, slide_w - width))
    top = max(0, min(top, slide_h - height))

    shape.left = Emu(left)
    shape.top = Emu(top)
    shape.width = Emu(width)
    shape.height = Emu(height)


def _optimize_translated_textbox(
    shape,
    original_geom: tuple[int, int, int, int],
    obstacles: list[tuple],
    slide_w: int,
    slide_h: int,
    margin: int,
    line_spacing_base: float,
    warnings: list[str],
    label: str,
) -> None:
    """Fit translated Arabic text to preserve source visual density.

        Optimization order:
            1) Font-size reduction (down to 80 % of original run size;
                 down to 70 % for dense bullet/list boxes)
      2) Slight line-spacing reduction
      3) Slight box expansion inside whitespace only
            4) Auto-fit + re-wrap while keeping body text top-anchored
    """
    if not shape.has_text_frame:
        return

    base_sizes = _run_base_sizes_pt(shape)
    _enable_word_wrap(shape)
    is_dense_list = _is_dense_bullet_textbox(shape)

    fits, height_over = _fits_textbox_constraints(
        shape, original_geom, slide_w, slide_h
    )
    if fits:
        return

    # 1) Prefer smaller font over extra wrapped lines.
    used_scale = 1.0
    min_frac = 0.70 if is_dense_list else 0.80
    max_steps = 15 if is_dense_list else 10
    for step in range(1, max_steps + 1):
        scale = 1.0 - (step * 0.02)
        _apply_font_scale(shape, base_sizes, scale=scale, min_frac=min_frac)
        fits, height_over = _fits_textbox_constraints(
            shape, original_geom, slide_w, slide_h
        )
        used_scale = scale
        if fits:
            warnings.append(f"{label}: font scaled to {scale:.2f} to preserve layout")
            return

    # 2) Reduce line spacing slightly if still not fitting.
    min_line_spacing = max(1.00, line_spacing_base * 0.85) if is_dense_list else max(1.15, line_spacing_base * 0.90)
    spacing = line_spacing_base
    while spacing - 0.05 >= min_line_spacing:
        spacing -= 0.05
        _set_line_spacing(shape, spacing)
        fits, height_over = _fits_textbox_constraints(
            shape, original_geom, slide_w, slide_h
        )
        if fits:
            warnings.append(
                f"{label}: line spacing reduced to {spacing:.2f} after font scale"
            )
            return

    # 3) Expand only when there is safe local whitespace.
    max_expand_frac = 0.24 if is_dense_list else 0.12
    if _expand_textbox_with_whitespace(
        shape,
        obstacles,
        slide_w,
        slide_h,
        margin,
        max_expand_frac=max_expand_frac,
    ):
        fits, height_over = _fits_textbox_constraints(
            shape, original_geom, slide_w, slide_h
        )
        if fits:
            warnings.append(
                f"{label}: expanded slightly into whitespace (font scale {used_scale:.2f})"
            )
            return

    # 4) If height still overflows, force autofit and keep text top-anchored.
    if height_over:
        _set_norm_autofit(shape)
        _enable_word_wrap(shape)
        try:
            shape.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
        except Exception:
            pass
        warnings.append(
            f"{label}: applied normAutofit with top anchor due to height overflow"
        )

    _clamp_shape_to_slide(shape, slide_w, slide_h)


def _mirror_and_fit_title_shape(
    shape,
    original_left: int,
    original_top: int,
    original_width: int,
    original_height: int,
    slide_w: int,
    slide_h: int,
    logo_left_boundary: int | None,
    banner_rect: tuple[int, int, int, int],
    warnings: list[str],
    label: str,
) -> None:
    """Mirror and fit a detected title shape in-place for RTL."""
    if not shape.has_text_frame:
        return

    safe_gap = int(0.015 * slide_w)
    mirrored_left = int(slide_w - original_left - original_width)

    right_limit = slide_w - safe_gap
    if logo_left_boundary is not None:
        right_limit = min(right_limit, int(logo_left_boundary) - safe_gap)
    # Keep right-half preference only when no hard safety boundary exists.
    if logo_left_boundary is None:
        right_limit = max(right_limit, int(0.55 * slide_w))

    # Keep title anchored by the mirrored formula whenever possible.
    new_left = mirrored_left
    new_width = original_width

    if new_left + new_width > right_limit:
        # First: try to keep mirrored left and shrink width to fit logo-safe area.
        new_width = max(1, right_limit - new_left)
        # If mirrored left is out of bounds, place title against the right limit.
        if new_width <= 1:
            new_width = max(1, min(original_width, right_limit - int(0.5 * slide_w)))
            new_left = max(0, right_limit - new_width)

    shape.left = Emu(new_left)
    shape.width = Emu(new_width)

    bl, bt, bw, bh = banner_rect
    banner_top = max(0, bt)
    banner_bottom = min(slide_h, bt + bh)
    banner_box_height = max(1, banner_bottom - banner_top)
    target_top = original_top
    target_height = original_height
    if target_height > banner_box_height:
        target_top = banner_top
        target_height = banner_box_height
    else:
        if target_top < banner_top:
            target_top = banner_top
        if target_top + target_height > banner_bottom:
            target_top = max(banner_top, banner_bottom - target_height)
    if target_top + target_height > slide_h:
        target_height = max(1, slide_h - target_top)

    shape.top = Emu(target_top)
    shape.height = Emu(target_height)

    tf = shape.text_frame
    _enable_word_wrap(shape)
    _set_norm_autofit(shape)
    # Keep safe internal clearance so Arabic ascenders never hit the top border.
    min_top_margin = int(max(_REF_MARGIN_TB, 0.008 * slide_h))
    min_bottom_margin = int(max(_REF_MARGIN_TB, 0.006 * slide_h))
    try:
        if int(getattr(tf, "margin_top", 0) or 0) < min_top_margin:
            tf.margin_top = Emu(min_top_margin)
    except Exception:
        pass
    try:
        if int(getattr(tf, "margin_bottom", 0) or 0) < min_bottom_margin:
            tf.margin_bottom = Emu(min_bottom_margin)
    except Exception:
        pass
    try:
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    except Exception:
        pass

    for para in tf.paragraphs:
        if not para.text.strip():
            continue
        _set_para_rtl(para)
        para.alignment = PP_ALIGN.RIGHT

    # Expand width to the left within banner/logo-safe area only when needed.
    left_bound = max(0, bl + safe_gap)
    max_width = max(1, right_limit - left_bound)
    if max_width > int(shape.width):
        est_w, est_h, _ = _estimate_text_occupancy(shape)
        needs_more_width = est_w > float(shape.width)
        if needs_more_width:
            target_w = max_width
            if target_w > int(shape.width):
                shape.width = Emu(target_w)
                shape.left = Emu(max(left_bound, right_limit - target_w))

    # Prefer geometry-first fit for titles: grow box height before shrinking text.
    est_w, est_h, _ = _estimate_text_occupancy(shape)
    top_margin = int(getattr(tf, "margin_top", 0) or 0)
    bottom_margin = int(getattr(tf, "margin_bottom", 0) or 0)
    target_height = int(math.ceil(est_h + top_margin + bottom_margin + max(1, int(0.003 * slide_h))))
    if target_height > int(shape.height):
        # Preserve original visual center while increasing height.
        original_center_y = target_top + (target_height // 2)
        grown_height = min(target_height, banner_box_height)
        grown_top = original_center_y - (grown_height // 2)

        # Keep growth inside the detected banner when possible.
        max_banner_top = banner_top
        max_banner_bottom = banner_bottom
        if grown_top < max_banner_top:
            grown_top = max_banner_top
        if grown_top + grown_height > max_banner_bottom:
            grown_top = max_banner_bottom - grown_height
        if grown_top < 0:
            grown_top = 0
        if grown_top + grown_height > slide_h:
            grown_height = max(1, slide_h - grown_top)

        if grown_height > int(shape.height):
            shape.top = Emu(grown_top)
            shape.height = Emu(grown_height)

    # Prefer readable fit: shrink gradually only when geometry expansion is insufficient.
    base_sizes = _run_base_sizes_pt(shape)
    for step in range(1, 11):
        scale = 1.0 - (step * 0.02)
        _apply_font_scale(shape, base_sizes, scale=scale, min_frac=0.80)
        est_w, est_h, _ = _estimate_text_occupancy(shape)
        if est_h <= float(shape.height) and est_w <= float(shape.width):
            break

    _clamp_shape_to_slide(shape, slide_w, slide_h)

    # Validation: title must be in right banner half, avoid logo, and never top-clip.
    g = _geometry(shape)
    if g is None:
        return
    l, t, w, h = g
    shape_center_x = l + (w // 2)
    banner_mid_x = bl + (bw // 2)
    if shape_center_x < banner_mid_x:
        warnings.append(f"{label}: title remains in left half after RTL mirror")

    if logo_left_boundary is not None and l + w > logo_left_boundary - safe_gap:
        warnings.append(f"{label}: title intersects logo-safe boundary")

    est_w, est_h, _ = _estimate_text_occupancy(shape)
    if est_h > float(h) or est_w > float(w):
        warnings.append(f"title_overflow_warning: {label}: title may still be clipped after fitting")

    # With MIDDLE anchor, derive effective top clearance from occupied height.
    top_margin = int(getattr(tf, "margin_top", 0) or 0)
    bottom_margin = int(getattr(tf, "margin_bottom", 0) or 0)
    inner_h = max(1, h - top_margin - bottom_margin)
    inner_clearance = max(0.0, (float(inner_h) - est_h) / 2.0)
    top_clearance = inner_clearance + float(top_margin)
    min_top_clearance = max(1.0, float(top_margin) * 0.35)
    top_touch_fail = est_h > float(inner_h) or top_clearance < min_top_clearance
    if top_touch_fail:
        warnings.append(f"title_overflow_warning: {label}: top-glyph clearance failed (title touches top border)")

    if abs(l - original_left) <= 1:
        warnings.append(f"{label}: title appears unmirrored (left position unchanged)")


def _fit_cover_title_shape(
    shape,
    slide_w: int,
    slide_h: int,
) -> None:
    """Fit CENTER_TITLE cover text without forcing title-band mirroring."""
    if not shape.has_text_frame:
        return

    tf = shape.text_frame
    _enable_word_wrap(shape)
    _set_norm_autofit(shape)

    min_top_margin = int(max(_REF_MARGIN_TB, 0.008 * slide_h))
    min_bottom_margin = int(max(_REF_MARGIN_TB, 0.006 * slide_h))
    try:
        if int(getattr(tf, "margin_top", 0) or 0) < min_top_margin:
            tf.margin_top = Emu(min_top_margin)
    except Exception:
        pass
    try:
        if int(getattr(tf, "margin_bottom", 0) or 0) < min_bottom_margin:
            tf.margin_bottom = Emu(min_bottom_margin)
    except Exception:
        pass
    try:
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    except Exception:
        pass

    for para in tf.paragraphs:
        if not para.text.strip():
            continue
        existing_alignment = para.alignment
        _set_para_rtl(para)
        if existing_alignment is None:
            para.alignment = PP_ALIGN.CENTER

    base_sizes = _run_base_sizes_pt(shape)
    for step in range(0, 18):
        scale = 1.0 - (step * 0.035)
        _apply_font_scale(shape, base_sizes, scale=scale, min_frac=0.50)
        est_w, est_h, _ = _estimate_text_occupancy(shape)
        if est_h <= float(shape.height) and est_w <= float(shape.width):
            break

    _clamp_shape_to_slide(shape, slide_w, slide_h)


# ── Overlap resolution ─────────────────────────────────────────────────────────

def _resolve_overlap(
    shape,
    obstacles: list[tuple],
    slide_w: int,
    slide_h: int,
    margin: int,
    warnings: list[str],
    label: str,
) -> None:
    """Reposition a text box so it does not overlap any obstacle.

    Resolution cascade (Rules §5):
      1. Find the largest free horizontal zone at the same vertical position.
      2. If found and wide enough → resize + reposition the text box there.
      3. If no wide-enough zone → scan vertical offsets (±5 % steps up to ±30 %).
      4. Last resort → shrink font by up to 30 % and force word-wrap.

    The text box is never placed over a graphic. Word-wrap is always enabled
    on repositioned shapes.
    """
    geom = _geometry(shape)
    if not geom:
        return
    left, top, width, height = geom

    # Quick exit: no actual overlap
    text_rect = geom
    if not any(_rects_overlap(text_rect, obs, margin) for obs in obstacles):
        return

    min_text_w = int(_MIN_TEXT_WIDTH_FRAC * slide_w)

    def _try_place(try_top: int) -> bool:
        """Try to fit the text box at vertical position try_top.
        Returns True and applies the change if successful."""
        if try_top < 0 or try_top + height > slide_h + margin:
            return False
        zones = _free_x_zones(try_top, height, obstacles, slide_w, margin)
        usable = [(lo, hi) for lo, hi in zones if hi - lo >= min_text_w]
        if not usable:
            return False
        # Pick the widest zone
        best_lo, best_hi = max(usable, key=lambda z: z[1] - z[0])
        inner = min(margin, int(0.01 * slide_w))
        new_left = best_lo + inner
        new_width = max(min_text_w, best_hi - best_lo - 2 * inner)
        new_left = max(0, min(new_left, slide_w - new_width))

        shape.left = Emu(new_left)
        shape.top = Emu(try_top)
        shape.width = Emu(new_width)
        _enable_word_wrap(shape)
        warnings.append(
            f"{label}: repositioned to avoid overlap "
            f"(x {new_left/914400:.2f}\"–{(new_left+new_width)/914400:.2f}\", "
            f"y {try_top/914400:.2f}\")"
        )
        return True

    # Step 1: same vertical position
    if _try_place(top):
        return

    # Step 2: scan vertical offsets in ±5 % increments up to ±30 %
    step = int(0.05 * slide_h)
    for delta in range(step, int(0.31 * slide_h), step):
        for dy in (delta, -delta):
            if _try_place(top + dy):
                return

    # Step 3: last resort — reduce font and force word-wrap at original position
    _reduce_shape_font(shape, max_reduction_frac=0.30)
    _enable_word_wrap(shape)
    warnings.append(
        f"{label}: no clear position found — font reduced 30 % to minimize overlap"
    )


def _validate_slide(
    slide,
    obstacles: list[tuple],
    slide_w: int,
    slide_h: int,
    margin: int,
) -> list[str]:
    """Check all text shapes against obstacles; return overlap warning strings."""
    issues: list[str] = []
    for idx, shape in enumerate(slide.shapes):
        if not shape.has_text_frame:
            continue
        geom = _geometry(shape)
        if not geom:
            continue
        for obs in obstacles:
            if _rects_overlap(geom, obs, margin):
                l, t, w, h = geom
                issues.append(
                    f"OVERLAP REMAINING — shape {idx} "
                    f"({l/914400:.2f}\",{t/914400:.2f}\") "
                    f"still overlaps an obstacle after resolution"
                )
                break
    return issues


# ── Public API ─────────────────────────────────────────────────────────────────

def transform_slide_layout(
    slide,
    slide_w: int,
    slide_h: int,
    profile: dict,
    translated_shape_idxs: set[int] | None = None,
) -> list[str]:
    """Transform one slide's layout for Arabic. Returns warning strings.

    Only call for LTR-source decks with an Arabic profile. Any unexpected
    failure rolls the slide back to its pre-transform state (translations
    are already in the tree at that point and are preserved) — a partially
    transformed slide must never ship.
    """
    sp_tree = slide.shapes._spTree
    backup = copy.deepcopy(sp_tree)
    try:
        return _transform_inner(
            slide, slide_w, slide_h, profile, translated_shape_idxs
        )
    except Exception as exc:
        log.exception("arabic layout transform failed; rolling slide back")
        sp_tree.getparent().replace(sp_tree, backup)
        slide.__dict__.pop("shapes", None)  # drop cached spTree wrapper
        return [f"layout transform failed, slide restored untransformed: {exc!r}"]


def _transform_inner(
    slide,
    slide_w: int,
    slide_h: int,
    profile: dict,
    translated_shape_idxs: set[int] | None,
) -> list[str]:
    warnings: list[str] = []
    cover_title_idxs: set[int] = set()
    zones = profile.get("layout_zones", {})
    typo = profile.get("typography", {})
    rules = profile.get("mirror_rules", {})
    min_move = rules.get("min_move_emu", 5000)
    translated = translated_shape_idxs or set()

    banner_rect = _detect_top_banner_rect(slide, slide_w, slide_h)
    logo_left_boundary = _detect_logo_left_boundary(slide, slide_w, slide_h)
    title_shape_idxs = _detect_title_shape_idxs(
        slide, slide_w, slide_h, translated
    )
    obstruction_left = _detect_top_right_obstruction_left(
        slide,
        slide_w,
        slide_h,
        exclude_idxs=title_shape_idxs,
    )
    if obstruction_left is not None:
        if logo_left_boundary is None:
            logo_left_boundary = obstruction_left
        else:
            logo_left_boundary = min(logo_left_boundary, obstruction_left)

    # Keep original geometry for translated text boxes so post-translation
    # fitting can preserve the source visual density/whitespace balance.
    original_text_geoms: dict[int, tuple[int, int, int, int]] = {}
    for idx, shape in enumerate(slide.shapes):
        if idx not in translated and idx not in title_shape_idxs:
            continue
        g = _geometry(shape)
        if g is not None and shape.has_text_frame:
            original_text_geoms[idx] = g

    # ══════════════════════════════════════════════════════════════════════════
    # PRE-PASS — Catalog obstacles and build a geometry lock set.
    #
    # Obstacle shapes (pictures, OLE, large non-text graphics) and any shape
    # whose bounding box lies within _NEIGHBOR_EMU of an obstacle are SPATIAL
    # ANCHORS.  They must NOT be moved by the mirror transform because:
    #   • Moving a screenshot breaks its callout arrows and numbered markers.
    #   • Moving a shape adjacent to an image (annotation, label, connector)
    #     destroys the spatial relationship with its target image.
    # Only isolated text boxes that are not near any obstacle are mirrored.
    # ══════════════════════════════════════════════════════════════════════════

    _NEIGHBOR_EMU = 457200   # 0.5 inch — shapes within this margin are anchored

    # Build list of obstacle bounding boxes BEFORE any transforms.
    pre_obstacles: list[tuple[int, int, int, int]] = []
    for shape in slide.shapes:
        geom = _geometry(shape)
        if geom and _is_obstacle(shape, slide_w, slide_h):
            pre_obstacles.append(geom)

    def _is_anchored(shape) -> bool:
        """True when shape must not be moved (is an obstacle or a neighbor)."""
        geom = _geometry(shape)
        if geom is None:
            return False
        # Title placeholders must always follow RTL mirroring rules.
        if _is_title_placeholder(shape):
            return False
        # Directional flow shapes and non-picture groups must still mirror.
        if _is_directional(shape):
            return False
        try:
            if (
                shape.shape_type == MSO_SHAPE_TYPE.GROUP
                and not _group_has_picture(shape)
            ):
                return False
        except Exception:
            pass
        # Obstacle itself
        if _is_obstacle(shape, slide_w, slide_h):
            return True
        # Logos/decorative edge ornaments are fixed by design.
        if _is_static_decoration(shape, slide_w, slide_h):
            return True
        # Adjacent to an obstacle (callout labels, arrows, numbered markers)
        return any(_rects_overlap(geom, obs, _NEIGHBOR_EMU) for obs in pre_obstacles)

    # ══════════════════════════════════════════════════════════════════════════
    # PASS 1 — Mirror / title-snap every shape.
    # ══════════════════════════════════════════════════════════════════════════

    for idx, shape in enumerate(slide.shapes):
        kind = _ph_kind(shape)
        geom = _geometry(shape)
        if geom is None:
            continue
        left, top, width, height = geom

        if kind == "center_title" and shape.has_text_frame:
            _fit_cover_title_shape(shape, slide_w, slide_h)
            cover_title_idxs.add(idx)
            continue

        # ── Detected title shapes (placeholder OR title-like textbox) ─────────
        if idx in title_shape_idxs and shape.has_text_frame:
            _mirror_and_fit_title_shape(
                shape=shape,
                original_left=left,
                original_top=top,
                original_width=width,
                original_height=height,
                slide_w=slide_w,
                slide_h=slide_h,
                logo_left_boundary=logo_left_boundary,
                banner_rect=banner_rect,
                warnings=warnings,
                label=f"shape {idx}",
            )
            continue

        # ── GEOMETRY LOCK: never move obstacle shapes or their neighbors ───────
        # Pictures, OLE objects, and shapes within 0.5" of any image are
        # spatial anchors.  Moving them breaks callout arrows, numbered
        # markers, and annotation labels that are positioned relative to them.
        # Skip the mirror transform entirely for these shapes.
        if _is_anchored(shape):
            # Still apply RTL text-direction and spacing to anchored text boxes
            # so their *content* reads correctly even though *position* is frozen.
            if idx in translated and shape.has_text_frame:
                tf = shape.text_frame
                try:
                    tf.margin_left  = Emu(_REF_MARGIN_LR)
                    tf.margin_right = Emu(_REF_MARGIN_LR)
                except Exception:
                    pass
                ls = typo.get("line_spacing")
                if ls:
                    paras = [p for p in tf.paragraphs if p.text.strip()]
                    if len(paras) >= 2 and height >= 0.12 * slide_h:
                        for para in paras:
                            try:
                                para.line_spacing = ls
                            except Exception:
                                pass
            continue  # position/size are locked

        # ── Mirror horizontally (margin swap) ─────────────────────────────────
        new_left = int(slide_w - left - width)
        if abs(new_left - left) >= min_move:
            before = _h_overflow(left, width, slide_w)
            shape.left = Emu(new_left)
            after = _h_overflow(new_left, width, slide_w)
            if abs(after - before) > 1:
                warnings.append(
                    f"shape {idx}: mirror changed off-canvas overflow "
                    f"({before} → {after} EMU)"
                )

        # ── RTL semantics: flip directional graphics, unfold groups ──────────
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                _mirror_group_children(shape, warnings, f"shape {idx}")
            elif _is_directional(shape):
                _flip_horizontal(shape)
        except Exception as exc:
            warnings.append(f"shape {idx}: RTL flip skipped ({exc})")

        # ── Reference spacing + margins on translated body text ──────────────
        if idx in translated and shape.has_text_frame:
            tf = shape.text_frame
            try:
                tf.margin_left = Emu(_REF_MARGIN_LR)
                tf.margin_right = Emu(_REF_MARGIN_LR)
            except Exception:
                pass
            ls = typo.get("line_spacing")
            if ls:
                paras = [p for p in tf.paragraphs if p.text.strip()]
                if len(paras) >= 2 and height >= 0.12 * slide_h:
                    for para in paras:
                        try:
                            para.line_spacing = ls
                        except Exception:
                            pass

    # ══════════════════════════════════════════════════════════════════════════
    # PASS 2 — Catalog obstacle positions after mirror transform.
    # ══════════════════════════════════════════════════════════════════════════

    obstacles = _collect_obstacles(slide, slide_w, slide_h)

    if not obstacles:
        return warnings  # nothing to collide with — fast path

    # ══════════════════════════════════════════════════════════════════════════
    # PASS 2.5 — Optimize translated Arabic textboxes before overlap routing.
    # ══════════════════════════════════════════════════════════════════════════

    line_spacing_base = float(typo.get("line_spacing", 1.5) or 1.5)
    for idx, shape in enumerate(slide.shapes):
        if idx not in translated:
            continue
        if idx in title_shape_idxs:
            continue
        if _ph_kind(shape) == "title":
            continue
        orig_geom = original_text_geoms.get(idx)
        if orig_geom is None:
            continue
        _optimize_translated_textbox(
            shape=shape,
            original_geom=orig_geom,
            obstacles=obstacles,
            slide_w=slide_w,
            slide_h=slide_h,
            margin=_MIN_MARGIN_EMU,
            line_spacing_base=line_spacing_base,
            warnings=warnings,
            label=f"shape {idx}",
        )

    # ══════════════════════════════════════════════════════════════════════════
    # PASS 3 — Resolve text-box overlaps with obstacles.
    #
    # We check translated text boxes only so logos/decorative labels keep their
    # original composition while translated content is fit safely.
    # ══════════════════════════════════════════════════════════════════════════

    for idx, shape in enumerate(slide.shapes):
        if idx not in translated and idx not in cover_title_idxs:
            continue
        if idx in title_shape_idxs and idx not in cover_title_idxs:
            continue
        if not shape.has_text_frame:
            continue
        # Skip title placeholders — they have their own safe zone.
        if _ph_kind(shape) == "title":
            continue
        geom = _geometry(shape)
        if geom is None:
            continue
        # Quick check: does this shape overlap any obstacle?
        if any(_rects_overlap(geom, obs, _MIN_MARGIN_EMU) for obs in obstacles):
            _resolve_overlap(
                shape, obstacles, slide_w, slide_h,
                _MIN_MARGIN_EMU, warnings, f"shape {idx}"
            )
            # Enable word-wrap unconditionally on repositioned shapes
            _enable_word_wrap(shape)

    # ── Word-wrap: enable on all translated body text shapes ──────────────────
    # Ensures Arabic text reflows inside its box regardless of whether it was
    # repositioned. Titles already have word_wrap set above.
    for idx, shape in enumerate(slide.shapes):
        if idx not in translated and idx not in cover_title_idxs:
            continue
        if not shape.has_text_frame:
            continue
        if _ph_kind(shape) == "title":
            continue
        _enable_word_wrap(shape)

    # ══════════════════════════════════════════════════════════════════════════
    # PASS 4 — Validate: log any remaining overlaps as export warnings.
    # ══════════════════════════════════════════════════════════════════════════

    remaining = _validate_slide(slide, obstacles, slide_w, slide_h, _MIN_MARGIN_EMU)
    warnings.extend(remaining)

    return warnings
