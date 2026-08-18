"""
Reference file processing for Education Studio.

Extracts text and detects structural elements (headings, warnings, procedures,
tables, figures, error codes, glossary terms) with page-level attribution.
"""
from __future__ import annotations

import io
import logging
import re
from typing import Any

log = logging.getLogger(__name__)

# ── Pattern library ────────────────────────────────────────────────────────────

_HEADING_PATTERNS = [
    re.compile(r'^(?:CHAPTER|SECTION|MODULE|UNIT|PART)\s+\d+[\s:.]+.{3,}', re.I),
    re.compile(r'^\d+(?:\.\d+){0,3}\s{1,4}[A-Z].{4,}$'),
    re.compile(r'^#{1,3}\s+.{3,}'),
    re.compile(r'^[A-Z][A-Z &\-/]{5,60}$'),
]

_WARNING_RE = re.compile(
    r'\b(WARNING|CAUTION|DANGER|NOTICE|NOTE|IMPORTANT|CRITICAL|ALERT)\b', re.I
)
_PROCEDURE_RE = re.compile(r'^\s*(\d{1,2})\.\s{1,4}[A-Z\(].{8,}')
_FIGURE_RE    = re.compile(r'\b(Figure|Fig\.?|FIGURE|Diagram|Illustration)\s+[\dA-Z]', re.I)
_TABLE_RE     = re.compile(r'\b(Table|TABLE)\s+[\dA-Z]', re.I)
_TABLE_ROW_RE = re.compile(r'.*\|.*\|')
_ERROR_RE     = re.compile(
    r'\b([EF][RR]?[-_]?\d{3,5}|ERR[-_]\d{3,5}|FAULT\s+\d{2,4}|ALARM\s+\d{2,4}|'
    r'CODE\s+[A-Z\d]{2,8}|E\d{3,4})\b'
)
_TROUBLESHOOT_RE = re.compile(
    r'\b(troubleshoot|fault diagnosis|corrective action|symptom|root cause|'
    r'probable cause|remedy|resolution)\b', re.I
)
_GLOSSARY_RE  = re.compile(r'^([A-Z][A-Za-z\s/\-()]{2,50}):\s{1,3}[A-Z].{10,}')
_MAINTENANCE_RE = re.compile(
    r'\b(preventive maintenance|PM\s+schedule|inspection interval|calibration|'
    r'service interval|lubrication|replacement interval)\b', re.I
)
_SPEC_RE = re.compile(
    r'\b(\d+(?:\.\d+)?\s*(?:mm|cm|m|kg|lb|psi|bar|kPa|MPa|A|mA|V|kV|W|kW|Hz|'
    r'rpm|°C|°F|dB|mSv|mR|mGy))\b'
)


# ── Public API ─────────────────────────────────────────────────────────────────

def process_reference_file(
    filename: str,
    file_bytes: bytes,
    file_type: str,
) -> dict[str, Any]:
    """
    Extract and structure a reference file.

    Returns:
        page_count, word_count, image_count, table_count, procedure_count,
        warning_count, section_count, figure_count, troubleshooting_count,
        doc_language, ocr_required, extracted_text, structure, error
    """
    ft = file_type.lower().strip(".")
    pages: list[dict] = []
    ocr_required = False
    error = None

    try:
        if ft == "pdf":
            pages, ocr_required = _extract_pdf(file_bytes)
        elif ft in ("docx", "doc"):
            pages = _extract_docx(file_bytes)
        elif ft in ("pptx", "ppt"):
            pages = _extract_pptx(file_bytes)
        elif ft in ("txt", "text", "md"):
            pages = _extract_txt(file_bytes)
        else:
            pages = _extract_txt(file_bytes)
    except Exception as exc:
        log.error("Reference extraction failed for %s: %s", filename, exc)
        return _empty(error=f"extraction_failed: {exc}")

    full_text = "\n\n".join(p["text"] for p in pages if p.get("text"))

    if not full_text.strip():
        err = "ocr_required" if ocr_required else "empty_reference"
        return _empty(page_count=len(pages), ocr_required=ocr_required, error=err)

    structure = _detect_structure(pages)
    counts    = _count_elements(full_text, structure)

    return {
        "page_count":            len(pages),
        "word_count":            counts["words"],
        "image_count":           counts["images"],
        "table_count":           counts["tables"],
        "procedure_count":       counts["procedures"],
        "warning_count":         counts["warnings"],
        "section_count":         counts["sections"],
        "figure_count":          counts["figures"],
        "troubleshooting_count": counts["troubleshooting"],
        "doc_language":          "en",
        "ocr_required":          ocr_required,
        "extracted_text":        full_text[:600_000],
        "structure":             structure[:3000],
        "error":                 None,
    }


# ── Text extractors ────────────────────────────────────────────────────────────

def _extract_pdf(file_bytes: bytes) -> tuple[list[dict], bool]:
    import fitz  # PyMuPDF
    pages: list[dict] = []
    total_chars = 0
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    for pno, page in enumerate(doc, 1):
        text = page.get_text("text", sort=True).strip()
        total_chars += len(text)
        pages.append({"page": pno, "text": text})
    doc.close()
    avg = total_chars / max(len(pages), 1)
    ocr_required = avg < 50 and len(pages) > 0
    return pages, ocr_required


def _extract_docx(file_bytes: bytes) -> list[dict]:
    from docx import Document  # type: ignore
    doc = Document(io.BytesIO(file_bytes))
    paras = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    chunk = 60
    pages: list[dict] = []
    for i in range(0, max(1, len(paras)), chunk):
        pages.append({"page": len(pages) + 1, "text": "\n".join(paras[i:i + chunk])})
    return pages or [{"page": 1, "text": "\n".join(paras[:300])}]


def _extract_pptx(file_bytes: bytes) -> list[dict]:
    from pptx import Presentation  # type: ignore
    prs = Presentation(io.BytesIO(file_bytes))
    pages: list[dict] = []
    for sno, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
        pages.append({"page": sno, "text": "\n".join(texts)})
    return pages or [{"page": 1, "text": ""}]


def _extract_txt(file_bytes: bytes) -> list[dict]:
    try:
        text = file_bytes.decode("utf-8", errors="replace")
    except Exception:
        text = file_bytes.decode("latin-1", errors="replace")
    lines = text.splitlines()
    chunk = 100
    pages: list[dict] = []
    for i in range(0, max(1, len(lines)), chunk):
        pages.append({"page": len(pages) + 1, "text": "\n".join(lines[i:i + chunk])})
    return pages or [{"page": 1, "text": text[:20_000]}]


# ── Structure detection ────────────────────────────────────────────────────────

def _detect_structure(pages: list[dict]) -> list[dict]:
    items: list[dict] = []
    for pi in pages:
        pno = pi.get("page", 1)
        for line in pi.get("text", "").splitlines():
            ls = line.strip()
            if not ls or len(ls) < 5:
                continue
            t = _classify_line(ls)
            if t:
                items.append({"page": pno, "type": t, "text": ls[:300]})
    return items


def _classify_line(ls: str) -> str | None:
    for pat in _HEADING_PATTERNS:
        if pat.match(ls):
            return "heading"
    if _WARNING_RE.search(ls):
        return "warning"
    if _PROCEDURE_RE.match(ls) and len(ls) > 15:
        return "procedure"
    if _FIGURE_RE.search(ls):
        return "figure"
    if _TABLE_RE.search(ls) or _TABLE_ROW_RE.match(ls):
        return "table"
    if _ERROR_RE.search(ls):
        return "error_code"
    if _TROUBLESHOOT_RE.search(ls):
        return "troubleshooting"
    if _MAINTENANCE_RE.search(ls):
        return "maintenance"
    if _GLOSSARY_RE.match(ls):
        return "glossary"
    return None


# ── Element counting ───────────────────────────────────────────────────────────

def _count_elements(full_text: str, structure: list[dict]) -> dict[str, int]:
    type_counts: dict[str, int] = {}
    for item in structure:
        t = item.get("type", "")
        type_counts[t] = type_counts.get(t, 0) + 1

    return {
        "words":          len(full_text.split()),
        "images":         max(len(_FIGURE_RE.findall(full_text)), type_counts.get("figure", 0)),
        "tables":         max(len(_TABLE_RE.findall(full_text)), type_counts.get("table", 0)),
        "warnings":       max(len(_WARNING_RE.findall(full_text)) // 2, type_counts.get("warning", 0)),
        "procedures":     max(len(_PROCEDURE_RE.findall(full_text)), type_counts.get("procedure", 0)),
        "sections":       type_counts.get("heading", 0),
        "figures":        max(len(_FIGURE_RE.findall(full_text)), type_counts.get("figure", 0)),
        "troubleshooting":max(len(_TROUBLESHOOT_RE.findall(full_text)) // 4, type_counts.get("troubleshooting", 0)),
    }


# ── Helpers ────────────────────────────────────────────────────────────────────

def _empty(
    page_count: int = 0,
    ocr_required: bool = False,
    error: str | None = None,
) -> dict[str, Any]:
    return {
        "page_count": page_count, "word_count": 0, "image_count": 0,
        "table_count": 0, "procedure_count": 0, "warning_count": 0,
        "section_count": 0, "figure_count": 0, "troubleshooting_count": 0,
        "doc_language": "en", "ocr_required": ocr_required,
        "extracted_text": "", "structure": [], "error": error,
    }
