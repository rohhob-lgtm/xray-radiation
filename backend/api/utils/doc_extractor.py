"""
Document extraction layer for the Translation Studio.

Extracts structured text segments from various file formats, preserving
identity information so translated text can be written back precisely.

Each segment is a dict:
  {
    "id":       str  — unique within the document,
    "source":   str  — original text,
    "seg_type": str  — "paragraph" | "table_cell" | "header" | "footer" |
                       "slide_title" | "slide_body" | "slide_notes" |
                       "sheet_cell" | "line" | "element",
    "loc":      dict — location metadata for rebuilder (format-specific),
  }
"""
from __future__ import annotations

import csv
import hashlib
import html as html_mod
import io
import json
import logging
import re
import uuid
from html.parser import HTMLParser
from typing import Any

log = logging.getLogger(__name__)


# ── Helpers ───────────────────────────────────────────────────────────────────

def _seg_id() -> str:
    return str(uuid.uuid4())[:8]


def _is_translatable(text: str) -> bool:
    """Skip empty, numeric-only, code-like, or very short strings."""
    t = text.strip()
    if not t or len(t) < 2:
        return False
    # Skip purely numeric / punctuation / symbols
    if re.fullmatch(r"[\d\s\.,;:!\?\-\+\=\(\)\[\]\{\}\/\\%&@#\*_~^<>|\"'`]+", t):
        return False
    # Skip lines that look like file paths, URLs, or code
    if t.startswith(("http://", "https://", "www.", "//", "/*", "#!")):
        return False
    # Must have at least one letter - Latin (incl. accented French/Spanish),
    # Arabic, or Cyrillic (Russian). Extend this set alongside api.languages
    # whenever a language using a different script is added.
    if not re.search(r"[a-zA-Z\u00C0-\u024F\u0400-\u04FF\u0600-\u06FF]", t):
        return False
    return True


# ── PDF extraction ─────────────────────────────────────────────────────────────

def extract_pdf(file_bytes: bytes, ocr_lang: str = "eng") -> list[dict]:
    """
    Extract text segments from a PDF.
    Tries PyMuPDF for text-selectable PDFs; falls back to pytesseract OCR.

    ocr_lang: tesseract language code for the scanned-page OCR fallback
        (e.g. "eng", "ara", "rus", "fra", "spa") — should match the
        document's *source* language, not the translation target.
    """
    segments = []
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        total_text = ""
        for page in doc:
            total_text += page.get_text("text") or ""
        doc.close()

        has_text = len(total_text.strip()) > 100
        if has_text:
            segments = _extract_pdf_text(file_bytes)
        else:
            segments = _extract_pdf_ocr(file_bytes, ocr_lang)
    except Exception as e:
        log.warning("PDF extraction error: %s", e)
        segments = _extract_pdf_ocr(file_bytes, ocr_lang)

    return segments


def _extract_pdf_text(file_bytes: bytes) -> list[dict]:
    import fitz
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    segments = []
    for page_num, page in enumerate(doc, 1):
        blocks = page.get_text("blocks") or []
        for block_idx, block in enumerate(blocks):
            # block: (x0, y0, x1, y1, text, block_no, block_type)
            if len(block) < 5:
                continue
            text = block[4].strip() if block[4] else ""
            # Split block into paragraphs
            for para_idx, para in enumerate(text.split("\n\n")):
                para = para.strip()
                if _is_translatable(para):
                    segments.append({
                        "id": _seg_id(),
                        "source": para,
                        "seg_type": "paragraph",
                        "loc": {
                            "format": "pdf",
                            "page": page_num,
                            "block": block_idx,
                            "para": para_idx,
                        },
                    })
    doc.close()
    return segments


def _extract_pdf_ocr(file_bytes: bytes, ocr_lang: str = "eng") -> list[dict]:
    """OCR fallback using pytesseract.

    ocr_lang: tesseract language code matching the document's source
    language (see api.languages.ocr_lang_code) — recognizing the page as
    English when the source is actually Russian/Arabic/etc. silently
    garbles or drops the text.
    """
    try:
        import fitz
        from PIL import Image
        import pytesseract

        doc = fitz.open(stream=file_bytes, filetype="pdf")
        segments = []
        for page_num, page in enumerate(doc, 1):
            mat = fitz.Matrix(2.0, 2.0)  # 2x scale for better OCR
            pix = page.get_pixmap(matrix=mat)
            img_bytes = pix.tobytes("png")
            img = Image.open(io.BytesIO(img_bytes))
            text = pytesseract.image_to_string(img, lang=ocr_lang)
            for line_idx, line in enumerate(text.split("\n")):
                line = line.strip()
                if _is_translatable(line):
                    segments.append({
                        "id": _seg_id(),
                        "source": line,
                        "seg_type": "line",
                        "loc": {"format": "pdf_ocr", "page": page_num, "line": line_idx},
                    })
        doc.close()
        return segments
    except ImportError:
        log.warning("pytesseract not available; returning OCR notice segment")
        return [{
            "id": _seg_id(),
            "source": "[OCR required — install pytesseract for scanned PDF support]",
            "seg_type": "line",
            "loc": {"format": "pdf_ocr", "page": 1, "line": 0},
        }]
    except Exception as e:
        log.warning("OCR failed: %s", e)
        return []


# ── DOCX extraction ────────────────────────────────────────────────────────────

def extract_docx(file_bytes: bytes) -> list[dict]:
    """Extract segments from a DOCX file preserving paragraph/table/header structure."""
    from docx import Document
    doc = Document(io.BytesIO(file_bytes))
    segments = []
    idx = 0

    def _add(text: str, seg_type: str, loc: dict):
        nonlocal idx
        if _is_translatable(text):
            segments.append({
                "id": _seg_id(),
                "source": text.strip(),
                "seg_type": seg_type,
                "loc": loc,
            })
            idx += 1

    # Body paragraphs
    # para.text joins soft line breaks (<w:br/>) as "\n" in one string.
    # We split on "\n" so each visible line becomes its own translatable segment
    # while preserving para_idx + line_idx so the rebuilder can reconstruct them.
    for para_idx, para in enumerate(doc.paragraphs):
        raw = para.text
        if not raw.strip():
            continue
        lines = raw.split("\n")
        total_lines = len(lines)
        if total_lines == 1:
            _add(raw.strip(), "paragraph", {
                "format": "docx",
                "para_idx": para_idx,
                "line_idx": 0,
                "total_lines": 1,
            })
        else:
            for line_idx, line in enumerate(lines):
                line = line.strip()
                if _is_translatable(line):
                    _add(line, "paragraph", {
                        "format": "docx",
                        "para_idx": para_idx,
                        "line_idx": line_idx,
                        "total_lines": total_lines,
                    })
                else:
                    # Keep non-translatable lines as passthrough so the rebuilder
                    # can reconstruct the full soft-return paragraph structure
                    if line:   # skip truly empty lines within the paragraph
                        segments.append({
                            "id": _seg_id(),
                            "source": line,
                            "target": line,   # passthrough — keep as-is
                            "seg_type": "paragraph",
                            "loc": {
                                "format": "docx",
                                "para_idx": para_idx,
                                "line_idx": line_idx,
                                "total_lines": total_lines,
                                "passthrough": True,
                            },
                        })
                        idx += 1

    # Tables
    for tbl_idx, table in enumerate(doc.tables):
        for row_idx, row in enumerate(table.rows):
            for col_idx, cell in enumerate(row.cells):
                text = cell.text.strip()
                if text:
                    _add(text, "table_cell", {
                        "format": "docx",
                        "tbl_idx": tbl_idx,
                        "row_idx": row_idx,
                        "col_idx": col_idx,
                    })

    # Headers and footers
    for sec_idx, section in enumerate(doc.sections):
        for hdr_para_idx, para in enumerate(section.header.paragraphs):
            if para.text.strip():
                _add(para.text, "header", {
                    "format": "docx",
                    "section": sec_idx,
                    "header_para": hdr_para_idx,
                })
        for ftr_para_idx, para in enumerate(section.footer.paragraphs):
            if para.text.strip():
                _add(para.text, "footer", {
                    "format": "docx",
                    "section": sec_idx,
                    "footer_para": ftr_para_idx,
                })

    return segments


# ── PPTX extraction ────────────────────────────────────────────────────────────

def _pptx_shape_is_offslide(shape, slide_w: int, slide_h: int) -> bool:
    """Return True if the shape lies completely outside the visible slide canvas.

    Off-slide shapes are used in PowerPoint for animation staging, hidden
    technical content, and diagram helpers.  Translating their text causes
    orphan broken-text artefacts in the output.

    A shape is "off-slide" when its bounding box does not intersect the
    slide rectangle [0, slide_w] × [0, slide_h] at all.
    """
    try:
        left   = shape.left   or 0
        top    = shape.top    or 0
        width  = shape.width  or 0
        height = shape.height or 0
        if left + width <= 0 or top + height <= 0:
            return True
        if slide_w > 0 and left >= slide_w:
            return True
        if slide_h > 0 and top >= slide_h:
            return True
        return False
    except Exception:
        return False


def _pptx_safe_name(shape) -> str:
    """Return the shape name as a string, or '' on any error."""
    try:
        return str(shape.name or "")
    except Exception:
        return ""


def _pptx_safe_shape_type(shape):
    """Return shape.shape_type or None on any error."""
    try:
        return shape.shape_type
    except Exception:
        return None


def _pptx_extract_table(
    shape, slide_idx: int, shape_idx: int,
    segments: list, diag: dict,
) -> None:
    """
    Safely extract text from a table shape.
    Each non-empty cell becomes one segment with all paragraph text joined.
    Malformed cells are skipped individually; the rest continue.
    """
    diag["tables"] += 1
    try:
        table = shape.table
        if table is None:
            log.warning("PPTX: slide %d shape %d — table is None, skipping",
                        slide_idx + 1, shape_idx)
            diag["skipped"] += 1
            return
        for row_idx, row in enumerate(table.rows):
            for col_idx, cell in enumerate(row.cells):
                try:
                    tf = cell.text_frame if cell is not None else None
                    if tf is None:
                        continue
                    parts = []
                    for p in tf.paragraphs:
                        if p is None:
                            continue
                        try:
                            t = p.text          # may be None for some cells
                            if t:
                                parts.append(t)
                        except Exception:
                            continue
                    text = " ".join(parts).strip()
                except Exception as exc:
                    log.debug("PPTX: slide %d table cell [%d,%d] skipped: %s",
                              slide_idx + 1, row_idx, col_idx, exc)
                    diag["skipped"] += 1
                    continue
                if _is_translatable(text):
                    segments.append({
                        "id": _seg_id(),
                        "source": text,
                        "seg_type": "table_cell",
                        "loc": {
                            "format": "pptx",
                            "slide_idx": slide_idx,
                            "shape_idx": shape_idx,
                            "table_cell": True,
                            "row_idx": row_idx,
                            "col_idx": col_idx,
                        },
                    })
    except Exception as exc:
        log.warning("PPTX: slide %d shape %d table extraction failed: %s",
                    slide_idx + 1, shape_idx, exc)
        diag["skipped"] += 1


def _pptx_extract_text_frame(
    shape, slide_idx: int, shape_idx: int,
    segments: list, diag: dict,
) -> None:
    """
    Safely extract paragraph text from a text-bearing shape.
    Returns without adding segments if the shape has no readable text.
    """
    # Determine segment type (slide_title vs slide_body)
    seg_type = "slide_body"
    try:
        if shape.is_placeholder:
            ph = shape.placeholder_format
            if ph is not None and getattr(ph, "idx", None) == 0:
                seg_type = "slide_title"
    except Exception:
        pass

    # Access the text frame — has_text_frame True does NOT guarantee
    # text_frame is non-None for all shape types in all PPTX variants
    tf = None
    try:
        tf = shape.text_frame
    except Exception as exc:
        log.debug("PPTX: slide %d shape %d text_frame inaccessible: %s",
                  slide_idx + 1, shape_idx, exc)
        diag["skipped"] += 1
        return
    if tf is None:
        diag["empty_placeholders"] += 1
        return

    paragraphs = []
    try:
        paragraphs = list(tf.paragraphs)
    except Exception as exc:
        log.warning("PPTX: slide %d shape %d paragraphs unreadable: %s",
                    slide_idx + 1, shape_idx, exc)
        diag["skipped"] += 1
        return

    found_text = False
    for para_idx, para in enumerate(paragraphs):
        if para is None:
            continue
        try:
            raw = para.text       # _Paragraph.text joins runs — may return None
            if raw is None:
                continue
            text = raw.strip()
        except Exception:
            continue
        if _is_translatable(text):
            segments.append({
                "id": _seg_id(),
                "source": text,
                "seg_type": seg_type,
                "loc": {
                    "format": "pptx",
                    "slide_idx": slide_idx,
                    "shape_idx": shape_idx,
                    "para_idx": para_idx,
                },
            })
            found_text = True

    if found_text:
        diag["text_shapes"] += 1
    else:
        diag["empty_placeholders"] += 1


def _pptx_dispatch_shape(
    shape, slide_idx: int, shape_idx: int,
    segments: list, diag: dict,
    slide_w: int, slide_h: int,
) -> None:
    """
    Route one shape (or group child) to the right extraction function.
    Skips off-slide shapes, classifies images/charts/OLE for diagnostics,
    and recurses into group shapes.

    All property reads are guarded so a single malformed shape never
    raises an uncaught exception.
    """
    diag["total_shapes"] += 1

    if _pptx_shape_is_offslide(shape, slide_w, slide_h):
        diag["skipped"] += 1
        return

    stype = _pptx_safe_shape_type(shape)

    # ── Group shape — recurse into children ───────────────────────────────────
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        _GROUP = MSO_SHAPE_TYPE.GROUP
        _PICTURE = MSO_SHAPE_TYPE.PICTURE
        _LINKED_PICTURE = MSO_SHAPE_TYPE.LINKED_PICTURE
        _CHART = MSO_SHAPE_TYPE.CHART
        _OLE = MSO_SHAPE_TYPE.OLE_OBJECT
        _MEDIA = MSO_SHAPE_TYPE.MEDIA
    except Exception:
        _GROUP = _PICTURE = _LINKED_PICTURE = _CHART = _OLE = _MEDIA = None

    if stype is not None and stype == _GROUP:
        diag["grouped_shapes"] += 1
        try:
            children = list(shape.shapes)
        except Exception:
            diag["skipped"] += 1
            return
        for child_idx, child in enumerate(children):
            if child is None:
                continue
            try:
                _pptx_dispatch_shape(
                    child,
                    slide_idx,
                    shape_idx * 10000 + child_idx,
                    segments, diag,
                    slide_w, slide_h,
                )
            except Exception as exc:
                log.debug(
                    "PPTX: slide %d group %d child %d skipped: %s",
                    slide_idx + 1, shape_idx, child_idx, exc,
                )
                diag["skipped"] += 1
        return

    # ── Table shape ───────────────────────────────────────────────────────────
    has_table = False
    try:
        has_table = bool(shape.has_table)
    except Exception:
        pass
    if has_table:
        _pptx_extract_table(shape, slide_idx, shape_idx, segments, diag)
        return

    # ── Non-text shapes — classify for diagnostics, preserve unchanged ────────
    has_tf = False
    try:
        has_tf = bool(shape.has_text_frame)
    except Exception:
        pass
    if not has_tf:
        if stype in (_PICTURE, _LINKED_PICTURE):
            diag["images"] += 1
        elif stype == _CHART:
            diag["charts"] += 1
        elif stype in (_OLE, _MEDIA):
            diag["unsupported"] += 1
        else:
            diag["unsupported"] += 1
        return   # preserve shape unchanged — do not attempt text read

    # ── Text-frame shape ──────────────────────────────────────────────────────
    _pptx_extract_text_frame(shape, slide_idx, shape_idx, segments, diag)


def extract_pptx(file_bytes: bytes) -> list[dict]:
    """
    Extract translatable text segments from a PPTX presentation.

    Robustness guarantees
    ─────────────────────
    • Slide-level fault isolation — one broken slide is logged and skipped;
      the remaining slides are still processed.
    • Shape-level fault isolation — one broken shape is logged and skipped;
      remaining shapes on the same slide are still processed.
    • Group shapes are traversed recursively.
    • Every .text access is guarded; None values are handled explicitly.
    • Unsupported objects (images, charts, OLE, media) are classified and
      preserved unchanged — no attempt is made to read .text from them.
    • Speaker notes are extracted safely; notes_text_frame None is handled.
    • Parser diagnostics are logged after extraction completes.

    Only raises if the file cannot be opened at all (corrupt / not a PPTX).
    """
    from pptx import Presentation
    prs = Presentation(io.BytesIO(file_bytes))
    segments: list[dict] = []
    skipped_slides: list[int] = []

    slide_w = int(prs.slide_width  or 0)
    slide_h = int(prs.slide_height or 0)

    diag: dict[str, int] = {
        "total_slides":       0,
        "total_shapes":       0,
        "text_shapes":        0,
        "tables":             0,
        "images":             0,
        "charts":             0,
        "grouped_shapes":     0,
        "smartart":           0,
        "unsupported":        0,
        "empty_placeholders": 0,
        "skipped":            0,
    }

    try:
        total_slides = len(prs.slides)
    except Exception:
        total_slides = 0
    diag["total_slides"] = total_slides

    for slide_idx, slide in enumerate(prs.slides):

        # ── Slide-level fault isolation ───────────────────────────────────────
        try:
            shapes_list = list(slide.shapes)
        except Exception as exc:
            log.warning(
                "PPTX: slide %d — cannot read shape list, slide preserved: %s",
                slide_idx + 1, exc,
            )
            skipped_slides.append(slide_idx + 1)
            diag["skipped"] += 1
            continue

        for shape_idx, shape in enumerate(shapes_list):
            if shape is None:
                diag["skipped"] += 1
                continue
            try:
                _pptx_dispatch_shape(
                    shape, slide_idx, shape_idx,
                    segments, diag,
                    slide_w, slide_h,
                )
            except Exception as exc:
                log.warning(
                    "PPTX: slide %d shape %d (%s) — unhandled error, shape preserved: %s",
                    slide_idx + 1, shape_idx, _pptx_safe_name(shape), exc,
                )
                diag["skipped"] += 1

        # ── Speaker notes ─────────────────────────────────────────────────────
        # CRASH GUARD: notes_text_frame explicitly returns None in python-pptx
        # when the notes slide has no body placeholder (common in imported or
        # converted PPTX files). Accessing .text on None raises
        # "AttributeError: 'NoneType' object has no attribute 'text'".
        try:
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                if notes_slide is not None:
                    ntf = notes_slide.notes_text_frame   # ← may be None
                    if ntf is not None:
                        raw = ntf.text                   # ← only safe when ntf is not None
                        if raw is not None:
                            notes_text = raw.strip()
                            if _is_translatable(notes_text):
                                segments.append({
                                    "id": _seg_id(),
                                    "source": notes_text,
                                    "seg_type": "slide_notes",
                                    "loc": {
                                        "format": "pptx",
                                        "slide_idx": slide_idx,
                                        "notes": True,
                                    },
                                })
        except Exception as exc:
            log.debug("PPTX: slide %d notes skipped: %s", slide_idx + 1, exc)

    if skipped_slides:
        log.warning(
            "PPTX extraction: %d slide(s) preserved without extraction: %s",
            len(skipped_slides), skipped_slides,
        )

    # ── Slide-master text ───────────────────────────────────────────────────
    # Some decks bake repeated content (e.g. a running course/deck title)
    # directly into the slide MASTER rather than overriding it on every
    # individual slide. That text is invisible to the per-slide loop above,
    # so without this it silently stays untranslated forever — and since it
    # renders in the same banner area as the (translated) slide title, the
    # two end up visually overlapping in mixed English/Arabic. Masters are
    # shared across many slides, so each is only walked once regardless of
    # how many slides use it.
    try:
        for master_idx, master in enumerate(prs.slide_masters):
            try:
                master_shapes = list(master.shapes)
            except Exception:
                continue
            for shape_idx, shape in enumerate(master_shapes):
                if not getattr(shape, "has_text_frame", False):
                    continue
                try:
                    tf = shape.text_frame
                    for para_idx, para in enumerate(tf.paragraphs):
                        if para is None:
                            continue
                        raw = para.text
                        if raw is None:
                            continue
                        text = raw.strip()
                        if _is_translatable(text):
                            segments.append({
                                "id": _seg_id(),
                                "source": text,
                                "seg_type": "master_text",
                                "loc": {
                                    "format": "pptx",
                                    "master_idx": master_idx,
                                    "shape_idx": shape_idx,
                                    "para_idx": para_idx,
                                },
                            })
                except Exception as exc:
                    log.debug("PPTX: master %d shape %d skipped: %s",
                              master_idx, shape_idx, exc)
    except Exception as exc:
        log.debug("PPTX master-text extraction skipped: %s", exc)

    log.info(
        "PPTX diagnostics — slides=%d shapes=%d text_shapes=%d tables=%d "
        "images=%d charts=%d groups=%d unsupported=%d "
        "empty=%d skipped=%d → %d segments extracted",
        diag["total_slides"], diag["total_shapes"],
        diag["text_shapes"], diag["tables"],
        diag["images"], diag["charts"],
        diag["grouped_shapes"], diag["unsupported"],
        diag["empty_placeholders"], diag["skipped"],
        len(segments),
    )

    return segments


# ── XLSX extraction ────────────────────────────────────────────────────────────

def extract_xlsx(file_bytes: bytes) -> list[dict]:
    """Extract cell text from all sheets (skip formulas)."""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True)
        segments = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            for row in ws.iter_rows():
                for cell in row:
                    if cell.value is not None and isinstance(cell.value, str):
                        if _is_translatable(cell.value):
                            segments.append({
                                "id": _seg_id(),
                                "source": cell.value.strip(),
                                "seg_type": "sheet_cell",
                                "loc": {
                                    "format": "xlsx",
                                    "sheet": sheet_name,
                                    "row": cell.row,
                                    "col": cell.column,
                                },
                            })
        return segments
    except Exception as e:
        log.warning("XLSX extraction failed: %s", e)
        return []


# ── Plain text / Markdown extraction ──────────────────────────────────────────

def extract_txt(file_bytes: bytes) -> list[dict]:
    """Extract non-empty lines from plain text or Markdown."""
    try:
        text = file_bytes.decode("utf-8", errors="replace")
    except Exception:
        text = file_bytes.decode("latin-1", errors="replace")

    segments = []
    for line_idx, line in enumerate(text.split("\n")):
        line = line.strip()
        # Skip Markdown headings markers but keep the text
        line_clean = re.sub(r"^#+\s*", "", line)
        if _is_translatable(line_clean):
            segments.append({
                "id": _seg_id(),
                "source": line_clean,
                "seg_type": "line",
                "loc": {"format": "txt", "line": line_idx},
            })
    return segments


# ── HTML extraction ────────────────────────────────────────────────────────────

class _HTMLTextExtractor(HTMLParser):
    """Minimal HTML → text list extractor."""

    _SKIP_TAGS = {"script", "style", "head", "meta", "link", "noscript"}

    def __init__(self):
        super().__init__()
        self._skip_depth = 0
        self._current_tag = ""
        self._segments: list[tuple[str, str]] = []  # (tag, text)

    def handle_starttag(self, tag, attrs):
        self._current_tag = tag
        if tag in self._SKIP_TAGS:
            self._skip_depth += 1

    def handle_endtag(self, tag):
        if tag in self._SKIP_TAGS and self._skip_depth > 0:
            self._skip_depth -= 1

    def handle_data(self, data):
        if self._skip_depth > 0:
            return
        text = data.strip()
        if text:
            self._segments.append((self._current_tag, text))


def extract_html(file_bytes: bytes) -> list[dict]:
    try:
        text = file_bytes.decode("utf-8", errors="replace")
    except Exception:
        text = file_bytes.decode("latin-1", errors="replace")

    parser = _HTMLTextExtractor()
    parser.feed(text)
    segments = []
    for elem_idx, (tag, content) in enumerate(parser._segments):
        content = html_mod.unescape(content).strip()
        if _is_translatable(content):
            seg_type = "element"
            if tag in ("h1", "h2", "h3", "h4", "h5", "h6"):
                seg_type = "slide_title"
            segments.append({
                "id": _seg_id(),
                "source": content,
                "seg_type": seg_type,
                "loc": {"format": "html", "elem": elem_idx, "tag": tag},
            })
    return segments


# ── CSV extraction ─────────────────────────────────────────────────────────────

def extract_csv(file_bytes: bytes) -> list[dict]:
    try:
        text = file_bytes.decode("utf-8", errors="replace")
    except Exception:
        text = file_bytes.decode("latin-1", errors="replace")

    segments = []
    reader = csv.reader(io.StringIO(text))
    for row_idx, row in enumerate(reader):
        for col_idx, cell in enumerate(row):
            cell = cell.strip()
            if _is_translatable(cell):
                segments.append({
                    "id": _seg_id(),
                    "source": cell,
                    "seg_type": "sheet_cell",
                    "loc": {"format": "csv", "row": row_idx, "col": col_idx},
                })
    return segments


# ── RTF extraction (basic) ─────────────────────────────────────────────────────

def extract_rtf(file_bytes: bytes) -> list[dict]:
    """Strip RTF control words and extract plain text."""
    try:
        text = file_bytes.decode("utf-8", errors="replace")
    except Exception:
        text = file_bytes.decode("latin-1", errors="replace")

    # Strip RTF control sequences
    text = re.sub(r"\\[a-z]+\d*\s?", " ", text)
    text = re.sub(r"[{}]", "", text)
    text = re.sub(r"\s+", " ", text).strip()

    segments = []
    for line_idx, line in enumerate(text.split(". ")):
        line = line.strip()
        if _is_translatable(line):
            segments.append({
                "id": _seg_id(),
                "source": line,
                "seg_type": "paragraph",
                "loc": {"format": "rtf", "line": line_idx},
            })
    return segments


# ── Dispatcher ────────────────────────────────────────────────────────────────

EXTRACTORS = {
    "pdf":  extract_pdf,
    "docx": extract_docx,
    "pptx": extract_pptx,
    "xlsx": extract_xlsx,
    "txt":  extract_txt,
    "md":   extract_txt,
    "html": extract_html,
    "htm":  extract_html,
    "csv":  extract_csv,
    "rtf":  extract_rtf,
    "xml":  extract_html,  # treat as HTML-like markup
}


def extract_document(file_bytes: bytes, file_type: str, ocr_lang: str = "eng") -> list[dict]:
    """
    Extract translatable segments from a document.

    Args:
        file_bytes: raw file content
        file_type:  lowercase extension without dot (e.g. "pdf", "docx")
        ocr_lang:   tesseract language code used only for the scanned-PDF
                    OCR fallback — should match the document's source
                    language (api.languages.ocr_lang_code(source_lang)).

    Returns:
        list of segment dicts with id, source, seg_type, loc
    """
    ftype = file_type.lower()
    extractor = EXTRACTORS.get(ftype)
    if extractor is None:
        # Fall back to plain text extraction
        log.warning("No extractor for '%s', falling back to txt", file_type)
        extractor = extract_txt

    try:
        if ftype == "pdf":
            segments = extract_pdf(file_bytes, ocr_lang)
        else:
            segments = extractor(file_bytes)
        log.info("Extracted %d segments from %s document", len(segments), file_type)
        return segments
    except Exception as e:
        log.error("Extraction failed for %s: %s", file_type, e)
        raise
