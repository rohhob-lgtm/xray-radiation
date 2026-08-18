"""
Single-canvas PPTX builder for AI-generated design specs — the Mode 2a
design-preserving path (see design_orchestrator.py). Unlike the Pillow PNG
renderer (design_renderer.py, Mode 2b), a PPTX's title/subtitle/bullets are
real text boxes, not flattened pixels — Canva's Design Import API
(POST /v1/imports) reconstructs them as editable text when the file is
imported, which a PNG import cannot do.

Shares canvas size/palette with design_renderer.py via design_specs.py so
both rendering paths agree on what a given design_type looks like. Mirrors
this repo's existing pptx_gen.py shape/textbox helper conventions and
arabic_layout_engine.py's `<a:pPr rtl="1"/>` convention for RTL paragraphs.
"""
from __future__ import annotations

import io
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from .design_specs import resolve_dimensions, resolve_palette

_PX_PER_INCH = 96


def _emu(px: int) -> Emu:
    return Inches(px / _PX_PER_INCH)


def _rgb(color: tuple[int, int, int]) -> RGBColor:
    return RGBColor(*color)


def _set_para_rtl(paragraph) -> None:
    """Force rtl="1" on paragraph properties — same convention as
    arabic_layout_engine.py's _set_para_rtl, reimplemented here rather than
    imported since that module's helper is private and tailored to a
    different (post-hoc slide transform) input shape."""
    try:
        ppr = paragraph._p.get_or_add_pPr()
        ppr.set("rtl", "1")
    except Exception:
        pass


def _fill(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _text_box(
    slide, text: str, left: Emu, top: Emu, width: Emu, height: Emu, *,
    size: int, bold: bool = False, color: RGBColor, is_rtl: bool, font: str = "Calibri",
) -> Any:
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT if is_rtl else PP_ALIGN.LEFT
    if is_rtl:
        _set_para_rtl(p)
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font
    return tb


def _bullet_box(
    slide, bullets: list[str], left: Emu, top: Emu, width: Emu, height: Emu, *,
    size: int, color: RGBColor, is_rtl: bool, font: str = "Calibri",
) -> Any:
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.RIGHT if is_rtl else PP_ALIGN.LEFT
        if is_rtl:
            _set_para_rtl(p)
        p.space_after = Pt(6)
        r = p.add_run()
        r.text = f"{bullet.strip()} •" if is_rtl else f"• {bullet.strip()}"
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.name = font
    return tb


def render_design_pptx(spec: dict, design_type: str) -> bytes:
    """Renders a structured design spec ({title, subtitle?, bullets?,
    palette?, direction?}) to a single-slide PPTX, sized to the design
    type's canvas. Returns raw .pptx bytes ready for canva.import_design.
    """
    width_px, height_px, layout = resolve_dimensions(design_type)
    palette = resolve_palette(spec.get("palette"))
    is_rtl = (spec.get("direction") or "").lower() == "rtl"

    prs = Presentation()
    prs.slide_width = _emu(width_px)
    prs.slide_height = _emu(height_px)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(palette["bg"])

    bar_h = _emu(max(8, int(height_px * 0.012)))
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, _emu(width_px), bar_h)
    _fill(top_bar, _rgb(palette["primary"]))
    bottom_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, _emu(height_px) - bar_h, _emu(width_px), bar_h)
    _fill(bottom_bar, _rgb(palette["accent"]))

    if layout == "certificate":
        border_px = max(6, int(min(width_px, height_px) * 0.012))
        border = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, _emu(border_px), _emu(border_px),
            _emu(width_px - 2 * border_px), _emu(height_px - 2 * border_px),
        )
        border.fill.background()
        border.line.color.rgb = _rgb(palette["accent"])
        border.line.width = Pt(max(2, border_px / 4))

    margin_px = int(width_px * 0.08)
    content_left = _emu(margin_px)
    content_width = _emu(width_px - 2 * margin_px)

    title = (spec.get("title") or "Untitled Design").strip()
    subtitle = (spec.get("subtitle") or "").strip()
    bullets = [b.strip() for b in (spec.get("bullets") or []) if isinstance(b, str) and b.strip()][:6]

    title_size = max(28, int(width_px * 0.045))
    subtitle_size = max(16, int(width_px * 0.022))
    bullet_size = max(14, int(width_px * 0.018))

    y_px = int(height_px * (0.20 if layout == "certificate" else 0.14))
    title_h_px = int(title_size * 2.2)
    _text_box(
        slide, title, content_left, _emu(y_px), content_width, _emu(title_h_px),
        size=title_size, bold=True, color=_rgb(palette["text"]), is_rtl=is_rtl,
    )
    y_px += title_h_px

    if subtitle:
        subtitle_h_px = int(subtitle_size * 3)
        _text_box(
            slide, subtitle, content_left, _emu(y_px), content_width, _emu(subtitle_h_px),
            size=subtitle_size, color=_rgb(palette["accent"]), is_rtl=is_rtl,
        )
        y_px += subtitle_h_px

    if bullets:
        y_px += int(height_px * 0.02)
        bullets_h_px = height_px - y_px - int(height_px * 0.05)
        _bullet_box(
            slide, bullets, content_left, _emu(y_px), content_width, _emu(max(bullets_h_px, 50)),
            size=bullet_size, color=_rgb(palette["text"]), is_rtl=is_rtl,
        )

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
