"""
Layout Quality Scorer — analyses a rebuilt PPTX against a style profile.

Computes:
  font_match_pct         — % of text runs using the target font family
  color_match_pct        — % of shapes with fills matching the theme palette
  overflow_count         — shapes whose text likely overflows the placeholder
  objects_adjusted       — placeholder shapes that were repositioned (informational)
  arabic_readability_pct — % of Arabic paragraphs with correct RTL + right alignment
  overall_score          — weighted composite 0-100
"""
from __future__ import annotations

import io
import logging

log = logging.getLogger(__name__)


def _safe(fn, default=None):
    try:
        return fn()
    except Exception:
        return default


def _text_has_arabic(text: str) -> bool:
    return any("\u0600" <= ch <= "\u06FF" for ch in text)


def score_pptx(
    pptx_bytes: bytes,
    style_profile: dict,
    target_lang: str = "ar",
) -> dict:
    """Return a quality score dict.  Never raises — returns partial results on error."""
    result: dict = {
        "font_match_pct": None,
        "color_match_pct": None,
        "overflow_count": 0,
        "objects_adjusted": 0,
        "arabic_readability_pct": None,
        "overall_score": None,
    }
    if not pptx_bytes or not style_profile:
        return result

    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(pptx_bytes))
    except Exception as exc:
        log.warning("layout_quality_scorer: cannot open pptx: %s", exc)
        return result

    is_rtl = target_lang in ("ar", "he", "fa", "ur")

    target_title_font = (style_profile.get("title_font_name") or "").lower()
    target_body_font  = (style_profile.get("body_font_name")  or "").lower()
    theme_colors: set[str] = {
        c.lstrip("#").upper()
        for c in (style_profile.get("theme_colors") or [])
        if c
    }

    total_runs    = 0
    matched_font  = 0
    total_shapes  = 0
    matched_color = 0
    overflow_count = 0
    ar_paras       = 0
    ar_correct     = 0

    try:
        from pptx.oxml.ns import qn as _qn
    except ImportError:
        _qn = None

    for slide in _safe(lambda: list(prs.slides), []):
        for shape in _safe(lambda s=slide: list(s.shapes), []):
            total_shapes += 1

            # ── Text analysis ─────────────────────────────────────────────────
            if _safe(lambda sh=shape: sh.has_text_frame, False):
                tf = _safe(lambda sh=shape: sh.text_frame, None)
                if tf:
                    shape_h = _safe(lambda sh=shape: sh.height, 0) or 0
                    total_chars = sum(len(_safe(lambda p=para: p.text, "") or "") for para in _safe(lambda t=tf: t.paragraphs, []))

                    # Overflow heuristic: very long text in a short box
                    if total_chars > 400 and shape_h < 2_000_000:
                        overflow_count += 1

                    for para in _safe(lambda t=tf: list(t.paragraphs), []):
                        para_text = _safe(lambda p=para: p.text, "") or ""

                        # Arabic readability
                        if is_rtl and _text_has_arabic(para_text):
                            ar_paras += 1
                            if _qn:
                                try:
                                    pPr = para._p.find(_qn("a:pPr"))
                                    rtl_ok  = pPr is not None and pPr.get("rtl") == "1"
                                    algn    = pPr.get("algn") if pPr is not None else None
                                    algn_ok = algn in ("r", "ctr", None)
                                    if rtl_ok and algn_ok:
                                        ar_correct += 1
                                except Exception:
                                    pass

                        for run in _safe(lambda p=para: list(p.runs), []):
                            total_runs += 1
                            fname = _safe(lambda r=run: (r.font.name or "").lower(), "") or ""
                            if fname and (
                                (target_body_font  and target_body_font  in fname)
                                or (target_title_font and target_title_font in fname)
                            ):
                                matched_font += 1

            # ── Fill color match ──────────────────────────────────────────────
            if theme_colors and _qn:
                try:
                    spPr = shape._element.find(_qn("p:spPr"))
                    if spPr is not None:
                        solidFill = spPr.find(f".//{_qn('a:solidFill')}")
                        if solidFill is not None:
                            srgb = solidFill.find(_qn("a:srgbClr"))
                            if srgb is not None:
                                val = (srgb.get("val") or "").upper()
                                if val in theme_colors:
                                    matched_color += 1
                except Exception:
                    pass

    # ── Compute individual metrics ────────────────────────────────────────────
    font_match_pct  = round(matched_font  / total_runs   * 100) if total_runs   > 0 else None
    color_match_pct = round(matched_color / total_shapes * 100) if total_shapes > 0 else None
    ar_readability  = round(ar_correct    / ar_paras     * 100) if ar_paras     > 0 else (100 if is_rtl else None)

    # ── Weighted composite ────────────────────────────────────────────────────
    # font 30% · color 20% · arabic_readability 40% · overflow penalty capped at 10%
    score_parts: list[tuple[float, float]] = []
    if font_match_pct  is not None: score_parts.append((float(font_match_pct),  30.0))
    if color_match_pct is not None: score_parts.append((float(color_match_pct), 20.0))
    if ar_readability  is not None: score_parts.append((float(ar_readability),  40.0))

    overflow_penalty = min(overflow_count * 3, 10)
    if score_parts:
        raw     = sum(s * w for s, w in score_parts) / sum(w for _, w in score_parts)
        overall = max(0, round(raw - overflow_penalty))
    else:
        overall = None

    return {
        "font_match_pct":         font_match_pct,
        "color_match_pct":        color_match_pct,
        "overflow_count":         overflow_count,
        "objects_adjusted":       0,
        "arabic_readability_pct": ar_readability,
        "overall_score":          overall,
    }
