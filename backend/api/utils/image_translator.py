"""
Image text detection, translation, and rendering pipeline.

Pipeline per image:
  1. detect_and_translate_regions() — GPT-4o Vision: detect text bounding
     boxes, extract text, translate to target language in one pass.
  2. render_translated_image()     — PIL: paint translated text over each
     bbox using arabic_reshaper + python-bidi for RTL languages.

Document image extraction:
  extract_document_images()  — pulls all raster images from PDF/DOCX/PPTX.
"""
from __future__ import annotations

import base64
import io
import json
import logging
import os
import uuid
from typing import Any

log = logging.getLogger(__name__)

# ── Font helpers ───────────────────────────────────────────────────────────────

_FONT_CACHE: dict[int, Any] = {}
_ARABIC_FONT_PATH: str | None = None


def _find_arabic_font() -> str | None:
    """Return path to an Arabic-capable TTF font, downloading if necessary."""
    global _ARABIC_FONT_PATH
    if _ARABIC_FONT_PATH:
        return _ARABIC_FONT_PATH

    # Candidates in order of preference
    candidates = [
        "/tmp/NotoNaskhArabic-Regular.ttf",
        "/usr/share/fonts/truetype/noto/NotoNaskhArabic-Regular.ttf",
        "/usr/share/fonts/truetype/noto/NotoSansArabic-Regular.ttf",
        "/usr/share/fonts/opentype/noto/NotoNaskhArabic-Regular.otf",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",  # partial Arabic
    ]
    for p in candidates:
        if os.path.exists(p):
            _ARABIC_FONT_PATH = p
            return p

    # Try to download Noto Naskh Arabic from Google Fonts
    target = "/tmp/NotoNaskhArabic-Regular.ttf"
    try:
        import urllib.request
        url = (
            "https://github.com/googlefonts/noto-fonts/raw/main/hinted/ttf/"
            "NotoNaskhArabic/NotoNaskhArabic-Regular.ttf"
        )
        log.info("Downloading Arabic font from %s", url)
        urllib.request.urlretrieve(url, target)
        _ARABIC_FONT_PATH = target
        return target
    except Exception as e:
        log.warning("Could not download Arabic font: %s — using DejaVu fallback", e)

    # Last resort: DejaVu
    dv = "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"
    if os.path.exists(dv):
        _ARABIC_FONT_PATH = dv
        return dv
    return None


def _get_font(size: int):
    """Return a PIL ImageFont at the given size."""
    from PIL import ImageFont
    if size in _FONT_CACHE:
        return _FONT_CACHE[size]
    path = _find_arabic_font()
    try:
        if path:
            font = ImageFont.truetype(path, size)
        else:
            font = ImageFont.load_default()
    except Exception:
        font = ImageFont.load_default()
    _FONT_CACHE[size] = font
    return font


# ── Arabic text shaping ────────────────────────────────────────────────────────

def _prepare_rtl_text(text: str, target_lang: str) -> str:
    """Shape and reorder Arabic/RTL text for correct PIL rendering."""
    if target_lang not in ("ar", "fa", "he", "ur"):
        return text
    try:
        import arabic_reshaper
        from bidi.algorithm import get_display
        reshaped = arabic_reshaper.reshape(text)
        return get_display(reshaped)
    except Exception as e:
        log.debug("arabic_reshaper/bidi failed: %s", e)
        return text


# ── GPT-4o Vision: detect + translate ─────────────────────────────────────────

_DETECT_SYSTEM = """\
You are an expert technical document analyst and Arabic translator.
Analyze the image and identify all visible text (labels, callouts, titles, annotations,
warnings, captions, measurements, button labels, UI text).

For EACH text region return:
  "id"               : sequential "r1", "r2", … string
  "bbox"             : {"x": %, "y": %, "w": %, "h": %} — top-left origin, percentages of image dims
  "source_text"      : exact text as it appears
  "translated_text"  : professional technical Arabic translation
  "confidence"       : float 0.0–1.0 (OCR confidence)
  "is_technical_code": true if model number / part number / connector ID / measurement value —
                       these must NOT be translated, keep source_text unchanged
  "font_size"        : estimated rendered font size in pixels (integer)
  "font_color"       : estimated hex color of the original text e.g. "#000000"

Rules:
- Do NOT translate: connector references (CN1, J2), part/serial numbers, software commands,
  measurement values (220V, 50Hz, M6x20), filenames, codes, URLs.
- Do NOT invent text that is unreadable — mark confidence < 0.5.
- Translate all other text to professional technical Arabic (engineering MSA register).
- Preserve numbers, units, and symbols unchanged inside the translation.
- Return ONLY a valid JSON object {"regions": [...]}.
"""


async def detect_and_translate_regions(
    image_bytes: bytes,
    source_lang: str,
    target_lang: str,
    style: str,
    client,
) -> list[dict]:
    """
    Call GPT-4o Vision to detect text regions and translate them in one pass.
    Returns a list of region dicts ready for storage.
    """
    b64 = base64.b64encode(image_bytes).decode()

    style_note = {
        "technical": "Use concise technical engineering Arabic (IEC/ISO register).",
        "formal":    "Use Modern Standard Arabic (MSA), professional and formal.",
        "bilingual": "Provide: Arabic translation (English original in parentheses).",
    }.get(style, "Use technical Arabic.")

    user_content = [
        {"type": "text", "text": f"Target language: {target_lang}. {style_note}\n\nAnalyze this image:"},
        {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{b64}", "detail": "high"}},
    ]

    try:
        resp = await client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": _DETECT_SYSTEM},
                {"role": "user",   "content": user_content},
            ],
            max_tokens=4096,
            temperature=0.05,
            response_format={"type": "json_object"},
        )
        try:
            from api.utils.usage_recorder import record_usage_from_response
            record_usage_from_response(
                "Image Translation", resp,
                sub_feature="image_translate_detect",
                meta={"target_lang": target_lang, "style": style},
            )
        except Exception:
            pass
        raw = resp.choices[0].message.content or "{}"
        data = json.loads(raw)
        regions = data.get("regions", [])
    except Exception as e:
        log.error("GPT-4o Vision failed: %s", e)
        return []

    # Normalise and fill defaults
    cleaned: list[dict] = []
    for r in regions:
        if not isinstance(r, dict):
            continue
        rid = str(r.get("id", uuid.uuid4().hex[:8]))
        cleaned.append({
            "id":               rid,
            "bbox":             r.get("bbox", {"x": 0, "y": 0, "w": 10, "h": 5}),
            "source_text":      str(r.get("source_text", "")).strip(),
            "translated_text":  str(r.get("translated_text", "")).strip(),
            "confidence":       float(r.get("confidence", 0.8)),
            "is_technical_code": bool(r.get("is_technical_code", False)),
            "font_size":        int(r.get("font_size", 14)),
            "font_color":       str(r.get("font_color", "#000000")),
            "edited":           False,
            "approved":         False,
            "keep_english":     bool(r.get("is_technical_code", False)),
        })
    return cleaned


# ── PIL rendering ──────────────────────────────────────────────────────────────

def render_translated_image(
    image_bytes: bytes,
    regions: list[dict],
    target_lang: str = "ar",
    bg_color: tuple = (255, 255, 255),
) -> bytes:
    """
    Paint translated Arabic text over each detected text region.

    Steps per region:
      1. Fill original text area with bg_color (erasing source text).
      2. Fit translated text to the bbox width.
      3. Draw RTL text right-aligned inside the bbox.

    Returns PNG bytes of the annotated image.
    """
    from PIL import Image, ImageDraw

    try:
        img = Image.open(io.BytesIO(image_bytes)).convert("RGBA")
    except Exception as e:
        log.error("Cannot open image for rendering: %s", e)
        return image_bytes  # return original unchanged

    W, H = img.size
    overlay = Image.new("RGBA", img.size, (0, 0, 0, 0))
    draw = ImageDraw.Draw(overlay)

    is_rtl = target_lang in ("ar", "fa", "he", "ur")

    for region in regions:
        if not isinstance(region, dict):
            continue
        # Technical codes are preserved unchanged
        if region.get("keep_english") or region.get("is_technical_code"):
            continue

        display_text = region.get("translated_text", "").strip()
        if not display_text:
            continue

        bbox = region.get("bbox", {})
        try:
            x = max(0, int(float(bbox.get("x", 0)) * W / 100))
            y = max(0, int(float(bbox.get("y", 0)) * H / 100))
            bw = max(10, int(float(bbox.get("w", 10)) * W / 100))
            bh = max(10, int(float(bbox.get("h", 5)) * H / 100))
        except (TypeError, ValueError):
            continue

        x2, y2 = min(W, x + bw), min(H, y + bh)

        # Parse hex color
        try:
            fc = region.get("font_color", "#000000").lstrip("#")
            fr, fg, fb = int(fc[0:2], 16), int(fc[2:4], 16), int(fc[4:6], 16)
        except Exception:
            fr, fg, fb = 0, 0, 0

        # Erase original text with a white/background rect
        draw.rectangle([x, y, x2, y2], fill=(bg_color[0], bg_color[1], bg_color[2], 230))

        # Shape Arabic text
        if is_rtl:
            display_text = _prepare_rtl_text(display_text, target_lang)

        # Auto-fit font size to bbox height
        hint_size = int(region.get("font_size", 14))
        font_size = max(8, min(hint_size, bh - 4))
        font = _get_font(font_size)

        # If text overflows width, reduce font size
        try:
            bbox_text = draw.textbbox((0, 0), display_text, font=font)
            text_w = bbox_text[2] - bbox_text[0]
            if text_w > bw and bw > 20:
                ratio = bw / text_w
                font_size = max(8, int(font_size * ratio * 0.95))
                font = _get_font(font_size)
                bbox_text = draw.textbbox((0, 0), display_text, font=font)
                text_w = bbox_text[2] - bbox_text[0]
        except Exception:
            text_w = bw

        # Vertical centering
        try:
            text_h = draw.textbbox((0, 0), display_text, font=font)[3]
        except Exception:
            text_h = font_size
        text_y = y + max(0, (bh - text_h) // 2)

        # Horizontal: right-align for RTL, left-align for LTR
        if is_rtl:
            text_x = x2  # anchor="rt" → right-top
            try:
                draw.text((text_x, text_y), display_text, font=font,
                          fill=(fr, fg, fb, 255), anchor="rt")
            except Exception:
                draw.text((x, text_y), display_text, font=font,
                          fill=(fr, fg, fb, 255))
        else:
            draw.text((x + 2, text_y), display_text, font=font,
                      fill=(fr, fg, fb, 255))

    # Composite overlay onto original
    result = Image.alpha_composite(img, overlay).convert("RGB")
    buf = io.BytesIO()
    result.save(buf, format="PNG", optimize=False)
    return buf.getvalue()


# ── Document image extraction ──────────────────────────────────────────────────

def _bytes_to_png(raw: bytes) -> bytes:
    """Normalise any image bytes to PNG via Pillow."""
    from PIL import Image as _PIL_Image
    try:
        img = _PIL_Image.open(io.BytesIO(raw)).convert("RGB")
        buf = io.BytesIO()
        img.save(buf, "PNG")
        return buf.getvalue()
    except Exception:
        return raw


def extract_document_images(
    file_bytes: bytes,
    file_type: str,
) -> list[dict]:
    """
    Extract all raster images from a document.

    Returns list of:
    {
      "doc_page":    int,   # page/slide index (1-based)
      "doc_type":    str,   # "pdf_page" | "docx_inline" | "pptx_slide"
      "image_index": int,   # nth image on that page (0-based)
      "image_bytes": bytes, # PNG bytes
      "mime":        str,   # "image/png"
      "width":       int,
      "height":      int,
    }
    """
    ft = file_type.lower()
    if ft == "pdf":
        return _extract_pdf_images(file_bytes)
    elif ft == "docx":
        return _extract_docx_images(file_bytes)
    elif ft == "pptx":
        return _extract_pptx_images(file_bytes)
    else:
        return []


def _extract_pdf_images(file_bytes: bytes) -> list[dict]:
    """Extract images from each PDF page via PyMuPDF."""
    try:
        import fitz
    except ImportError:
        log.warning("PyMuPDF not available — skipping PDF image extraction")
        return []

    results = []
    try:
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        for page_num, page in enumerate(doc, start=1):
            img_list = page.get_images(full=True)
            for img_idx, img_info in enumerate(img_list):
                xref = img_info[0]
                try:
                    base_img = doc.extract_image(xref)
                    raw = base_img.get("image", b"")
                    if not raw or len(raw) < 100:
                        continue
                    png = _bytes_to_png(raw)
                    from PIL import Image as _PIL
                    with _PIL.open(io.BytesIO(png)) as im:
                        w, h = im.size
                    results.append({
                        "doc_page":    page_num,
                        "doc_type":    "pdf_page",
                        "image_index": img_idx,
                        "image_bytes": png,
                        "mime":        "image/png",
                        "width":       w,
                        "height":      h,
                    })
                except Exception as e:
                    log.debug("PDF image xref=%d extract error: %s", xref, e)
        doc.close()
    except Exception as e:
        log.error("PDF image extraction failed: %s", e)
    return results


def _extract_docx_images(file_bytes: bytes) -> list[dict]:
    """Extract inline images from a DOCX file."""
    results = []
    try:
        import zipfile
        with zipfile.ZipFile(io.BytesIO(file_bytes)) as z:
            media_files = [n for n in z.namelist() if n.startswith("word/media/")]
            for idx, name in enumerate(sorted(media_files)):
                try:
                    raw = z.read(name)
                    if len(raw) < 100:
                        continue
                    png = _bytes_to_png(raw)
                    from PIL import Image as _PIL
                    with _PIL.open(io.BytesIO(png)) as im:
                        w, h = im.size
                    # Skip tiny icons (< 50px in either dimension)
                    if w < 50 or h < 50:
                        continue
                    results.append({
                        "doc_page":    1,
                        "doc_type":    "docx_inline",
                        "image_index": idx,
                        "image_bytes": png,
                        "mime":        "image/png",
                        "width":       w,
                        "height":      h,
                    })
                except Exception as e:
                    log.debug("DOCX media %s error: %s", name, e)
    except Exception as e:
        log.error("DOCX image extraction failed: %s", e)
    return results


def _extract_pptx_images(file_bytes: bytes) -> list[dict]:
    """Extract images from PPTX slides."""
    results = []
    try:
        from pptx import Presentation
        from pptx.util import Pt
        prs = Presentation(io.BytesIO(file_bytes))
        for slide_idx, slide in enumerate(prs.slides, start=1):
            img_idx = 0
            for shape in slide.shapes:
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    try:
                        raw = shape.image.blob
                        if len(raw) < 100:
                            continue
                        png = _bytes_to_png(raw)
                        from PIL import Image as _PIL
                        with _PIL.open(io.BytesIO(png)) as im:
                            w, h = im.size
                        if w < 50 or h < 50:
                            continue
                        results.append({
                            "doc_page":    slide_idx,
                            "doc_type":    "pptx_slide",
                            "image_index": img_idx,
                            "image_bytes": png,
                            "mime":        "image/png",
                            "width":       w,
                            "height":      h,
                        })
                        img_idx += 1
                    except Exception as e:
                        log.debug("PPTX slide %d shape error: %s", slide_idx, e)
    except Exception as e:
        log.error("PPTX image extraction failed: %s", e)
    return results


# ── Image-level quality checks ─────────────────────────────────────────────────

def check_image_quality(regions: list[dict], W: int, H: int) -> list[dict]:
    """
    Run automatic quality checks on translated regions.
    Returns list of issue dicts: {region_id, type, severity, message}.
    """
    issues: list[dict] = []
    for r in regions:
        if r.get("is_technical_code") or r.get("keep_english"):
            continue
        rid = r.get("id", "?")

        if not r.get("translated_text", "").strip():
            issues.append({
                "region_id": rid,
                "type": "missing_translation",
                "severity": "error",
                "message": f'Region "{r.get("source_text","")[:40]}" was not translated.',
            })

        conf = float(r.get("confidence", 1.0))
        if conf < 0.5:
            issues.append({
                "region_id": rid,
                "type": "low_ocr_confidence",
                "severity": "warning",
                "message": f'Low OCR confidence ({conf:.0%}) — verify: "{r.get("source_text","")[:40]}"',
            })

        # Check if bbox is within image bounds
        bbox = r.get("bbox", {})
        try:
            rx = float(bbox.get("x", 0))
            ry = float(bbox.get("y", 0))
            rw = float(bbox.get("w", 10))
            rh = float(bbox.get("h", 5))
            if rx + rw > 105 or ry + rh > 105:
                issues.append({
                    "region_id": rid,
                    "type": "out_of_bounds",
                    "severity": "warning",
                    "message": f'Region "{r.get("source_text","")[:30]}" bbox extends outside image.',
                })
        except (TypeError, ValueError):
            pass

        # Check number preservation
        import re
        src_nums = set(re.findall(r"\b\d+(?:[.,]\d+)*\b", r.get("source_text", "")))
        tgt_nums = set(re.findall(r"\b\d+(?:[.,]\d+)*\b", r.get("translated_text", "")))
        missing = src_nums - tgt_nums
        if missing:
            issues.append({
                "region_id": rid,
                "type": "number_mismatch",
                "severity": "warning",
                "message": f'Numbers missing in translation: {", ".join(sorted(missing))}',
            })

    return issues
