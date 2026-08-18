"""
Professional PPTX generator — Rapiscan Training Style.

Reverse-engineered from the LZBV ZBV Level 3 Field Service Course reference deck.
All dimensions, colors, and layout positions faithfully match the reference.

Slide format: 10.00" × 7.50"  (4:3 standard — matches reference deck exactly)

Layout language:
  ┌─────────────────────────────────────────────────────────────────┐
  │ WHITE title bar (0.38,0.38) 7.40×0.84 │  branding top-right    │  ← row 1
  ├─────────────────────────────────────────────────────────────────┤
  │  DARK GRAY content area (0.38,1.33) 9.38×5.43 — white text     │  ← rows 2-8
  ├─────────────────────────────────────────────────────────────────┤
  │  slide background (dark charcoal)  │  date  │  slide num        │  ← footer
  └─────────────────────────────────────────────────────────────────┘

Color palette (extracted from reference):
  FFFFFF  — title bar bg, white text on dark
  303030  — content area bg (dark gray)
  141414  — slide background
  000000  — alternate content area (image slides)
  CC0000  — warnings / danger
  FFFFFF  — all body text on dark backgrounds
"""
from __future__ import annotations
import io
import re
from datetime import datetime, timezone
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree


# ── Color palette ──────────────────────────────────────────────────────────────
# Extracted from reference LZBV ZBV Level 3 Field Service Course

SLIDE_BG        = RGBColor(0xFF, 0xFF, 0xFF)  # white background
TITLE_BAR_BG    = RGBColor(0x0F, 0x2A, 0x4D)  # dark navy title bar
TITLE_TEXT      = RGBColor(0xFF, 0xFF, 0xFF)  # white text on navy bar
CONTENT_BG      = RGBColor(0xF3, 0xF5, 0xF8)  # very light gray-blue content area
CONTENT_TEXT    = RGBColor(0x1A, 0x1A, 0x1A)  # near-black body text on light background
SUBTEXT         = RGBColor(0x44, 0x44, 0x44)  # dark gray for sub-bullets
CITATION        = RGBColor(0x70, 0x70, 0x70)  # muted gray for (p.N) references
FOOTER_TEXT     = RGBColor(0x70, 0x70, 0x70)  # date + slide number

# Accent colors for special slide types
SECTION_BG      = RGBColor(0x0F, 0x2A, 0x4D)  # dark navy for section dividers
SECTION_RULE    = RGBColor(0x4A, 0x90, 0xD9)  # horizontal rule on section slides
SECTION_LABEL   = RGBColor(0x8A, 0xC0, 0xEE)  # "SECTION N" label (light blue on navy)

KQ_BANNER_BG    = RGBColor(0x1A, 0x4A, 0x80)  # knowledge-check header strip
KQ_BANNER_TEXT  = RGBColor(0xFF, 0xFF, 0xFF)
KQ_OPT_BG       = RGBColor(0xF0, 0xF2, 0xF6)  # answer option box (light)
KQ_OPT_BORDER   = RGBColor(0xC7, 0xCE, 0xD8)
KQ_CORRECT_BG   = RGBColor(0x22, 0x8B, 0x4E)  # correct answer (instructor)
KQ_CORRECT_TEXT = RGBColor(0xFF, 0xFF, 0xFF)

PRAC_BANNER_BG  = RGBColor(0x9A, 0x44, 0x00)  # practical-exercise header strip
PRAC_BANNER_TEXT= RGBColor(0xFF, 0xFF, 0xFF)
PRAC_STEP_BG    = RGBColor(0xF3, 0xF0, 0xEC)  # step box (light warm gray)
PRAC_NUM_BG     = RGBColor(0x9A, 0x44, 0x00)  # step-number circle bg

WARN_BG         = RGBColor(0x8B, 0x00, 0x00)   # WARNING box
WARN_TEXT       = RGBColor(0xFF, 0xFF, 0xFF)

OBJ_NUM_BG      = RGBColor(0x1A, 0x4A, 0x80)  # objective number box
OBJ_NUM_TEXT    = RGBColor(0xFF, 0xFF, 0xFF)
WHITE           = RGBColor(0xFF, 0xFF, 0xFF)
BLACK           = RGBColor(0x00, 0x00, 0x00)
COVER_BAND_BG   = RGBColor(0x0F, 0x2A, 0x4D)  # dark navy band behind the cover title (white text)


# ── Layout constants (10" × 7.5" — matches reference exactly) ─────────────────
W      = Inches(10.0)
H      = Inches(7.5)

# Title bar — white rectangle, top-left (reference positions)
TB_L   = Inches(0.38)
TB_T   = Inches(0.38)
TB_W   = Inches(7.40)   # intentionally less than full width; right area = branding
TB_H   = Inches(0.84)

# Right-side branding strip (top-right of title row)
BR_L   = Inches(7.90)   # 0.38 + 7.40 + 0.12 gap
BR_T   = Inches(0.38)
BR_W   = Inches(1.92)   # fills to right margin
BR_H   = Inches(0.84)

# Content area — dark gray rectangle (reference positions)
CA_L   = Inches(0.38)
CA_T   = Inches(1.33)
CA_W   = Inches(9.38)   # wider than title bar
CA_H   = Inches(5.43)

# Inner text margins within content area
CT_PAD_X = Inches(0.25)   # horizontal text padding inside content rect
CT_PAD_Y = Inches(0.18)   # vertical text padding inside content rect
CT_L     = CA_L + CT_PAD_X
CT_T     = CA_T + CT_PAD_Y
CT_W     = CA_W - CT_PAD_X * 2
CT_H     = CA_H - CT_PAD_Y * 2

# Footer positions (below content area, on bare slide background)
DATE_L   = Inches(0.38)
DATE_T   = Inches(6.82)
DATE_W   = Inches(3.0)
DATE_H   = Inches(0.35)

SNUM_L   = Inches(9.10)
SNUM_T   = Inches(7.08)
SNUM_W   = Inches(0.75)
SNUM_H   = Inches(0.28)


# ── Primitive helpers ──────────────────────────────────────────────────────────

def _fill(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def _no_line(shape) -> None:
    shape.line.fill.background()


def _line_color(shape, color: RGBColor, width_pt: float = 0.75) -> None:
    shape.line.color.rgb = color
    shape.line.width = Pt(width_pt)


def _set_bg(slide, color: RGBColor) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _rect(slide, left, top, width, height,
          fill: RGBColor, border: RGBColor | None = None, border_pt: float = 0.75):
    sh = slide.shapes.add_shape(1, left, top, width, height)
    _fill(sh, fill)
    if border:
        _line_color(sh, border, border_pt)
    else:
        _no_line(sh)
    return sh


def _text_box(slide, text: str, left, top, width, height,
              size: int = 18, bold: bool = False, italic: bool = False,
              color: RGBColor = WHITE, align=PP_ALIGN.LEFT,
              wrap: bool = True, font: str = "Calibri") -> Any:
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = font
    return tb


def _bullet_frame(slide, bullets: list[str], left, top, width, height,
                  size: int = 18, color: RGBColor = CONTENT_TEXT,
                  sub_size: int = 16, sub_color: RGBColor = SUBTEXT,
                  font: str = "Calibri") -> Any:
    """
    Add a bulleted text frame. Bullets starting with spaces are treated as
    sub-bullets (level-1, smaller, dimmer).
    Citation patterns like (p.N) are rendered in muted gray.
    """
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True

    for raw in bullets:
        is_sub = raw.startswith("  ") or raw.startswith("\t")
        text = raw.strip()
        if not text:
            continue

        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()

        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(3)

        # Choose bullet character and size
        if is_sub:
            bullet_char = "  –  "
            fsize = sub_size
            fcolor = sub_color
        else:
            bullet_char = "•  "
            fsize = size
            fcolor = color

        # Check for citation pattern — split and dim the citation portion
        cite_match = re.search(r'(\s*\(p\.[\d,\s]+\))', text)
        if cite_match:
            main_part = text[:cite_match.start()]
            cite_part = text[cite_match.start():]
            r = p.add_run()
            r.text = bullet_char + main_part
            r.font.size = Pt(fsize)
            r.font.color.rgb = fcolor
            r.font.name = font
            r2 = p.add_run()
            r2.text = cite_part
            r2.font.size = Pt(fsize - 2)
            r2.font.color.rgb = CITATION
            r2.font.name = font
        else:
            r = p.add_run()
            r.text = bullet_char + text
            r.font.size = Pt(fsize)
            r.font.color.rgb = fcolor
            r.font.name = font

    return tb


def _notes(slide, text: str) -> None:
    if text:
        slide.notes_slide.notes_text_frame.text = text


# ── Shared chrome elements ─────────────────────────────────────────────────────

def _title_bar(slide, title: str) -> None:
    """White title bar with dark bold text — matches reference Rectangle 2."""
    _rect(slide, TB_L, TB_T, TB_W, TB_H, TITLE_BAR_BG)
    _text_box(slide, title,
              TB_L + Inches(0.18), TB_T + Pt(6),
              TB_W - Inches(0.36), TB_H - Pt(12),
              size=26, bold=True, color=TITLE_TEXT,
              align=PP_ALIGN.LEFT)


def _branding_strip(slide, label: str) -> None:
    """Top-right strip — course/equipment branding where the master logo lives."""
    _rect(slide, BR_L, BR_T, BR_W, BR_H, CONTENT_BG)
    _text_box(slide, label,
              BR_L + Pt(4), BR_T + Pt(4), BR_W - Pt(8), BR_H - Pt(8),
              size=9, bold=False, color=SUBTEXT,
              align=PP_ALIGN.CENTER, wrap=True)


def _content_rect(slide) -> None:
    """Dark gray content rectangle — matches reference Rectangle 3."""
    _rect(slide, CA_L, CA_T, CA_W, CA_H, CONTENT_BG)


def _footer(slide, course_title: str, slide_num: int, total: int,
            source_pages: list[int] | None = None) -> None:
    """Date left + manual citation + slide number right, in footer zone."""
    now_utc = datetime.now(timezone.utc)
    # %-d is not supported on Windows strftime; build a portable date string.
    date_str = f"{now_utc.day} {now_utc.strftime('%B %Y')}"
    cite_str = ""
    if source_pages:
        pages = ", ".join(str(p) for p in sorted(set(source_pages))[:5])
        cite_str = f"  |  Manual p.{pages}"

    footer_text = f"{date_str}{cite_str}"
    _text_box(slide, footer_text,
              DATE_L, DATE_T, DATE_W + Inches(3), DATE_H,
              size=9, color=FOOTER_TEXT, align=PP_ALIGN.LEFT)

    _text_box(slide, f"{slide_num} / {total}",
              SNUM_L, SNUM_T, SNUM_W, SNUM_H,
              size=9, color=FOOTER_TEXT, align=PP_ALIGN.RIGHT)


def _standard_chrome(slide, title: str, course_title: str,
                     slide_num: int, total: int,
                     source_pages: list[int] | None = None) -> None:
    """Full chrome: slide bg + title bar + branding + content rect + footer."""
    _set_bg(slide, SLIDE_BG)
    _title_bar(slide, title)
    _branding_strip(slide, course_title)
    _content_rect(slide)
    _footer(slide, course_title, slide_num, total, source_pages)


# ── Slide builders ─────────────────────────────────────────────────────────────

def _title_slide(prs: Presentation, s: dict, project: dict, num: int, total: int) -> None:
    """
    Cover slide. Large title in content area, equipment/audience details below.
    Matches reference S1 (Title Slide layout).
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, SLIDE_BG)

    # Full-width dark band at top (title area)
    _rect(slide, Inches(0), Inches(0), W, Inches(2.80), COVER_BAND_BG)

    # Horizontal rule below band
    _rect(slide, Inches(0), Inches(2.80), W, Pt(3), SECTION_RULE)

    # Course / equipment label in top band
    equip = f"{project.get('manufacturer','')} {project.get('equipment_model','')}".strip()
    if equip:
        _text_box(slide, equip.upper(),
                  TB_L, Inches(0.42), W - TB_L * 2, Inches(0.50),
                  size=13, bold=True, color=SECTION_LABEL,
                  align=PP_ALIGN.LEFT)

    # Main title
    title = s.get("title", project.get("course_title", "Training Course"))
    _text_box(slide, title,
              TB_L, Inches(1.00), W - TB_L * 2, Inches(1.60),
              size=36, bold=True, color=WHITE,
              align=PP_ALIGN.LEFT)

    # Details block (audience, type, date)
    bullets = s.get("bullets", [])
    if not bullets:
        bullets = [
            project.get("training_type", ""),
            f"Audience: {project.get('audience', '')}",
            datetime.now(timezone.utc).strftime("%B %Y"),
        ]
    details = "\n".join(b for b in bullets if b)
    _text_box(slide, details,
              TB_L, Inches(3.20), W - TB_L * 2, Inches(1.60),
              size=18, color=SUBTEXT, align=PP_ALIGN.LEFT)

    # Export control / confidentiality label (matches reference)
    _text_box(slide, "PROPRIETARY — FOR AUTHORIZED PERSONNEL ONLY",
              TB_L, H - Inches(0.55), W - TB_L * 2, Inches(0.40),
              size=10, bold=True, color=FOOTER_TEXT, align=PP_ALIGN.LEFT)

    _footer(slide, project.get("course_title", ""), num, total)
    _notes(slide, s.get("speaker_notes", ""))


def _agenda_slide(prs: Presentation, s: dict, project: dict, num: int, total: int) -> None:
    """Course agenda — two-column layout for longer lists."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _standard_chrome(slide, s.get("title", "Course Agenda"),
                     project.get("course_title", ""), num, total)

    bullets = s.get("bullets", [])
    mid = (len(bullets) + 1) // 2
    col_w = (CT_W - Inches(0.4)) / 2

    _bullet_frame(slide, bullets[:mid],
                  CT_L, CT_T, col_w, CT_H, size=17)
    if bullets[mid:]:
        _bullet_frame(slide, bullets[mid:],
                      CT_L + col_w + Inches(0.4), CT_T, col_w, CT_H, size=17)

    _notes(slide, s.get("speaker_notes", ""))


def _section_slide(prs: Presentation, s: dict, project: dict, num: int, total: int) -> None:
    """
    Section divider. Dark navy background, centred large title, rule above.
    Matches reference S4/S13 (Section Title layout).
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, SECTION_BG)

    # Horizontal rule above title
    _rect(slide, Inches(0.50), Inches(2.90), Inches(9.0), Pt(2), SECTION_RULE)

    # Section number label
    sec_num = s.get("section_num", "")
    if sec_num:
        _text_box(slide, f"SECTION  {sec_num}",
                  Inches(0.79), Inches(2.45), Inches(8.50), Inches(0.45),
                  size=12, bold=True, italic=True,
                  color=SECTION_LABEL, align=PP_ALIGN.LEFT)

    # Section title (matches reference: text at (0.79,3.18) 8.50×1.07)
    _text_box(slide, s.get("title", ""),
              Inches(0.79), Inches(3.10), Inches(8.50), Inches(1.40),
              size=38, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # Section intro (first bullet, if present)
    desc = (s.get("bullets") or [""])
    if desc[0]:
        _text_box(slide, desc[0],
                  Inches(0.79), Inches(4.80), Inches(8.50), Inches(0.80),
                  size=16, italic=True, color=SUBTEXT, align=PP_ALIGN.LEFT)

    # Footer
    _footer(slide, project.get("course_title", ""), num, total)
    _notes(slide, s.get("speaker_notes", ""))


def _objectives_slide(prs: Presentation, s: dict, project: dict, num: int, total: int) -> None:
    """
    Lesson Objective slide.
    Matches reference S5 (Title and Content layout with Lesson Objective).
    Opener: "By the end of this lesson, utilizing student guide materials,
             you will be able to [verb]:"
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _standard_chrome(slide, s.get("title", "Lesson Objective"),
                     project.get("course_title", ""), num, total,
                     s.get("source_pages"))

    bullets = s.get("bullets", [])

    # Opener sentence (first bullet if it looks like an opener, otherwise add one)
    opener = ""
    obj_bullets = bullets
    if bullets and any(kw in bullets[0].lower() for kw in
                       ["by the end", "you will be able", "able to"]):
        opener = bullets[0]
        obj_bullets = bullets[1:]
    else:
        opener = ("By the end of this lesson, utilizing student guide materials, "
                  "you will be able to:")

    # Opener text (white, slightly smaller)
    _text_box(slide, opener,
              CT_L, CT_T, CT_W, Inches(0.65),
              size=15, italic=True, color=SUBTEXT, align=PP_ALIGN.LEFT)

    # Numbered objectives with blue number badges
    top = CT_T + Inches(0.72)
    for i, bullet in enumerate(obj_bullets[:8]):
        if top + Inches(0.55) > CA_T + CA_H - Inches(0.15):
            break
        badge_size = Inches(0.36)
        # Number badge (blue box)
        _rect(slide, CT_L, top + Pt(2), badge_size, badge_size, OBJ_NUM_BG)
        _text_box(slide, str(i + 1),
                  CT_L, top + Pt(2), badge_size, badge_size,
                  size=13, bold=True, color=OBJ_NUM_TEXT, align=PP_ALIGN.CENTER)
        # Objective text
        _text_box(slide, bullet.strip().lstrip("•- "),
                  CT_L + badge_size + Inches(0.12), top,
                  CT_W - badge_size - Inches(0.12), Inches(0.55),
                  size=16, color=CONTENT_TEXT, align=PP_ALIGN.LEFT)
        top += Inches(0.60)

    _notes(slide, s.get("speaker_notes", ""))


def _content_slide(prs: Presentation, s: dict, project: dict, num: int, total: int) -> None:
    """
    Standard content slide. White title bar + dark content area.
    Matches reference S2/S3/S5 (Title and Content layout).
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _standard_chrome(slide, s.get("title", ""),
                     project.get("course_title", ""), num, total,
                     s.get("source_pages"))

    bullets = s.get("bullets", [])
    if bullets:
        _bullet_frame(slide, bullets,
                      CT_L, CT_T, CT_W, CT_H,
                      size=18, sub_size=16)

    _notes(slide, s.get("speaker_notes", ""))


def _quiz_slide(prs: Presentation, s: dict, project: dict, num: int, total: int,
                show_answers: bool = False) -> None:
    """
    Knowledge Check slide.
    Banner header "KNOWLEDGE CHECK" in blue, question text, then A/B/C/D option boxes.
    Instructor version highlights the correct answer green.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, SLIDE_BG)

    # White title bar — "Knowledge Check" as the title
    _rect(slide, TB_L, TB_T, TB_W, TB_H, TITLE_BAR_BG)
    _text_box(slide, "KNOWLEDGE CHECK",
              TB_L + Inches(0.18), TB_T + Pt(6),
              TB_W - Inches(0.36), TB_H - Pt(12),
              size=24, bold=True, color=TITLE_TEXT, align=PP_ALIGN.LEFT)

    # KQ accent strip under title bar
    _rect(slide, TB_L, TB_T + TB_H, TB_W, Pt(4), KQ_BANNER_BG)

    _branding_strip(slide, project.get("course_title", ""))

    # Content rectangle
    _content_rect(slide)

    # Question text (top of content area)
    question = s.get("title", "Question")
    _text_box(slide, question,
              CT_L, CT_T, CT_W, Inches(0.90),
              size=20, bold=True, color=CONTENT_TEXT, align=PP_ALIGN.LEFT)

    # Answer option boxes
    options = s.get("bullets", [])
    answer_idx = s.get("answer_index", -1)
    top = CT_T + Inches(1.05)
    box_h = Inches(0.64)
    gap = Inches(0.08)

    for i, opt in enumerate(options[:6]):
        label = chr(65 + i)  # A, B, C …
        is_correct = show_answers and (i == answer_idx)
        bg   = KQ_CORRECT_BG   if is_correct else KQ_OPT_BG
        txt  = KQ_CORRECT_TEXT if is_correct else CONTENT_TEXT
        border = KQ_BANNER_BG if is_correct else KQ_OPT_BORDER

        # Strip ONLY leading option labels (A. B) C. etc.) — never bare letters
        # Using regex prevents stripping first letter of answers like "Curie" or "Continue"
        opt_clean = re.sub(r'^[A-Da-d]\s*[.)]\s*', '', opt.strip())

        _rect(slide, CT_L, top, CT_W, box_h, bg, border, 0.75)
        _text_box(slide, f"  {label}.   {opt_clean}",
                  CT_L + Pt(4), top + Pt(6),
                  CT_W - Pt(8), box_h - Pt(8),
                  size=16, bold=is_correct, color=txt, align=PP_ALIGN.LEFT)
        top += box_h + gap

    _footer(slide, project.get("course_title", ""), num, total, s.get("source_pages"))
    _notes(slide, s.get("speaker_notes", ""))


def _practical_slide(prs: Presentation, s: dict, project: dict, num: int, total: int) -> None:
    """
    Practical Exercise slide.
    Orange banner, numbered step boxes in the content area.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, SLIDE_BG)

    # White title bar
    _rect(slide, TB_L, TB_T, TB_W, TB_H, TITLE_BAR_BG)
    _text_box(slide, "PRACTICAL EXERCISE",
              TB_L + Inches(0.18), TB_T + Pt(6),
              TB_W - Inches(0.36), TB_H - Pt(12),
              size=24, bold=True, color=TITLE_TEXT, align=PP_ALIGN.LEFT)

    # Orange accent strip
    _rect(slide, TB_L, TB_T + TB_H, TB_W, Pt(4), PRAC_BANNER_BG)

    _branding_strip(slide, project.get("course_title", ""))
    _content_rect(slide)

    # Practical title / task description
    task_title = s.get("title", "Practical Exercise")
    _text_box(slide, task_title,
              CT_L, CT_T, CT_W, Inches(0.60),
              size=18, bold=True, color=CONTENT_TEXT, align=PP_ALIGN.LEFT)

    # Step boxes
    steps = s.get("bullets", [])
    top = CT_T + Inches(0.70)
    step_h = Inches(0.58)
    num_w = Inches(0.44)
    gap = Pt(5)

    for i, step in enumerate(steps[:7]):
        if top + step_h > CA_T + CA_H - Inches(0.15):
            break
        # Step number box
        _rect(slide, CT_L, top, num_w, step_h, PRAC_NUM_BG)
        _text_box(slide, str(i + 1),
                  CT_L, top, num_w, step_h,
                  size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # Step text
        _text_box(slide, step.strip().lstrip("0123456789). "),
                  CT_L + num_w + Inches(0.10), top + Pt(4),
                  CT_W - num_w - Inches(0.10), step_h - Pt(4),
                  size=15, color=CONTENT_TEXT, align=PP_ALIGN.LEFT)
        top += step_h + gap

    _footer(slide, project.get("course_title", ""), num, total, s.get("source_pages"))
    _notes(slide, s.get("speaker_notes", ""))


def _summary_slide(prs: Presentation, s: dict, project: dict, num: int, total: int) -> None:
    """Course summary — content slide with a 'COURSE SUMMARY' blue accent strip."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _standard_chrome(slide, s.get("title", "Course Summary"),
                     project.get("course_title", ""), num, total)

    # Thin blue accent line below title bar
    _rect(slide, TB_L, TB_T + TB_H, TB_W, Pt(4), KQ_BANNER_BG)

    bullets = s.get("bullets", [])
    if bullets:
        _bullet_frame(slide, bullets, CT_L, CT_T, CT_W, CT_H, size=17)

    _notes(slide, s.get("speaker_notes", ""))


def _references_slide(prs: Presentation, s: dict, project: dict, num: int, total: int) -> None:
    """References & Sources slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _standard_chrome(slide, s.get("title", "References & Sources"),
                     project.get("course_title", ""), num, total)

    bullets = s.get("bullets", [])
    mfn = project.get("manual_filename", "Uploaded Manual")
    if not bullets:
        bullets = [f"[Manual]  {mfn}"]

    _bullet_frame(slide, bullets, CT_L, CT_T, CT_W, CT_H, size=16, color=SUBTEXT)
    _notes(slide, s.get("speaker_notes", ""))


def _tracker_slide(prs: Presentation, s: dict, project: dict, num: int, total: int) -> None:
    """
    Operating Sequence Tracker slide.
    Lists all steps in the system's operating sequence; the current step is
    highlighted in accent blue. Used after every step-section divider.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, SECTION_BG)

    # Title bar — slightly different from section (dark, not white)
    _rect(slide, TB_L, TB_T, TB_W, TB_H, RGBColor(0x10, 0x18, 0x2C))
    _text_box(slide, s.get("title", "System Operating Sequence"),
              TB_L + Inches(0.18), TB_T + Pt(6),
              TB_W - Inches(0.36), TB_H - Pt(12),
              size=22, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # Thin accent rule under title bar
    _rect(slide, TB_L, TB_T + TB_H, TB_W, Pt(3), SECTION_RULE)
    _branding_strip(slide, project.get("course_title", ""))
    _footer(slide, project.get("course_title", ""), num, total)

    steps = s.get("bullets", [])
    n = len(steps)
    if n == 0:
        _notes(slide, s.get("speaker_notes", ""))
        return

    # Layout: single column (≤8 steps) or two columns (>8)
    two_col = n > 8
    rows = (n + 1) // 2 if two_col else n
    available_h = CA_H - Inches(0.20)
    row_h = min(Inches(0.56), available_h / max(rows, 1))
    col_w = (CA_W - Inches(0.40)) / 2 if two_col else CA_W - Inches(0.20)
    NUM_W = Inches(0.38)
    GAP   = Pt(4)

    import re as _re
    for idx, step_text in enumerate(steps):
        col = (idx // rows) if two_col else 0
        row = idx % rows

        x = CA_L + Inches(0.10) + col * (col_w + Inches(0.20))
        y = CA_T + Inches(0.10) + row * (row_h + GAP)

        is_current = "→" in step_text or "CURRENT" in step_text.upper()
        clean = _re.sub(r'→|←\s*CURRENT', '', step_text).strip()

        row_bg   = SECTION_RULE         if is_current else RGBColor(0x1A, 0x24, 0x3C)
        num_bg   = RGBColor(0x00, 0x50, 0xA0) if is_current else RGBColor(0x0A, 0x14, 0x28)
        txt_col  = WHITE

        _rect(slide, x, y, col_w - Inches(0.05), row_h, row_bg)
        _rect(slide, x, y, NUM_W, row_h, num_bg)

        m = _re.match(r'Step\s*(\d+)', clean, _re.I)
        lbl = m.group(1) if m else str(idx + 1)
        _text_box(slide, lbl,
                  x, y, NUM_W, row_h,
                  size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        _text_box(slide, clean,
                  x + NUM_W + Pt(5), y + Pt(4),
                  col_w - NUM_W - Inches(0.10), row_h - Pt(4),
                  size=11 if two_col else 13,
                  bold=is_current, color=txt_col, align=PP_ALIGN.LEFT)

    _notes(slide, s.get("speaker_notes", ""))


def _lab_slide(prs: Presentation, s: dict, project: dict, num: int, total: int) -> None:
    """
    Laboratory Exercise slide — orange chrome (matches Practical), activity checklist.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, SLIDE_BG)

    # White title bar
    _rect(slide, TB_L, TB_T, TB_W, TB_H, TITLE_BAR_BG)
    _text_box(slide, "LABORATORY EXERCISE",
              TB_L + Inches(0.18), TB_T + Pt(6),
              TB_W - Inches(0.36), TB_H - Pt(12),
              size=24, bold=True, color=TITLE_TEXT, align=PP_ALIGN.LEFT)

    # Orange accent strip
    _rect(slide, TB_L, TB_T + TB_H, TB_W, Pt(4), PRAC_BANNER_BG)
    _branding_strip(slide, project.get("course_title", ""))
    _content_rect(slide)

    # Lab title
    lab_title = s.get("title", "Laboratory Exercise")
    _text_box(slide, lab_title,
              CT_L, CT_T, CT_W, Inches(0.62),
              size=19, bold=True, color=CONTENT_TEXT, align=PP_ALIGN.LEFT)

    # Thin divider
    _rect(slide, CT_L, CT_T + Inches(0.66), CT_W, Pt(1), RGBColor(0xC7, 0xCE, 0xD8))

    activities = s.get("bullets", [])
    ref_items  = [a for a in activities if
                  a.lower().startswith("reference:") or a.lower().startswith("ref:")]
    main_items = [a for a in activities if a not in ref_items]

    top    = CT_T + Inches(0.78)
    act_h  = Inches(0.52)
    chk_w  = Inches(0.34)
    gap    = Pt(5)

    for i, act in enumerate(main_items[:7]):
        if top + act_h > CA_T + CA_H - Inches(0.48):
            break
        # Checkmark badge
        _rect(slide, CT_L, top + Pt(2), chk_w, act_h - Pt(4), PRAC_NUM_BG)
        _text_box(slide, "\u2713",
                  CT_L, top + Pt(2), chk_w, act_h - Pt(4),
                  size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # Activity text
        _text_box(slide, act.strip(),
                  CT_L + chk_w + Inches(0.10), top + Pt(4),
                  CT_W - chk_w - Inches(0.10), act_h - Pt(4),
                  size=15, color=CONTENT_TEXT, align=PP_ALIGN.LEFT)
        top += act_h + gap

    # Reference handouts in footnote area
    if ref_items:
        ref_text = "  \u2014  ".join(r.strip() for r in ref_items)
        _text_box(slide, ref_text,
                  CT_L, CA_T + CA_H - Inches(0.38), CT_W, Inches(0.34),
                  size=11, italic=True, color=CITATION, align=PP_ALIGN.LEFT)

    _footer(slide, project.get("course_title", ""), num, total, s.get("source_pages"))
    _notes(slide, s.get("speaker_notes", ""))


def _image_content_slide(
    prs: Presentation, s: dict, project: dict, num: int, total: int,
    images: dict[str, bytes] | None = None,
) -> None:
    """
    Split-layout slide: figure on the left (~55 %), bullet points on the right.
    Falls back to a standard content slide when the referenced image is missing.
    """
    image_tag = s.get("image_tag")
    img_bytes  = (images or {}).get(image_tag or "") if image_tag else None

    if not img_bytes:
        # No image available — render as plain content slide
        _content_slide(prs, s, project, num, total)
        return

    import io as _io
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, SLIDE_BG)
    _title_bar(slide, s.get("title", ""))
    _branding_strip(slide, project.get("course_title", ""))
    _footer(slide, project.get("course_title", ""), num, total, s.get("source_pages"))

    # Dark content background (full width, same as normal slides)
    _rect(slide, CA_L, CA_T, CA_W, CA_H, CONTENT_BG)

    # ── Image column (left 54 % of content area) ──────────────────────────────
    IMG_W = CA_W * 0.54
    IMG_H = CA_H - Inches(0.12)
    img_stream = _io.BytesIO(img_bytes)
    try:
        slide.shapes.add_picture(
            img_stream,
            CA_L + Inches(0.10), CA_T + Inches(0.06),
            IMG_W - Inches(0.10), IMG_H,
        )
    except Exception:
        pass  # malformed image — content still renders

    # ── Vertical divider ──────────────────────────────────────────────────────
    DIV_L = CA_L + IMG_W
    _rect(slide, DIV_L, CA_T + Inches(0.12), Pt(2), CA_H - Inches(0.24),
          RGBColor(0x60, 0x60, 0x60))

    # ── Text column (right 46 %) ──────────────────────────────────────────────
    TEXT_L = DIV_L + Inches(0.18)
    TEXT_W = CA_L + CA_W - TEXT_L - Inches(0.10)
    bullets = s.get("bullets", [])
    if bullets:
        _bullet_frame(slide, bullets,
                      TEXT_L, CA_T + Inches(0.12),
                      TEXT_W, CA_H - Inches(0.24),
                      size=15, sub_size=13)

    # Caption under image (image tag as dim label)
    if image_tag:
        _text_box(slide, image_tag.lower().replace("_", " "),
                  CA_L + Inches(0.12), CA_T + IMG_H - Inches(0.02),
                  IMG_W - Inches(0.14), Inches(0.28),
                  size=9, color=CITATION, align=PP_ALIGN.LEFT)

    _notes(slide, s.get("speaker_notes", ""))


def _full_image_slide(
    prs: Presentation, s: dict, project: dict, num: int, total: int,
    images: dict[str, bytes] | None = None,
) -> None:
    """
    Full content-area image slide with title bar and caption bar at the bottom.
    Falls back to content slide when image is missing.
    """
    image_tag = s.get("image_tag")
    img_bytes  = (images or {}).get(image_tag or "") if image_tag else None

    if not img_bytes:
        _content_slide(prs, s, project, num, total)
        return

    import io as _io
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, SLIDE_BG)
    _title_bar(slide, s.get("title", ""))
    _branding_strip(slide, project.get("course_title", ""))
    _footer(slide, project.get("course_title", ""), num, total, s.get("source_pages"))

    # Dark frame for image
    _rect(slide, CA_L, CA_T, CA_W, CA_H, RGBColor(0x00, 0x00, 0x00))

    # Caption bar at very bottom of content area
    CAPTION_H = Inches(0.42)
    IMG_H     = CA_H - CAPTION_H - Inches(0.04)

    img_stream = _io.BytesIO(img_bytes)
    try:
        slide.shapes.add_picture(
            img_stream,
            CA_L + Inches(0.04), CA_T + Inches(0.04),
            CA_W - Inches(0.08), IMG_H,
        )
    except Exception:
        pass

    # Caption
    caption_text = (s.get("bullets") or [""])[0]
    if caption_text:
        _rect(slide, CA_L, CA_T + IMG_H + Inches(0.02), CA_W, CAPTION_H,
              RGBColor(0x1A, 0x1A, 0x1A))
        _text_box(slide, caption_text,
                  CA_L + Inches(0.15), CA_T + IMG_H + Inches(0.06),
                  CA_W - Inches(0.30), CAPTION_H - Inches(0.08),
                  size=13, color=WHITE, align=PP_ALIGN.LEFT)

    _notes(slide, s.get("speaker_notes", ""))


# ── Slide type registry ────────────────────────────────────────────────────────

SLIDE_BUILDERS = {
    "title":         _title_slide,
    "agenda":        _agenda_slide,
    "section":       _section_slide,
    "objectives":    _objectives_slide,
    "content":       _content_slide,
    "image_content": _image_content_slide,
    "full_image":    _full_image_slide,
    "quiz":          _quiz_slide,
    "practical":     _practical_slide,
    "procedure":     _practical_slide,   # alias
    "tracker":       _tracker_slide,
    "lab":           _lab_slide,
    "summary":       _summary_slide,
    "references":    _references_slide,
}


# ── Master build function ──────────────────────────────────────────────────────

def build_pptx(
    project: dict,
    slides: list[dict],
    version: str = "combined",
    images: dict[str, bytes] | None = None,
) -> bytes:
    """
    Build a .pptx in the Rapiscan training style.

    project: dict with keys: course_title, manufacturer, equipment_model,
             audience, training_type, language, manual_filename
    slides:  list of dicts: {type, title, bullets, speaker_notes,
                              source_pages, is_visible, image_tag, [answer_index]}
    version: combined | instructor | student | operator | engineer | executive
             instructor → shows correct answers on quiz slides
             student    → strips speaker notes
    images:  optional {FIGURE_N: png_bytes} map for image slide types
    """
    show_answers  = version in ("instructor",)
    include_notes = version not in ("student",)
    imgs = images or {}

    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    visible = [s for s in slides if s.get("is_visible", True)]
    total   = len(visible)

    for idx, s in enumerate(visible, 1):
        stype   = s.get("type", "content")

        # Downgrade image slides to content when no image is available
        if stype in ("image_content", "full_image") and not imgs.get(s.get("image_tag", "")):
            stype = "content"

        builder = SLIDE_BUILDERS.get(stype, _content_slide)
        kwargs: dict = {}
        if stype == "quiz":
            kwargs["show_answers"] = show_answers
        if stype in ("image_content", "full_image"):
            kwargs["images"] = imgs
        builder(prs, s, project, idx, total, **kwargs)

        # Strip speaker notes for student version
        if not include_notes:
            try:
                prs.slides[idx - 1].notes_slide.notes_text_frame.text = ""
            except Exception:
                pass

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
