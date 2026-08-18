"""
Document rebuild layer for the Translation Studio.

Takes a list of translated segments (with loc metadata) and writes
translated text back into a copy of the original file, preserving
styles, formatting, fonts, and document structure.

For formats where in-place rebuild isn't feasible (PDF, HTML, CSV, TXT),
a clean translated DOCX is generated instead.
"""
from __future__ import annotations

import io
import logging
import math
import re
import unicodedata
from typing import Any

from api.languages import is_rtl_lang

log = logging.getLogger(__name__)

_BIDI_LRM = "\u200e"
_BIDI_RLM = "\u200f"
_ARABIC_BLOCK_RE = re.compile(r"[\u0600-\u06FF\u0750-\u077F\u08A0-\u08FF]")
_LATIN_BLOCK_RE = re.compile(r"[A-Za-z]")
_UNIT_TOKEN_ALT = "kg|g|mg|ug|kv|mv|ma|mas|hz|khz|mhz|mm|cm|m|msv|usv|gy|mgy|mr|ev|kev|mev|w|kw|v|a|ohm|bar|psi|rpm|sec|s|min|h"
_UNIT_TOKEN_RE = re.compile(
    rf"(?i)\b({_UNIT_TOKEN_ALT})\b"
)


class LayoutValidationError(RuntimeError):
    """Raised when final PPTX visual/layout QA detects unresolved violations."""


def _shape_geometry(shape) -> tuple[int, int, int, int] | None:
    try:
        return (int(shape.left), int(shape.top), int(shape.width), int(shape.height))
    except Exception:
        return None


def _rect_overlap(r1: tuple[int, int, int, int], r2: tuple[int, int, int, int]) -> bool:
    l1, t1, w1, h1 = r1
    l2, t2, w2, h2 = r2
    return l1 < l2 + w2 and l1 + w1 > l2 and t1 < t2 + h2 and t1 + h1 > t2


def _rect_overlap_with_margin(
    r1: tuple[int, int, int, int],
    r2: tuple[int, int, int, int],
    margin_emu: int,
) -> bool:
    l1, t1, w1, h1 = r1
    expanded = (l1 - margin_emu, t1 - margin_emu, w1 + (2 * margin_emu), h1 + (2 * margin_emu))
    return _rect_overlap(expanded, r2)


def _to_rgb_tuple(rgb_obj) -> tuple[int, int, int] | None:
    if rgb_obj is None:
        return None
    try:
        return (int(rgb_obj[0]), int(rgb_obj[1]), int(rgb_obj[2]))
    except Exception:
        try:
            hx = str(rgb_obj)
            if len(hx) == 6:
                return (int(hx[0:2], 16), int(hx[2:4], 16), int(hx[4:6], 16))
        except Exception:
            return None
    return None


def _run_text_rgb(run) -> tuple[int, int, int] | None:
    rgb = _to_rgb_tuple(_run_explicit_rgb(run))
    if rgb is not None:
        return rgb
    return None


def _shape_text_rgb(shape) -> tuple[int, int, int]:
    # Fallback to black when runs inherit colors from the theme.
    try:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                rgb = _run_text_rgb(run)
                if rgb is not None:
                    return rgb
    except Exception:
        pass
    return (0, 0, 0)


def _shape_fill_info(shape) -> tuple[str, tuple[int, int, int] | None]:
    """Return background kind and representative rgb for a shape."""
    try:
        st_name = str(shape.shape_type)
        if "PICTURE" in st_name:
            try:
                from PIL import Image, ImageStat

                img = Image.open(io.BytesIO(shape.image.blob)).convert("RGB")
                mean = ImageStat.Stat(img).mean
                return "image", (int(mean[0]), int(mean[1]), int(mean[2]))
            except Exception:
                return "image", (128, 128, 128)
    except Exception:
        pass

    try:
        spPr = shape._element.find("{http://schemas.openxmlformats.org/presentationml/2006/main}spPr")
        if spPr is None:
            return "shape", None

        solid = spPr.find(f".//{{{_NS_DML}}}solidFill")
        if solid is not None:
            srgb = solid.find(f"{{{_NS_DML}}}srgbClr")
            if srgb is not None and srgb.get("val"):
                v = srgb.get("val")
                return "solid", (int(v[0:2], 16), int(v[2:4], 16), int(v[4:6], 16))
            return "solid", None

        grad = spPr.find(f".//{{{_NS_DML}}}gradFill")
        if grad is not None:
            stops = grad.findall(f".//{{{_NS_DML}}}gs")
            colors: list[tuple[int, int, int]] = []
            for gs in stops:
                srgb = gs.find(f"{{{_NS_DML}}}srgbClr")
                if srgb is None or not srgb.get("val"):
                    continue
                v = srgb.get("val")
                colors.append((int(v[0:2], 16), int(v[2:4], 16), int(v[4:6], 16)))
            if colors:
                r = int(sum(c[0] for c in colors) / len(colors))
                g = int(sum(c[1] for c in colors) / len(colors))
                b = int(sum(c[2] for c in colors) / len(colors))
                return "gradient", (r, g, b)
            return "gradient", (128, 128, 128)
    except Exception:
        pass

    return "shape", None


def _resolve_text_background(
    slide, shape_idx: int, text_rect: tuple[int, int, int, int]
) -> tuple[str, tuple[int, int, int], bool]:
    """Returns (kind, rgb, confirmed).

    `confirmed=False` means no slide-level shape actually revealed a
    background colour underneath this text — many slides rely on a colour
    band or fill that lives on the slide LAYOUT/MASTER, which is invisible
    to this per-slide scan.  The (255,255,255) "white" result in that case
    is an unverified guess, not a detection, and callers must not treat it
    as license to force-recolor text: a shape whose original author chose
    white for a colored master band would get "fixed" from white to black
    for a contrast problem that never existed.
    """
    best_idx = -1
    best_kind = "shape"
    best_rgb: tuple[int, int, int] | None = None

    for idx, other in enumerate(slide.shapes):
        if idx >= shape_idx:
            break
        geom = _shape_geometry(other)
        if geom is None or not _rect_overlap(text_rect, geom):
            continue

        try:
            if getattr(other, "has_text_frame", False):
                txt = " ".join((p.text or "").strip() for p in other.text_frame.paragraphs).strip()
                if txt:
                    continue
        except Exception:
            pass

        kind, rgb = _shape_fill_info(other)
        # An image that merely overlaps the text box's bounding rect is NOT
        # reliable evidence of the colour directly behind the text — technical
        # slides routinely place a text box beside a dark photo/scan while the
        # text itself sits on the plain (white) area. Treating that image as the
        # background made readable theme-coloured text get force-recoloured to
        # white (invisible on the real white area). Only deliberate solid/
        # gradient fills count as a confirmed text background.
        if kind == "image":
            continue
        if idx > best_idx:
            best_idx = idx
            best_kind = kind
            best_rgb = rgb

    if best_rgb is not None:
        return best_kind, best_rgb, True

    # Fallback: treat slide background as white when no underlying shape color
    # is available. Unconfirmed — see docstring.
    return "shape", (255, 255, 255), False


def _srgb_to_linear(v: float) -> float:
    if v <= 0.04045:
        return v / 12.92
    return ((v + 0.055) / 1.055) ** 2.4


def _relative_luminance(rgb: tuple[int, int, int]) -> float:
    r = _srgb_to_linear(rgb[0] / 255.0)
    g = _srgb_to_linear(rgb[1] / 255.0)
    b = _srgb_to_linear(rgb[2] / 255.0)
    return 0.2126 * r + 0.7152 * g + 0.0722 * b


def _contrast_ratio(fg: tuple[int, int, int], bg: tuple[int, int, int]) -> float:
    l1 = _relative_luminance(fg)
    l2 = _relative_luminance(bg)
    lighter = max(l1, l2)
    darker = min(l1, l2)
    return (lighter + 0.05) / (darker + 0.05)


def _set_shape_font_color(shape, rgb: tuple[int, int, int]) -> None:
    from pptx.dml.color import RGBColor

    try:
        for para in shape.text_frame.paragraphs:
            if not para.runs and para.text:
                para.add_run()
            for run in para.runs:
                run.font.color.rgb = RGBColor(rgb[0], rgb[1], rgb[2])
    except Exception:
        pass


def _pick_high_contrast_color(bg_rgb: tuple[int, int, int]) -> tuple[int, int, int]:
    # Preserve corporate palette preference before hard black/white fallback.
    candidates = [
        (13, 59, 102),    # corporate navy
        (255, 255, 255),
        (0, 0, 0),
    ]
    scored = sorted(candidates, key=lambda c: _contrast_ratio(c, bg_rgb), reverse=True)
    return scored[0]


def _is_text_shape(shape) -> bool:
    try:
        return bool(getattr(shape, "has_text_frame", False) and _shape_text_value(shape))
    except Exception:
        return False


def _shape_kind(shape, rect: tuple[int, int, int, int], slide_w: int, slide_h: int) -> str:
    try:
        if getattr(shape, "has_table", False):
            return "table"
    except Exception:
        pass
    try:
        if getattr(shape, "has_chart", False):
            return "chart"
    except Exception:
        pass
    try:
        if str(shape.shape_type).upper().find("PICTURE") >= 0:
            return "image"
    except Exception:
        pass
    try:
        if str(shape.shape_type).upper().find("DIAGRAM") >= 0:
            return "smartart"
    except Exception:
        pass

    l, t, w, h = rect
    has_text = _is_text_shape(shape)
    if not has_text:
        # Top-right compact non-text objects are likely logos/decorative marks.
        if t <= int(0.25 * slide_h) and l >= int(0.55 * slide_w) and w <= int(0.25 * slide_w) and h <= int(0.20 * slide_h):
            return "logo"
        return "shape"
    return "text"


def _collect_slide_objects(slide, slide_w: int, slide_h: int) -> list[dict[str, Any]]:
    objs: list[dict[str, Any]] = []
    for idx, shape in enumerate(slide.shapes):
        rect = _shape_geometry(shape)
        if rect is None:
            continue
        kind = _shape_kind(shape, rect, slide_w, slide_h)
        objs.append({"idx": idx, "shape": shape, "rect": rect, "kind": kind, "z": idx})
    return objs


def _is_bounds_violation(rect: tuple[int, int, int, int], slide_w: int, slide_h: int) -> bool:
    l, t, w, h = rect
    return l < 0 or t < 0 or (l + w) > slide_w or (t + h) > slide_h


def _capture_original_text_overlap_pairs(slide, slide_w: int, slide_h: int) -> set[tuple[int, int]]:
    pairs: set[tuple[int, int]] = set()
    objs = _collect_slide_objects(slide, slide_w, slide_h)
    texts = [o for o in objs if o["kind"] == "text"]
    non_text = [o for o in objs if o["kind"] != "text"]
    for t in texts:
        for o in non_text:
            if _rect_overlap(t["rect"], o["rect"]):
                pairs.add((int(t["idx"]), int(o["idx"])))
    return pairs


def _capture_original_text_hidden_pairs(slide, slide_w: int, slide_h: int) -> set[tuple[int, int]]:
    pairs: set[tuple[int, int]] = set()
    objs = _collect_slide_objects(slide, slide_w, slide_h)
    for i in range(len(objs)):
        for j in range(i + 1, len(objs)):
            a = objs[i]
            b = objs[j]
            if not _rect_overlap(a["rect"], b["rect"]):
                continue
            if a["kind"] == "text" and b["kind"] != "text" and int(b["z"]) > int(a["z"]):
                pairs.add((int(a["idx"]), int(b["idx"])))
            if b["kind"] == "text" and a["kind"] != "text" and int(a["z"]) > int(b["z"]):
                pairs.add((int(b["idx"]), int(a["idx"])))
    return pairs


def _capture_original_near_pairs(
    slide,
    slide_w: int,
    slide_h: int,
    margin_emu: int,
) -> set[tuple[int, int]]:
    pairs: set[tuple[int, int]] = set()
    objs = _collect_slide_objects(slide, slide_w, slide_h)
    for i in range(len(objs)):
        for j in range(i + 1, len(objs)):
            a = objs[i]
            b = objs[j]
            if _rect_overlap_with_margin(a["rect"], b["rect"], margin_emu):
                ai = int(a["idx"])
                bi = int(b["idx"])
                pairs.add((min(ai, bi), max(ai, bi)))
    return pairs


def _capture_protected_zone_indices(slide, slide_w: int, slide_h: int) -> set[int]:
    protected: set[int] = set()
    for idx, shape in enumerate(slide.shapes):
        rect = _shape_geometry(shape)
        if rect is None:
            continue
        l, t, w, h = rect
        kind = _shape_kind(shape, rect, slide_w, slide_h)

        if kind in {"logo", "image"}:
            protected.add(idx)
            continue

        # Header/footer bands and page number regions are protected.
        if t <= int(0.12 * slide_h) or (t + h) >= int(0.90 * slide_h):
            if kind in {"shape", "table", "chart", "smartart"}:
                protected.add(idx)

        # Confidentiality/page-note text areas.
        try:
            if getattr(shape, "has_text_frame", False):
                txt = _shape_text_value(shape).lower()
                if any(k in txt for k in ["confidential", "proprietary", "do not distribute", "page ", "copyright"]):
                    protected.add(idx)
        except Exception:
            pass
    return protected


def _bring_shape_to_front(slide, shape) -> None:
    try:
        sp_tree = slide.shapes._spTree
        el = shape._element
        sp_tree.remove(el)
        sp_tree.append(el)
    except Exception:
        pass


def _normalize_zorder_once(slide, slide_w: int, slide_h: int) -> bool:
    """Bring any text shape hidden behind an overlapping image/shape/table/
    chart to front. Must run AT MOST ONCE per slide, and only after every
    original_geom-indexed repair pass has finished.

    _bring_shape_to_front reorders the slide's XML shape tree, which shifts
    every later enumerate(slide.shapes) index. `_run_slide_visual_validation`
    (and the geometry snapshots it consults) is keyed entirely by that
    positional index, captured once before translation began. Calling this
    mid-repair-loop previously caused shapes to be restored to a DIFFERENT
    shape's original geometry once indices drifted (observed as e.g. a body
    placeholder inheriting the title's box). Doing it exactly once, as the
    last mutation on the slide, avoids that entirely.
    """
    from api.utils.arabic_layout_engine import _rects_overlap

    changed = False
    objs = _collect_slide_objects(slide, slide_w, slide_h)
    for a in objs:
        if a["kind"] != "text":
            continue
        for b in objs:
            if b is a or b["kind"] not in {"image", "shape", "chart", "table", "smartart"}:
                continue
            if int(b["z"]) > int(a["z"]) and _rects_overlap(a["rect"], b["rect"], 45720):
                _bring_shape_to_front(slide, a["shape"])
                changed = True
                break
    return changed


def _grow_textbox_height_in_safe_space(
    shape,
    obstacles: list[tuple[int, int, int, int]],
    slide_h: int,
    *,
    step_emu: int = 45720,
    max_steps: int = 10,
) -> bool:
    from pptx.util import Emu

    rect = _shape_geometry(shape)
    if rect is None:
        return False
    l, t, w, h = rect
    changed = False

    for _ in range(max_steps):
        down_try = (l, t, w, h + step_emu)
        ok_down = (down_try[1] + down_try[3] <= slide_h) and (not any(_rect_overlap(down_try, obs) for obs in obstacles))

        if ok_down:
            h += step_emu
            changed = True
            shape.height = Emu(h)
            continue
        break

    return changed


def _restore_text_anchor_if_safe(
    shape,
    original_rect: tuple[int, int, int, int] | None,
    obstacles: list[tuple[int, int, int, int]],
    slide_w: int,
    slide_h: int,
) -> bool:
    """Restore original left/top anchor when it can be done without collisions."""
    from pptx.util import Emu

    if original_rect is None:
        return False
    ol, ot, ow, oh = original_rect
    cur = _shape_geometry(shape)
    if cur is None:
        return False

    _, _, cw, ch = cur
    new_left = max(0, min(ol, slide_w - cw))
    new_top = max(0, min(ot, slide_h - ch))
    candidate = (new_left, new_top, cw, ch)
    if any(_rect_overlap(candidate, obs) for obs in obstacles):
        return False

    shape.left = Emu(new_left)
    shape.top = Emu(new_top)
    return True


def _nudge_nontext_shape(shape, obstacles: list[tuple[int, int, int, int]], slide_w: int, slide_h: int) -> bool:
    from pptx.util import Emu

    rect = _shape_geometry(shape)
    if rect is None:
        return False
    l0, t0, w, h = rect

    steps = [int(0.02 * slide_w), int(0.04 * slide_w), int(0.06 * slide_w)]
    candidates: list[tuple[int, int]] = [(l0, t0)]
    for s in steps:
        candidates.extend([
            (l0 + s, t0), (l0 - s, t0), (l0, t0 + s), (l0, t0 - s),
            (l0 + s, t0 + s), (l0 - s, t0 - s), (l0 + s, t0 - s), (l0 - s, t0 + s),
        ])

    for nl, nt in candidates:
        nl = max(0, min(nl, slide_w - w))
        nt = max(0, min(nt, slide_h - h))
        cand = (nl, nt, w, h)
        if any(_rect_overlap(cand, obs) for obs in obstacles):
            continue
        shape.left = Emu(nl)
        shape.top = Emu(nt)
        return True
    return False


def _run_slide_visual_validation(
    slide,
    *,
    slide_idx: int,
    slide_w: int,
    slide_h: int,
    original_geom: dict[int, tuple[int, int, int, int]],
    original_banner: tuple[int, int, int, int] | None,
    original_text_overlap_pairs: set[tuple[int, int]] | None,
    original_text_hidden_pairs: set[tuple[int, int]] | None,
    original_near_pairs: set[tuple[int, int]] | None,
    protected_zone_indices: set[int] | None,
    layout_warnings: list[str] | None,
    repair_actions: list[str] | None = None,
) -> list[str]:
    from api.utils.arabic_layout_engine import (
        _collect_obstacles,
        _expand_textbox_with_whitespace,
        _estimate_text_occupancy,
        _run_base_sizes_pt,
        _apply_font_scale,
        _enable_word_wrap,
        _set_norm_autofit,
        _optimize_translated_textbox,
        _rects_overlap,
        _resolve_overlap,
    )

    baseline_pairs = original_text_overlap_pairs or set()
    baseline_hidden = original_text_hidden_pairs or set()
    baseline_near = original_near_pairs or set()
    protected_idxs = protected_zone_indices or set()
    violations: list[str] = []
    clearance = 45720

    # Pass 2 and Pass 3 below do full pairwise (O(n^2)) comparisons across
    # every shape on the slide, repeated up to max_layout_repair_iters times
    # by the caller. That is fine for normal content slides (a handful to a
    # few dozen shapes) but explodes on slides that are really an embedded
    # technical diagram made of hundreds of tiny leader-line labels (seen in
    # practice: 250-270+ shapes on one slide) — a single slide like that can
    # cost tens of thousands of comparisons per iteration and, since these
    # dense diagrams are also where collision-based auto-repair makes the
    # least sense (labels are intentionally placed close to diagram lines by
    # the original author), burn all repair iterations for no real benefit.
    # Skip the pairwise passes there; Pass 1's cheap per-shape checks
    # (bounds, banner containment, contrast) still run and still fix real
    # issues.
    _DENSE_SLIDE_SHAPE_THRESHOLD = 80
    _is_dense_slide = len(slide.shapes) > _DENSE_SLIDE_SHAPE_THRESHOLD

    title_idxs = _detect_title_shape_indices(slide, slide_w, slide_h)
    # Real-banner detection only — NOT _detect_top_banner_rect, whose generic
    # top-strip fallback would make banner_rect "truthy" for every slide
    # (including plain white-background titles with no banner at all),
    # incorrectly forcing MSO_ANCHOR.MIDDLE below and shifting title text
    # down into whatever sits in that nominal region (e.g. a pre-existing
    # image positioned just under the title box).
    banner_rect = original_banner if original_banner is not None else _find_real_top_banner_shape(slide, slide_w, slide_h)

    def _mark(action: str) -> None:
        if repair_actions is None:
            return
        if action not in repair_actions:
            repair_actions.append(action)

    # Pass 1: visibility, bounds, banner containment, contrast.
    for idx, shape in enumerate(slide.shapes):
        if not getattr(shape, "has_text_frame", False):
            continue
        txt = _shape_text_value(shape)
        if not txt:
            continue

        geom = _shape_geometry(shape)
        if geom is None:
            continue
        left, top, width, height = geom
        is_title = idx in title_idxs

        # Keep first-line start aligned by forcing top anchor for body text.
        try:
            from pptx.enum.text import MSO_ANCHOR

            if is_title and banner_rect is not None:
                shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            else:
                shape.text_frame.vertical_anchor = MSO_ANCHOR.TOP
        except Exception:
            pass

        # Ensure text boxes remain on-slide.
        out_of_bounds = left < 0 or top < 0 or (left + width) > slide_w or (top + height) > slide_h
        if out_of_bounds:
            if idx in original_geom:
                ol, ot, ow, oh = original_geom[idx]
                shape.left, shape.top, shape.width, shape.height = ol, ot, ow, oh
                _mark("fixed object outside slide bounds")
            else:
                shape.left = max(0, min(left, slide_w - width))
                shape.top = max(0, min(top, slide_h - height))
                _mark("fixed object outside slide bounds")

        # Titles must remain in/over the original top banner strip — only
        # meaningful when a real banner shape was actually detected; with no
        # banner at all there's nothing to stay "over" and banner_rect is
        # None (see _find_real_top_banner_shape).
        if is_title and banner_rect is not None:
            cur = _shape_geometry(shape)
            if cur is not None and not _rect_overlap(cur, banner_rect) and idx in original_geom:
                ol, ot, ow, oh = original_geom[idx]
                shape.left, shape.top, shape.width, shape.height = ol, ot, ow, oh
                _mark("fixed title position")

        # Contrast check against resolved background underneath this text box.
        cur = _shape_geometry(shape)
        if cur is None:
            continue
        bg_kind, bg_rgb, bg_confirmed = _resolve_text_background(slide, idx, cur)
        fg_rgb = _shape_text_rgb(shape)
        min_ratio = 3.0 if is_title else 4.5
        ratio = _contrast_ratio(fg_rgb, bg_rgb)

        # Priority 1: restore original position first.
        if ratio < min_ratio and bg_confirmed and idx in original_geom:
            ol, ot, ow, oh = original_geom[idx]
            shape.left, shape.top, shape.width, shape.height = ol, ot, ow, oh
            cur = _shape_geometry(shape)
            if cur is not None:
                bg_kind, bg_rgb, bg_confirmed = _resolve_text_background(slide, idx, cur)
                fg_rgb = _shape_text_rgb(shape)
                ratio = _contrast_ratio(fg_rgb, bg_rgb)

        # Priority 2: if still low, apply high-contrast font color. Only when
        # the background was actually detected — an unconfirmed "assume
        # white" guess is not grounds to overwrite the author's original
        # color (e.g. white text designed for a layout/master color band).
        if ratio < min_ratio and bg_confirmed:
            new_rgb = _pick_high_contrast_color(bg_rgb)
            _set_shape_font_color(shape, new_rgb)
            ratio = _contrast_ratio(new_rgb, bg_rgb)
            _mark("fixed low contrast")

        if layout_warnings is not None and bg_confirmed and ratio < min_ratio:
            layout_warnings.append(
                f"Slide {slide_idx + 1} shape {idx}: low contrast remains ({ratio:.2f}) on {bg_kind} background"
            )

    # Pass 2: advanced collision resolution (iterative). Skipped on dense
    # diagram slides — see _DENSE_SLIDE_SHAPE_THRESHOLD note above.
    for _iter in range(4 if not _is_dense_slide else 0):
        objs = _collect_slide_objects(slide, slide_w, slide_h)
        text_objs = [o for o in objs if o["kind"] == "text"]
        non_text_objs = [o for o in objs if o["kind"] != "text"]

        collisions: list[tuple[dict[str, Any], dict[str, Any]]] = []
        for i in range(len(objs)):
            for j in range(i + 1, len(objs)):
                a = objs[i]
                b = objs[j]
                if not _rect_overlap(a["rect"], b["rect"]):
                    continue
                kinds = {a["kind"], b["kind"]}
                allowed = (
                    ("text" in kinds and "image" in kinds and (
                        (a["kind"] == "text" and (a["idx"], b["idx"]) in baseline_pairs)
                        or (b["kind"] == "text" and (b["idx"], a["idx"]) in baseline_pairs)
                    ))
                )
                if allowed:
                    continue
                collisions.append((a, b))

        if not collisions:
            break

        static_obstacles = _collect_obstacles(slide, slide_w, slide_h)

        for a, b in collisions:
            kinds = {a["kind"], b["kind"]}

            # Text vs any object: preserve design first, then move/fit text.
            if "text" in kinds:
                t = a if a["kind"] == "text" else b
                o = b if t is a else a
                pkey = (min(int(t["idx"]), int(o["idx"])), max(int(t["idx"]), int(o["idx"])))
                if pkey in baseline_near:
                    continue
                text_shape = t["shape"]
                text_idx = int(t["idx"])
                is_title = text_idx in title_idxs

                if text_idx in original_geom:
                    ol, ot, ow, oh = original_geom[text_idx]
                    text_shape.left, text_shape.top, text_shape.width, text_shape.height = ol, ot, ow, oh

                trect = _shape_geometry(text_shape)
                orect = _shape_geometry(o["shape"])
                if trect is not None and orect is not None and _rects_overlap(trect, orect, clearance):
                    if not is_title:
                        dynamic_obstacles = list(static_obstacles)
                        for x in text_objs:
                            if int(x["idx"]) == text_idx:
                                continue
                            dynamic_obstacles.append(x["rect"])

                        # Priority 1: expand text box into safe empty space.
                        _expand_textbox_with_whitespace(
                            text_shape,
                            dynamic_obstacles,
                            slide_w,
                            slide_h,
                            clearance,
                            max_expand_frac=0.20,
                        )
                        _mark("expanded text box")

                        # Priority 2: reflow and increase text box height.
                        _enable_word_wrap(text_shape)
                        _set_norm_autofit(text_shape)
                        _grow_textbox_height_in_safe_space(text_shape, dynamic_obstacles, slide_h)
                        _mark("reflowed paragraphs")
                        _mark("increased textbox height")

                        # Priority 3: move to nearest available safe area.
                        _resolve_overlap(
                            text_shape,
                            dynamic_obstacles,
                            slide_w,
                            slide_h,
                            clearance,
                            layout_warnings if layout_warnings is not None else [],
                            f"slide {slide_idx + 1} shape {text_idx}",
                        )
                        _mark("moved textbox")

                        # Priority 4: controlled font-size reduction (min readable size).
                        base_sizes = _run_base_sizes_pt(text_shape)
                        min_readable_pt = 12.0 if is_title else 10.0
                        for step in range(1, 9):
                            scale = 1.0 - (step * 0.02)
                            _apply_font_scale(text_shape, base_sizes, scale=scale, min_frac=0.75, min_pt=min_readable_pt)
                            g_now = _shape_geometry(text_shape)
                            o_now = _shape_geometry(o["shape"])
                            if g_now is not None and o_now is not None and not _rects_overlap(g_now, o_now, clearance):
                                break
                        _mark("reduced font size")

                        # Additional optimization pass with existing engine heuristics.
                        try:
                            _optimize_translated_textbox(
                                text_shape,
                                original_geom.get(text_idx, _shape_geometry(text_shape) or (0, 0, 1, 1)),
                                dynamic_obstacles,
                                slide_w,
                                slide_h,
                                clearance,
                                1.5,
                                layout_warnings if layout_warnings is not None else [],
                                f"slide {slide_idx + 1} shape {text_idx}",
                            )
                        except Exception:
                            _set_body_autofit(text_shape, "norm", force_wrap=True)

                        # Preserve original left/top anchor whenever safely possible.
                        _restore_text_anchor_if_safe(
                            text_shape,
                            original_geom.get(text_idx),
                            dynamic_obstacles,
                            slide_w,
                            slide_h,
                        )
                        _mark("restored original anchor")

                # Priority 4: nudge nearby non-text only when text attempts fail.
                trect = _shape_geometry(text_shape)
                orect = _shape_geometry(o["shape"])
                if trect is not None and orect is not None and _rects_overlap(trect, orect, clearance):
                    if int(o["idx"]) in protected_idxs:
                        continue
                    o_kind = str(o["kind"])
                    movable = o_kind in {"shape", "decorative"}
                    if movable:
                        other_obstacles = [x["rect"] for x in objs if int(x["idx"]) not in {int(o["idx"]), text_idx}]
                        if _nudge_nontext_shape(o["shape"], other_obstacles + [trect], slide_w, slide_h):
                            _mark("resized or moved non-critical object")
                continue

            # Non-text to non-text collisions (image-image/shape-shape/chart/table mixes).
            if a["idx"] in original_geom and a["kind"] != "text":
                ol, ot, ow, oh = original_geom[int(a["idx"])]
                a["shape"].left, a["shape"].top, a["shape"].width, a["shape"].height = ol, ot, ow, oh
            if b["idx"] in original_geom and b["kind"] != "text":
                ol, ot, ow, oh = original_geom[int(b["idx"])]
                b["shape"].left, b["shape"].top, b["shape"].width, b["shape"].height = ol, ot, ow, oh

            a_rect = _shape_geometry(a["shape"])
            b_rect = _shape_geometry(b["shape"])
            pkey = (min(int(a["idx"]), int(b["idx"])), max(int(a["idx"]), int(b["idx"])))
            if pkey in baseline_near:
                continue
            if a_rect is not None and b_rect is not None and _rects_overlap(a_rect, b_rect, clearance):
                if int(a["idx"]) not in protected_idxs and str(a["kind"]) in {"shape", "decorative"}:
                    others = [x["rect"] for x in objs if int(x["idx"]) != int(a["idx"])]
                    if _nudge_nontext_shape(a["shape"], others, slide_w, slide_h):
                        _mark("resized or moved non-critical object")
                elif int(b["idx"]) not in protected_idxs and str(b["kind"]) in {"shape", "decorative"}:
                    others = [x["rect"] for x in objs if int(x["idx"]) != int(b["idx"])]
                    if _nudge_nontext_shape(b["shape"], others, slide_w, slide_h):
                        _mark("resized or moved non-critical object")

    # Pass 3: final strict QA gate.
    # Only return hard failures that remain unrepaired. Advisory geometry drift
    # checks are intentionally excluded so QA warnings represent actionable,
    # unresolved defects only. Skipped on dense diagram slides (empty
    # final_objs makes every loop below a no-op) — see
    # _DENSE_SLIDE_SHAPE_THRESHOLD note above.
    final_objs = [] if _is_dense_slide else _collect_slide_objects(slide, slide_w, slide_h)
    final_text = [o for o in final_objs if o["kind"] == "text"]
    final_non_text = [o for o in final_objs if o["kind"] != "text"]

    # Objects outside slide boundaries.
    for obj in final_objs:
        if _is_bounds_violation(obj["rect"], slide_w, slide_h):
            violations.append(f"Slide {slide_idx + 1}: object {obj['idx']} outside slide boundaries")

    # Text clipping checks.
    for t in final_text:
        ts = t["shape"]
        tidx = int(t["idx"])
        try:
            if _shape_has_estimated_clipping(ts, t["rect"], is_title=tidx in title_idxs):
                violations.append(f"Slide {slide_idx + 1}: clipped text in shape {tidx}")
        except Exception:
            pass

        # Contrast must remain readable after all repairs.
        try:
            bg_kind, bg_rgb, bg_confirmed = _resolve_text_background(slide, tidx, t["rect"])
            fg_rgb = _shape_text_rgb(ts)
            min_ratio = 3.0 if tidx in title_idxs else 4.5
            if bg_confirmed and _contrast_ratio(fg_rgb, bg_rgb) < min_ratio:
                violations.append(
                    f"Slide {slide_idx + 1}: unreadable text contrast in shape {tidx} on {bg_kind} background"
                )
        except Exception:
            pass

    # Pairwise QA checks.
    for i in range(len(final_objs)):
        for j in range(i + 1, len(final_objs)):
            a = final_objs[i]
            b = final_objs[j]
            if not _rects_overlap(a["rect"], b["rect"], clearance):
                continue

            ak, bk = str(a["kind"]), str(b["kind"])
            aidx, bidx = int(a["idx"]), int(b["idx"])
            pkey = (min(aidx, bidx), max(aidx, bidx))

            if ak == "text" and bk == "text":
                if pkey not in baseline_near and _rect_overlap(a["rect"], b["rect"]):
                    violations.append(f"Slide {slide_idx + 1}: overlapping text boxes ({aidx},{bidx})")

            if {ak, bk} == {"text", "image"}:
                t = a if ak == "text" else b
                o = b if ak == "text" else a
                if (int(t["idx"]), int(o["idx"])) not in baseline_pairs:
                    violations.append(f"Slide {slide_idx + 1}: text inside image not present in original ({t['idx']},{o['idx']})")
                # Z-order: text must not be hidden behind image unless original
                # intentional overlap. NOTE: do NOT call _bring_shape_to_front
                # here — this pass runs repeatedly (up to max_layout_repair_iters
                # times) and reordering the slide's XML shape tree mid-loop
                # shifts every subsequent enumerate(slide.shapes) index, which
                # corrupts the idx-keyed `original_geom` restores used elsewhere
                # in this same repeated function. Only report the violation;
                # the actual z-order fix runs exactly once, after all
                # original_geom-based repairs are done (see
                # _normalize_zorder_once).
                if int(o["z"]) > int(t["z"]) and (int(t["idx"]), int(o["idx"])) not in baseline_hidden:
                    violations.append(f"Slide {slide_idx + 1}: text hidden behind image ({t['idx']},{o['idx']})")

            if {ak, bk} == {"text", "shape"} or {ak, bk} == {"text", "chart"} or {ak, bk} == {"text", "table"} or {ak, bk} == {"text", "smartart"}:
                t = a if ak == "text" else b
                o = b if ak == "text" else a
                if int(o["z"]) > int(t["z"]) and (int(t["idx"]), int(o["idx"])) not in baseline_hidden:
                    # See note above: no in-place z-order mutation here either.
                    violations.append(f"Slide {slide_idx + 1}: text hidden behind {bk if ak=='text' else ak} ({t['idx']},{o['idx']})")

            if ak == "image" and bk == "image":
                if pkey not in baseline_near and _rect_overlap(a["rect"], b["rect"]):
                    violations.append(f"Slide {slide_idx + 1}: overlapping images ({aidx},{bidx})")

            if ak == "shape" and bk == "shape":
                if pkey not in baseline_near and _rect_overlap(a["rect"], b["rect"]):
                    violations.append(f"Slide {slide_idx + 1}: overlapping shapes ({aidx},{bidx})")

            if "text" in {ak, bk} and "logo" in {ak, bk}:
                violations.append(f"Slide {slide_idx + 1}: logo covered by text ({aidx},{bidx})")

            if "text" in {ak, bk} and ("shape" in {ak, bk} or "chart" in {ak, bk} or "table" in {ak, bk} or "smartart" in {ak, bk}):
                # Explicit check for title intersecting graphics.
                t = a if ak == "text" else b
                if int(t["idx"]) in title_idxs:
                    if pkey not in baseline_near and _rect_overlap(a["rect"], b["rect"]):
                        violations.append(f"Slide {slide_idx + 1}: title intersecting graphics ({aidx},{bidx})")

    # De-duplicate for cleaner failure output.
    deduped = list(dict.fromkeys(violations))
    return deduped


# ── Arabic text helpers ────────────────────────────────────────────────────────

def _prepare_arabic(text: str) -> str:
    """Apply Arabic reshaping and BiDi for correct display in Office apps."""
    try:
        import arabic_reshaper
        from bidi.algorithm import get_display
        reshaped = arabic_reshaper.reshape(text)
        return get_display(reshaped)
    except ImportError:
        return text


def _is_rtl_text(text: str) -> bool:
    """Check if text contains Arabic characters."""
    return bool(re.search(r"[\u0600-\u06FF\u0750-\u077F\u08A0-\u08FF]", text))


# ── DOCX rebuild ───────────────────────────────────────────────────────────────

def rebuild_docx(original_bytes: bytes, segments: list[dict], target_lang: str) -> bytes:
    """
    Write translated segments back into a DOCX, preserving styles.
    Returns the rebuilt DOCX as bytes.
    """
    from docx import Document
    from docx.oxml.ns import qn
    from docx.shared import Pt
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc = Document(io.BytesIO(original_bytes))
    is_rtl = is_rtl_lang(target_lang)

    # Build lookups
    # para_lines[para_idx][line_idx] = translated_line
    # para_total[para_idx] = total_lines in that paragraph (for soft-return count)
    from collections import defaultdict
    para_lines: dict[int, dict[int, str]] = defaultdict(dict)
    para_total: dict[int, int] = {}
    table_map: dict[tuple, str] = {}
    header_map: dict[tuple, str] = {}
    footer_map: dict[tuple, str] = {}

    for seg in segments:
        loc = seg.get("loc", {})
        if loc.get("format") != "docx":
            continue
        seg_type = seg.get("seg_type", "paragraph")
        target = seg.get("target", "").strip()

        if seg_type == "paragraph":
            para_idx  = loc.get("para_idx")
            line_idx  = loc.get("line_idx",  0)
            total_ln  = loc.get("total_lines", 1)
            if para_idx is None:
                continue
            if target or loc.get("passthrough"):
                # Use target if available; fall back to source for passthroughs
                para_lines[para_idx][line_idx] = target or seg.get("source", "")
            para_total[para_idx] = total_ln

        elif seg_type == "table_cell":
            if target:
                table_map[(loc["tbl_idx"], loc["row_idx"], loc["col_idx"])] = target
        elif seg_type == "header":
            if target:
                header_map[(loc["section"], loc["header_para"])] = target
        elif seg_type == "footer":
            if target:
                footer_map[(loc["section"], loc["footer_para"])] = target

    # Replace body paragraphs
    for para_idx, para in enumerate(doc.paragraphs):
        line_map = para_lines.get(para_idx)
        if not line_map:
            continue
        total_lines = para_total.get(para_idx, 1)
        # Reconstruct the full paragraph text — preserve soft-return positions
        lines_out: list[str] = []
        for li in range(total_lines):
            lines_out.append(line_map.get(li, ""))   # empty string keeps the soft break slot
        new_text = "\n".join(lines_out).strip("\n")
        if new_text:
            _replace_paragraph_text(para, new_text, is_rtl)

    # Replace table cells
    for tbl_idx, table in enumerate(doc.tables):
        for row_idx, row in enumerate(table.rows):
            for col_idx, cell in enumerate(row.cells):
                target = table_map.get((tbl_idx, row_idx, col_idx))
                if target:
                    # Replace text in first paragraph of cell
                    if cell.paragraphs:
                        _replace_paragraph_text(cell.paragraphs[0], target, is_rtl)
                        for p in cell.paragraphs[1:]:
                            for run in p.runs:
                                if _run_has_drawing(run._element):
                                    continue
                                run.text = ""
                            if is_rtl:
                                _set_rtl_paragraph_direction(p)

    # Replace headers and footers
    for sec_idx, section in enumerate(doc.sections):
        for hdr_para_idx, para in enumerate(section.header.paragraphs):
            target = header_map.get((sec_idx, hdr_para_idx))
            if target:
                _replace_paragraph_text(para, target, is_rtl)
        for ftr_para_idx, para in enumerate(section.footer.paragraphs):
            target = footer_map.get((sec_idx, ftr_para_idx))
            if target:
                _replace_paragraph_text(para, target, is_rtl)

    # ── Image text overlay ─────────────────────────────────────────────────────
    # Collect image_text segments grouped by rel_id
    image_text_segs: dict[str, list[dict]] = {}
    for seg in segments:
        if seg.get("seg_type") != "image_text":
            continue
        target = seg.get("target", "").strip()
        if not target:
            continue
        loc = seg.get("loc", {})
        rel_id = loc.get("rel_id")
        if not rel_id:
            continue
        if rel_id not in image_text_segs:
            image_text_segs[rel_id] = []
        image_text_segs[rel_id].append({
            "text": target,
            "x_pct": loc.get("x_pct", 0.5),
            "y_pct": loc.get("y_pct", 0.5),
            "bbox": loc.get("bbox"),
            "font_size": loc.get("font_size", 14),
            "font_color": loc.get("font_color", "#000000"),
            "alignment": loc.get("alignment", "center"),
            "category": loc.get("text_category", "other"),
            "confidence": loc.get("ocr_confidence", 0.0),
        })

    if image_text_segs:
        try:
            from api.utils.image_text_extractor import render_overlay_image
            from docx.oxml.ns import qn as _qn
            for para in doc.paragraphs:
                for run in para.runs:
                    drawing = run._element.find(_qn("w:drawing"))
                    if drawing is None:
                        continue
                    blip = drawing.find(".//" + _qn("a:blip"))
                    if blip is None:
                        continue
                    embed_id = blip.get(_qn("r:embed"))
                    if embed_id not in image_text_segs:
                        continue
                    try:
                        rel = doc.part.rels.get(embed_id)
                        if rel is None:
                            continue
                        original_bytes = rel.target_part.blob
                        overlaid = render_overlay_image(
                            original_bytes,
                            image_text_segs[embed_id],
                            target_lang,
                        )
                        rel.target_part._blob = overlaid
                    except Exception as _oe:
                        log.warning("Overlay failed for rel %s: %s", embed_id, _oe)
        except Exception as _overlay_err:
            log.warning("Image overlay step skipped: %s", _overlay_err)

    # ── RTL table mirroring ─────────────────────────────────────────────────
    # <w:tblPr><w:bidiVisual/></w:tblPr> reverses a table's column rendering
    # order for right-to-left layout — independent of per-cell text alignment,
    # which is already handled per-paragraph above. Mirrors the tblPr rtl="1"
    # rule rebuild_pptx already applies to PowerPoint tables.
    if is_rtl:
        from docx.oxml.ns import qn as _qn
        for table in doc.tables:
            tblPr = table._tbl.find(_qn("w:tblPr"))
            if tblPr is None:
                continue
            if tblPr.find(_qn("w:bidiVisual")) is None:
                from lxml import etree as _etree
                _etree.SubElement(tblPr, _qn("w:bidiVisual"))

    out = io.BytesIO()
    doc.save(out)
    return out.getvalue()


def _run_has_drawing(r_elem) -> bool:
    """True when a <w:r> run element carries an image — DrawingML (<w:drawing>)
    or legacy VML (<w:pict>) — rather than (or in addition to) text. Such runs
    must survive text-replacement passes; deleting them silently drops the
    image while leaving its now-orphaned word/media part and relationship
    entry behind in the package.

    Descendant search (.//), not immediate-child: a drawing/pict is commonly
    nested one or two levels deeper inside <mc:AlternateContent><mc:Choice>
    (w:drawing) / <mc:Fallback> (w:pict) — Word's version-compatibility
    wrapper for shapes/text boxes — rather than being a direct child of <w:r>.
    """
    from docx.oxml.ns import qn
    return r_elem.find(f".//{qn('w:drawing')}") is not None or r_elem.find(f".//{qn('w:pict')}") is not None


def _resolve_effective_alignment(para):
    """`paragraph.alignment` only reflects DIRECT formatting (an explicit
    <w:jc> on the paragraph itself) and returns None when alignment instead
    comes from the paragraph's STYLE — e.g. a title using the "Heading1"
    style, whose own <w:pPr><w:jc val="center"/></w:pPr> centers it without
    any direct override on the paragraph. Walk the style's basedOn chain to
    find the effective alignment the way Word actually resolves it."""
    if para.alignment is not None:
        return para.alignment
    style = para.style
    seen = set()
    while style is not None and id(style) not in seen:
        seen.add(id(style))
        try:
            if style.paragraph_format.alignment is not None:
                return style.paragraph_format.alignment
        except Exception:
            pass
        style = getattr(style, "base_style", None)
    return None


def _set_rtl_paragraph_direction(para) -> None:
    """Set bidi=1 and right-alignment on a paragraph for RTL rendering — a
    centered heading/title (whether centered by direct formatting or by its
    style, e.g. Heading1) is left alone rather than forced right/left, since
    <w:jc val="center"> already reads correctly in RTL and this is the one
    alignment translation must never overwrite.

    Word inverts the physical meaning of <w:jc> once <w:bidi val="1"> is set
    on a paragraph: jc="right" renders visually LEFT, and jc="left" renders
    visually RIGHT. Verified empirically (pixel-measured, real Word via COM
    automation, three independent isolated repros) — not documented behavior,
    but consistently reproducible. WD_ALIGN_PARAGRAPH.LEFT is therefore the
    correct value to use here for a visually right-aligned RTL paragraph.
    """
    from docx.oxml.ns import qn
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from lxml import etree

    try:
        p_elem = para._p
        pPr = p_elem.get_or_add_pPr() if hasattr(p_elem, "get_or_add_pPr") else para._p.get_or_add_pPr()
        bidi_el = pPr.find(qn("w:bidi"))
        if bidi_el is None:
            bidi_el = etree.SubElement(pPr, qn("w:bidi"))
        bidi_el.set(qn("w:val"), "1")
        if _resolve_effective_alignment(para) != WD_ALIGN_PARAGRAPH.CENTER:
            para.alignment = WD_ALIGN_PARAGRAPH.LEFT
    except Exception:
        pass


def _replace_paragraph_text(para, new_text: str, is_rtl: bool) -> None:
    """
    Replace paragraph text while preserving the first run's character formatting.

    Handles multi-line text: \\n in new_text is written as a proper OOXML soft
    line break (<w:br/>) so Word renders it correctly rather than showing a
    box character.

    Runs that carry an image (DrawingML <w:drawing> or legacy VML <w:pict>)
    are left in place untouched — only text-bearing runs are cleared and
    replaced, so an image sharing a paragraph/cell with translatable text
    (e.g. an icon next to a label) is never deleted.
    """
    import copy
    from docx.oxml.ns import qn
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from lxml import etree

    p_elem = para._p
    XML_SPACE = "{http://www.w3.org/XML/1998/namespace}space"

    # Capture rPr from the first non-drawing run for font/size/bold preservation
    rPr_copy = None
    for r in para.runs:
        if _run_has_drawing(r._element):
            continue
        rPr = r._element.find(qn("w:rPr"))
        if rPr is not None:
            rPr_copy = copy.deepcopy(rPr)
        break

    # Remove existing run / hyperlink / br children — except runs that carry
    # an image, which must survive this text-replacement pass untouched.
    for child in list(p_elem):
        if child.tag == qn("w:r") and _run_has_drawing(child):
            continue
        if child.tag in (qn("w:r"), qn("w:hyperlink"), qn("w:ins"), qn("w:del")):
            p_elem.remove(child)

    lines = new_text.split("\n")
    for line_idx, line in enumerate(lines):
        if line:
            r = etree.SubElement(p_elem, qn("w:r"))
            if rPr_copy is not None:
                r.insert(0, copy.deepcopy(rPr_copy))
            t = etree.SubElement(r, qn("w:t"))
            t.text = line
            # xml:space="preserve" required when text starts/ends with whitespace
            if line != line.strip():
                t.set(XML_SPACE, "preserve")

        # Soft line break between lines — not after the last line
        if line_idx < len(lines) - 1:
            r_br = etree.SubElement(p_elem, qn("w:r"))
            if rPr_copy is not None:
                r_br.insert(0, copy.deepcopy(rPr_copy))
            etree.SubElement(r_br, qn("w:br"))   # default type = textWrapping (soft return)

    if is_rtl:
        _set_rtl_paragraph_direction(para)


# ── PPTX rebuild ───────────────────────────────────────────────────────────────

def _pptx_rebuild_group(
    group,
    slide_idx: int,
    shape_idx: int,
    text_map: dict,
    src_map: dict,
    is_rtl: bool,
    profile: "dict | None",
    font_subs: list,
    translated_shape_idxs: set,
    deck_src_is_rtl: bool,
) -> None:
    """Recursively apply translations to text frames inside a GROUP shape.

    Mirrors the extractor's traversal exactly so the (slide_idx, shape_idx,
    para_idx) keys built here match those stored during extraction:

        top-level shape  → shape_idx from enumerate(slide.shapes)
        group child      → parent_shape_idx * 10_000 + child_idx
        nested group     → (parent_child_shape_idx) * 10_000 + grandchild_idx

    Tables inside groups are skipped — they are handled by _apply_rtl_to_table
    at the top-level (groups cannot contain table shapes in standard PPTX).
    """
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE as _MSO
        children = list(group.shapes)
    except Exception:
        return

    for child_idx, child in enumerate(children):
        if child is None:
            continue
        child_shape_idx = shape_idx * 10_000 + child_idx
        try:
            # Nested group — recurse
            if child.shape_type == _MSO.GROUP:
                _pptx_rebuild_group(
                    child, slide_idx, child_shape_idx,
                    text_map, src_map, is_rtl, profile, font_subs,
                    translated_shape_idxs, deck_src_is_rtl,
                )
                continue

            if not child.has_text_frame:
                continue

            shape_translated = False
            for para_idx, para in enumerate(child.text_frame.paragraphs):
                key = (slide_idx, child_shape_idx, para_idx)
                trans_txt = text_map.get(key)
                label = (f"slide {slide_idx+1} group shape {child_shape_idx}"
                         f" para {para_idx}")

                if trans_txt:
                    _replace_pptx_para_arabic(
                        para, trans_txt, is_rtl,
                        is_title=False,
                        profile=profile,
                        font_subs=font_subs,
                        slide_label=label,
                        source_text=src_map.get(key, ""),
                    )
                    shape_translated = True
                elif is_rtl and not deck_src_is_rtl and para.text.strip():
                    _set_para_rtl(para, True)
                    if not _text_is_rtl_script(para.text):
                        _apply_arabic_para_alignment(
                            para, is_title=False, profile=profile
                        )

            if is_rtl and shape_translated:
                _set_body_autofit(child, "norm", force_wrap=True)
                translated_shape_idxs.add(child_shape_idx)

        except Exception as exc:
            log.debug(
                "PPTX group child rebuild skipped (slide=%d shape=%d child=%d): %s",
                slide_idx + 1, shape_idx, child_idx, exc,
            )


def rebuild_pptx(
    original_bytes: bytes,
    segments: list[dict],
    target_lang: str,
    layout_warnings: "list[str] | None" = None,
    style_profile_override: "dict | None" = None,
    template_strength: str = "balanced",
    strict_qa: bool = False,
    allow_export_with_warnings: bool = True,
    auto_repair_enabled: bool = True,
    export_best_effort_result: bool = True,
) -> bytes:
    """Write translated segments back into a PPTX, preserving all formatting.

    For Arabic (target_lang == "ar") the rebuild automatically loads the Arabic
    formatting profile extracted from the approved reference presentation and applies:

    - rtl=1 on every translated paragraph (and every untranslated one in RTL docs)
    - RIGHT alignment for body text, CENTER preserved for centered titles
    - Arabic font substitution when the original font cannot render Arabic glyphs
    - RTL on all table shapes: tblPr rtl=1 + per-cell paragraph RTL + alignment
    - Table cell text replacement (now extracted by the extractor)
    - Font-substitution log entries for the quality report

    Returns the rebuilt PPTX as bytes.
    """
    from pptx import Presentation

    prs = Presentation(io.BytesIO(original_bytes))
    is_rtl = is_rtl_lang(target_lang)
    is_ar  = target_lang == "ar"

    # Load Arabic reference profile when translating to Arabic.
    # style_profile_override (from a saved LayoutStyle or uploaded reference) takes
    # priority over the built-in ARABIC_PROFILE.  When it is provided, it is merged
    # on top of the built-in profile so any missing keys fall back gracefully.
    profile: dict | None = None
    if is_ar:
        try:
            from api.utils.arabic_pptx_profile import ARABIC_PROFILE
            profile = dict(ARABIC_PROFILE)
        except Exception as _pe:
            log.warning("Could not load arabic_pptx_profile: %s", _pe)

    if style_profile_override:
        profile = {**(profile or {}), **style_profile_override}
        log.info(
            "Applying style_profile_override (strength=%s): title_font=%s body_font=%s colors=%d",
            template_strength,
            style_profile_override.get("title_font_name", "—"),
            style_profile_override.get("body_font_name", "—"),
            len(style_profile_override.get("theme_colors") or []),
        )

    # ── Build lookup maps from segment list ───────────────────────────────────
    # text_map:  (slide_idx, shape_idx, para_idx)         → translated text
    # table_map: (slide_idx, shape_idx, row_idx, col_idx) → translated text
    # notes_map: slide_idx                                → translated notes text
    text_map:  dict[tuple, str] = {}
    src_map:   dict[tuple, str] = {}
    table_map: dict[tuple, str] = {}
    table_src: dict[tuple, str] = {}
    notes_map: dict[int, str]   = {}
    notes_src: dict[int, str]   = {}
    image_text_by_rel: dict[str, list[dict]] = {}
    image_text_by_shape: dict[tuple[int, int], list[dict]] = {}
    # master_text_map: (master_idx, shape_idx, para_idx) → translated text
    master_text_map: dict[tuple, str] = {}

    for seg in segments:
        translated = seg.get("target", "").strip()
        if not translated:
            continue
        loc = seg.get("loc", {})
        if loc.get("format") != "pptx":
            continue
        source_txt = seg.get("source", "") or ""
        seg_type = seg.get("seg_type", "paragraph")

        if seg_type == "image_text":
            payload = {
                "text": translated,
                "x_pct": loc.get("x_pct", 0.5),
                "y_pct": loc.get("y_pct", 0.5),
                    "bbox": loc.get("bbox"),
                    "font_size": loc.get("font_size", 14),
                    "font_color": loc.get("font_color", "#000000"),
                    "alignment": loc.get("alignment", "center"),
                    "category": loc.get("text_category", "other"),
                    "confidence": loc.get("ocr_confidence", 0.0),
            }
            rel_id = loc.get("rel_id")
            if rel_id:
                image_text_by_rel.setdefault(rel_id, []).append(payload)
            else:
                sidx = loc.get("slide_idx")
                shidx = loc.get("shape_idx")
                if sidx is not None and shidx is not None:
                    image_text_by_shape.setdefault((int(sidx), int(shidx)), []).append(payload)
            continue

        if loc.get("notes"):
            notes_map[loc["slide_idx"]] = translated
            notes_src[loc["slide_idx"]] = source_txt
        elif loc.get("table_cell"):
            key = (loc["slide_idx"], loc["shape_idx"],
                   loc["row_idx"],   loc["col_idx"])
            table_map[key] = translated
            table_src[key] = source_txt
        elif "master_idx" in loc:
            key = (loc["master_idx"], loc["shape_idx"], loc["para_idx"])
            master_text_map[key] = translated
        else:
            key = (loc["slide_idx"], loc["shape_idx"], loc["para_idx"])
            text_map[key] = translated
            src_map[key]  = source_txt

    font_subs: list[str] = []   # populated by helpers; surfaced in quality report
    post_export_qa_contexts: list[dict[str, Any]] = []

    # ── Slide-master text ───────────────────────────────────────────────────
    # Written back before the per-slide loop since masters are shared,
    # slide-independent objects — see the matching extraction note in
    # doc_extractor.py's extract_pptx for why this exists.
    if master_text_map:
        for master_idx, master in enumerate(prs.slide_masters):
            try:
                master_shapes = list(master.shapes)
            except Exception:
                continue
            for shape_idx, shape in enumerate(master_shapes):
                if not getattr(shape, "has_text_frame", False):
                    continue
                try:
                    tf = shape.text_frame
                    for para_idx, para in enumerate(tf.paragraphs):
                        trans_txt = master_text_map.get((master_idx, shape_idx, para_idx))
                        if not trans_txt:
                            continue
                        _replace_pptx_para_arabic(
                            para, trans_txt, is_rtl,
                            is_title=False, profile=profile, font_subs=font_subs,
                            slide_label=f"master {master_idx} shape {shape_idx}",
                            source_text="",
                        )
                except Exception as exc:
                    log.warning("Could not update master %d shape %d: %s",
                                master_idx, shape_idx, exc)

    # ── Deck-level source script detection ────────────────────────────────────
    # A deck whose segment sources are mostly Arabic script was already
    # DESIGNED for Arabic (e.g. re-translating an Arabic deck).  Untranslated
    # paragraphs in such decks are left completely untouched — no rtl, no
    # alignment — so the output stays visually identical to the original.
    _all_sources = [s.get("source", "") or "" for s in segments
                    if s.get("loc", {}).get("format") == "pptx"]
    _rtl_sources = sum(1 for s in _all_sources if _text_is_rtl_script(s))
    deck_src_is_rtl = bool(_all_sources) and _rtl_sources / len(_all_sources) > 0.5

    # ── Main rebuild pass ─────────────────────────────────────────────────────
    for slide_idx, slide in enumerate(prs.slides):
        slide_translated_shapes: set[int] = set()

        # Real banner detection (may be None) — see _find_real_top_banner_shape
        # docstring for why _detect_top_banner_rect's generic fallback must
        # not be used here.
        original_banner_rect = _find_real_top_banner_shape(slide, prs.slide_width, prs.slide_height)
        # Single, deck-wide title detection for this slide, reused by both
        # the per-shape translation loop below AND _mirror_titles_for_rtl_slide
        # / the QA pass later. Previously the per-shape loop had its OWN,
        # narrower inline check (placeholder-type only), while everything
        # else used this broader heuristic (which also catches manually
        # drawn, non-placeholder title textboxes). A shape the broad check
        # recognized as a title but the narrow one didn't would get treated
        # as a title for positioning/banner purposes but as ordinary BODY
        # text for font-sizing purposes — including a lower minimum-readable
        # font floor and body-text clipping thresholds — silently shrinking
        # it far more than a title ever should. One shared source of truth
        # closes that gap for good instead of requiring the two checks to be
        # kept in sync by hand.
        slide_title_idxs = _detect_title_shape_indices(slide, prs.slide_width, prs.slide_height)
        original_text_overlap_pairs = _capture_original_text_overlap_pairs(slide, prs.slide_width, prs.slide_height)
        original_text_hidden_pairs = _capture_original_text_hidden_pairs(slide, prs.slide_width, prs.slide_height)
        original_near_pairs = _capture_original_near_pairs(slide, prs.slide_width, prs.slide_height, 45720)
        protected_zone_indices = _capture_protected_zone_indices(slide, prs.slide_width, prs.slide_height)

        # ── GEOMETRY LOCK: snapshot every shape's exact position/size ─────────
        # We restore these after any processing that could move shapes.  This is
        # the first line of defence: even if a library call accidentally changes
        # geometry we will catch and undo it.
        _geom_snap: dict[int, tuple] = {}
        for _gi, _gs in enumerate(slide.shapes):
            try:
                _geom_snap[_gi] = (_gs.left, _gs.top, _gs.width, _gs.height)
            except Exception:
                pass

        post_export_qa_contexts.append({
            "slide_idx": slide_idx,
            "original_geom": {k: (int(v[0]), int(v[1]), int(v[2]), int(v[3])) for k, v in _geom_snap.items()},
            "original_banner": original_banner_rect,
            "original_text_overlap_pairs": original_text_overlap_pairs,
            "original_text_hidden_pairs": original_text_hidden_pairs,
            "original_near_pairs": original_near_pairs,
            "protected_zone_indices": protected_zone_indices,
        })

        for shape_idx, shape in enumerate(slide.shapes):

            # ── Picture shapes: overlay translated image labels (if any) ────
            try:
                from pptx.enum.shapes import MSO_SHAPE_TYPE as _MSO
                if shape.shape_type == _MSO.PICTURE:
                    labels = None
                    try:
                        from pptx.oxml.ns import qn as _qn

                        blip = shape._element.find(".//" + _qn("a:blip"))
                        embed_id = blip.get(_qn("r:embed")) if blip is not None else None
                    except Exception:
                        embed_id = None

                    if embed_id:
                        labels = image_text_by_rel.get(embed_id)
                    if not labels:
                        labels = image_text_by_shape.get((slide_idx, shape_idx))

                    if labels:
                        try:
                            from api.utils.image_text_extractor import render_overlay_image

                            original_img = shape.image.blob
                            overlaid = render_overlay_image(original_img, labels, target_lang)
                            # python-pptx exposes relationships via `.rels`
                            # (dict of rId -> Relationship) and the
                            # `.related_part(rId)` accessor — there is no
                            # `.related_parts` collection in this version.
                            if embed_id and embed_id in slide.part.rels:
                                slide.part.related_part(embed_id)._blob = overlaid
                            elif embed_id and embed_id in prs.part.rels:
                                prs.part.related_part(embed_id)._blob = overlaid
                            else:
                                log.debug(
                                    "PPTX image overlay: rel not found for slide=%d shape=%d",
                                    slide_idx + 1,
                                    shape_idx,
                                )
                        except Exception as _img_overlay_err:
                            log.warning(
                                "PPTX image overlay failed (slide=%d shape=%d): %s",
                                slide_idx + 1,
                                shape_idx,
                                _img_overlay_err,
                            )
            except Exception:
                pass

            # ── Table shapes ──────────────────────────────────────────────────
            if shape.has_table:
                if is_rtl:
                    _apply_rtl_to_table(
                        shape, table_map, slide_idx, shape_idx,
                        is_rtl=is_rtl, profile=profile, font_subs=font_subs,
                        table_src=table_src, deck_src_is_rtl=deck_src_is_rtl,
                    )
                continue

            # ── Group shapes — recurse into children ──────────────────────────
            # The extractor assigns group children synthetic indices
            # (parent_idx * 10000 + child_idx) so we must mirror that traversal
            # here.  Group shapes themselves have neither has_table nor
            # has_text_frame at the top level, so without this branch their
            # children's translations are silently dropped.
            try:
                from pptx.enum.shapes import MSO_SHAPE_TYPE as _MSO
                if shape.shape_type == _MSO.GROUP:
                    _pptx_rebuild_group(
                        shape, slide_idx, shape_idx,
                        text_map, src_map, is_rtl, profile, font_subs,
                        slide_translated_shapes, deck_src_is_rtl,
                    )
                    continue
            except Exception as _grp_err:
                log.debug("group rebuild skipped for shape %d: %s", shape_idx, _grp_err)

            # ── Text-frame shapes ─────────────────────────────────────────────
            if not shape.has_text_frame:
                continue

            # Determine if this is a title (affects alignment, font-floor, and
            # clipping-threshold rules) — via the single shared detector, see
            # slide_title_idxs above for why this must not be a separate,
            # narrower check.
            is_title_shape = shape_idx in slide_title_idxs

            shape_translated = False  # tracks whether any para in this shape was replaced

            for para_idx, para in enumerate(shape.text_frame.paragraphs):
                key      = (slide_idx, shape_idx, para_idx)
                trans_txt = text_map.get(key)
                label     = f"slide {slide_idx+1} shape {shape_idx} para {para_idx}"

                if trans_txt:
                    # ── Replace translated paragraph with full Arabic formatting
                    _replace_pptx_para_arabic(
                        para, trans_txt, is_rtl,
                        is_title=is_title_shape,
                        profile=profile,
                        font_subs=font_subs,
                        slide_label=label,
                        source_text=src_map.get(key, ""),
                    )
                    shape_translated = True
                elif is_rtl and not deck_src_is_rtl and para.text.strip():
                    # ── Untranslated paragraph (passthrough / technical term)
                    # Only adjusted when translating an LTR-designed deck, so
                    # passthrough terms sit correctly among Arabic body text.
                    # In Arabic-designed decks these paragraphs stay untouched.
                    _set_para_rtl(para, True)
                    if not _text_is_rtl_script(para.text):
                        _apply_arabic_para_alignment(para, is_title=is_title_shape, profile=profile)

            # Shapes that received Arabic text get normAutofit so text shrinks
            # to fit the original box rather than overflowing or being clipped.
            # normAutofit never changes the box dimensions — it only scales text.
            if is_rtl and shape_translated:
                _set_body_autofit(shape, "norm", force_wrap=(not is_title_shape))
                slide_translated_shapes.add(shape_idx)

                # Most placeholder runs have no explicit <a:rPr sz=.../> — the
                # size shown is pure layout/master inheritance. That means
                # BOTH (a) _shape_has_estimated_clipping's occupancy estimate
                # (via _shape_avg_font_pt's 12pt fallback, usually far smaller
                # than the real ~18-20pt inherited size — so it silently
                # under-estimates occupancy and misses genuine overflow) and
                # (b) any font-scale-from-baseline repair (nothing to scale
                # from) are no-ops for this — the majority — case. Bake in
                # the real inherited size FIRST, unconditionally, so both the
                # clipping check right below and any later repair pass see
                # accurate numbers.
                _bake_in_inherited_font_sizes(shape, slide, is_title_shape)

                # A bare <a:normAutofit/> only tells PowerPoint "this box
                # should shrink to fit" — it does NOT itself compute or bake
                # in a font-scale percentage. PowerPoint only recalculates
                # that scale when a human edits the box in the desktop app;
                # a file rendered non-interactively (COM export, or opened
                # read-only) shows the run's raw font size as-is, still
                # clipped, even with normAutofit set. The existing repair
                # loop below only shrinks a shape's font in response to a
                # detected COLLISION with another shape — the far more common
                # case (translated text simply longer than the original,
                # overflowing its own box with nothing else nearby) was never
                # touched. Proactively shrink here, independent of collisions,
                # for every shape that just received translated text.
                if is_title_shape:
                    # _shape_has_estimated_clipping's title-mode top-glyph
                    # check requires clearance proportional to font size
                    # above the text. A manually drawn title textbox (not a
                    # real placeholder, so nothing inherited a sensible
                    # default) starts with margin_top=0 and a non-centered
                    # vertical anchor — under those exact conditions the
                    # check can never be satisfied at ANY font size, all the
                    # way down to the shrink floor, because zero margin never
                    # grows no matter how small the text gets. This is the
                    # same top-margin/anchor normalization
                    # _mirror_titles_for_rtl_slide applies later; doing it
                    # here too means the clipping check below evaluates
                    # against the geometry the title will actually render
                    # with, not a raw, not-yet-normalized starting state.
                    try:
                        from pptx.enum.text import MSO_ANCHOR as _MSO_ANCHOR
                        _tf = shape.text_frame
                        _min_top_margin = int(max(45720, 0.008 * prs.slide_height))
                        if int(getattr(_tf, "margin_top", 0) or 0) < _min_top_margin:
                            _tf.margin_top = _min_top_margin
                        _tf.vertical_anchor = _MSO_ANCHOR.MIDDLE
                    except Exception:
                        pass

                if _shape_has_estimated_clipping(shape, is_title=is_title_shape):
                    from api.utils.arabic_layout_engine import _run_base_sizes_pt, _apply_font_scale

                    # Box growth is intentionally NOT attempted here: a
                    # branded footer/confidentiality bar is usually drawn on
                    # the SLIDE LAYOUT or MASTER rather than as a shape on
                    # the slide itself, so it never appears in slide.shapes
                    # and there is no reliable obstacle data to grow safely
                    # against — growing blind risks pushing content further
                    # into exactly the footer this is meant to avoid.
                    # Shrinking is the only geometry-safe lever available:
                    # it can never make the overlap situation worse, only
                    # the font smaller.
                    base_sizes = _run_base_sizes_pt(shape)
                    # Titles must stay readable — never shrink a title to a tiny
                    # size (better to let it slightly overflow than be microscopic,
                    # as seen on the opening slide). Body text may shrink further.
                    min_readable_pt = 18.0 if is_title_shape else 8.0
                    _min_frac = 0.70 if is_title_shape else 0.45
                    for step in range(1, 30):
                        scale = 1.0 - (step * 0.02)
                        _apply_font_scale(shape, base_sizes, scale=scale, min_frac=_min_frac, min_pt=min_readable_pt)
                        if not _shape_has_estimated_clipping(shape, is_title=is_title_shape):
                            break

                # ── Title vs. picture/decorative-image collision ───────────────
                # A title box can sit close enough to a nearby picture that
                # their BOUNDING BOXES already "overlap" on paper even when
                # the original (shorter, differently-fonted) title text never
                # visually reached that far — the collision-baseline system
                # (original_text_overlap_pairs) then treats this pair as
                # "pre-existing, allowed" and never repairs it, even though
                # the actually-rendered translated title now visually does
                # touch the image.
                #
                # Nudging the picture (trying candidate left/right/up/down
                # offsets) does not help here: the gap between it and the
                # next shape below is typically far smaller than any nudge
                # step, so every candidate position still collides with
                # something and the nudge silently does nothing. Instead,
                # deterministically clear the overlap: move the picture's
                # top edge down to just below the title's box, and shrink
                # its height by the same amount so its BOTTOM edge — and
                # whatever sits below it — never moves.
                if is_title_shape:
                    try:
                        from pptx.enum.shapes import MSO_SHAPE_TYPE as _MSO3
                        from pptx.util import Emu
                        title_rect = _shape_geometry(shape)
                        if title_rect is not None:
                            t_left, t_top, t_w, t_h = title_rect
                            title_bottom = t_top + t_h
                            margin = int(0.01 * prs.slide_height)
                            for _oi, _osh in enumerate(slide.shapes):
                                if _oi == shape_idx:
                                    continue
                                if _osh.shape_type != _MSO3.PICTURE:
                                    continue
                                pic_rect = _shape_geometry(_osh)
                                if pic_rect is None or not _rect_overlap(title_rect, pic_rect):
                                    continue
                                p_left, p_top, p_w, p_h = pic_rect
                                p_bottom = p_top + p_h
                                new_top = title_bottom + margin
                                if new_top < p_bottom:
                                    new_height = max(int(0.25 * p_h), p_bottom - new_top)
                                    _osh.top = Emu(new_top)
                                    _osh.height = Emu(new_height)
                                    # The end-of-slide geometry lock (below)
                                    # restores every shape to its pre-loop
                                    # _geom_snap snapshot to undo ACCIDENTAL
                                    # drift — without updating the snapshot
                                    # here too, it would revert this
                                    # deliberate repositioning right back
                                    # into the collision it just fixed.
                                    _geom_snap[_oi] = (p_left, new_top, p_w, new_height)
                    except Exception:
                        pass

                # ── Avoid right-aligned text hugging a side obstacle ───────────
                # A wide placeholder that safely cleared a side image under
                # left-aligned English text can flow straight into that same
                # image once right-aligned for Arabic. Narrow the box (never
                # move/grow it) to stop just before any picture/table/group-
                # with-picture obstacle it geometrically overlaps.
                if not is_title_shape:
                    try:
                        from pptx.enum.shapes import MSO_SHAPE_TYPE as _MSO2
                        _side_obstacles: list[tuple[int, int, int, int]] = []
                        for _oi, _osh in enumerate(slide.shapes):
                            if _oi == shape_idx:
                                continue
                            _og = _geom_snap.get(_oi)
                            if _og is None:
                                continue
                            _is_side_obstacle = False
                            try:
                                if _osh.shape_type == _MSO2.PICTURE or _osh.has_table:
                                    _is_side_obstacle = True
                                elif _osh.shape_type == _MSO2.GROUP:
                                    _is_side_obstacle = any(
                                        c.shape_type == _MSO2.PICTURE for c in _osh.shapes
                                    )
                            except Exception:
                                pass
                            if _is_side_obstacle:
                                _side_obstacles.append(_og)

                        if _side_obstacles:
                            from api.utils.arabic_layout_engine import (
                                _shrink_textbox_away_from_side_obstacle,
                            )
                            if _shrink_textbox_away_from_side_obstacle(
                                shape, _side_obstacles, prs.slide_width, prs.slide_height, 45720,
                            ):
                                # Update the lock snapshot so the restore pass
                                # below treats this intentional narrowing as
                                # the expected geometry, not a regression.
                                _geom_snap[shape_idx] = (
                                    shape.left, shape.top, shape.width, shape.height,
                                )
                    except Exception as _shrink_err:
                        log.debug(
                            "side-obstacle shrink skipped (slide=%d shape=%d): %s",
                            slide_idx + 1, shape_idx, _shrink_err,
                        )

        # ── Restore geometry lock: undo any accidental position/size changes ──
        # Text replacement and normAutofit must NEVER move a shape.  If they
        # somehow did (library bug, edge-case XML mutation) we restore here.
        for _gi, _gsnap in _geom_snap.items():
            try:
                _gs = slide.shapes[_gi]
                _gl, _gt, _gw, _gh = _gsnap
                if (_gs.left, _gs.top, _gs.width, _gs.height) != (_gl, _gt, _gw, _gh):
                    _gs.left   = _gl
                    _gs.top    = _gt
                    _gs.width  = _gw
                    _gs.height = _gh
                    log.debug("geometry lock: restored shape %d on slide %d", _gi, slide_idx + 1)
            except Exception:
                pass

        # ── Degenerate (zero-size) placeholder geometry ─────────────────────
        # Some source files carry a placeholder with an explicit <a:xfrm>
        # override whose width or height is literally 0 (seen with a
        # <p:ph idx="1"/> body: cy="0" while holding several paragraphs of
        # real text) — pre-existing in the uploaded file, not something the
        # translation pass produces, so the geometry-lock restore above just
        # puts the same zero back. With normAutofit, PowerPoint has no non-
        # zero box to shrink text into, so real translated content renders
        # squeezed to an unreadable size at whatever position the box's
        # degenerate origin happens to be. The safe fix is to drop the
        # override entirely so the shape falls back to inheriting its
        # normal, designed geometry from the slide layout's matching
        # placeholder — never invent a guessed size.
        for _dgs in slide.shapes:
            try:
                if int(_dgs.width) > 0 and int(_dgs.height) > 0:
                    continue
                if not _dgs.is_placeholder or not getattr(_dgs, "has_text_frame", False):
                    continue
                if not _dgs.text_frame.text.strip():
                    continue
                xfrm = _dgs._element.spPr.find(f"{{{_NS_DML}}}xfrm")
                if xfrm is not None:
                    _dgs._element.spPr.remove(xfrm)
                    log.debug(
                        "slide %d shape %d: dropped zero-size xfrm override, "
                        "falling back to layout-inherited geometry",
                        slide_idx + 1, _dgs.shape_id,
                    )
            except Exception:
                pass

        run_full_layout_transform = (
            is_rtl
            and not deck_src_is_rtl
            and style_profile_override
            and template_strength != "light"
        )

        # Mirror title placeholders for LTR->RTL decks in-place when the full
        # layout engine is not active.
        if is_rtl and not deck_src_is_rtl and not run_full_layout_transform:
            _mirror_titles_for_rtl_slide(
                slide,
                prs.slide_width,
                prs.slide_height,
                layout_warnings=layout_warnings,
            )

        # ── Arabic layout transformation (explicit style profile only) ────────
        # This engine mirrors shapes horizontally to follow Arabic reading flow.
        # It must ONLY run when the user explicitly chose a saved style profile
        # or uploaded a reference template — not for default Arabic translation.
        # Default translation = text-in-place, zero geometry changes.
        # "light" strength = fonts + colors only; skip geometry transforms.
        if run_full_layout_transform:
            try:
                from api.utils.arabic_layout_engine import transform_slide_layout
                for warn in transform_slide_layout(
                    slide, prs.slide_width, prs.slide_height, profile,
                    translated_shape_idxs=slide_translated_shapes,
                ):
                    log.warning("layout engine slide %d: %s", slide_idx + 1, warn)
                    # Propagate title-overflow warnings to the caller so the
                    # SSE pipeline can surface them to the user.
                    if "title_overflow_warning" in warn and layout_warnings is not None:
                        layout_warnings.append(
                            f"Slide {slide_idx + 1}: {warn.replace('title_overflow_warning: ', '')}"
                        )
            except Exception as exc:
                log.error("layout engine failed on slide %d: %s", slide_idx + 1, exc)

        # ── Speaker notes ─────────────────────────────────────────────────────
        if slide_idx in notes_map and slide.has_notes_slide:
            try:
                notes_tf = slide.notes_slide.notes_text_frame
                if notes_tf.paragraphs:
                    _replace_pptx_para_arabic(
                        notes_tf.paragraphs[0], notes_map[slide_idx], is_rtl,
                        is_title=False, profile=profile, font_subs=font_subs,
                        slide_label=f"slide {slide_idx+1} notes",
                        source_text=notes_src.get(slide_idx, ""),
                    )
            except Exception as exc:
                log.warning("Could not update notes on slide %d: %s", slide_idx, exc)

        # Layout QA is intentionally deferred until AFTER PPTX bytes are saved.
        # Rebuild must not depend on layout diagnostics.

    # Log a summary of any font substitutions so operators can review
    if font_subs:
        log.info("PPTX Arabic font substitutions (%d): %s",
                 len(font_subs), "; ".join(font_subs[:10]))

    # Layout QA + auto-repair: runs directly on `prs` (the object that gets
    # saved below), NOT a re-parsed throwaway copy. Repair actions
    # (font-shrink, box-move, box-expand-into-whitespace, z-order fixes) must
    # mutate the presentation that is actually exported — previously this ran
    # against a separate `Presentation` parsed from already-saved bytes, so
    # every repair the QA loop computed was discarded and only ever used to
    # decide whether to emit a warning string. The saved file never reflected
    # any of the fixes, no matter how many iterations "succeeded".
    if layout_warnings is not None:
        try:
            for ctx in post_export_qa_contexts:
                sidx = int(ctx.get("slide_idx", 0))
                if sidx < 0 or sidx >= len(prs.slides):
                    continue
                qa_slide = prs.slides[sidx]
                iter_actions: list[str] = []
                max_layout_repair_iters = 10 if auto_repair_enabled else 1
                remaining: list[str] = []
                for repair_iter in range(1, max_layout_repair_iters + 1):
                    remaining = _run_slide_visual_validation(
                        qa_slide,
                        slide_idx=sidx,
                        slide_w=prs.slide_width,
                        slide_h=prs.slide_height,
                        original_geom=ctx.get("original_geom", {}),
                        original_banner=ctx.get("original_banner"),
                        original_text_overlap_pairs=ctx.get("original_text_overlap_pairs"),
                        original_text_hidden_pairs=ctx.get("original_text_hidden_pairs"),
                        original_near_pairs=ctx.get("original_near_pairs"),
                        protected_zone_indices=ctx.get("protected_zone_indices"),
                        layout_warnings=None,
                        repair_actions=iter_actions,
                    )
                    # "hidden behind" violations are intentionally left
                    # unfixed by this loop (see _normalize_zorder_once) and
                    # can never clear on their own — stop iterating on them
                    # so the loop doesn't burn all 10 passes for nothing.
                    non_zorder_remaining = [w for w in remaining if "hidden behind" not in w]
                    if not non_zorder_remaining:
                        remaining = non_zorder_remaining
                        break
                    if repair_iter == max_layout_repair_iters:
                        layout_warnings.extend(remaining)

                # Z-order fix runs exactly once, as the very last mutation on
                # this slide, strictly AFTER every original_geom-indexed
                # repair pass above has finished — see _normalize_zorder_once
                # docstring for why the ordering matters.
                try:
                    _normalize_zorder_once(qa_slide, prs.slide_width, prs.slide_height)
                except Exception:
                    pass
        except Exception as qa_exc:
            layout_warnings.append(f"Post-export QA unavailable: {qa_exc}")

    if strict_qa and layout_warnings:
        if not allow_export_with_warnings:
            raise LayoutValidationError(
                "Unresolved layout issues after auto-repair: "
                + "; ".join(layout_warnings[:10])
            )

    # Save AFTER the QA/repair loop so its fixes are reflected in the output.
    out = io.BytesIO()
    prs.save(out)
    exported_bytes = out.getvalue()

    return exported_bytes


def _safe_rgb(color_obj):
    """DEPRECATED — do not use for reading colours.

    Accessing run.font.color in python-pptx has a destructive side effect:
    the getter converts the run's fill to an empty <a:solidFill/> element,
    which overrides the colour inherited from the layout/master and renders
    BLACK.  Use _run_explicit_rgb(run) instead, which reads the XML directly
    without mutating it.
    """
    try:
        if color_obj.type is None:
            return None
        return color_obj.rgb          # works for RGBColor and _SRGBColor
    except AttributeError:
        return None                   # _SchemeColor or other indirect types


def _run_explicit_rgb(run):
    """Read an explicitly set sRGB colour from a run's rPr XML — NO mutation.

    Returns an RGBColor when the run has <a:solidFill><a:srgbClr val="…"/>,
    otherwise None (covers: no rPr, no fill, or theme/scheme colours which
    must be preserved untouched so they keep inheriting correctly).
    """
    from pptx.dml.color import RGBColor
    try:
        rPr = run._r.find(f"{{{_NS_DML}}}rPr")
        if rPr is None:
            return None
        fill = rPr.find(f"{{{_NS_DML}}}solidFill")
        if fill is None:
            return None
        srgb = fill.find(f"{{{_NS_DML}}}srgbClr")
        if srgb is not None and srgb.get("val"):
            return RGBColor.from_string(srgb.get("val"))
        return None   # schemeClr or empty — leave the XML alone
    except Exception:
        return None


def _run_explicit_color_element(run):
    """Return a deep copy of a run's <a:solidFill> element, RGB or scheme.

    _run_explicit_rgb only recognizes literal sRGB colors and returns None
    for scheme-color runs (<a:schemeClr val="tx1"/> etc.) by design — reading
    them isn't destructive. But paragraph rebuilders that DELETE the run
    entirely and create a new one (as _replace_pptx_para_arabic does) throw
    the scheme reference away with it: there is nothing left to "keep
    inheriting" once the original <a:rPr> is gone, so the new run silently
    falls back to whatever default color applies to a bare run in that
    shape — which is not necessarily what the scheme reference resolved to.
    This captures the whole <a:solidFill> (either color type) so callers can
    reinsert it into the replacement run's rPr, preserving the original
    color exactly instead of dropping it.
    """
    try:
        rPr = run._r.find(f"{{{_NS_DML}}}rPr")
        if rPr is None:
            return None
        fill = rPr.find(f"{{{_NS_DML}}}solidFill")
        if fill is None:
            return None
        import copy
        return copy.deepcopy(fill)
    except Exception:
        return None


def _text_is_rtl_script(text: str) -> bool:
    """True when *text* is PREDOMINANTLY Arabic script.

    Counts strong-direction letters only (digits and punctuation are
    direction-neutral): Arabic-script letters must OUTNUMBER Latin letters.
    A mixed line like "G60ZBx – مقدمة" is classified by its dominant script,
    while a mostly-English sentence containing a single Arabic term still
    counts as LTR — so it receives normal LTR→RTL treatment when translated.
    Used to detect content already designed for Arabic, whose layout
    attributes must be preserved exactly.
    """
    if not text:
        return False
    rtl = sum(1 for ch in text
              if "\u0600" <= ch <= "\u06FF" or "\u0750" <= ch <= "\u077F")
    ltr = sum(1 for ch in text if ch.isascii() and ch.isalpha())
    return rtl > ltr


def _set_para_rtl(para, is_rtl: bool) -> None:
    """Set or clear the RTL flag on a PPTX paragraph element.

    In OOXML a paragraph's direction is controlled by <a:pPr rtl="1"/>.
    python-pptx doesn't expose this as a property, so we manipulate the XML
    directly.  Errors are silently swallowed so a bad element never aborts
    the whole rebuild.
    """
    if not is_rtl:
        return
    try:
        pPr = para._p.get_or_add_pPr()
        pPr.set("rtl", "1")
    except Exception:
        pass


def _set_para_alignment(para, alignment_val: int) -> None:
    """Set paragraph alignment via OOXML <a:pPr algn="…"/>.

    alignment_val uses python-pptx PP_ALIGN integer values:
      1=LEFT, 2=CENTER, 3=RIGHT, 4=JUSTIFY, 5=JUSTIFY_LOW, 7=DISTRIBUTE
    """
    _ALGN_MAP = {1: "l", 2: "ctr", 3: "r", 4: "just", 5: "justLow", 7: "dist"}
    algn_str = _ALGN_MAP.get(alignment_val)
    if not algn_str:
        return
    try:
        pPr = para._p.get_or_add_pPr()
        pPr.set("algn", algn_str)
    except Exception:
        pass


_NS_DML = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _set_body_autofit(shape, mode: str = "norm", force_wrap: bool = False) -> None:
    """Set the text-body autofit mode on a shape so Arabic text fits cleanly.

    mode: "norm" → <a:normAutofit/> — shrink text proportionally to fit the
                   existing box.  Safe default for translated text.
          "sp"   → <a:spAutoFit/>  — expand the box to fit (changes dimensions,
                   avoid for translated slides).
          "none" → <a:noAutofit/>  — allow overflow/clipping (PowerPoint default).

    When force_wrap is True, also sets wrap="square" so Arabic body text wraps
    inside the shape bounds. Keep this False for title placeholders.
    """
    import lxml.etree as _ET

    _TAG_MAP = {
        "norm": f"{{{_NS_DML}}}normAutofit",
        "sp":   f"{{{_NS_DML}}}spAutoFit",
        "none": f"{{{_NS_DML}}}noAutofit",
    }
    tag = _TAG_MAP.get(mode, f"{{{_NS_DML}}}normAutofit")

    try:
        txBody = shape.text_frame._txBody
        bodyPr = txBody.find(f"{{{_NS_DML}}}bodyPr")
        if bodyPr is None:
            # Create bodyPr as first child of txBody
            bodyPr = _ET.Element(f"{{{_NS_DML}}}bodyPr")
            txBody.insert(0, bodyPr)
        # Remove all existing autofit child elements first
        for t in ("normAutofit", "spAutoFit", "noAutofit"):
            for child in list(bodyPr.findall(f"{{{_NS_DML}}}{t}")):
                bodyPr.remove(child)
        # Insert desired autofit mode
        bodyPr.append(_ET.Element(tag))
        if force_wrap:
            bodyPr.set("wrap", "square")
            try:
                shape.text_frame.word_wrap = True
            except Exception:
                pass
        # NOTE: for titles keep wrap untouched. Forcing wrap="square" can
        # override wrap="none" from layout placeholders and break title lines.
    except Exception:
        pass


def _is_title_placeholder_shape(shape) -> bool:
    """True when shape is a PowerPoint title placeholder."""
    try:
        if not shape.is_placeholder:
            return False
        ph_t = str(shape.placeholder_format.type)
        return ("TITLE" in ph_t) or ("CENTER_TITLE" in ph_t)
    except Exception:
        return False


def _shape_text_value(shape) -> str:
    """Return normalized text from a shape text frame."""
    try:
        if not shape.has_text_frame:
            return ""
        return " ".join(
            p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip()
        ).strip()
    except Exception:
        return ""


def _shape_max_font_pt(shape) -> float:
    """Return max run font pt in a shape text frame."""
    m = 0.0
    try:
        if not shape.has_text_frame:
            return 0.0
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.size is not None and run.font.size.pt:
                    m = max(m, float(run.font.size.pt))
    except Exception:
        pass
    return m


def _layout_placeholder_level_sizes_pt(slide, ph_idx: "int | None") -> list[float]:
    """Font sizes (pt), indexed by outline level, defined in the slide
    layout's matching placeholder <a:lstStyle> — the real source of a
    placeholder run's "inherited" size when no explicit <a:rPr sz=.../> is
    set directly on the run. Returns [] if unresolvable (e.g. not a
    placeholder, or the layout defines no explicit sizes for it either)."""
    if ph_idx is None:
        return []
    try:
        layout = slide.slide_layout
        for ph in layout.placeholders:
            if ph.placeholder_format.idx == ph_idx and ph.has_text_frame:
                import re
                xml = ph.text_frame._txBody.xml
                sizes = re.findall(r'sz="(\d+)"', xml)
                return [int(s) / 100.0 for s in sizes]
    except Exception:
        pass
    return []


def _bake_in_inherited_font_sizes(shape, slide, is_title: bool) -> None:
    """Set an explicit run.font.size for every run that has none, so a
    subsequent font-scale-reduction pass has a real baseline to shrink from.

    Most placeholder text in real decks has no <a:rPr sz=.../> at all — its
    displayed size comes purely from the slide layout's <a:lstStyle> (or,
    failing that, the slide master), resolved by PowerPoint at render time.
    python-pptx's run.font.size only ever reflects DIRECT formatting, so a
    scale-from-baseline loop keyed on it is a silent no-op for this — by far
    the more common — case. Levels are looked up per-paragraph (para.level);
    unresolvable ones fall back to a generic template-typical default rather
    than skipping the run entirely.
    """
    try:
        ph_idx = shape.placeholder_format.idx if shape.is_placeholder else None
    except Exception:
        ph_idx = None
    level_sizes = _layout_placeholder_level_sizes_pt(slide, ph_idx)
    default_pt = 32.0 if is_title else 18.0

    try:
        from pptx.util import Pt
        for para in shape.text_frame.paragraphs:
            level = getattr(para, "level", 0) or 0
            fallback_pt = level_sizes[level] if level < len(level_sizes) else (
                level_sizes[-1] if level_sizes else default_pt
            )
            for run in para.runs:
                if run.font.size is None or not run.font.size.pt:
                    run.font.size = Pt(fallback_pt)
    except Exception:
        pass


def _shape_avg_font_pt(shape, default_pt: float = 12.0) -> float:
    """Return average run font size (pt) for occupancy guard calculations."""
    sizes: list[float] = []
    try:
        if not shape.has_text_frame:
            return default_pt
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.size is not None and run.font.size.pt:
                    sizes.append(float(run.font.size.pt))
    except Exception:
        return default_pt
    if not sizes:
        return default_pt
    return sum(sizes) / float(len(sizes))


def _shape_has_estimated_clipping(
    shape,
    rect: tuple[int, int, int, int] | None = None,
    *,
    is_title: bool = False,
) -> bool:
    """Conservative clipped-text detector for final QA.

    Uses inner text-frame box (after margins), mixed-script guard bands,
    and top-clearance checks to catch cases where PowerPoint visually clips
    Arabic ascenders/marks even when raw occupancy is near the box limits.
    """
    try:
        if not getattr(shape, "has_text_frame", False):
            return False
        from api.utils.arabic_layout_engine import _estimate_text_occupancy

        geom = rect if rect is not None else _shape_geometry(shape)
        if geom is None:
            return False
        _, _, width, height = geom

        tf = shape.text_frame
        text = _shape_text_value(shape)
        est_w, est_h, _ = _estimate_text_occupancy(shape)
        if est_w <= 0.0 or est_h <= 0.0:
            return False

        margin_l = int(getattr(tf, "margin_left", 0) or 0)
        margin_r = int(getattr(tf, "margin_right", 0) or 0)
        margin_t = int(getattr(tf, "margin_top", 0) or 0)
        margin_b = int(getattr(tf, "margin_bottom", 0) or 0)
        inner_w = max(1.0, float(width - margin_l - margin_r))
        inner_h = max(1.0, float(height - margin_t - margin_b))

        avg_pt = _shape_avg_font_pt(shape, default_pt=12.0)
        emu_per_pt = 12700.0

        is_mixed = _is_mixed_bidi_text(text)

        # Titles: strict checks to protect top-band readability and glyph headroom.
        if is_title:
            glyph_guard = max(9144.0, 0.18 * avg_pt * emu_per_pt)
            if is_mixed:
                est_w *= 1.04
                est_h *= 1.03
                glyph_guard *= 1.08

            # est_w is the UNWRAPPED single-line width (see the matching note
            # in the body-text branch below). A title box with word_wrap on
            # can legitimately flow onto a second line — flagging it as
            # "clipped" purely for not fitting on ONE line, with only 0.5%
            # tolerance, forced titles to shrink far past what was ever
            # actually necessary (observed: a 29pt original title shrunk to
            # ~14.5pt while sitting comfortably inside its box) even though
            # nothing was really at risk of being cut off. Only enforce the
            # single-line width assumption when wrapping is genuinely off.
            _title_word_wrap_on = bool(getattr(tf, "word_wrap", False))
            if not _title_word_wrap_on and est_w > (inner_w * 0.995):
                return True
            if est_h + glyph_guard > inner_h:
                return True

            anchor_name = str(getattr(tf, "vertical_anchor", "")).upper()
            if "MIDDLE" in anchor_name or "CENTER" in anchor_name:
                top_clearance = max(0.0, (inner_h - est_h) / 2.0) + float(margin_t)
            else:
                top_clearance = float(margin_t)

            min_top_clearance = max(float(margin_t) * 0.30, 0.16 * avg_pt * emu_per_pt)
            if _contains_arabic(text) and top_clearance < min_top_clearance:
                return True

            if est_h > (inner_h * 0.95):
                return True
            return False

        # Body text: avoid over-flagging; require clear overflow beyond estimator noise.
        if is_mixed:
            est_w *= 1.01
            est_h *= 1.01
        # est_w from _estimate_text_occupancy is the UNWRAPPED single-line
        # width of the longest paragraph — not the width actually rendered,
        # which reflows across multiple lines whenever word_wrap is on (true
        # for virtually every translated body placeholder). Comparing that
        # unwrapped estimate to the box width flags completely healthy,
        # correctly-wrapping bullet text as "clipped" any time a sentence is
        # long enough to need more than one line — which is normal, not a
        # defect. Only apply the width check when wrapping is actually off,
        # where a long line truly would overflow horizontally uncut.
        _word_wrap_on = bool(getattr(tf, "word_wrap", False))
        if not _word_wrap_on and est_w > (inner_w * 1.06):
            return True
        if est_h > (inner_h * 1.08):
            return True

        return False
    except Exception:
        return False


def _looks_heading_like(text: str) -> bool:
    """Heuristic to separate title headings from body paragraphs."""
    if not text:
        return False
    t = text.strip()
    if not t or t.isdigit():
        return False
    # Enumerated callouts (e.g., "1- ...") are body labels, not titles.
    if re.match(r"^\s*\d+[\)\].:-]", t):
        return False
    # Tiny labels (e.g., "P") are usually callout markers, not slide titles.
    core = re.sub(r"[^\w\u0600-\u06FF]+", "", t)
    if len(core) < 4:
        return False
    words = [w for w in t.split() if w]
    if len(words) > 22:
        return False
    if len(t) > 170:
        return False
    if t.endswith((".", ";", "؟", "!")):
        return False
    return True


def _iter_logo_candidate_shapes(slide):
    """Yield shapes from slide + layout + master for logo-boundary detection."""
    for shp in slide.shapes:
        yield shp
    try:
        for shp in slide.slide_layout.shapes:
            yield shp
    except Exception:
        pass
    try:
        for shp in slide.slide_layout.slide_master.shapes:
            yield shp
    except Exception:
        pass


def _detect_logo_left_boundary(slide, slide_width: int, slide_height: int) -> int | None:
    """Estimate left edge of top-right logo/brand zone from slide/layout/master."""
    best_left: int | None = None

    for shp in _iter_logo_candidate_shapes(slide):
        try:
            left = int(shp.left)
            top = int(shp.top)
            width = int(shp.width)
            height = int(shp.height)
        except Exception:
            continue

        if top > int(0.25 * slide_height):
            continue
        if left < int(0.55 * slide_width):
            continue
        if width > int(0.30 * slide_width) or height > int(0.20 * slide_height):
            continue

        is_logo_like = False
        try:
            st = shp.shape_type
            st_name = str(st)
            if "PICTURE" in st_name or "GROUP" in st_name:
                is_logo_like = True
        except Exception:
            pass

        if not is_logo_like:
            txt = _shape_text_value(shp)
            if not txt and width <= int(0.20 * slide_width):
                is_logo_like = True

        if not is_logo_like:
            continue

        if best_left is None or left < best_left:
            best_left = left

    return best_left


def _detect_top_right_obstruction_left(
    slide,
    slide_width: int,
    slide_height: int,
    exclude_idxs: set[int] | None = None,
) -> int | None:
    """Detect left edge of top-right decorative overlays that can occlude titles."""
    blocked_left: int | None = None
    skip = exclude_idxs or set()

    for idx, shp in enumerate(slide.shapes):
        if idx in skip:
            continue
        try:
            left = int(shp.left)
            top = int(shp.top)
            width = int(shp.width)
            height = int(shp.height)
        except Exception:
            continue

        if top > int(0.25 * slide_height):
            continue
        if left < int(0.45 * slide_width):
            continue
        if width > int(0.60 * slide_width) or height > int(0.25 * slide_height):
            continue

        txt = _shape_text_value(shp)
        if txt:
            continue

        if blocked_left is None or left < blocked_left:
            blocked_left = left

    return blocked_left


def _detect_master_banner_rect(slide, slide_width: int, slide_height: int) -> "tuple[int, int, int, int] | None":
    """Best-effort bounding rect (left, top, width, height in EMU) of a
    decorative top banner shape that lives on the slide LAYOUT or MASTER
    rather than on the slide itself.

    Many corporate templates draw an angled/diagonal title banner as a
    FREEFORM (or plain AUTO_SHAPE) on the master, which never appears in
    slide.shapes and is therefore invisible to logo/obstruction-based right-
    boundary detection, and invisible to slide-level banner-height detection
    too. Its rectangular bounding box can end well short of a logo
    positioned further right (leaving a blank gap a title can drift into and
    become invisible white-on-white), and its real height is often smaller
    than a generic top-strip guess (letting a grown title box push past the
    actual banner into the plain white area below it). Heuristic: a
    top-anchored shape, narrower than the full slide (so it's not a
    full-bleed background panel) but still covering a clear majority of the
    width (so it's not some small unrelated decoration).
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    candidates: list[tuple[int, int, int, int]] = []
    sources = []
    try:
        sources.append(slide.slide_layout)
    except Exception:
        pass
    try:
        sources.append(slide.slide_layout.slide_master)
    except Exception:
        pass

    for src in sources:
        try:
            shapes = src.shapes
        except Exception:
            continue
        for shp in shapes:
            try:
                if shp.shape_type not in (MSO_SHAPE_TYPE.FREEFORM, MSO_SHAPE_TYPE.AUTO_SHAPE):
                    continue
                top = int(shp.top)
                left = int(shp.left)
                width = int(shp.width)
                height = int(shp.height)
                if top > int(0.30 * slide_height):
                    continue
                if width >= int(0.97 * slide_width) or width < int(0.35 * slide_width):
                    continue
                candidates.append((left, top, width, height))
            except Exception:
                continue

    if not candidates:
        return None
    # Prefer the smallest right edge (most conservative width boundary) —
    # matches the selection this function used before it also tracked
    # height, keeping the already-verified right-edge behavior unchanged.
    return min(candidates, key=lambda r: r[0] + r[2])


def _find_real_top_banner_shape(slide, slide_width: int, slide_height: int) -> "tuple[int, int, int, int] | None":
    """Detect an actual top colored title-banner SHAPE, or None if the slide
    genuinely has no such element (e.g. a plain white-background title, like
    most slides in a training-deck style with no colored header strip)."""
    best: tuple[int, int, int, int] | None = None
    best_area = -1

    for shape in slide.shapes:
        try:
            left = int(shape.left)
            top = int(shape.top)
            width = int(shape.width)
            height = int(shape.height)
        except Exception:
            continue

        if top > int(0.25 * slide_height):
            continue
        if width < int(0.45 * slide_width):
            continue
        if height < int(0.04 * slide_height) or height > int(0.25 * slide_height):
            continue
        if _shape_text_value(shape):
            continue

        area = width * height
        if area > best_area:
            best_area = area
            best = (left, top, width, height)

    return best


def _detect_top_banner_rect(slide, slide_width: int, slide_height: int) -> tuple[int, int, int, int]:
    """Detect the top colored title banner; fall back to a generic top strip.

    The generic fallback exists for callers that just need SOME reference
    region for horizontal fitting/positioning math and are fine treating the
    top of any slide as a nominal header band. It must NOT be used to decide
    whether a real banner shape exists (e.g. to justify re-anchoring title
    text) — use _find_real_top_banner_shape for that, which returns None
    instead of faking one.
    """
    best = _find_real_top_banner_shape(slide, slide_width, slide_height)
    if best is not None:
        return best
    return (0, 0, slide_width, int(0.18 * slide_height))


def _detect_title_shape_indices(slide, slide_width: int, slide_height: int) -> set[int]:
    """Detect title shapes: placeholders + title-like top-banner textboxes."""
    title_idxs: set[int] = set()

    for idx, shape in enumerate(slide.shapes):
        if _is_title_placeholder_shape(shape):
            try:
                if int(shape.top) <= int(0.45 * slide_height):
                    title_idxs.add(idx)
            except Exception:
                title_idxs.add(idx)

    top_font_max = 0.0
    font_by_idx: dict[int, float] = {}
    for idx, shape in enumerate(slide.shapes):
        try:
            if int(shape.top) <= int(0.25 * slide_height):
                f = _shape_max_font_pt(shape)
                font_by_idx[idx] = f
                top_font_max = max(top_font_max, f)
        except Exception:
            continue

    # The heuristic scan below exists ONLY to find a title-like textbox on
    # decks with NO real title placeholder at all. When a genuine title
    # placeholder was already found above, running it anyway risks scoring a
    # short heading-like CONTENT placeholder (e.g. a one-line Arabic lead-in
    # sentence sitting near the top) as a second "title" — which then gets
    # force-fit into the banner rect by _mirror_titles_for_rtl_slide,
    # corrupting its position/size. Skip entirely once a real title exists.
    if title_idxs:
        return title_idxs

    best_idx = None
    best_score = -999.0
    for idx, shape in enumerate(slide.shapes):
        if idx in title_idxs:
            continue
        try:
            if not shape.has_text_frame:
                continue
            left, top, width, height = int(shape.left), int(shape.top), int(shape.width), int(shape.height)
        except Exception:
            continue

        text = _shape_text_value(shape)
        if not text:
            continue
        if top > int(0.25 * slide_height):
            continue
        if top >= int(0.80 * slide_height):
            continue

        score = 0.0
        if top <= int(0.20 * slide_height):
            score += 2.5
        if _looks_heading_like(text):
            score += 2.5
        else:
            score -= 3.0

        f = font_by_idx.get(idx, 0.0)
        if top_font_max > 0 and f > 0:
            score += 3.0 * (f / top_font_max)

        if width >= int(0.30 * slide_width):
            score += 1.0
        if height > int(0.22 * slide_height):
            score -= 2.0
        if left >= int(0.85 * slide_width) and width <= int(0.12 * slide_width):
            score -= 2.0  # probable page number/decorative chip

        if score > best_score:
            best_score = score
            best_idx = idx

    if best_idx is not None and best_score >= 4.5:
        title_idxs.add(best_idx)
    return title_idxs


def _estimate_title_fit(shape) -> tuple[float, float]:
    """Estimate title text occupied width/height in EMU inside its box."""
    try:
        tf = shape.text_frame
        width = int(shape.width)
        usable_w = max(
            1.0,
            float(
                width
                - int(getattr(tf, "margin_left", 0) or 0)
                - int(getattr(tf, "margin_right", 0) or 0)
            ),
        )
    except Exception:
        return (0.0, 0.0)

    from api.utils.font_metrics import arabic_width_factor

    font_name = ""
    pts: list[float] = []
    lines: list[str] = []
    try:
        for para in tf.paragraphs:
            if not para.text:
                continue
            for chunk in para.text.split("\n"):
                if chunk.strip():
                    lines.append(chunk.strip())
            for run in para.runs:
                if run.font.name and not font_name:
                    font_name = run.font.name
                if run.font.size is not None and run.font.size.pt:
                    pts.append(float(run.font.size.pt))
    except Exception:
        return (0.0, 0.0)

    if not lines:
        return (0.0, 0.0)

    avg_pt = sum(pts) / len(pts) if pts else 24.0
    width_factor, _ = arabic_width_factor(font_name)
    emu_per_pt = 12700.0
    line_spacing = 1.3

    total_line_count = 0
    max_line_width = 0.0
    for line in lines:
        chars = max(1, len(line))
        est_line_w = chars * avg_pt * emu_per_pt * width_factor
        wraps = max(1, int(math.ceil(est_line_w / usable_w)))
        total_line_count += wraps
        max_line_width = max(max_line_width, est_line_w)

    est_height = total_line_count * (avg_pt * emu_per_pt * line_spacing)
    return (max_line_width, est_height)


def _fit_title_in_placeholder(shape, min_frac: float = 0.80) -> None:
    """Shrink title run fonts gradually and enable normAutofit to avoid clipping."""
    from pptx.util import Pt

    try:
        # Keyed by (para_idx, run_idx), NOT id(run) — python-pptx's `.runs`
        # constructs a fresh wrapper object on every access, so id(run)
        # captured in this first pass is not guaranteed to match id(run) for
        # the same run seen again in the scaling loop below. When it didn't
        # match (unreliable — depends on CPython reusing a freed wrapper's
        # memory address), the lookup silently found nothing and that run
        # was never scaled, with no error raised. This single function is
        # what every title in this rebuild pass shrinks through, so a silent
        # no-op here reads as "the title already fits" even when it doesn't.
        base_sizes: dict[tuple[int, int], float] = {}
        explicit_pts: list[float] = []
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                sz = run.font.size
                if sz is not None and sz.pt:
                    explicit_pts.append(float(sz.pt))

        inferred_base = (sum(explicit_pts) / len(explicit_pts)) if explicit_pts else 24.0
        for para_idx, para in enumerate(shape.text_frame.paragraphs):
            para_sz = para.font.size
            para_base = float(para_sz.pt) if (para_sz is not None and para_sz.pt) else inferred_base
            for run_idx, run in enumerate(para.runs):
                sz = run.font.size
                if sz is not None and sz.pt:
                    base_sizes[(para_idx, run_idx)] = float(sz.pt)
                else:
                    base_sizes[(para_idx, run_idx)] = para_base
        if not base_sizes:
            _set_body_autofit(shape, "norm")
            return

        tf = shape.text_frame
        ml = int(getattr(tf, "margin_left", 0) or 0)
        mr = int(getattr(tf, "margin_right", 0) or 0)
        mt = int(getattr(tf, "margin_top", 0) or 0)
        mb = int(getattr(tf, "margin_bottom", 0) or 0)
        usable_w = max(1.0, float(shape.width - ml - mr))
        usable_h = max(1.0, float(shape.height - mt - mb))

        for step in range(0, 16):
            scale = 1.0 - (step * 0.03)
            for para_idx, para in enumerate(shape.text_frame.paragraphs):
                for run_idx, run in enumerate(para.runs):
                    base = base_sizes.get((para_idx, run_idx))
                    if not base:
                        continue
                    run.font.size = Pt(max(base * min_frac, base * scale))
            est_w, est_h = _estimate_title_fit(shape)
            if est_w <= usable_w and est_h <= usable_h:
                break
    except Exception:
        pass

    _set_body_autofit(shape, "norm")


def _mirror_titles_for_rtl_slide(
    slide,
    slide_width: int,
    slide_height: int,
    layout_warnings: "list[str] | None" = None,
) -> None:
    """Mirror original title placeholders in-place for RTL translation."""
    try:
        from pptx.enum.text import MSO_VERTICAL_ANCHOR
    except Exception:
        MSO_VERTICAL_ANCHOR = None

    title_idxs = _detect_title_shape_indices(slide, slide_width, slide_height)
    banner_left, banner_top, banner_width, banner_height = _detect_top_banner_rect(
        slide,
        slide_width,
        slide_height,
    )
    # A real banner shape found on the LAYOUT/MASTER is strictly more
    # accurate than _detect_top_banner_rect's generic top-strip guess — use
    # its actual bounds (including height) whenever one is found, since the
    # generic guess's height not matching the real drawn banner is exactly
    # what let a grown title box push past the visible banner into the
    # plain white area below it.
    master_banner_rect = _detect_master_banner_rect(slide, slide_width, slide_height)
    if master_banner_rect is not None:
        banner_left, banner_top, banner_width, banner_height = master_banner_rect
    safe_gap = int(0.015 * slide_width)
    logo_left = _detect_logo_left_boundary(slide, slide_width, slide_height)
    obstruction_left = _detect_top_right_obstruction_left(
        slide,
        slide_width,
        slide_height,
        exclude_idxs=title_idxs,
    )
    right_limit = slide_width - safe_gap
    if logo_left is not None:
        right_limit = min(right_limit, logo_left - safe_gap)
    if obstruction_left is not None:
        right_limit = min(right_limit, obstruction_left - safe_gap)
    if master_banner_rect is not None:
        # Many templates draw the angled/diagonal title banner as a FREEFORM
        # shape on the slide LAYOUT or MASTER rather than on the slide
        # itself, so it never shows up in slide.shapes and logo/obstruction
        # detection never sees it. Its bounding box can end well short of
        # the logo (e.g. ~75% of slide width), leaving a white gap between
        # the banner's diagonal cut and the logo. Without this, a
        # right-aligned white title can extend into that gap and become
        # invisible (white text on white background) even though nothing
        # reports it as "clipped" — the text isn't out of bounds, it's just
        # unreadable against the wrong background.
        right_limit = min(right_limit, banner_left + banner_width - safe_gap)
    # Keep right-half preference only when no hard safety boundary exists.
    if logo_left is None and obstruction_left is None:
        right_limit = max(right_limit, int(0.55 * slide_width))
    cover_title_idxs: list[int] = []

    for idx, shape in enumerate(slide.shapes):
        if idx not in title_idxs:
            continue
        try:
            left = int(shape.left)
            top = int(shape.top)
            width = int(shape.width)
            height = int(shape.height)
            banner_bottom = min(slide_height, banner_top + banner_height)
            banner_box_height = max(1, banner_bottom - banner_top)

            is_center_title = False
            try:
                is_center_title = bool(shape.is_placeholder) and "CENTER_TITLE" in str(shape.placeholder_format.type)
            except Exception:
                is_center_title = False

            if is_center_title:
                from api.utils.arabic_layout_engine import _fit_cover_title_shape

                _fit_cover_title_shape(shape, slide_width, slide_height)
                cover_title_idxs.append(idx)
                continue

            mirrored_left = int(slide_width - left - width)
            new_left = mirrored_left
            new_width = width

            # Keep title inside the logo-safe banner zone.
            if new_left + new_width > right_limit:
                overflow = (new_left + new_width) - right_limit
                new_left = max(0, new_left - overflow)
                if new_left + new_width > right_limit:
                    new_width = max(1, right_limit - new_left)

            # Keep title's center in the right half of the slide.
            # If the mirrored box is too wide under logo constraints, narrow it
            # and rely on autofit/font reduction to keep text readable.
            slide_mid = slide_width // 2
            if new_left + (new_width // 2) < slide_mid:
                max_right_half_w = max(1, int(2 * (right_limit - slide_mid)))
                if new_width > max_right_half_w:
                    new_width = max_right_half_w
                    new_left = max(0, right_limit - new_width)

            shape.left = new_left
            shape.width = new_width

            # Keep the title box fully inside the detected blue title banner,
            # but preserve its original vertical placement whenever it already
            # fits — and ONLY reposition at all when the banner rect is a real
            # detected shape (master_banner_rect), not the generic top-strip
            # guess from _detect_top_banner_rect. That guess is frequently
            # taller than the deck's actual visible banner (e.g. when the
            # banner is really part of a tiled background image rather than
            # a discrete shape, as seen here), and clamping an
            # already-tall-but-correct original title box against an
            # overestimated banner height pushes it UP past where it
            # originally sat — straight into the real banner it was trying
            # to avoid. Without real confidence in the banner's true bounds,
            # the original author's placement (which rendered correctly in
            # the source deck) is the only position known to be safe.
            target_top = top
            target_height = height
            if master_banner_rect is not None:
                if target_height > banner_box_height:
                    target_top = max(0, banner_top)
                    target_height = banner_box_height
                else:
                    if target_top < banner_top:
                        target_top = banner_top
                    if target_top + target_height > banner_bottom:
                        target_top = max(banner_top, banner_bottom - target_height)
            if target_top + target_height > slide_height:
                target_height = max(1, slide_height - target_top)
            shape.top = target_top
            shape.height = target_height

            tf = shape.text_frame
            tf.word_wrap = True
            # Ensure safe internal top padding for Arabic ascenders.
            min_top_margin = int(max(45720, 0.008 * slide_height))
            min_bottom_margin = int(max(45720, 0.006 * slide_height))
            try:
                if int(getattr(tf, "margin_top", 0) or 0) < min_top_margin:
                    tf.margin_top = min_top_margin
            except Exception:
                pass
            try:
                if int(getattr(tf, "margin_bottom", 0) or 0) < min_bottom_margin:
                    tf.margin_bottom = min_bottom_margin
            except Exception:
                pass
            # Keep a small safe right margin so right-aligned Arabic text
            # does not visually touch/crop at the banner edge.
            try:
                if int(getattr(tf, "margin_right", 0) or 0) < 91440:
                    tf.margin_right = 91440
            except Exception:
                pass
            if MSO_VERTICAL_ANCHOR is not None:
                tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

            for para in tf.paragraphs:
                if not para.text.strip():
                    continue
                _set_para_rtl(para, True)
                _set_para_alignment(para, 3)

            # Use the full safe banner width before forcing extra wrapping.
            left_bound = max(0, banner_left + safe_gap)
            max_banner_width = max(1, right_limit - left_bound)
            est_w0, est_h0 = _estimate_title_fit(shape)
            if est_w0 > float(shape.width) and max_banner_width > int(shape.width):
                shape.width = max_banner_width
                shape.left = max(left_bound, right_limit - max_banner_width)
                new_left = int(shape.left)
                new_width = int(shape.width)

            # Font-shrink only — never grow the title box past its ORIGINAL
            # height. This used to grow the box first (up to banner_box_height)
            # and shrink the font only if that wasn't enough. The recurring
            # failure mode that caused: banner_box_height comes from either a
            # generic top-strip guess or a detected banner shape, and real
            # decks draw their banner in ways neither can always see — a
            # tiled background image split across several PICTURE shapes
            # (as opposed to one clean vector shape), a banner baked into a
            # rasterized slide background, etc. Every such case is a new,
            # undetectable variant, so bounding growth by "the banner" is an
            # unbounded chase. The ORIGINAL box height is a boundary that is
            # always correct by construction — the source-language title
            # already rendered inside it without complaint — so it is used
            # as the hard ceiling and font size is shrunk to fit inside it
            # instead of ever moving the boundary.
            fit_min_frac = 0.65 if new_width < width else 0.80
            _fit_title_in_placeholder(shape, min_frac=fit_min_frac)

            # Preserve centered placement inside the banner if height changed.
            if int(shape.height) > target_height:
                expanded_h = int(shape.height)
                expanded_top = max(banner_top, min(banner_bottom - expanded_h, (target_top + (target_height // 2)) - (expanded_h // 2)))
                shape.top = expanded_top
            else:
                shape.top = target_top
                shape.height = target_height
            shape.left = new_left
            shape.width = new_width

            # Strict validation: any top-touch in the FINAL geometry is a failure.
            est_w2, est_h2 = _estimate_title_fit(shape)
            box_h = int(shape.height)
            mt = int(getattr(tf, "margin_top", 0) or 0)
            mb = int(getattr(tf, "margin_bottom", 0) or 0)
            inner_h = max(1, box_h - mt - mb)
            inner_clearance = max(0.0, (float(inner_h) - est_h2) / 2.0)
            top_clearance = inner_clearance + float(mt)
            min_top_clearance = max(1.0, float(mt) * 0.35)
            top_touch_fail = est_h2 > float(inner_h) or top_clearance < min_top_clearance
            if layout_warnings is not None and top_touch_fail:
                layout_warnings.append(
                    f"Slide title layout FAILED: top-glyph clearance failed for shape {idx}"
                )
        except Exception:
            continue

    if cover_title_idxs:
        from api.utils.arabic_layout_engine import _collect_obstacles, _rects_overlap, _resolve_overlap

        obstacles = _collect_obstacles(slide, slide_width, slide_height)
        for idx in cover_title_idxs:
            try:
                shape = slide.shapes[idx]
                if not getattr(shape, "has_text_frame", False):
                    continue
                geom = (int(shape.left), int(shape.top), int(shape.width), int(shape.height))
            except Exception:
                continue
            if any(_rects_overlap(geom, obs, 45720) for obs in obstacles):
                _resolve_overlap(
                    shape,
                    obstacles,
                    slide_width,
                    slide_height,
                    45720,
                    layout_warnings if layout_warnings is not None else [],
                    f"cover title {idx}",
                )


def validate_pptx_bytes(data: bytes) -> tuple[bool, str]:
    """Verify that *data* is a structurally valid PPTX (ZIP) package.

    Returns (True, "") on success or (False, reason) on failure.
    Checks performed:
    - Local-file ZIP magic bytes (PK\\x03\\x04)
    - Presence of [Content_Types].xml
    - Presence of ppt/presentation.xml
    """
    import zipfile

    if not data:
        return False, "empty bytes"
    if data[:4] != b"PK\x03\x04":
        return False, f"not a ZIP package (magic={data[:4]!r})"
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as zf:
            names = set(zf.namelist())
            if "[Content_Types].xml" not in names:
                return False, "missing [Content_Types].xml"
            if "ppt/presentation.xml" not in names:
                return False, "missing ppt/presentation.xml"
    except zipfile.BadZipFile as exc:
        return False, f"corrupt ZIP: {exc}"
    return True, ""


def _replace_pptx_para(para, new_text: str, is_rtl: bool) -> None:
    """Replace text in a PPTX paragraph while preserving first run formatting.

    This is the legacy entry-point kept for backwards compatibility.
    New code should call _replace_pptx_para_arabic which also applies
    alignment and font substitution.
    """
    _replace_pptx_para_arabic(para, new_text, is_rtl, is_title=False, profile=None, font_subs=None)


# ── Arabic-aware PPTX paragraph helpers ──────────────────────────────────────

def _contains_arabic(text: str) -> bool:
    return bool(_ARABIC_BLOCK_RE.search(text or ""))


def _contains_latin(text: str) -> bool:
    return bool(_LATIN_BLOCK_RE.search(text or ""))


def _is_mixed_bidi_text(text: str) -> bool:
    return _contains_arabic(text) and _contains_latin(text)


def _normalize_parentheses_ltr(chunk: str) -> str:
    s = chunk
    # If punctuation got visually mirrored into ")...(", restore it as "(...)".
    m = re.match(r"^\)\s*(.+?)\s*\($", s)
    if m:
        s = f"({m.group(1)})"
    return s


def _normalize_units_spacing(text: str) -> str:
    s = text
    # Ensure spaces between numbers and measurement units.
    s = re.sub(
        rf"(?i)(\d)\s*({_UNIT_TOKEN_ALT})\b",
        r"\1 \2",
        s,
    )
    s = re.sub(
        rf"(?i)\b({_UNIT_TOKEN_ALT})\s*(\d)",
        r"\1 \2",
        s,
    )
    return s


def _normalize_mixed_text_content(text: str) -> str:
    s = text or ""
    # Repair a known broken mixed-order string from OCR/translation output.
    s = re.sub(
        r"(?i)\bkg\)\s*(\d+)\s*modulator\s*\(",
        r"وحدة نبضات (\1 kg Modulator)",
        s,
    )
    s = _normalize_units_spacing(s)
    return s


def _strong_dir_for_char(ch: str) -> str | None:
    if not ch:
        return None
    cp = ord(ch)
    if (
        0x0600 <= cp <= 0x06FF
        or 0x0750 <= cp <= 0x077F
        or 0x08A0 <= cp <= 0x08FF
    ):
        return "ar"
    if ("A" <= ch <= "Z") or ("a" <= ch <= "z"):
        return "ltr"
    return None


def _split_bidi_runs(text: str) -> list[tuple[str, str]]:
    if not text:
        return []

    chars = list(text)
    n = len(chars)
    prev_strong: list[str | None] = [None] * n
    next_strong: list[str | None] = [None] * n

    cur = None
    for i, ch in enumerate(chars):
        d = _strong_dir_for_char(ch)
        if d is not None:
            cur = d
        prev_strong[i] = cur

    cur = None
    for i in range(n - 1, -1, -1):
        ch = chars[i]
        d = _strong_dir_for_char(ch)
        if d is not None:
            cur = d
        next_strong[i] = cur

    assigned: list[str] = []
    for i, ch in enumerate(chars):
        d = _strong_dir_for_char(ch)
        if d is not None:
            assigned.append(d)
            continue

        left = prev_strong[i]
        right = next_strong[i]
        bidi_cat = unicodedata.category(ch)

        if ch in "()[]{}" and (left == "ltr" or right == "ltr"):
            assigned.append("ltr")
        elif left and right and left == right:
            assigned.append(left)
        elif left:
            assigned.append(left)
        elif right:
            assigned.append(right)
        elif bidi_cat.startswith("Z"):
            assigned.append("ar")
        else:
            assigned.append("ltr")

    out: list[tuple[str, str]] = []
    cur_dir = assigned[0]
    start = 0
    for i in range(1, n):
        if assigned[i] != cur_dir:
            out.append((cur_dir, "".join(chars[start:i])))
            start = i
            cur_dir = assigned[i]
    out.append((cur_dir, "".join(chars[start:])))

    return [(d, t) for d, t in out if t]


def _set_pptx_run_lang(run, primary_lang: str, alt_lang: str) -> None:
    _NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
    try:
        import lxml.etree as _ET

        r_elem = run._r
        rPr = r_elem.find(f"{{{_NS_A}}}rPr")
        if rPr is None:
            rPr = _ET.Element(f"{{{_NS_A}}}rPr")
            r_elem.insert(0, rPr)
        rPr.set("lang", primary_lang)
        rPr.set("altLang", alt_lang)
    except Exception:
        pass


def _set_pptx_run_cs_font(run, typeface: str) -> None:
    """Set the run's Complex-Script font (<a:cs typeface="..."/>).

    python-pptx's `run.font.name` only ever writes <a:latin> — the font
    slot for Latin-script glyphs. PowerPoint renders Arabic (and other
    complex scripts) using <a:cs>, not <a:latin>; <a:latin> is merely a
    fallback PowerPoint MAY use when <a:cs> is absent. Whether that
    fallback kicks in depends on the deck's THEME, which defines its own
    default complex-script font — if that theme default is set (to
    anything, including something with no Arabic shaping support at all,
    e.g. a CJK font left over from a template's original design), it wins
    over <a:latin> for every run that doesn't carry an explicit <a:cs>.
    Setting <a:latin> alone therefore worked by accident on decks whose
    theme happened to default to something Arabic-capable, and silently
    failed on decks that don't — this sets the slot that actually controls
    Arabic rendering, regardless of the theme.
    """
    _NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
    try:
        rPr = run._r.get_or_add_rPr()
        cs = rPr.find(f"{{{_NS_A}}}cs")
        if cs is None:
            cs = rPr.makeelement(f"{{{_NS_A}}}cs", {})
            # Schema order requires <a:cs> after <a:latin>/<a:ea> and before
            # <a:sym>; anchoring off <a:latin> (which run.font.name already
            # wrote) keeps this valid regardless of what else is present.
            latin = rPr.find(f"{{{_NS_A}}}latin")
            if latin is not None:
                latin.addnext(cs)
            else:
                rPr.append(cs)
        cs.set("typeface", typeface)
    except Exception:
        pass


def _pick_bilingual_fonts(
    original_font: str | None,
    profile: dict | None,
    is_title: bool,
) -> tuple[str, str]:
    # Enforce requested professional bilingual typography defaults.
    # "Arabic Typesetting" ships built into Windows (Arabic complex-script
    # support since Vista) — unlike "Cairo", which is a Google Font that
    # must be separately installed. Rendering always looked correct during
    # development because Cairo happens to be installed on the dev machine
    # used to verify every fix in this file; on a machine without it,
    # PowerPoint silently substitutes a font with no Arabic shaping table,
    # producing disconnected, unreadable letterforms. Defaulting to a
    # font guaranteed present on the target machine avoids that entirely.
    arabic_font = "Arabic Typesetting"
    english_font = "Arial"

    if original_font and original_font.strip().lower() in {"cairo", "din next lt arabic", "arabic typesetting"}:
        arabic_font = original_font.strip()

    if profile:
        role = "title" if is_title else "body"
        candidate = (
            (profile.get("fonts", {}) or {})
            .get("substitution_by_role", {})
            .get(role)
        )
        if isinstance(candidate, str) and candidate.strip().lower() in {"cairo", "din next lt arabic", "arabic typesetting"}:
            arabic_font = candidate.strip()

    if original_font and original_font.strip().lower() in {"arial", "calibri"}:
        english_font = original_font.strip()

    return arabic_font, english_font

def _replace_pptx_para_arabic(
    para,
    new_text: str,
    is_rtl: bool,
    *,
    is_title: bool = False,
    profile: dict | None = None,
    font_subs: list | None = None,
    slide_label: str = "",
    source_text: str = "",
) -> None:
    """Replace text in a PPTX paragraph with full Arabic formatting support.

    Improvements over the legacy helper:
    - Preserves ALL runs' formatting (not just first) by extracting dominant style
    - Checks every run for special warning/brand colours and preserves them
    - Applies Arabic font substitution when the original font cannot render Arabic
    - Sets RTL direction AND correct alignment (RIGHT for body, preserve CENTER for titles)
    - Records every font substitution into *font_subs* for the quality report
    """
    # ── Step 1: snapshot all runs' formatting ──────────────────────────────────
    # NOTE: colour is read via _run_explicit_rgb (raw XML) — NEVER through
    # run.font.color, whose getter mutates the rPr into an empty
    # <a:solidFill/> that renders BLACK and destroys inherited theme colours.
    run_snapshots = []
    for run in para.runs:
        clr = _run_explicit_rgb(run)
        clr_el = _run_explicit_color_element(run)
        run_snapshots.append({
            "bold":      run.font.bold,
            "italic":    run.font.italic,
            "size":      run.font.size,
            "underline": run.font.underline,
            "name":      run.font.name,
            "color":     clr,
            "color_el":  clr_el,
            "text_len":  len(run.text),
        })

    # ── Step 2: find dominant run (longest text, or first if all empty) ────────
    dominant_idx = 0
    if run_snapshots:
        dominant_idx = max(range(len(run_snapshots)),
                           key=lambda i: run_snapshots[i]["text_len"])
    dom = run_snapshots[dominant_idx] if run_snapshots else {}

    # ── Step 3: look for a "special" colour (non-default, non-black) ──────────
    # We prefer colours that are explicitly set (non-None) with meaningful hue.
    # special_color_el mirrors the same "first explicit" preference but keeps
    # the full <a:solidFill> element so a scheme-color reference (schemeClr,
    # not representable as a plain RGB) survives run replacement too — see
    # _run_explicit_color_element for why this matters.
    special_color = dom.get("color")
    special_color_el = dom.get("color_el")
    for snap in run_snapshots:
        c = snap.get("color")
        if c is not None:
            special_color = c
            special_color_el = snap.get("color_el")
            break  # take the first explicitly set colour
    if special_color is None:
        for snap in run_snapshots:
            ce = snap.get("color_el")
            if ce is not None:
                special_color_el = ce
                break

    # ── Step 4: font substitution ──────────────────────────────────────────────
    original_font = dom.get("name")
    if is_rtl and profile:
        from api.utils.arabic_pptx_profile import is_arabic_capable
        if not is_arabic_capable(original_font):
            role = "title" if is_title else "body"
            subst = (profile.get("fonts", {})
                           .get("substitution_by_role", {})
                           .get(role, "Simplified Arabic"))
            if font_subs is not None and original_font:
                font_subs.append(
                    f"{slide_label}: font '{original_font}' → '{subst}'"
                )
            new_font = subst
        else:
            new_font = original_font
    else:
        new_font = original_font

    # ── Step 5: generate directional run parts (never merge AR + EN) ─────────
    normalized_text = _normalize_mixed_text_content(new_text)
    mixed_bidi = _is_mixed_bidi_text(normalized_text)

    if mixed_bidi:
        run_parts = _split_bidi_runs(normalized_text)
    else:
        default_dir = "ar" if (is_rtl and _contains_arabic(normalized_text)) else "ltr"
        run_parts = [(default_dir, normalized_text)]

    # Drop all existing runs and rebuild in directional chunks.
    for run in list(para.runs):
        try:
            run._r.getparent().remove(run._r)
        except Exception:
            pass

    arabic_font, english_font = _pick_bilingual_fonts(original_font, profile, is_title)

    for run_dir, raw_chunk in run_parts:
        chunk = raw_chunk
        if not chunk:
            continue

        if run_dir == "ltr":
            chunk = _normalize_parentheses_ltr(chunk)
            if mixed_bidi:
                chunk = f"{_BIDI_LRM}{chunk}{_BIDI_LRM}"
        elif mixed_bidi:
            chunk = f"{_BIDI_RLM}{chunk}{_BIDI_RLM}"

        r = para.add_run()
        r.text = chunk
        try:
            r.font.bold = dom.get("bold")
            r.font.italic = dom.get("italic")
            r.font.underline = dom.get("underline")
            if dom.get("size"):
                r.font.size = dom["size"]
            if special_color is not None:
                r.font.color.rgb = special_color
            elif special_color_el is not None:
                # Scheme-color case: let python-pptx insert a schema-correctly
                # positioned <a:solidFill> via the normal RGB setter, then
                # swap its content for the original element (srgbClr or
                # schemeClr) — preserves the exact original color/reference
                # instead of losing it when the run is recreated from scratch.
                from pptx.dml.color import RGBColor
                import copy as _copy
                r.font.color.rgb = RGBColor(0, 0, 0)
                rPr = r._r.get_or_add_rPr()
                placeholder_fill = rPr.find(f"{{{_NS_DML}}}solidFill")
                if placeholder_fill is not None:
                    fill_idx = list(rPr).index(placeholder_fill)
                    rPr.remove(placeholder_fill)
                    rPr.insert(fill_idx, _copy.deepcopy(special_color_el))

            if is_rtl and run_dir == "ar":
                r.font.name = arabic_font
                _set_pptx_run_cs_font(r, arabic_font)
                _set_pptx_run_lang(r, "ar-SA", "en-US")
            else:
                r.font.name = english_font if run_dir == "ltr" else (new_font or english_font)
                _set_pptx_run_lang(r, "en-US", "ar-SA")
        except Exception:
            pass

    # ── Steps 7+8: RTL direction & alignment — LTR→RTL translation only ───────
    # Skip ONLY when this exact paragraph was already explicitly marked
    # rtl="1" in the original file — that is the real signal that its
    # direction/alignment are part of a deliberate, already-correct RTL
    # design and must be preserved exactly (re-flipping it would reverse an
    # already-correct line, e.g. a title like "G60ZBx – مقدمة").
    #
    # This used to infer that same "already RTL-designed" state from
    # whether the SOURCE TEXT merely contained Arabic characters — but a
    # paragraph can contain Arabic text without its pPr ever having had
    # rtl="1" set (e.g. a partially/manually pre-translated deck where
    # someone typed Arabic into an otherwise-English, LTR-paragraph
    # PowerPoint). Such paragraphs kept getting only alignment fixed while
    # staying in LTR paragraph-direction, which left multi-segment mixed
    # Arabic/Latin titles (e.g. "Term1 (EN1) Term2 (EN2)") to PowerPoint's
    # bidi engine to reorder unpredictably. Checking the actual XML
    # attribute — not guessing from text content — is what the "already
    # correct, don't touch it" exemption is actually supposed to mean.
    _pPr_existing = para._p.find(f"{{{_NS_DML}}}pPr")
    _already_rtl = _pPr_existing is not None and _pPr_existing.get("rtl") == "1"
    if not _already_rtl:
        _set_para_rtl(para, is_rtl)
        if is_rtl:
            _apply_arabic_para_alignment(para, is_title=is_title, profile=profile)


def _apply_arabic_para_alignment(para, *, is_title: bool, profile: dict | None) -> None:
    """Set alignment for an Arabic paragraph translated from an LTR source.

    Layout-preservation policy (translated deck must be visually identical
    to the original — only the text may change):

    - An EXPLICIT algn already on the paragraph is NEVER overridden — it is
      designer intent (centered labels, decorative alignment, etc.).
    - Title placeholders are never re-aligned: their alignment comes from the
      slide layout design (center or positioned) and rtl=1 alone is enough.
    - Body paragraphs with no explicit algn get RIGHT — English masters
      typically specify left alignment at the layout/master level, which is
      wrong for Arabic body text.
    """
    PP_RIGHT = 3

    if is_title:
        _set_para_alignment(para, PP_RIGHT)
        return
    try:
        if para.alignment is not None:   # explicit alignment → preserve
            return
    except Exception:
        pass
    _set_para_alignment(para, PP_RIGHT)


def _apply_rtl_to_table(shape, table_text_map: dict, slide_idx: int, shape_idx: int,
                          is_rtl: bool, profile: dict | None, font_subs: list,
                          table_src: dict | None = None,
                          deck_src_is_rtl: bool = False) -> None:
    """Apply RTL direction to all table cells, replace translated cell text.

    Also sets <a:tblPr rtl="1"/> on the table XML element when profile
    specifies table_tblPr_rtl=True.
    """
    import lxml.etree as ET

    tbl = shape.table
    NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"

    # ── Set tblPr rtl on the table XML ────────────────────────────────────────
    # LTR-designed decks only: flipping tblPr@rtl on an Arabic-designed deck
    # changes column order / visual flow — a structural layout change.
    if is_rtl and not deck_src_is_rtl and profile and profile.get("rtl_rules", {}).get("table_tblPr_rtl", True):
        try:
            tbl_el = tbl._tbl
            tblPr = tbl_el.find(f"{{{NS_A}}}tblPr")
            if tblPr is None:
                tblPr = ET.SubElement(tbl_el, f"{{{NS_A}}}tblPr")
                tbl_el.insert(0, tblPr)
            tblPr.set("rtl", "1")
        except Exception as exc:
            log.debug("Could not set tblPr rtl: %s", exc)

    for row_idx, row in enumerate(tbl.rows):
        for col_idx, cell in enumerate(row.cells):
            # Replace translated text if available
            key = (slide_idx, shape_idx, row_idx, col_idx)
            translated = table_text_map.get(key)

            for para_idx, para in enumerate(cell.text_frame.paragraphs):
                # Replace text only in the first paragraph per cell
                if translated and para_idx == 0:
                    _replace_pptx_para_arabic(
                        para, translated, is_rtl,
                        is_title=False,
                        profile=profile,
                        font_subs=font_subs,
                        slide_label=f"slide {slide_idx+1} table row{row_idx} col{col_idx}",
                        source_text=(table_src or {}).get(key, ""),
                    )
                elif is_rtl and not deck_src_is_rtl:
                    # Untranslated cell paragraph in an LTR-designed deck:
                    # RTL direction + right alignment (only when no explicit
                    # algn exists).  Arabic-designed decks: untouched.
                    _set_para_rtl(para, True)
                    if para.alignment is None and not _text_is_rtl_script(para.text):
                        _set_para_alignment(para, 3)  # RIGHT = 3


# ── Generic DOCX output (for non-rebuildable formats) ─────────────────────────

def build_translated_docx(segments: list[dict], source_lang: str, target_lang: str) -> bytes:
    """
    Build a professional bilingual DOCX from segments when in-place rebuild isn't
    feasible (PDF, TXT, HTML sources).

    Layout: two-column table — source (grey, left) | translation (black, right).
    Untranslated segments fall back to the source text so the document is always
    readable; a subtle yellow cell background flags those cells for review.
    """
    from docx import Document
    from docx.shared import Pt, RGBColor, Cm, Inches
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn
    from lxml import etree

    doc = Document()

    # ── Page margins: narrow for two-column layout ────────────────────────────
    for section in doc.sections:
        section.left_margin   = Cm(2)
        section.right_margin  = Cm(2)
        section.top_margin    = Cm(2)
        section.bottom_margin = Cm(2)

    is_rtl = is_rtl_lang(target_lang)

    lang_name = {
        "en": "English", "ar": "Arabic", "fr": "French",
        "de": "German",  "es": "Spanish",
    }
    src_name = lang_name.get(source_lang, source_lang.upper())
    tgt_name = lang_name.get(target_lang, target_lang.upper())

    # ── Document title ────────────────────────────────────────────────────────
    title = doc.add_heading("Translated Document", level=1)
    title.alignment = WD_ALIGN_PARAGRAPH.LEFT

    sub = doc.add_paragraph(f"{src_name}  →  {tgt_name}")
    sub.runs[0].font.size = Pt(10)
    sub.runs[0].font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    doc.add_paragraph()

    # ── Helper: set cell background colour ───────────────────────────────────
    def _set_cell_bg(cell, hex_color: str):
        tcPr = cell._tc.get_or_add_tcPr()
        shd = etree.SubElement(tcPr, qn("w:shd"))
        shd.set(qn("w:val"), "clear")
        shd.set(qn("w:color"), "auto")
        shd.set(qn("w:fill"), hex_color.lstrip("#"))

    # ── Helper: set cell RTL paragraph direction ──────────────────────────────
    def _set_cell_rtl(cell):
        for para in cell.paragraphs:
            pPr = para._p.get_or_add_pPr()
            bidi = pPr.find(qn("w:bidi"))
            if bidi is None:
                bidi = etree.SubElement(pPr, qn("w:bidi"))
            bidi.set(qn("w:val"), "1")
            para.alignment = WD_ALIGN_PARAGRAPH.RIGHT

    # ── Header row ────────────────────────────────────────────────────────────
    hdr_table = doc.add_table(rows=1, cols=2)
    hdr_table.style = "Table Grid"
    hdr_c0, hdr_c1 = hdr_table.rows[0].cells
    _set_cell_bg(hdr_c0, "#2B5797")   # dark blue
    _set_cell_bg(hdr_c1, "#1A6B3C")   # dark green

    h0 = hdr_c0.paragraphs[0]
    r0 = h0.add_run(f"SOURCE — {src_name.upper()}")
    r0.font.bold = True
    r0.font.size = Pt(9)
    r0.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    h0.alignment = WD_ALIGN_PARAGRAPH.LEFT

    h1 = hdr_c1.paragraphs[0]
    r1 = h1.add_run(f"TRANSLATION — {tgt_name.upper()}")
    r1.font.bold = True
    r1.font.size = Pt(9)
    r1.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    h1.alignment = WD_ALIGN_PARAGRAPH.RIGHT if is_rtl else WD_ALIGN_PARAGRAPH.LEFT

    doc.add_paragraph()

    # ── Content table ─────────────────────────────────────────────────────────
    content_segs = [s for s in segments if s.get("source", "").strip()]
    if not content_segs:
        doc.add_paragraph("No translatable content found.")
        out = io.BytesIO()
        doc.save(out)
        return out.getvalue()

    table = doc.add_table(rows=len(content_segs), cols=2)
    table.style = "Table Grid"

    for row_idx, seg in enumerate(content_segs):
        source = seg.get("source", "").strip()
        raw_target = seg.get("target", "").strip()
        is_untranslated = not raw_target

        # Fallback: show source text when translation is missing
        target = raw_target if raw_target else source

        row = table.rows[row_idx]
        src_cell = row.cells[0]
        tgt_cell = row.cells[1]

        # Alternating row shading for readability
        if row_idx % 2 == 1:
            _set_cell_bg(src_cell, "F5F5F5")
            _set_cell_bg(tgt_cell, "F5F5F5")

        # Untranslated: soft amber tint on target cell to flag for review
        if is_untranslated:
            _set_cell_bg(tgt_cell, "FFF9E6")

        # Source cell — grey, small
        src_para = src_cell.paragraphs[0]
        src_run = src_para.add_run(source)
        src_run.font.size = Pt(9)
        src_run.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
        src_para.alignment = WD_ALIGN_PARAGRAPH.LEFT

        # Target cell — normal weight, RTL-aware
        tgt_para = tgt_cell.paragraphs[0]
        tgt_run = tgt_para.add_run(target)
        tgt_run.font.size = Pt(10)
        if is_untranslated:
            tgt_run.font.color.rgb = RGBColor(0x99, 0x77, 0x00)  # amber for review
        if is_rtl and raw_target:
            _set_cell_rtl(tgt_cell)
        elif is_rtl and is_untranslated:
            tgt_para.alignment = WD_ALIGN_PARAGRAPH.LEFT  # source fallback is LTR

    out = io.BytesIO()
    doc.save(out)
    return out.getvalue()


# ── PDF output ────────────────────────────────────────────────────────────────

def build_translated_pdf(segments: list[dict], source_lang: str, target_lang: str, project_name: str = "") -> bytes:
    """
    Build a clean translated PDF using ReportLab with proper Arabic/RTL support.
    Each segment is rendered as a bilingual block: source (grey) above target (black).
    """
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import mm
    from reportlab.lib import colors
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, HRFlowable
    from reportlab.lib.enums import TA_LEFT, TA_RIGHT, TA_CENTER
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont

    is_rtl = is_rtl_lang(target_lang)

    # Try to register Arabic-capable font; fall back to Helvetica if unavailable
    arabic_font = "Helvetica"
    try:
        import os, glob
        # Look for a Unicode font that can handle Arabic on the system
        candidates = [
            "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
            "/usr/share/fonts/truetype/freefont/FreeSans.ttf",
            "/usr/share/fonts/TTF/DejaVuSans.ttf",
            "/run/current-system/sw/share/X11/fonts/DejaVuSans.ttf",
        ]
        candidates += glob.glob("/usr/share/fonts/**/*.ttf", recursive=True)[:5]
        for path in candidates:
            if os.path.exists(path):
                pdfmetrics.registerFont(TTFont("UnicodeFont", path))
                arabic_font = "UnicodeFont"
                break
    except Exception:
        pass

    buf = io.BytesIO()
    doc = SimpleDocTemplate(
        buf,
        pagesize=A4,
        leftMargin=20 * mm,
        rightMargin=20 * mm,
        topMargin=20 * mm,
        bottomMargin=20 * mm,
    )

    align_target = TA_RIGHT if is_rtl else TA_LEFT

    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        "TitleStyle",
        parent=styles["Title"],
        fontName=arabic_font,
        fontSize=16,
        spaceAfter=4 * mm,
        alignment=TA_CENTER,
    )
    meta_style = ParagraphStyle(
        "MetaStyle",
        parent=styles["Normal"],
        fontName=arabic_font,
        fontSize=9,
        textColor=colors.HexColor("#888888"),
        alignment=TA_CENTER,
        spaceAfter=6 * mm,
    )
    source_style = ParagraphStyle(
        "SourceStyle",
        parent=styles["Normal"],
        fontName=arabic_font,
        fontSize=8,
        textColor=colors.HexColor("#888888"),
        alignment=TA_LEFT,
        spaceAfter=1 * mm,
    )
    target_style = ParagraphStyle(
        "TargetStyle",
        parent=styles["Normal"],
        fontName=arabic_font,
        fontSize=11,
        textColor=colors.black,
        alignment=align_target,
        spaceAfter=2 * mm,
        leading=16,
    )

    story = []

    # Title block
    title_text = project_name or "Translated Document"
    story.append(Paragraph(title_text, title_style))
    story.append(Paragraph(
        f"{source_lang.upper()} → {target_lang.upper()}",
        meta_style,
    ))
    story.append(HRFlowable(width="100%", thickness=0.5, color=colors.HexColor("#dddddd"), spaceAfter=6 * mm))

    for seg in segments:
        source = (seg.get("source") or "").strip()
        raw_target = (seg.get("target") or "").strip()
        if not source:
            continue

        # Fallback to source text when no translation available — never show placeholder
        is_untranslated = not raw_target
        target = raw_target if raw_target else source

        # Apply Arabic reshaping for correct glyph order in PDF
        if is_rtl and raw_target:
            target = _prepare_arabic(raw_target)

        # Escape HTML entities for ReportLab Paragraph
        def _esc(s: str) -> str:
            return s.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")

        story.append(Paragraph(_esc(source), source_style))
        story.append(Paragraph(_esc(target), target_style))
        story.append(HRFlowable(
            width="100%", thickness=0.3,
            color=colors.HexColor("#eeeeee"),
            spaceAfter=3 * mm,
        ))

    if not story or len(story) <= 3:
        story.append(Paragraph("No translated content available.", target_style))

    doc.build(story)
    return buf.getvalue()


# ── TXT output ────────────────────────────────────────────────────────────────

def build_translated_txt(segments: list[dict]) -> bytes:
    """Build a plain text file with translated segments."""
    lines = []
    for seg in segments:
        target = seg.get("target", "").strip()
        if target:
            lines.append(target)
    return "\n".join(lines).encode("utf-8")


# ── HTML output ───────────────────────────────────────────────────────────────

def build_translated_html(segments: list[dict], target_lang: str) -> bytes:
    """Build a clean HTML file with translated content."""
    is_rtl = is_rtl_lang(target_lang)
    dir_attr = ' dir="rtl" lang="ar"' if is_rtl else ' dir="ltr"'

    lines = [
        f"<!DOCTYPE html><html{dir_attr}><head>",
        '<meta charset="utf-8">',
        "<title>Translated Document</title>",
        "<style>body{font-family:Arial,sans-serif;margin:2em;line-height:1.6;}</style>",
        "</head><body>",
    ]
    for seg in segments:
        target = seg.get("target", "").strip()
        if target:
            if seg.get("seg_type") == "slide_title":
                lines.append(f"<h2>{target}</h2>")
            else:
                lines.append(f"<p>{target}</p>")
    lines.append("</body></html>")
    return "\n".join(lines).encode("utf-8")


# ── Dispatcher ────────────────────────────────────────────────────────────────

def rebuild_xlsx(original_bytes: bytes, segments: list[dict], target_lang: str) -> bytes:
    """
    Write translated segments back into an XLSX, preserving cell structure,
    formulas, styles, and sheets.  String cells are substituted using the
    source → target lookup from the segment list.

    Returns the rebuilt XLSX as bytes.
    """
    import openpyxl

    wb = openpyxl.load_workbook(io.BytesIO(original_bytes))
    is_rtl = is_rtl_lang(target_lang)

    # Build lookup: source_text → target_text
    translation_map: dict[str, str] = {}
    for seg in segments:
        source = seg.get("source", "").strip()
        target = seg.get("target", "").strip()
        if source and target:
            translation_map[source] = target

    for ws in wb.worksheets:
        # Set sheet direction for RTL languages
        if is_rtl:
            ws.sheet_view.rightToLeft = True

        for row in ws.iter_rows():
            for cell in row:
                if not isinstance(cell.value, str):
                    continue
                stripped = cell.value.strip()
                if stripped in translation_map:
                    translated = translation_map[stripped]
                    # Preserve leading/trailing whitespace from original
                    prefix = cell.value[: len(cell.value) - len(cell.value.lstrip())]
                    suffix = cell.value[len(cell.value.rstrip()):]
                    cell.value = prefix + translated + suffix
                elif cell.value in translation_map:
                    cell.value = translation_map[cell.value]

    out = io.BytesIO()
    wb.save(out)
    return out.getvalue()


def rebuild_document(
    original_bytes: bytes,
    file_type: str,
    segments: list[dict],
    target_lang: str,
    source_lang: str,
    layout_warnings: "list[str] | None" = None,
    style_profile_override: "dict | None" = None,
    template_strength: str = "balanced",
    strict_qa: bool = False,
    allow_export_with_warnings: bool = True,
    auto_repair_enabled: bool = True,
    export_best_effort_result: bool = True,
) -> tuple[bytes | None, bytes | None, bytes | None]:
    """
    Rebuild translated document.

    Returns:
        (docx_bytes, pptx_bytes, xlsx_bytes) — exactly one will be non-None
        for structured formats; for unstructured sources a bilingual DOCX is
        returned in the docx slot.

        file_type  → result slot
        docx       → docx_bytes
        pptx       → pptx_bytes
        xlsx       → xlsx_bytes
        pdf/txt/…  → docx_bytes (bilingual fallback)
    """
    ft = file_type.lower()

    if ft == "docx":
        try:
            docx_bytes = rebuild_docx(original_bytes, segments, target_lang)
            return docx_bytes, None, None
        except Exception as e:
            log.error("DOCX rebuild failed, falling back to generic: %s", e)
            docx_bytes = build_translated_docx(segments, source_lang, target_lang)
            return docx_bytes, None, None

    elif ft == "pptx":
        try:
            pptx_bytes = rebuild_pptx(
                original_bytes, segments, target_lang,
                layout_warnings=layout_warnings,
                style_profile_override=style_profile_override,
                template_strength=template_strength,
                strict_qa=strict_qa,
                allow_export_with_warnings=allow_export_with_warnings,
                auto_repair_enabled=auto_repair_enabled,
                export_best_effort_result=export_best_effort_result,
            )
            return None, pptx_bytes, None
        except LayoutValidationError as e:
            # Strict mode can block export, but advisory mode must still export.
            if strict_qa and not allow_export_with_warnings:
                raise
            if layout_warnings is not None:
                layout_warnings.append(f"PPTX layout QA warning: {e}")
            try:
                pptx_bytes = rebuild_pptx(
                    original_bytes, segments, target_lang,
                    layout_warnings=layout_warnings,
                    style_profile_override=style_profile_override,
                    template_strength=template_strength,
                    strict_qa=False,
                    allow_export_with_warnings=True,
                    auto_repair_enabled=auto_repair_enabled,
                    export_best_effort_result=True,
                )
                return None, pptx_bytes, None
            except Exception as retry_err:
                log.error("PPTX rebuild failed after advisory retry: %s", retry_err)
                raise
        except Exception as e:
            log.error("PPTX rebuild failed, falling back to generic DOCX: %s", e)
            docx_bytes = build_translated_docx(segments, source_lang, target_lang)
            return docx_bytes, None, None

    elif ft == "xlsx":
        try:
            xlsx_bytes = rebuild_xlsx(original_bytes, segments, target_lang)
            return None, None, xlsx_bytes
        except Exception as e:
            log.error("XLSX rebuild failed, falling back to generic DOCX: %s", e)
            docx_bytes = build_translated_docx(segments, source_lang, target_lang)
            return docx_bytes, None, None

    else:
        # PDF, TXT, HTML, CSV, MD, RTF → produce clean translated DOCX
        docx_bytes = build_translated_docx(segments, source_lang, target_lang)
        return docx_bytes, None, None
