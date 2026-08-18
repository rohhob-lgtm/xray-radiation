"""
Image text extraction using GPT-4o Vision.

Scans images embedded in documents for text labels (callouts, diagram
annotations, schematic labels, etc.) and returns their approximate
positions so translated text can be overlaid.

Public API
----------
extract_image_labels(image_bytes, client) -> list[dict]
    Each dict: {"text": str, "x_pct": float, "y_pct": float}
    Positions are expressed as fractions of image width / height (0.0–1.0).

extract_pdf_images(file_bytes) -> list[dict]
    Each dict: {"page": int, "image_index": int, "image_bytes": bytes, "xref": int}

extract_docx_images(file_bytes) -> list[dict]
    Each dict: {"shape_idx": int, "rel_id": str, "image_bytes": bytes}

extract_pptx_images(file_bytes) -> list[dict]
    Each dict: {"slide_idx": int, "shape_idx": int, "rel_id": str, "image_bytes": bytes}
"""
from __future__ import annotations

import base64
import io
import json
import logging
import re
from typing import Any

log = logging.getLogger(__name__)

# Maximum images to scan per document to bound API cost
MAX_IMAGES_PER_DOC = 15
# Minimum pixel dimension (width AND height) for an image to be scanned
MIN_IMAGE_DIM = 80
OCR_LOW_CONFIDENCE_THRESHOLD = 0.62

_TEXT_CATEGORIES = {
    "equipment_label",
    "control_panel_label",
    "button",
    "indicator",
    "warning_label",
    "technical_annotation",
    "callout_box",
    "caption",
    "other",
}

_IMAGE_TYPES = {
    "screenshot",
    "photograph",
    "scanned_page",
    "diagram",
    "control_panel",
    "manual_illustration",
    "other",
}


# ── Vision extraction ──────────────────────────────────────────────────────────

_VISION_PROMPT = """You are analysing a technical training image for OCR translation.
Detect visible English text and classify each text item and image type.

Return ONLY valid JSON in this exact shape:
{
  "image_type": "screenshot|photograph|scanned_page|diagram|control_panel|manual_illustration|other",
  "regions": [
    {
      "text": "exact source text",
      "x_pct": 0.0,
      "y_pct": 0.0,
      "bbox": {"x": 0.0, "y": 0.0, "w": 0.0, "h": 0.0},
      "category": "equipment_label|control_panel_label|button|indicator|warning_label|technical_annotation|callout_box|caption|other",
      "confidence": 0.0,
      "font_size": 12,
      "font_color": "#000000",
      "alignment": "left|center|right",
      "translate": true
    }
  ]
}

Rules:
- Include all readable UI labels, callouts, warning labels, technical annotations, and captions.
- Skip watermark lines, page numbers, and unreadable noise.
- If no text is found return {"image_type":"other","regions":[]}.
- confidence must be 0.0-1.0 and reflect OCR certainty.
- Keep translate=false for technical IDs/codes that should stay in English.
- bbox/x_pct/y_pct are percentages in range 0..1.
"""


def _enhance_image_for_ocr(image_bytes: bytes) -> bytes:
    """Increase OCR readability for low-confidence images."""
    try:
        from PIL import Image as PILImage, ImageEnhance, ImageFilter

        img = PILImage.open(io.BytesIO(image_bytes)).convert("RGB")
        img = img.resize((max(1, img.width * 2), max(1, img.height * 2)), resample=PILImage.Resampling.LANCZOS)
        img = ImageEnhance.Contrast(img).enhance(1.45)
        img = ImageEnhance.Sharpness(img).enhance(1.35)
        img = img.filter(ImageFilter.MedianFilter(size=3))
        buf = io.BytesIO()
        img.save(buf, format="PNG")
        return buf.getvalue()
    except Exception:
        return image_bytes


def _parse_vision_labels(raw: str) -> tuple[str, list[dict]]:
    # Strip markdown code fences if present
    raw = re.sub(r"^```(?:json)?\s*", "", raw.strip(), flags=re.MULTILINE)
    raw = re.sub(r"\s*```$", "", raw.strip(), flags=re.MULTILINE)
    data = json.loads(raw or "{}")
    if isinstance(data, list):
        data = {"image_type": "other", "regions": data}

    image_type = str(data.get("image_type", "other")).strip().lower()
    if image_type not in _IMAGE_TYPES:
        image_type = "other"

    regions = data.get("regions", [])
    if not isinstance(regions, list):
        regions = []

    cleaned: list[dict] = []
    for item in regions:
        if not isinstance(item, dict):
            continue
        text = str(item.get("text", "")).strip()
        if not text:
            continue

        category = str(item.get("category", "other")).strip().lower()
        if category not in _TEXT_CATEGORIES:
            category = "other"

        alignment = str(item.get("alignment", "center")).strip().lower()
        if alignment not in {"left", "center", "right"}:
            alignment = "center"

        conf = float(item.get("confidence", 0.70))
        conf = max(0.0, min(1.0, conf))

        bbox_in = item.get("bbox") or {}
        bx = max(0.0, min(1.0, float(bbox_in.get("x", item.get("x_pct", 0.5)))))
        by = max(0.0, min(1.0, float(bbox_in.get("y", item.get("y_pct", 0.5)))))
        bw = max(0.01, min(1.0, float(bbox_in.get("w", 0.10))))
        bh = max(0.01, min(1.0, float(bbox_in.get("h", 0.05))))

        if bx + bw > 1.0:
            bw = max(0.01, 1.0 - bx)
        if by + bh > 1.0:
            bh = max(0.01, 1.0 - by)

        color = str(item.get("font_color", "#000000")).strip()
        if not re.fullmatch(r"#[0-9A-Fa-f]{6}", color):
            color = "#000000"

        cleaned.append({
            "text": text,
            "x_pct": bx,
            "y_pct": by,
            "bbox": {"x": bx, "y": by, "w": bw, "h": bh},
            "category": category,
            "confidence": conf,
            "font_size": max(8, int(item.get("font_size", 14))),
            "font_color": color,
            "alignment": alignment,
            "translate": bool(item.get("translate", True)),
            "image_type": image_type,
        })
    return image_type, cleaned


async def extract_image_labels(
    image_bytes: bytes,
    client: Any,  # openai.AsyncOpenAI
    model: str = "gpt-4o-mini",
) -> list[dict]:
    """
    Call GPT-4o Vision on a single image and return detected text labels
    with approximate (x_pct, y_pct) positions.

    Uses gpt-4o-mini by default to preserve the gpt-4o rate-limit budget
    for the primary translation pipeline.  Pass model="gpt-4o" when higher
    accuracy is needed.

    Returns [] on any error so callers can treat it as non-fatal.
    """
    if not image_bytes:
        return []

    b64 = base64.b64encode(image_bytes).decode()
    # Detect format from magic bytes
    mime = _detect_mime(image_bytes)

    try:
        response = await client.chat.completions.create(
            model=model,
            max_tokens=1024,
            messages=[
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "image_url",
                            "image_url": {"url": f"data:{mime};base64,{b64}"},
                        },
                        {"type": "text", "text": _VISION_PROMPT},
                    ],
                }
            ],
        )
        try:
            from api.utils.usage_recorder import record_usage_from_response
            record_usage_from_response(
                "Image Translation", response,
                sub_feature="image_label_extraction",
                meta={"model": model},
            )
        except Exception:
            pass
        raw = response.choices[0].message.content or "{}"
        _, labels = _parse_vision_labels(raw)

        # OCR confidence validation + retry with enhanced preprocessing.
        avg_conf = (sum(float(l.get("confidence", 0.0)) for l in labels) / len(labels)) if labels else 1.0
        if labels and avg_conf < OCR_LOW_CONFIDENCE_THRESHOLD:
            enhanced = _enhance_image_for_ocr(image_bytes)
            if enhanced != image_bytes:
                retry_b64 = base64.b64encode(enhanced).decode()
                retry = await client.chat.completions.create(
                    model=model,
                    max_tokens=1024,
                    messages=[
                        {
                            "role": "user",
                            "content": [
                                {
                                    "type": "image_url",
                                    "image_url": {"url": f"data:image/png;base64,{retry_b64}"},
                                },
                                {"type": "text", "text": _VISION_PROMPT},
                            ],
                        }
                    ],
                )
                retry_raw = retry.choices[0].message.content or "{}"
                _, retry_labels = _parse_vision_labels(retry_raw)
                retry_avg = (sum(float(l.get("confidence", 0.0)) for l in retry_labels) / len(retry_labels)) if retry_labels else 0.0
                if retry_labels and retry_avg >= avg_conf:
                    labels = retry_labels

        return labels
    except Exception as e:
        log.warning("GPT-4o Vision call failed: %s", e)
        return []


def _detect_mime(image_bytes: bytes) -> str:
    if image_bytes[:4] == b"\x89PNG":
        return "image/png"
    if image_bytes[:2] in (b"\xff\xd8",):
        return "image/jpeg"
    if image_bytes[:4] in (b"GIF8",):
        return "image/gif"
    if image_bytes[:4] == b"RIFF" and image_bytes[8:12] == b"WEBP":
        return "image/webp"
    return "image/png"  # safe default


# ── PDF image extraction ───────────────────────────────────────────────────────

def extract_pdf_images(file_bytes: bytes) -> list[dict]:
    """
    Extract embedded images from a PDF using PyMuPDF.

    Returns list of:
        {
          "page": int (1-indexed),
          "image_index": int,
          "image_bytes": bytes,   # PNG
          "xref": int,            # PyMuPDF xref for replacement
        }
    Limited to MAX_IMAGES_PER_DOC entries.
    """
    try:
        import fitz  # PyMuPDF
    except ImportError:
        log.warning("PyMuPDF not available — cannot extract PDF images")
        return []

    results = []
    try:
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        for page_num, page in enumerate(doc, 1):
            img_list = page.get_images(full=True)
            for img_idx, img_info in enumerate(img_list):
                if len(results) >= MAX_IMAGES_PER_DOC:
                    break
                xref = img_info[0]
                try:
                    base_img = doc.extract_image(xref)
                    img_bytes = base_img.get("image", b"")
                    width = base_img.get("width", 0)
                    height = base_img.get("height", 0)
                    if width < MIN_IMAGE_DIM or height < MIN_IMAGE_DIM:
                        continue
                    # Convert to PNG if needed for consistent handling
                    ext = base_img.get("ext", "png").lower()
                    if ext != "png":
                        try:
                            from PIL import Image as PILImage
                            buf = io.BytesIO()
                            PILImage.open(io.BytesIO(img_bytes)).save(buf, format="PNG")
                            img_bytes = buf.getvalue()
                        except Exception:
                            pass
                    results.append({
                        "page": page_num,
                        "image_index": img_idx,
                        "image_bytes": img_bytes,
                        "xref": xref,
                    })
                except Exception as e:
                    log.debug("Skipping PDF image xref=%d: %s", xref, e)
            if len(results) >= MAX_IMAGES_PER_DOC:
                break
        doc.close()
    except Exception as e:
        log.warning("PDF image extraction failed: %s", e)

    return results


# ── DOCX image extraction ──────────────────────────────────────────────────────

def extract_docx_images(file_bytes: bytes) -> list[dict]:
    """
    Extract embedded images from a DOCX file.

    Returns list of:
        {
          "shape_idx": int,
          "rel_id": str,
          "image_bytes": bytes,   # PNG
        }
    Limited to MAX_IMAGES_PER_DOC entries.
    """
    try:
        from docx import Document
        from docx.oxml.ns import qn
    except ImportError:
        log.warning("python-docx not available — cannot extract DOCX images")
        return []

    results = []
    try:
        doc = Document(io.BytesIO(file_bytes))
        shape_idx = 0

        for para in doc.paragraphs:
            for run in para.runs:
                if len(results) >= MAX_IMAGES_PER_DOC:
                    break
                # Check for inline image (w:drawing element)
                drawing = run._element.find(qn("w:drawing"))
                if drawing is None:
                    continue
                # Try to locate the blip relationship id
                blip = drawing.find(".//" + qn("a:blip"))
                if blip is None:
                    shape_idx += 1
                    continue
                embed_id = blip.get(qn("r:embed"))
                if not embed_id:
                    shape_idx += 1
                    continue
                try:
                    rel = doc.part.rels.get(embed_id)
                    if rel is None:
                        shape_idx += 1
                        continue
                    img_bytes = rel.target_part.blob
                    if not img_bytes or len(img_bytes) < 64:
                        shape_idx += 1
                        continue
                    # Check dimensions
                    from PIL import Image as PILImage
                    im = PILImage.open(io.BytesIO(img_bytes))
                    w, h = im.size
                    if w < MIN_IMAGE_DIM or h < MIN_IMAGE_DIM:
                        shape_idx += 1
                        continue
                    # Normalise to PNG
                    buf = io.BytesIO()
                    im.save(buf, format="PNG")
                    results.append({
                        "shape_idx": shape_idx,
                        "rel_id": embed_id,
                        "image_bytes": buf.getvalue(),
                    })
                except Exception as e:
                    log.debug("Skipping DOCX image rel=%s: %s", embed_id, e)
                shape_idx += 1

    except Exception as e:
        log.warning("DOCX image extraction failed: %s", e)

    return results


# ── PPTX image extraction ─────────────────────────────────────────────────────

def extract_pptx_images(file_bytes: bytes) -> list[dict]:
    """
    Extract picture shapes from a PPTX file.

    Returns list of:
        {
          "slide_idx": int,      # 0-based slide index
          "shape_idx": int,      # enumerate(slide.shapes) index
          "rel_id": str,         # image relationship id (r:embed)
          "image_bytes": bytes,  # PNG
        }
    Limited to MAX_IMAGES_PER_DOC entries.
    """
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from pptx.oxml.ns import qn
    except ImportError:
        log.warning("python-pptx not available — cannot extract PPTX images")
        return []

    results = []
    try:
        prs = Presentation(io.BytesIO(file_bytes))
        for slide_idx, slide in enumerate(prs.slides):
            for shape_idx, shape in enumerate(slide.shapes):
                if len(results) >= MAX_IMAGES_PER_DOC:
                    break
                try:
                    if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                        continue

                    blip = shape._element.find(".//" + qn("a:blip"))
                    if blip is None:
                        continue
                    rel_id = blip.get(qn("r:embed"))
                    if not rel_id:
                        continue

                    img_bytes = shape.image.blob
                    if not img_bytes or len(img_bytes) < 64:
                        continue

                    from PIL import Image as PILImage

                    im = PILImage.open(io.BytesIO(img_bytes))
                    w, h = im.size
                    if w < MIN_IMAGE_DIM or h < MIN_IMAGE_DIM:
                        continue

                    buf = io.BytesIO()
                    im.save(buf, format="PNG")
                    results.append({
                        "slide_idx": slide_idx,
                        "shape_idx": shape_idx,
                        "rel_id": rel_id,
                        "image_bytes": buf.getvalue(),
                    })
                except Exception as e:
                    log.debug(
                        "Skipping PPTX image slide=%d shape=%d: %s",
                        slide_idx,
                        shape_idx,
                        e,
                    )
            if len(results) >= MAX_IMAGES_PER_DOC:
                break
    except Exception as e:
        log.warning("PPTX image extraction failed: %s", e)

    return results


# ── Overlay rendering ──────────────────────────────────────────────────────────

def render_overlay_image(
    image_bytes: bytes,
    labels: list[dict],  # [{"text": str, "x_pct": float, "y_pct": float}]
    target_lang: str = "ar",
) -> bytes:
    """
    Composite translated labels onto an image.

    Draws a semi-transparent dark pill behind each label and the translated
    text in white on top.  Returns PNG bytes.

    `labels` must already contain the *translated* text in the `text` field.
    """
    try:
        from PIL import Image as PILImage, ImageDraw, ImageFont, ImageStat
        import arabic_reshaper
        from bidi.algorithm import get_display

        img = PILImage.open(io.BytesIO(image_bytes)).convert("RGBA")
        w, h = img.size

        # Guard against destructive overlays. Each label is masked with an opaque
        # rectangle before the translation is drawn. On photos / X-ray scans /
        # diagrams the text detector can return large or spurious regions, and
        # masking them paints big grey boxes over the artwork (a real regression
        # seen on scan slides). If the detected regions cover a large share of
        # the image, the detection is untrustworthy — return the ORIGINAL image
        # untouched rather than damage it. Small embedded labels still overlay.
        _total_frac = 0.0
        _max_frac = 0.0
        for _lbl in labels:
            _bb = _lbl.get("bbox") if isinstance(_lbl.get("bbox"), dict) else None
            if _bb:
                _fw = max(0.0, min(1.0, float(_bb.get("w", 0.10) or 0.10)))
                _fh = max(0.0, min(1.0, float(_bb.get("h", 0.05) or 0.05)))
            else:
                _fw, _fh = 0.12, 0.06
            _a = _fw * _fh
            _total_frac += _a
            _max_frac = max(_max_frac, _a)
        if _max_frac > 0.35 or _total_frac > 0.30:
            log.info(
                "image overlay skipped (regions too large: max=%.2f total=%.2f) — keeping original image",
                _max_frac, _total_frac,
            )
            return image_bytes

        # Overlay layer
        overlay = PILImage.new("RGBA", img.size, (0, 0, 0, 0))
        draw = ImageDraw.Draw(overlay)

        # Try to get a font — fall back to default
        font_size_default = max(12, min(22, h // 30))
        try:
            font_default = ImageFont.load_default()
        except Exception:
            font_default = None

        def _font_for(size_hint: int):
            candidates = [
                "C:/Windows/Fonts/tahoma.ttf",
                "C:/Windows/Fonts/arial.ttf",
                "C:/Windows/Fonts/trado.ttf",
                "/usr/share/fonts/truetype/noto/NotoNaskhArabic-Regular.ttf",
                "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
            ]
            for fp in candidates:
                try:
                    return ImageFont.truetype(fp, max(9, int(size_hint)))
                except Exception:
                    continue
            try:
                return ImageFont.load_default()
            except Exception:
                return None

        def _sample_bg(x0: int, y0: int, x1: int, y1: int) -> tuple[int, int, int]:
            try:
                pad = 2
                sx0 = max(0, x0 - pad)
                sy0 = max(0, y0 - pad)
                sx1 = min(w, x1 + pad)
                sy1 = min(h, y1 + pad)
                crop = img.crop((sx0, sy0, sx1, sy1)).convert("RGB")
                mean = ImageStat.Stat(crop).mean
                return (int(mean[0]), int(mean[1]), int(mean[2]))
            except Exception:
                return (255, 255, 255)

        for lbl in labels:
            text = lbl.get("text", "")
            if not text:
                continue

            # Reshape Arabic if needed
            if target_lang == "ar" or any("\u0600" <= c <= "\u06FF" for c in text):
                try:
                    reshaped = arabic_reshaper.reshape(text)
                    text = get_display(reshaped)
                except Exception:
                    pass

            # Prefer bbox-based placement for true replacement; fallback to center point.
            bbox = lbl.get("bbox") if isinstance(lbl.get("bbox"), dict) else None
            if bbox:
                rx0 = int(max(0.0, min(1.0, float(bbox.get("x", lbl.get("x_pct", 0.5))))) * w)
                ry0 = int(max(0.0, min(1.0, float(bbox.get("y", lbl.get("y_pct", 0.5))))) * h)
                rw = int(max(0.01, min(1.0, float(bbox.get("w", 0.10)))) * w)
                rh = int(max(0.01, min(1.0, float(bbox.get("h", 0.05)))) * h)
                rx1 = min(w, rx0 + max(8, rw))
                ry1 = min(h, ry0 + max(8, rh))
            else:
                cx = int(float(lbl.get("x_pct", 0.5)) * w)
                cy = int(float(lbl.get("y_pct", 0.5)) * h)
                rx0 = max(0, cx - 40)
                ry0 = max(0, cy - 12)
                rx1 = min(w, cx + 40)
                ry1 = min(h, cy + 12)

            bw = max(8, rx1 - rx0)
            bh = max(8, ry1 - ry0)

            # Measure text
            font_size = int(lbl.get("font_size", font_size_default))
            font_size = max(9, min(64, font_size))
            font = _font_for(font_size)
            if font:
                try:
                    bb = draw.textbbox((0, 0), text, font=font)
                    tw = bb[2] - bb[0]
                    th = bb[3] - bb[1]
                except AttributeError:
                    tw, th = draw.textsize(text, font=font)
            else:
                tw, th = len(text) * 7, 14

            # Fit text within replacement box width.
            while tw > bw - 4 and font_size > 9:
                font_size -= 1
                font = _font_for(font_size)
                if font:
                    try:
                        bb = draw.textbbox((0, 0), text, font=font)
                        tw = bb[2] - bb[0]
                        th = bb[3] - bb[1]
                    except Exception:
                        tw = len(text) * 7
                        th = font_size

            bg = _sample_bg(rx0, ry0, rx1, ry1)
            draw.rectangle([rx0, ry0, rx1, ry1], fill=(bg[0], bg[1], bg[2], 250))

            # Font color hint with fallback.
            fhex = str(lbl.get("font_color", "#000000")).lstrip("#")
            try:
                fill = (int(fhex[0:2], 16), int(fhex[2:4], 16), int(fhex[4:6], 16), 255)
            except Exception:
                fill = (0, 0, 0, 255)

            align = str(lbl.get("alignment", "center")).lower()
            ty = ry0 + max(0, (bh - th) // 2)
            if align == "right" or target_lang == "ar":
                tx = max(rx0 + 1, rx1 - tw - 1)
            elif align == "left":
                tx = rx0 + 1
            else:
                tx = rx0 + max(0, (bw - tw) // 2)

            if font:
                draw.text((tx, ty), text, fill=fill, font=font)
            else:
                draw.text((tx, ty), text, fill=fill)

        composite = PILImage.alpha_composite(img, overlay).convert("RGB")
        buf = io.BytesIO()
        composite.save(buf, format="PNG")
        return buf.getvalue()

    except Exception as e:
        log.warning("Overlay rendering failed: %s — returning original", e)
        return image_bytes
